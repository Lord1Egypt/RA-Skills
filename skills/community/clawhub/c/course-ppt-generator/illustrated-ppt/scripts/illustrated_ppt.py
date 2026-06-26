#!/usr/bin/env python3
"""
图文并茂PPT生成器
根据文字内容生成匹配的图片并嵌入PPT
"""

import os
import sys
import json
import requests
import datetime
from pathlib import Path
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# API配置
API_URL = "https://api.minimaxi.com/v1/image_generation"
MODEL = "image-01"

# 场景风格配置
SCENE_STYLES = {
    "education": {
        "name": "蓝色科技风",
        "primary": "1E3A5F",
        "accent": "4A90D9",
        "light_bg": "F0F5FA",
        "text": "FFFFFF",
        "prompt_suffix": "教育科技风格，蓝色系PPT幻灯片，高清"
    },
    "business": {
        "name": "专业商务风",
        "primary": "2C3E50",
        "accent": "3498DB",
        "light_bg": "ECF0F1",
        "text": "FFFFFF",
        "prompt_suffix": "专业商务风格PPT幻灯片，高清"
    },
    "tech": {
        "name": "技术简约风",
        "primary": "1A1A2E",
        "accent": "16213E",
        "light_bg": "EAEAEA",
        "text": "FFFFFF",
        "prompt_suffix": "技术简约风格PPT幻灯片，高清"
    },
    "creative": {
        "name": "创意活力风",
        "primary": "6C5CE7",
        "accent": "A29BFE",
        "light_bg": "DFE6E9",
        "text": "FFFFFF",
        "prompt_suffix": "创意活力风格PPT幻灯片，高清"
    }
}


class MiniMaxClient:
    """MiniMax API客户端"""
    
    def __init__(self, api_key: str):
        self.api_key = api_key
        self.headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}"
        }
    
    def generate_image(self, prompt: str, aspect_ratio: str = "16:9") -> Optional[str]:
        """生成单张图片"""
        payload = {
            "model": MODEL,
            "prompt": prompt,
            "aspect_ratio": aspect_ratio,
            "response_format": "url",
            "n": 1,
            "prompt_optimizer": True
        }
        
        try:
            resp = requests.post(API_URL, json=payload, headers=self.headers, timeout=120)
            result = resp.json()
            
            if "data" in result and "image_urls" in result["data"]:
                return result["data"]["image_urls"][0]
        except Exception as e:
            print(f"❌ 生成失败: {e}")
        
        return None
    
    def download_image(self, url: str, save_path: str) -> bool:
        """下载图片"""
        try:
            resp = requests.get(url, timeout=60)
            if resp.status_code == 200:
                with open(save_path, "wb") as f:
                    f.write(resp.content)
                return True
        except Exception as e:
            print(f"❌ 下载失败: {e}")
        return False


class IllustratedPPTGenerator:
    """图文并茂PPT生成器"""
    
    def __init__(self, api_key: str = None):
        self.api_key = api_key or os.environ.get("MINIMAX_API_KEY", "")
        if not self.api_key:
            raise ValueError("请提供MINIMAX_API_KEY")
        
        self.client = MiniMaxClient(self.api_key)
        self.scene = "education"
        self.style = "blue_tech"
    
    def set_scene(self, scene: str):
        """设置场景风格"""
        if scene in SCENE_STYLES:
            self.scene = scene
    
    def set_style(self, style: str):
        """设置视觉风格"""
        self.style = style
    
    def _get_image_prompt(self, slide: Dict) -> str:
        """根据幻灯片内容生成图片描述"""
        slide_type = slide.get("type", "content")
        title = slide.get("title", "")
        content = slide.get("content", "")
        bullets = slide.get("bullets", [])
        chapters = slide.get("chapters", [])
        
        style_cfg = SCENE_STYLES.get(self.scene, SCENE_STYLES["education"])
        suffix = style_cfg["prompt_suffix"]
        
        # 根据类型生成匹配的图片描述
        if slide_type == "cover":
            prompt = f"PPT封面，标题《{title}》，副标题《{content}》，深蓝色背景，几何线条装饰，现代简约教育风格，{suffix}"
        
        elif slide_type == "toc":
            chapter_text = "、".join([f"{i+1} {ch}" for i, ch in enumerate(chapters)])
            prompt = f"PPT目录页，标题《{title}》，章节列表：{chapter_text}，蓝色编号卡片，浅色背景，{suffix}"
        
        elif slide_type == "content":
            bullet_text = "、".join(bullets[:4])
            prompt = f"PPT内容页，标题《{title}》，要点：{bullet_text}，蓝色标题栏，图标+文字布局，浅色背景，{suffix}"
        
        elif slide_type == "code":
            prompt = f"PPT代码演示页，标题《{title}》，深色代码编辑器背景，语法高亮，右侧：相关示意图，{suffix}"
        
        elif slide_type == "summary":
            bullet_text = "、".join([f"✓ {b}" for b in bullets[:4]])
            prompt = f"PPT总结页，标题《{title}》，要点：{bullet_text}，深色背景，白色文字，{suffix}"
        
        elif slide_type == "homework":
            bullet_text = "、".join([f"{i+1}. {b}" for i, b in enumerate(bullets[:4])])
            prompt = f"PPT作业页，标题《{title}》，作业：{bullet_text}，浅色背景，编号卡片，{suffix}"
        
        elif slide_type == "closing":
            prompt = f"PPT结束页，标题《{title}》，副标题《{content}》，深蓝背景，居中文字，装饰线条，{suffix}"
        
        else:
            prompt = f"PPT幻灯片，标题《{title}》，{suffix}"
        
        return prompt
    
    def _create_ppt(self, slides: List[Dict], output_path: str, image_dir: str):
        """创建PPT并嵌入图片"""
        style_cfg = SCENE_STYLES.get(self.scene, SCENE_STYLES["education"])
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for i, slide in enumerate(slides, 1):
            slide_type = slide.get("type", "content")
            title = slide.get("title", "")
            content = slide.get("content", "")
            bullets = slide.get("bullets", [])
            chapters = slide.get("chapters", [])
            
            # 加载图片
            img_path = os.path.join(image_dir, f"slide_{i:02d}.png")
            
            if slide_type == "cover":
                self._add_cover_slide(prs, title, content, img_path, style_cfg)
            elif slide_type == "toc":
                self._add_toc_slide(prs, title, chapters, img_path, style_cfg, i)
            elif slide_type == "content":
                self._add_content_slide(prs, title, bullets, img_path, style_cfg, i)
            elif slide_type == "code":
                self._add_code_slide(prs, title, content, img_path, style_cfg, i)
            elif slide_type == "summary":
                self._add_summary_slide(prs, title, bullets, img_path, style_cfg, i)
            elif slide_type == "homework":
                self._add_homework_slide(prs, title, bullets, img_path, style_cfg, i)
            elif slide_type == "closing":
                self._add_closing_slide(prs, title, content, img_path, style_cfg)
        
        prs.save(output_path)
        return output_path
    
    def _add_cover_slide(self, prs, title, subtitle, img_path, style_cfg):
        """添加封面页"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                     width=Inches(13.333), height=Inches(7.5))
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(13.333), Inches(1.2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(style_cfg["text"])
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(0), Inches(3.8), Inches(13.333), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor.from_string(style_cfg["accent"])
        p.alignment = PP_ALIGN.CENTER
    
    def _add_toc_slide(self, prs, title, chapters, img_path, style_cfg, page_num):
        """添加目录页"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                     width=Inches(13.333), height=Inches(7.5))
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(13.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(style_cfg["primary"])
        p.alignment = PP_ALIGN.CENTER
        
        # 章节列表
        y_start = 2.0
        for j, ch in enumerate(chapters):
            ch_box = slide.shapes.add_textbox(Inches(3), Inches(y_start + j * 1.2), Inches(7), Inches(0.8))
            tf = ch_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{j+1:02d}  {ch}"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor.from_string(style_cfg["primary"])
    
    def _add_content_slide(self, prs, title, bullets, img_path, style_cfg, page_num):
        """添加内容页"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # 背景图片
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(6.5), Inches(0),
                                     width=Inches(6.833), height=Inches(7.5))
        
        # 标题栏
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(6.5), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(style_cfg["primary"])
        shape.line.fill.background()
        
        # 标题文字
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5.5), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(style_cfg["text"])
        
        # 内容区
        content_bg = slide.shapes.add_shape(1, Inches(0), Inches(1.2), Inches(6.5), Inches(6.3))
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = RGBColor.from_string(style_cfg["light_bg"])
        content_bg.line.fill.background()
        
        # 要点列表
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for j, bullet in enumerate(bullets):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor.from_string(style_cfg["primary"])
            p.space_after = Pt(18)
    
    def _add_code_slide(self, prs, title, code, img_path, style_cfg, page_num):
        """添加代码页"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # 右侧图片
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(6.5), Inches(0),
                                     width=Inches(6.833), height=Inches(7.5))
        
        # 左侧深色背景
        left_bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(6.5), Inches(7.5))
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = RGBColor(30, 41, 59)
        left_bg.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(6), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(style_cfg["text"])
        
        # 代码内容
        code_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(5.9), Inches(6))
        tf = code_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code
        p.font.size = Pt(12)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(226, 232, 240)
    
    def _add_summary_slide(self, prs, title, bullets, img_path, style_cfg, page_num):
        """添加总结页"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                     width=Inches(13.333), height=Inches(7.5))
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(13.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(style_cfg["text"])
        p.alignment = PP_ALIGN.CENTER
        
        # 要点
        y_start = 2.5
        for j, bullet in enumerate(bullets):
            pt_box = slide.shapes.add_textbox(Inches(2), Inches(y_start + j * 1.1), Inches(9), Inches(0.8))
            tf = pt_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"✓  {bullet}"
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor.from_string(style_cfg["text"])
            p.alignment = PP_ALIGN.CENTER
    
    def _add_homework_slide(self, prs, title, bullets, img_path, style_cfg, page_num):
        """添加作业页"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                     width=Inches(13.333), height=Inches(7.5))
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(13.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(style_cfg["primary"])
        p.alignment = PP_ALIGN.CENTER
        
        # 作业卡片
        for j, bullet in enumerate(bullets):
            hw_box = slide.shapes.add_textbox(Inches(2), Inches(2 + j * 1.2), Inches(9), Inches(0.8))
            tf = hw_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{j+1}. {bullet}"
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor.from_string(style_cfg["primary"])
    
    def _add_closing_slide(self, prs, title, subtitle, img_path, style_cfg):
        """添加结束页"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                     width=Inches(13.333), height=Inches(7.5))
        
        # 主标题
        title_box = slide.shapes.add_textbox(Inches(0), Inches(3), Inches(13.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(style_cfg["text"])
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), Inches(13.333), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(176, 196, 222)
        p.alignment = PP_ALIGN.CENTER
    
    def generate(self, topic: str, slides: List[Dict], scene: str = "education",
                 output_dir: str = "/tmp", output_name: str = None) -> Dict:
        """
        生成图文并茂PPT
        
        Args:
            topic: PPT主题
            slides: 幻灯片内容列表
            scene: 场景风格 (education/business/tech/creative)
            output_dir: 输出目录
            output_name: 输出文件名
        
        Returns:
            dict: 包含output_path等信息
        """
        self.set_scene(scene)
        
        if not output_name:
            safe_topic = topic.replace(" ", "_").replace("/", "-")
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_name = f"{safe_topic}_{timestamp}.pptx"
        
        output_path = os.path.join(output_dir, output_name)
        image_dir = os.path.join(output_dir, "illustrated_images")
        os.makedirs(image_dir, exist_ok=True)
        
        print(f"📚 开始生成 {topic} PPT...")
        print(f"🎨 场景风格: {SCENE_STYLES[scene]['name']}")
        print(f"📄 共 {len(slides)} 页")
        
        # 生成图片
        for i, slide in enumerate(slides, 1):
            prompt = self._get_image_prompt(slide)
            print(f"\n[{i}/{len(slides)}] 生成图片...")
            print(f"   描述: {prompt[:50]}...")
            
            img_url = self.client.generate_image(prompt)
            if img_url:
                img_path = os.path.join(image_dir, f"slide_{i:02d}.png")
                if self.client.download_image(img_url, img_path):
                    print(f"   ✅ 已保存: {img_path}")
                else:
                    print(f"   ⚠️ 下载失败，使用占位符")
            else:
                print(f"   ❌ 图片生成失败")
        
        # 创建PPT
        print(f"\n📦 生成PPT文件...")
        self._create_ppt(slides, output_path, image_dir)
        print(f"✅ 完成: {output_path}")
        
        return {
            "output_path": output_path,
            "image_dir": image_dir,
            "topic": topic,
            "scene": scene,
            "slide_count": len(slides)
        }


def main():
    if len(sys.argv) < 2:
        print("用法: python illustrated_ppt.py <主题> [场景] [输出路径]")
        print("场景选项: education / business / tech / creative")
        print("示例: python illustrated_ppt.py '深度学习课程' education /tmp/output.pptx")
        sys.exit(1)
    
    topic = sys.argv[1]
    scene = sys.argv[2] if len(sys.argv) > 2 else "education"
    output_path = sys.argv[3] if len(sys.argv) > 3 else None
    
    api_key = os.environ.get("MINIMAX_API_KEY", "")
    if not api_key:
        print("❌ 请设置环境变量 MINIMAX_API_KEY")
        sys.exit(1)
    
    output_dir = os.path.dirname(output_path) if output_path else "/tmp"
    output_name = os.path.basename(output_path) if output_path else None
    
    # 默认幻灯片结构
    default_slides = [
        {"type": "cover", "title": topic, "content": "图文并茂演示"},
        {"type": "toc", "title": "目 录", "chapters": ["第1章", "第2章", "第3章"]},
        {"type": "content", "title": "核心概念", "bullets": ["概念一", "概念二", "概念三"]},
        {"type": "content", "title": "应用实践", "bullets": ["实践一", "实践二", "实践三"]},
        {"type": "closing", "title": "本文由AI辅助创作", "content": "作者：TJMtaotao"},
    ]
    
    generator = IllustratedPPTGenerator(api_key)
    result = generator.generate(topic, default_slides, scene, output_dir, output_name)
    
    print(f"\n📋 结果:")
    print(f"   主题: {result['topic']}")
    print(f"   场景: {result['scene']}")
    print(f"   页数: {result['slide_count']}")
    print(f"   路径: {result['output_path']}")


if __name__ == "__main__":
    main()