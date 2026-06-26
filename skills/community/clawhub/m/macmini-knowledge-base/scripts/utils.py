#!/usr/bin/env python3
"""
知识库共享工具函数 v2.0
集成 kreuzberg（统一提取层）+ antiword（.doc 极速提取）+ soffice（兜底）
"""
import os
import json
import shutil
import subprocess
import re as re_module
import uuid
import pymupdf
import docx
import openpyxl
import pptx

try:
    from kreuzberg import extract_file_sync as kreuzberg_extract
    KREUZBERG_AVAILABLE = True
except Exception:
    KREUZBERG_AVAILABLE = False

_DEFAULT_KNOWLEDGE_DIR = os.path.expanduser("~/.openclaw/workspace/knowledge")
_STATE_FILE    = os.path.join(_DEFAULT_KNOWLEDGE_DIR, ".analysis/analysis_state.json")
_CACHE_FILE    = os.path.join(_DEFAULT_KNOWLEDGE_DIR, ".analysis/.catalog_cache.json")
_PROGRESS_FILE = os.path.join(_DEFAULT_KNOWLEDGE_DIR, ".analysis/.catalog_progress.json")

def load_json(path=None):
    if path is None:
        path = _STATE_FILE
    if os.path.exists(path):
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except (json.JSONDecodeError, IOError):
            pass
    return None

def save_json(path, data):
    if path is None:
        path = _STATE_FILE
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def load_state():
    data = load_json(_STATE_FILE)
    return data if data else {}

def save_state(state):
    save_json(_STATE_FILE, state)

def load_cache():
    data = load_json(_CACHE_FILE)
    return data if data else {}

def save_cache(cache):
    save_json(_CACHE_FILE, cache)

def load_progress():
    data = load_json(_PROGRESS_FILE)
    return data if data else {}

def save_progress(progress):
    save_json(_PROGRESS_FILE, progress)

def is_gibberish(text, strict=False):
    if not text or len(text.strip()) < 20:
        return True
    alpha_count   = sum(1 for c in text if c.isalpha())
    upper_count   = sum(1 for c in text if c.isupper())
    space_count   = sum(1 for c in text if c.isspace())
    chinese_count = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
    chinese_ratio = chinese_count / len(text) if len(text) > 0 else 0
    if alpha_count > 10:
        upper_ratio = upper_count / alpha_count
        space_ratio = space_count / len(text)
        if upper_ratio > 0.45 and space_ratio < 0.05:
            return True
    if chinese_ratio > 0.3:
        return False
    cleaned = re_module.sub(r'[\s\n]', '', text)
    if len(cleaned) == 0:
        return True
    normal_count = len(re_module.findall(r'[A-Za-z0-9 .,;:\'"!?()-]', cleaned))
    ratio = normal_count / len(cleaned) if len(cleaned) > 0 else 0
    return ratio < (0.5 if strict else 0.3)

def extract_via_kreuzberg(filepath):
    if not KREUZBERG_AVAILABLE:
        return None
    try:
        result = kreuzberg_extract(filepath)
        if result and result.content and result.content.strip():
            return result.content
    except Exception:
        pass
    return None

def extract_doc_via_antiword(filepath):
    try:
        result = subprocess.run(['antiword', filepath], capture_output=True, timeout=10)
        if result.returncode == 0:
            text = result.stdout.decode('utf-8', errors='replace')
            if len(text) > 100:
                return text
    except Exception:
        pass
    return None

def _kill_proc_tree(pid):
    try:
        os.system(f"pkill -P {pid} 2>/dev/null")
    except:
        pass
    try:
        os.kill(pid, 9)
    except:
        pass

def convert_old_office(filepath, ext):
    tmp_dir = "/tmp/office_convert"
    os.makedirs(tmp_dir, exist_ok=True)
    convert_map = {".doc": "docx", ".xls": "xlsx", ".ppt": "pptx"}
    new_ext = convert_map.get(ext)
    if not new_ext:
        return (False, "")
    expected_name = os.path.basename(filepath).replace(ext, '.' + new_ext)
    tmp_out = os.path.join(tmp_dir, expected_name)
    if os.path.exists(tmp_out):
        try:
            os.remove(tmp_out)
        except:
            pass
    proc = subprocess.Popen(
        ['soffice', '--headless', '--convert-to', new_ext, '--outdir', tmp_dir, filepath],
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
    )
    try:
        outs, errs = proc.communicate(timeout=60)
        if proc.returncode == 0 and os.path.exists(tmp_out):
            return (True, tmp_out)
    except subprocess.TimeoutExpired:
        try:
            proc.kill()
        except:
            pass
        _kill_proc_tree(proc.pid)
    return (False, "")

def extract_doc_text(filepath):
    text = extract_doc_via_antiword(filepath)
    if text:
        return text
    success, converted = convert_old_office(filepath, '.doc')
    if success:
        try:
            doc = docx.Document(converted)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text[:8000] if text.strip() else None
        finally:
            try:
                os.remove(converted)
            except:
                pass
    return "【DOC提取失败】"

def extract_pdf_text(filepath):
    text = extract_via_kreuzberg(filepath)
    if text and not is_gibberish(text, strict=True):
        return text
    try:
        doc = pymupdf.open(filepath)
        text = ""
        for page_num in range(min(len(doc), 10)):
            t = doc[page_num].get_text()
            if t.strip():
                text += t + "\n"
        doc.close()
        if text.strip() and not is_gibberish(text, strict=True):
            return text
    except:
        pass
    return extract_pdf_via_ocr(filepath)

def extract_docx_text(filepath):
    text = extract_via_kreuzberg(filepath)
    if text:
        return text
    try:
        doc = docx.Document(filepath)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return text[:8000] if text.strip() else "【Word文档无文字内容】"
    except Exception as e:
        return "【DOCX提取失败】" + str(e)

def extract_xlsx_text(filepath):
    text = extract_via_kreuzberg(filepath)
    if text:
        return text
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        sheets_text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    sheets_text.append(row_text)
        return "\n".join(sheets_text)[:8000] if sheets_text else "【Excel无内容】"
    except Exception as e:
        return "【XLSX提取失败】" + str(e)

def extract_pptx_text(filepath):
    text = extract_via_kreuzberg(filepath)
    if text:
        return text
    try:
        prs = pptx.Presentation(filepath)
        text_parts = []
        for i, slide in enumerate(prs.slides[:20]):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append(f"[Page {i+1}]\n" + "\n".join(slide_text))
        return "\n".join(text_parts)[:8000] if text_parts else "【PPT无文字内容】"
    except Exception as e:
        return "【PPTX提取失败】" + str(e)

def extract_pdf_via_ocr(filepath):
    tmp_dir = "/tmp/ocrmypdf_work"
    os.makedirs(tmp_dir, exist_ok=True)
    tmp_pdf = os.path.join(tmp_dir, uuid.uuid4().hex + ".pdf")
    try:
        shutil.copy2(filepath, tmp_pdf)
        result = subprocess.run(
            ["ocrmypdf", "-l", "chi_sim+eng", "--skip-text",
             "--pages", "1-10", tmp_pdf, tmp_pdf],
            capture_output=True, timeout=40
        )
    except Exception:
        pass
    finally:
        if os.path.exists(tmp_pdf):
            try:
                os.remove(tmp_pdf)
            except:
                pass
    try:
        doc = pymupdf.open(filepath)
        text = ""
        for page_num in range(min(len(doc), 10)):
            t = doc[page_num].get_text()
            if t.strip():
                text += t + "\n"
        doc.close()
        if text.strip():
            return text
    except:
        pass
    return "【PDF文字提取失败】"
