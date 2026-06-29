#!/usr/bin/env python3
"""
fix_fonts.py — 批量修复 PPTX 字号 + 边界
用法: python fix_fonts.py <输入pptx> <输出pptx> [--min-size 16]

功能:
  1. 将所有 < min_size (pt) 的 run.font.size 调整为 min_size
  2. 启用 text_frame.word_wrap = True 防止横向溢出
  3. 修正超出幻灯片边界的形状（右溢出、下溢出、负偏移）
  4. 输出修复统计

原则: 只改格式(run.font.size + word_wrap + shape bounds)，不动文字内容(run.text)
"""
import sys
import os
import argparse

def fix_fonts(input_path, output_path, min_size=16):
    from pptx import Presentation
    from pptx.util import Pt

    if not os.path.exists(input_path):
        print(f"错误：输入文件不存在 — {input_path}")
        sys.exit(1)

    prs = Presentation(input_path)
    MIN_FONT_SIZE = Pt(min_size)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    print(f"文件：{os.path.basename(input_path)}")
    print(f"页数：{len(prs.slides)}")
    print(f"最小字号阈值：{min_size}pt")
    print(f"幻灯片尺寸：{slide_w} x {slide_h} EMU")
    print()

    stats = {
        "total_runs": 0,
        "fixed_runs": 0,
        "already_ok": 0,
        "empty_or_none": 0,
        "boundary_fixes": 0,
    }

    for slide_idx, slide in enumerate(prs.slides):
        slide_fixes = 0
        for shape in slide.shapes:
            shape_modified = False

            if shape.has_text_frame:
                tf = shape.text_frame

                # 启用自动换行
                if tf.word_wrap is not True:
                    tf.word_wrap = True
                    shape_modified = True

                for para in tf.paragraphs:
                    for run in para.runs:
                        stats["total_runs"] += 1
                        current_size = run.font.size

                        if current_size is None:
                            stats["empty_or_none"] += 1
                            continue

                        if current_size < MIN_FONT_SIZE:
                            run.font.size = MIN_FONT_SIZE
                            stats["fixed_runs"] += 1
                            slide_fixes += 1
                            shape_modified = True
                        else:
                            stats["already_ok"] += 1

        if slide_fixes > 0:
            print(f"  第{slide_idx + 1}页：修复 {slide_fixes} 个文本运行")

    # 第二阶段：独立边界修正（覆盖所有形状，不限于字号被修改的）
    print()
    print("第二遍：修正超出幻灯片边界的形状...")
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            sl = shape.left or 0
            st = shape.top or 0
            sw = shape.width or 0
            sh = shape.height or 0
            boundary_fixed = False

            if sl < 0:
                shape.left = 0
                boundary_fixed = True

            if st < 0:
                shape.top = 0
                boundary_fixed = True

            # 重新读取修正后的值
            sl = shape.left or 0
            sw = shape.width or 0

            if sl + sw > slide_w and sw > 0:
                new_w = slide_w - sl
                if new_w > 0:
                    shape.width = new_w
                    boundary_fixed = True

            st = shape.top or 0
            sh = shape.height or 0

            if st + sh > slide_h and sh > 0:
                new_h = slide_h - st
                if new_h > 0:
                    shape.height = new_h
                    boundary_fixed = True

            if boundary_fixed:
                stats["boundary_fixes"] += 1

    print(f"\n{'=' * 60}")
    print(f"修复摘要：")
    print(f"  文本运行总数：  {stats['total_runs']}")
    print(f"  已 ≥{min_size}pt：   {stats['already_ok']}")
    print(f"  已修复至 {min_size}pt：{stats['fixed_runs']}")
    print(f"  无字号/空文本： {stats['empty_or_none']}")
    print(f"  边界修正形状：  {stats['boundary_fixes']}")

    # 保存
    try:
        prs.save(output_path)
        print(f"\n已保存到：{output_path}")
    except PermissionError:
        print(f"\n错误：无法写入 {output_path}（文件被占用）")
        sys.exit(1)

    # 验证
    print("\n验证中...")
    prs2 = Presentation(output_path)
    bad_font = 0
    bad_bound = 0
    total = 0
    for si, s in enumerate(prs2.slides):
        for sh in s.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        total += 1
                        if r.font.size and r.font.size < MIN_FONT_SIZE:
                            bad_font += 1
            sl2 = sh.left or 0
            st2 = sh.top or 0
            sr = sl2 + (sh.width or 0)
            sb = st2 + (sh.height or 0)
            if sl2 < -1000 or st2 < -1000 or sr > slide_w + 1000 or sb > slide_h + 1000:
                bad_bound += 1
    print(f"  文本运行总数：{total}，字号 <{min_size}pt：{bad_font}，边界问题：{bad_bound}")
    if bad_font == 0 and bad_bound == 0:
        print("  [PASS] 全部通过")
    else:
        print(f"  [FAIL] 仍有问题：字号 {bad_font} + 边界 {bad_bound}")


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='修复 PPTX 字号和边界')
    parser.add_argument('input', help='输入 PPTX 路径')
    parser.add_argument('output', help='输出 PPTX 路径')
    parser.add_argument('--min-size', type=int, default=16, help='最小字号(pt)，默认 16')
    args = parser.parse_args()
    fix_fonts(args.input, args.output, args.min_size)
