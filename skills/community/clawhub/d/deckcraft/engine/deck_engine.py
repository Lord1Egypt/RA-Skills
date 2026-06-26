"""
DeckCraft v6 — DeckEngine
High-level layout API for creating professional PPTX presentations.

## Role Modes (v6.0+)

DeckEngine supports two role modes:

- **single** (default): The classic one-pass workflow. LLM generates complete
  content.json in a single call. Identical to v5.3.0 behavior. No changes
  to existing API — fully backward compatible.

- **multi**: An optional two-phase workflow inspired by "Strategist → Executor"
  role separation. In this mode, the LLM first produces a strategic plan
  (page count, key points per page, visual style) via ``strategist_plan()``,
  then generates actual slides via ``execute_plan()``.

  Note: The ``multi`` mode only provides the API scaffolding. Users are
  responsible for calling their LLM to produce the plan JSON. DeckEngine
  does NOT call any LLM internally.

Example::

    # single mode (default, unchanged)
    eng = DeckEngine(theme_name="business", canvas="16:9")
    eng.cover(title="Hello")
    eng.save("out.pptx")

    # multi mode (strategist first, then execute)
    eng = DeckEngine(theme_name="business", canvas="16:9", role_mode="multi")
    plan = eng.strategist_plan(brief)     # returns plan template
    # ... user calls LLM with plan template, gets filled plan.json ...
    deck = eng.execute_plan(plan)         # returns self for chaining
    eng.save("out.pptx")
"""

import os
import json
from typing import Optional, List, Tuple, Any, Literal, Union
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from .constants import (
    GAP_SMALL, GAP_MEDIUM, GAP_LARGE,
    get_canvas, list_canvases, CANVAS_PRESETS, CANVAS_ALIASES,
    TITLE_SIZE, SUBTITLE_SIZE, BODY_SIZE, BODY_SMALL, CAPTION_SIZE,
    TABLE_HEADER_SIZE, TABLE_BODY_SIZE, STAT_NUMBER_SIZE, STAT_LABEL_SIZE,
    TOP_BAR_HEIGHT, ACCENT_LINE_HEIGHT,
    FONT_CN, FONT_EN_BODY, FONT_EN_TITLE,
    c, get_theme, THEMES
)
from .core import (
    set_bg, add_rect, add_rounded_rect, add_textbox,
    add_multiline_textbox, add_line, add_circle, add_image_safe,
    clean_shape, full_cleanup, auto_cjk_font
)
from .chart_engine import (
    bar_chart, pie_chart, line_chart, gauge_chart,
    chart_funnel, chart_gantt, chart_swot, chart_porter, chart_sankey,
    chart_heatmap, chart_radar, chart_treemap, chart_waterfall,
)


# ── Input validation helpers ──────────────────────────────────────────────────

def _validate_text(value: Any, name: str, allow_none: bool = True, max_len: int = 500) -> Optional[str]:
    """Validate a text parameter. Returns the (possibly stripped) string or None.

    Raises TypeError if value is not str/None.
    Raises ValueError if empty (when allow_none=False) or too long.
    """
    if value is None:
        if not allow_none:
            raise ValueError(f"{name} is required (got None)")
        return None
    if not isinstance(value, str):
        raise TypeError(f"{name} must be str, got {type(value).__name__}: {value!r}")
    stripped = value.strip() if value else ""
    if not stripped:
        if not allow_none:
            raise ValueError(f"{name} is required (got empty string)")
        return None
    if max_len and len(stripped) > max_len:
        raise ValueError(f"{name} is too long ({len(stripped)} > {max_len} chars). Truncate or split.")
    return stripped


def _validate_int(value: Any, name: str, min_val: int = 0) -> int:
    """Validate a non-negative integer parameter."""
    if not isinstance(value, int) or isinstance(value, bool):
        raise TypeError(f"{name} must be int, got {type(value).__name__}: {value!r}")
    if value < min_val:
        raise ValueError(f"{name} must be >= {min_val}, got {value}")
    return value


def _validate_list(value: Any, name: str, min_items: int = 0, max_items: int = 100) -> List:
    """Validate a list parameter with item count bounds."""
    if not isinstance(value, (list, tuple)):
        raise TypeError(f"{name} must be list/tuple, got {type(value).__name__}")
    if len(value) < min_items:
        raise ValueError(f"{name} must have at least {min_items} items, got {len(value)}")
    if len(value) > max_items:
        raise ValueError(f"{name} must have at most {max_items} items, got {len(value)}")
    return list(value)


def _validate_image_path(value: Any, name: str) -> Optional[str]:
    """Validate image_path. Returns path string or None. Warns if file missing but doesn't fail (for graceful render)."""
    if value is None or value == "":
        return None
    if not isinstance(value, str):
        raise TypeError(f"{name} must be str, got {type(value).__name__}")
    if not os.path.isfile(value):
        # Warn but don't raise — image may be added later in some workflows
        import warnings
        warnings.warn(f"{name} file not found: {value!r}. Slide will be created without image.", UserWarning, stacklevel=3)
    return value


class DeckEngine:
    """Create professional PPTX presentations with high-level layout methods.

    Args:
        theme_name: One of the 10 built-in themes (see THEMES keys).
                    Default: "business". Use THEMES dict for the full list.
        canvas: One of 6 canvas presets (16:9, 9:16, 1:1, 4:3, A4, A4-portrait)
                plus aliases (mobile, square, ppt, ppt-16x9).
                Default: "16:9". Use list_canvases() for the full list.
        role_mode: "single" (default) for classic one-pass workflow,
                   "multi" for optional strategist→executor two-phase workflow.
                   v6.0+ feature; does not affect single-mode behavior.

    Raises:
        ValueError: If theme_name, canvas, or role_mode is unknown.
    """

    def __init__(self, theme_name: str = "business", canvas: str = "16:9",
                 role_mode: Literal["single", "multi"] = "single"):
        # Validate role_mode
        if role_mode not in ("single", "multi"):
            raise ValueError(
                f"Unknown role_mode: {role_mode!r}. Must be 'single' or 'multi'."
            )
        self.role_mode: Literal["single", "multi"] = role_mode

        # Validate theme
        if theme_name not in THEMES:
            raise ValueError(
                f"Unknown theme: {theme_name!r}. Available: {list(THEMES.keys())}. "
                f"Use list_canvases() to inspect."
            )
        # Validate canvas (v6.0+: resolve aliases first, e.g. xiaohongshu → 9:16)
        resolved_canvas = CANVAS_ALIASES.get(canvas, canvas)
        if resolved_canvas not in CANVAS_PRESETS:
            raise ValueError(
                f"Unknown canvas: {canvas!r}. Available: {list(CANVAS_PRESETS.keys())} "
                f"+ aliases {list(CANVAS_ALIASES.keys())}. Use list_canvases() to inspect."
            )

        self.theme_name: str = theme_name
        self.theme: dict = THEMES[theme_name]

        # Canvas configuration (v5.2+, alias resolved in v6.0+)
        self.canvas_name: str = resolved_canvas
        self.canvas_alias: Optional[str] = canvas if canvas in CANVAS_ALIASES else None
        self.canvas: dict = get_canvas(resolved_canvas)
        self.cw = self.canvas["width"]          # slide width (Emu)
        self.ch = self.canvas["height"]         # slide height (Emu)
        self.ml = self.canvas["margin_left"]    # margin left
        self.mr = self.canvas["margin_right"]   # margin right
        self.mt = self.canvas["margin_top"]     # margin top
        self.mb = self.canvas["margin_bottom"]  # margin bottom
        self.content_w = self.canvas["content_width"]
        self.content_top = self.canvas["content_top"]
        # 缩放因子（相对 16:9 参考值，供可选比例缩放使用）
        self.scale_x = self.canvas["width_in"] / 10.0
        self.scale_y = self.canvas["height_in"] / 5.625
        # 创建 Presentation
        self.prs: Presentation = Presentation()
        self.prs.slide_width = self.cw
        self.prs.slide_height = self.ch
        self._slide_count: int = 0

        # Multi-mode state
        self._strategist_plan: Optional[dict] = None

    # ── Multi-mode API (v6.0+) ────────────────────────────────────

    def strategist_plan(self, brief: dict) -> dict:
        """Generate a strategist plan template from a brief.

        In ``role_mode="multi"``, call this first to get a plan template.
        Then feed the template to your LLM to fill in the strategic decisions
        (page count, key points per page, visual style, etc.).

        The returned plan has placeholder values — it's a schema, not content.
        Your LLM should replace the placeholders with actual strategy.

        Args:
            brief: Dict with keys like 'title', 'audience', 'goal', 'duration'.
                   Minimal validation — just needs to be a dict.

        Returns:
            A plan dict (template) ready for LLM completion.

        Raises:
            RuntimeError: If role_mode is not "multi".
            TypeError: If brief is not a dict.
        """
        if self.role_mode != "multi":
            raise RuntimeError(
                "strategist_plan() requires role_mode='multi'. "
                f"Current mode: {self.role_mode!r}"
            )
        if not isinstance(brief, dict):
            raise TypeError(f"brief must be dict, got {type(brief).__name__}")

        plan = {
            "meta": {
                "source_brief": brief.get("title", "Untitled"),
                "theme": self.theme_name,
                "canvas": self.canvas_name,
            },
            "strategy": {
                "total_pages": None,         # int: target page count
                "page_allocation": None,     # list of dicts: [{"page": 1, "type": "cover", "key_points": []}]
                "visual_style": None,        # str: e.g. "professional", "creative"
                "emphasis_pages": [],        # list of page numbers to emphasize
            },
            "content_plan": None,            # list of page content dicts (filled by LLM)
            "notes": "Fill strategy.* and content_plan, then call execute_plan().",
        }
        self._strategist_plan = plan
        return plan

    def execute_plan(self, plan: dict) -> "DeckEngine":
        """Execute a filled strategist plan to generate slides.

        In ``role_mode="multi"``, call this after your LLM has filled the plan
        template returned by ``strategist_plan()``.

        The plan dict's ``content_plan`` should be a list of page dicts, each
        with ``"type"`` and type-specific keys matching DeckEngine method names
        (e.g. ``cover``, ``content``, ``closing``).

        Args:
            plan: A filled plan dict (from strategist_plan, completed by LLM).

        Returns:
            self (for chaining).

        Raises:
            RuntimeError: If role_mode is not "multi".
            TypeError: If plan is not a dict or content_plan is missing/empty.
        """
        if self.role_mode != "multi":
            raise RuntimeError(
                "execute_plan() requires role_mode='multi'. "
                f"Current mode: {self.role_mode!r}"
            )
        if not isinstance(plan, dict):
            raise TypeError(f"plan must be dict, got {type(plan).__name__}")

        content_plan = plan.get("content_plan")
        if not content_plan or not isinstance(content_plan, list):
            raise TypeError(
                "plan['content_plan'] must be a non-empty list of page dicts. "
                "Have your LLM fill this before calling execute_plan()."
            )

        # Map page type strings to DeckEngine methods
        _METHOD_MAP = {
            "cover": self.cover,
            "toc": self.toc,
            "section_divider": self.section_divider,
            "content": self.content,
            "content_with_icon": self.content_with_icon,
            "two_col": self.two_col,
            "vs_compare": self.vs_compare,
            "table": self.table,
            "stat_cards": self.stat_cards,
            "chart_bar": self.chart_bar,
            "chart_pie": self.chart_pie,
            "chart_line": self.chart_line,
            "chart_gauge": self.chart_gauge,
            "timeline": self.timeline,
            "process_flow": self.process_flow,
            "matrix_2x2": self.matrix_2x2,
            "quote": self.quote,
            "image_full": self.image_full,
            "image_split": self.image_split,
            "kpi_dashboard": self.kpi_dashboard,
            "team_grid": self.team_grid,
            "checklist": self.checklist,
            "summary": self.summary,
            "closing": self.closing,
        }

        for page in content_plan:
            page_type = page.get("type")
            if not page_type or page_type not in _METHOD_MAP:
                import warnings
                warnings.warn(
                    f"Unknown page type {page_type!r} in plan, skipping.",
                    UserWarning, stacklevel=2,
                )
                continue
            method = _METHOD_MAP[page_type]
            # Pass through all keys except "type" as kwargs
            kwargs = {k: v for k, v in page.items() if k != "type"}
            method(**kwargs)

        return self

    @property
    def slide_count(self) -> int:
        """Number of slides created so far."""
        return self._slide_count

    def _new_slide(self):
        """Add a blank slide and return it."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._slide_count += 1
        return slide

    def _add_top_bar(self, slide):
        """Add accent bar at top of content slides."""
        add_rect(slide, 0, 0, self.cw, TOP_BAR_HEIGHT, self.theme["accent"])

    def _add_title_area(self, slide, title, subtitle=None):
        """Add title + optional subtitle + accent line."""
        add_textbox(slide, self.ml, Inches(0.35), self.content_w, Inches(0.45),
                   title, size=Pt(26), bold=True, color=self.theme["primary"])
        add_line(slide, self.ml, Inches(0.82), Inches(1.2), self.theme["accent"], 0.035)
        if subtitle:
            add_textbox(slide, self.ml, Inches(0.9), self.content_w, Inches(0.25),
                       subtitle, size=Pt(11), color=self.theme.get("text_muted", (120,120,120)))

    def _add_page_number(self, slide, num, total=None):
        """Add page number at bottom right."""
        text = f"{num}" if total is None else f"{num}/{total}"
        add_textbox(slide, self.cw - Inches(1.0), self.ch - Inches(0.525), Inches(0.7), Inches(0.3),
                   text, size=CAPTION_SIZE, color=self.theme.get("text_muted", (120,120,120)),
                   align=PP_ALIGN.RIGHT)

    # ── Cover & Closing ───────────────────────────────────────────

    def cover(self, title: str, subtitle: str = "", author: str = "",
              date: str = "", image_path: Optional[str] = None):
        """Title slide with balanced center-left layout.

        Args:
            title: Main title text (required, max 200 chars).
            subtitle: Subtitle line (optional).
            author: Author/byline (optional).
            date: Date or version (optional).
            image_path: Background image path (optional, e.g. logo).

        Raises:
            TypeError: If any string param is not str.
            ValueError: If title is empty/too long.
        """
        # Input validation
        title = _validate_text(title, "title", allow_none=False, max_len=200)
        subtitle = _validate_text(subtitle, "subtitle") or ""
        author = _validate_text(author, "author") or ""
        date = _validate_text(date, "date") or ""
        image_path = _validate_image_path(image_path, "image_path")

        slide = self._new_slide()
        set_bg(slide, self.theme["bg_dark"])

        if image_path:
            add_image_safe(slide, image_path, 0, 0, self.cw, self.ch)
            add_rect(slide, 0, 0, self.cw, self.ch, (0, 0, 0), alpha=0.55)
        
        # Accent line — tight to title group
        add_line(slide, self.ml, Inches(1.8), Inches(1.8), self.theme["accent"], 0.05)
        
        # Title — centered vertically in upper-middle zone
        add_textbox(slide, self.ml, Inches(1.95), self.content_w, Inches(1.2),
                   title, size=Pt(36), bold=True, color=self.theme["text_light"])
        
        # Subtitle — close below title
        if subtitle:
            add_textbox(slide, self.ml, Inches(3.2), self.content_w, Inches(0.5),
                       subtitle, size=Pt(16), color=self.theme["secondary"])
        
        # Author/date — bottom area, higher contrast
        info_parts = [p for p in [author, date] if p]
        if info_parts:
            add_textbox(slide, self.ml, Inches(4.6), self.content_w, Inches(0.4),
                       "  |  ".join(info_parts), size=BODY_SMALL, color=self.theme["text_light"])
        
        # Decorative side accent
        add_rect(slide, 0, 0, Inches(0.12), self.ch, self.theme["accent"])
        return slide

    def closing(self, title: str = "Thank You", message: str = "", contact: str = ""):
        """Closing slide with centered layout.

        Args:
            title: Closing title (default: "Thank You").
            message: Sub-message (optional).
            contact: Contact info (optional).
        """
        title = _validate_text(title, "title", allow_none=False, max_len=100) or "Thank You"
        message = _validate_text(message, "message") or ""
        contact = _validate_text(contact, "contact") or ""

        slide = self._new_slide()
        set_bg(slide, self.theme["bg_dark"])
        
        add_line(slide, self.cw / 2 - Inches(1.5), Inches(1.8), Inches(3), self.theme["accent"], 0.05)
        add_textbox(slide, self.ml, Inches(2.0), self.content_w, Inches(1),
                   title, size=Pt(36), bold=True, color=self.theme["text_light"],
                   align=PP_ALIGN.CENTER)
        
        if message:
            add_textbox(slide, self.ml, Inches(3.2), self.content_w, Inches(0.5),
                       message, size=BODY_SIZE, color=self.theme["secondary"],
                       align=PP_ALIGN.CENTER)
        if contact:
            add_textbox(slide, self.ml, Inches(4.0), self.content_w, Inches(0.4),
                       contact, size=BODY_SMALL, color=self.theme["text_light"],
                       align=PP_ALIGN.CENTER)
        
        # Side accent
        add_rect(slide, 0, 0, Inches(0.12), self.ch, self.theme["accent"])
        return slide

    # ── Table of Contents ──────────────────────────────────────────

    def toc(self, items):
        """
        Table of contents slide.
        items: list of (number, title, description) tuples
        """
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, "Agenda")
        
        y = self.content_top
        for num, title, desc in items:
            # Number circle
            add_circle(slide, Inches(1.1), y + Inches(0.2), Inches(0.2), self.theme["primary"])
            add_textbox(slide, Inches(0.85), y + Inches(0.02), Inches(0.5), Inches(0.35),
                       str(num), size=Pt(13), bold=True, color=self.theme["text_light"],
                       align=PP_ALIGN.CENTER)
            # Title + description
            add_textbox(slide, self.ml + Inches(1.0), y, self.content_w - Inches(1.3), Inches(0.3),
                       title, size=Pt(16), bold=True, color=self.theme["primary"])
            if desc:
                add_textbox(slide, self.ml + Inches(1.0), y + Inches(0.3), self.content_w - Inches(1.3), Inches(0.25),
                           desc, size=BODY_SMALL, color=self.theme.get("text_muted", (120,120,120)))
            y += Inches(0.7) if desc else Inches(0.5)
        
        return slide

    # ── Content Slides ─────────────────────────────────────────────

    def content(self, title: str, bullets: List[str], key_point: str = "",
               image_path: Optional[str] = None, page_num: Optional[int] = None):
        """Standard content slide with bullets, optional image on right.

        Args:
            title: Slide title (required).
            bullets: List of bullet strings (1–20 items).
            key_point: Highlighted bottom card text (optional).
            image_path: Optional right-side image path.
            page_num: Page number to display (optional).

        Raises:
            ValueError: If bullets is empty or exceeds 20.
        """
        title = _validate_text(title, "title", allow_none=False, max_len=200)
        bullets = _validate_list(bullets, "bullets", min_items=1, max_items=20)
        for i, b in enumerate(bullets):
            if not isinstance(b, str):
                raise TypeError(f"bullets[{i}] must be str, got {type(b).__name__}: {b!r}")
        key_point = _validate_text(key_point, "key_point") or ""
        image_path = _validate_image_path(image_path, "image_path")
        if page_num is not None:
            page_num = _validate_int(page_num, "page_num", min_val=1)

        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        text_width = Inches(4.8) if image_path else self.content_w
        bullet_y = Inches(1.4)
        
        add_multiline_textbox(slide, self.ml, bullet_y, text_width, Inches(2.5),
                             bullets, size=Pt(16), color=self.theme["text"],
                             bullet=True, spacing_pt=18)
        
        if image_path:
            add_image_safe(slide, image_path, Inches(5.8), bullet_y, Inches(3.6), Inches(3.0))
        
        if key_point:
            key_y = Inches(3.8)
            add_rounded_rect(slide, self.ml, key_y, self.content_w, Inches(0.5),
                            self.theme["bg_dark"])
            add_textbox(slide, self.ml + Inches(0.15), key_y + Inches(0.08),
                       self.content_w - Inches(0.3), Inches(0.4),
                       key_point, size=Pt(12), bold=True, color=self.theme["accent"])
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def content_with_icon(self, title, items, page_num=None):
        """
        Content slide with icon-style rows.
        items: list of (icon_text, heading, description)
        icon_text: short text/emoji shown in colored circle
        """
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        y = self.content_top
        for icon, heading, desc in items:
            # Icon circle
            add_circle(slide, Inches(1.1), y + Inches(0.2), Inches(0.22), self.theme["primary"])
            add_textbox(slide, Inches(0.85), y + Inches(0.02), Inches(0.5), Inches(0.35),
                       icon[:2], size=Pt(10), bold=True, color=self.theme["text_light"],
                       align=PP_ALIGN.CENTER)
            add_textbox(slide, self.ml + Inches(1.0), y, self.content_w - Inches(1.3), Inches(0.3),
                       heading, size=Pt(15), bold=True, color=self.theme["primary"])
            if desc:
                add_textbox(slide, self.ml + Inches(1.0), y + Inches(0.3), self.content_w - Inches(1.3), Inches(0.4),
                           desc, size=BODY_SMALL, color=self.theme.get("text_muted", (120,120,120)))
            y += Inches(0.8)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── Comparison ─────────────────────────────────────────────────

    def two_col(self, title, left_title, left_items, right_title, right_items, page_num=None):
        """Two-column comparison slide."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        col_width = (self.content_w - Inches(0.3)) / 2  # 2 cols + gap
        col_gap = Inches(0.2)
        left_x = self.ml
        right_x = self.ml + col_width + col_gap
        
        card_h = max(Inches(2.4), Inches(0.5) + Inches(0.4) + max(len(left_items), len(right_items)) * Inches(0.38))
        
        # Left card
        add_rounded_rect(slide, left_x, self.content_top, col_width, card_h, self.theme["card_bg"])
        add_textbox(slide, left_x + Inches(0.2), self.content_top + Inches(0.15),
                   col_width - Inches(0.4), Inches(0.35),
                   left_title, size=Pt(15), bold=True, color=self.theme["primary"],
                   align=PP_ALIGN.CENTER)
        add_multiline_textbox(slide, left_x + Inches(0.3), self.content_top + Inches(0.6),
                             col_width - Inches(0.6), Inches(2.7),
                             left_items, size=Pt(12), color=self.theme["text"], bullet=True, spacing_pt=6)
        
        # Right card
        add_rounded_rect(slide, right_x, self.content_top, col_width, card_h, self.theme["card_bg"])
        add_textbox(slide, right_x + Inches(0.2), self.content_top + Inches(0.15),
                   col_width - Inches(0.4), Inches(0.35),
                   right_title, size=Pt(15), bold=True, color=self.theme["secondary"],
                   align=PP_ALIGN.CENTER)
        add_multiline_textbox(slide, right_x + Inches(0.3), self.content_top + Inches(0.6),
                             col_width - Inches(0.6), Inches(2.7),
                             right_items, size=Pt(12), color=self.theme["text"], bullet=True, spacing_pt=6)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def vs_compare(self, title, left_title, right_title, rows, page_num=None):
        """
        Side-by-side comparison table.
        rows: list of (dimension, left_value, right_value) tuples
        """
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        y = self.content_top
        # 按 content_w 分配三列（dim + left + right = content_w）
        dim_w = self.content_w * 0.22   # 维度列
        col_w = self.content_w * 0.39   # 值列
        
        # Headers
        add_rect(slide, self.ml, y, dim_w, Inches(0.35), self.theme["primary"])
        add_textbox(slide, self.ml + Inches(0.1), y, dim_w - Inches(0.2), Inches(0.35),
                   "", size=TABLE_HEADER_SIZE, bold=True, color=self.theme["text_light"],
                   align=PP_ALIGN.CENTER)
        
        add_rect(slide, self.ml + dim_w, y, col_w, Inches(0.35), self.theme["primary"])
        add_textbox(slide, self.ml + dim_w + Inches(0.1), y, col_w - Inches(0.2), Inches(0.35),
                   left_title, size=TABLE_HEADER_SIZE, bold=True, color=self.theme["text_light"],
                   align=PP_ALIGN.CENTER)
        
        add_rect(slide, self.ml + dim_w + col_w, y, col_w, Inches(0.35), self.theme["secondary"])
        add_textbox(slide, self.ml + dim_w + col_w + Inches(0.1), y, col_w - Inches(0.2), Inches(0.35),
                   right_title, size=TABLE_HEADER_SIZE, bold=True, color=self.theme["text_light"],
                   align=PP_ALIGN.CENTER)
        
        y += Inches(0.35)
        for i, (dim, left_val, right_val) in enumerate(rows):
            bg_color = self.theme["card_bg"] if i % 2 == 0 else self.theme["bg"]
            row_h = Inches(0.45)
            
            add_rect(slide, self.ml, y, dim_w + col_w * 2, row_h, bg_color)
            add_textbox(slide, self.ml + Inches(0.1), y + Inches(0.05), dim_w - Inches(0.2), row_h,
                       dim, size=TABLE_BODY_SIZE, bold=True, color=self.theme["primary"])
            add_textbox(slide, self.ml + dim_w + Inches(0.1), y + Inches(0.05), col_w - Inches(0.2), row_h,
                       str(left_val), size=TABLE_BODY_SIZE, color=self.theme["text"],
                       align=PP_ALIGN.CENTER)
            add_textbox(slide, self.ml + dim_w + col_w + Inches(0.1), y + Inches(0.05), col_w - Inches(0.2), row_h,
                       str(right_val), size=TABLE_BODY_SIZE, color=self.theme["text"],
                       align=PP_ALIGN.CENTER)
            y += row_h
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── Data & Tables ──────────────────────────────────────────────

    def table(self, title, headers, rows, insights=None, page_num=None):
        """
        Data table slide with optional insight bullets below.
        insights: list of key takeaway strings
        """
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        num_cols = len(headers)
        num_rows = len(rows) + 1
        table_width = self.content_w
        col_width = int(table_width / num_cols)
        row_height = int(Inches(0.4))
        
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            int(self.ml), int(self.content_top),
            int(table_width), int(row_height * num_rows)
        )
        table = table_shape.table
        
        # Set column widths
        for i in range(num_cols):
            table.columns[i].width = col_width
        
        # Header row
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = h
            run.font.size = TABLE_HEADER_SIZE
            run.font.bold = True
            run.font.color.rgb = c(self.theme["text_light"])
            run.font.name = FONT_EN_BODY
            p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = c(self.theme["primary"])
        
        # Data rows
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = table.cell(ri + 1, ci)
                cell.text = ""
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = str(val)
                run.font.size = TABLE_BODY_SIZE
                run.font.color.rgb = c(self.theme["text"])
                run.font.name = FONT_EN_BODY
                p.alignment = PP_ALIGN.CENTER
                if ri % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = c(self.theme["card_bg"])
        
        # Insights
        if insights:
            insight_y = self.content_top + row_height * num_rows + Inches(0.2)
            for ins in insights:
                add_textbox(slide, self.ml, insight_y, self.content_w, Inches(0.3),
                           f"→ {ins}", size=BODY_SMALL, bold=True, color=self.theme["accent"],
                           font_name=FONT_CN)
                insight_y += Inches(0.3)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def stat_cards(self, title, stats, page_num=None):
        """
        Big number stat cards.
        stats: list of (number_text, label, optional_color) tuples
        Max 4 cards.
        """
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        num_stats = min(len(stats), 4)
        card_width = Inches(2.0)
        total_width = card_width * num_stats + GAP_MEDIUM * (num_stats - 1)
        start_x = (self.cw - total_width) / 2
        
        y = Inches(1.8)
        for i, stat in enumerate(stats):
            num_text = stat[0]
            label = stat[1]
            stat_color = stat[2] if len(stat) > 2 else self.theme["accent"]
            
            x = start_x + (card_width + GAP_MEDIUM) * i
            add_rounded_rect(slide, x, y, card_width, Inches(2.8), self.theme["card_bg"])
            # Number — smaller to prevent overflow, single line
            add_textbox(slide, x + Inches(0.1), y + Inches(0.4), card_width - Inches(0.2), Inches(0.7),
                       num_text, size=Pt(36), bold=True, color=stat_color,
                       align=PP_ALIGN.CENTER, font_name=FONT_EN_TITLE)
            # Label — use contrasting color (white on dark card, dark on light card)
            luma = 0.2126 * self.theme["card_bg"][0] + 0.7152 * self.theme["card_bg"][1] + 0.0722 * self.theme["card_bg"][2]
            label_color = (240, 250, 255) if luma < 128 else (51, 51, 51)
            add_textbox(slide, x + Inches(0.1), y + Inches(1.2), card_width - Inches(0.2), Inches(1.2),
                       label, size=STAT_LABEL_SIZE, color=label_color,
                       align=PP_ALIGN.CENTER, bold=True)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── Chart Slides ───────────────────────────────────────────────

    def chart_bar(self, title, data, labels, series_names=None,
                  orientation="vertical", page_num=None):
        """Native bar/column chart slide — fully editable in PowerPoint."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        bar_chart(slide, self.ml, self.content_top, self.content_w, Inches(3.6),
                  data, labels, title="", theme_name=self.theme_name,
                  orientation=orientation, series_names=series_names)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def chart_pie(self, title, data, labels, donut=True, page_num=None):
        """Native pie/donut chart slide — fully editable in PowerPoint."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        pie_chart(slide, self.ml, self.content_top, self.content_w, Inches(3.8),
                  data, labels, title="", theme_name=self.theme_name, donut=donut)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def chart_line(self, title, data, labels, series_names=None,
                   fill_area=False, page_num=None):
        """Native line/area chart slide — fully editable in PowerPoint."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        line_chart(slide, self.ml, self.content_top, self.content_w, Inches(3.6),
                  data, labels, title="", theme_name=self.theme_name,
                  series_names=series_names, fill_area=fill_area)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def chart_gauge(self, title, value, max_value=100, label="", page_num=None):
        """Native gauge chart slide — fully editable in PowerPoint."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        gauge_chart(slide, self.ml, self.content_top, self.content_w, Inches(3.5),
                    value, max_value=max_value, title="",
                    theme_name=self.theme_name, label=label)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── v6.0.0 Extended Charts ──────────────────────────────────

    def chart_funnel(self, title, stages, values, page_num=None):
        """Funnel chart slide with horizontal bars."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        chart_funnel(slide, title="", stages=stages, values=values,
                     left=self.ml, top=self.content_top + Inches(0.2),
                     width=self.content_w, height=Inches(3.3),
                     theme_name=self.theme_name)
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def chart_gantt(self, title, tasks, page_num=None):
        """Gantt chart slide with horizontal task bars."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        chart_gantt(slide, title="", tasks=tasks,
                    left=self.ml, top=self.content_top + Inches(0.2),
                    width=self.content_w, height=Inches(3.3),
                    theme_name=self.theme_name)
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def chart_swot(self, title, strengths, weaknesses, opportunities, threats, page_num=None):
        """SWOT 2x2 matrix slide."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        chart_swot(slide, title="", strengths=strengths, weaknesses=weaknesses,
                   opportunities=opportunities, threats=threats,
                   left=self.ml, top=self.content_top + Inches(0.2),
                   width=self.content_w, height=Inches(3.3),
                   theme_name=self.theme_name)
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def chart_porter(self, title, forces, page_num=None):
        """Porter's Five Forces diagram slide."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        chart_porter(slide, title="", forces=forces,
                     left=self.ml, top=self.content_top + Inches(0.2),
                     width=self.content_w, height=Inches(3.3),
                     theme_name=self.theme_name)
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def chart_sankey(self, title, nodes, links, page_num=None):
        """Simplified Sankey/flow diagram slide."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        chart_sankey(slide, title="", nodes=nodes, links=links,
                     left=self.ml, top=self.content_top + Inches(0.2),
                     width=self.content_w, height=Inches(3.3),
                     theme_name=self.theme_name)
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def chart_heatmap(self, title, rows, cols, values, color_scale=None, page_num=None):
        """Heatmap grid slide with color-coded cells."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        chart_heatmap(slide, title="", rows=rows, cols=cols, values=values,
                      color_scale=color_scale,
                      left=self.ml, top=self.content_top + Inches(0.2),
                      width=self.content_w, height=Inches(3.3),
                      theme_name=self.theme_name)
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def chart_radar(self, title, axes, series_data, series_names=None, page_num=None):
        """Radar/spider chart slide."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        chart_radar(slide, title="", axes=axes, series_data=series_data,
                    series_names=series_names,
                    left=self.ml, top=self.content_top + Inches(0.2),
                    width=self.content_w, height=Inches(3.3),
                    theme_name=self.theme_name)
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def chart_treemap(self, title, items, page_num=None):
        """Treemap chart slide with squarified layout."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        chart_treemap(slide, title="", items=items,
                      left=self.ml, top=self.content_top + Inches(0.2),
                      width=self.content_w, height=Inches(3.3),
                      theme_name=self.theme_name)
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def chart_waterfall(self, title, categories, values, is_total=None, page_num=None):
        """Waterfall chart slide."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        chart_waterfall(slide, title="", categories=categories, values=values,
                        is_total=is_total,
                        left=self.ml, top=self.content_top + Inches(0.2),
                        width=self.content_w, height=Inches(3.3),
                        theme_name=self.theme_name)
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── Timeline & Process ─────────────────────────────────────────

    def timeline(self, title, milestones, page_num=None):
        """
        Timeline slide.
        milestones: list of (period, event_text) tuples
        """
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        num = len(milestones)
        if num == 0:
            return slide
        
        line_y = Inches(2.8)
        start_x = Inches(1.0)
        end_x = self.cw - Inches(1.0)
        span = end_x - start_x
        
        # Horizontal line
        add_line(slide, start_x, line_y, span, self.theme["secondary"], 0.03)
        
        for i, (period, event) in enumerate(milestones):
            x = start_x + span * i / max(num - 1, 1)
            
            # Node circle
            add_circle(slide, x, line_y, Inches(0.12), self.theme["accent"])
            
            # Period label above
            add_textbox(slide, x - Inches(0.6), line_y - Inches(0.6), Inches(1.2), Inches(0.3),
                       period, size=BODY_SMALL, bold=True, color=self.theme["primary"],
                       align=PP_ALIGN.CENTER)
            
            # Event text below
            add_textbox(slide, x - Inches(0.8), line_y + Inches(0.3), Inches(1.6), Inches(0.8),
                       event, size=Pt(10), color=self.theme["text"],
                       align=PP_ALIGN.CENTER)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def process_flow(self, title, steps, page_num=None):
        """
        Process flow slide with chevron-style steps.
        steps: list of step strings (max 5)
        """
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        num = min(len(steps), 5)
        gap = Inches(0.15)
        step_width = (self.content_w - gap * (num - 1)) / num  # 自适应画布宽度
        total = step_width * num + gap * (num - 1)
        start_x = (self.cw - total) / 2
        y = Inches(2.2)
        
        for i in range(num):
            x = start_x + (step_width + gap) * i
            colors = self.theme.get("chart_colors", [self.theme["primary"]])
            color = colors[i % len(colors)]
            
            add_rounded_rect(slide, x, y, step_width, Inches(1.8), color)
            add_textbox(slide, x + Inches(0.1), y + Inches(0.1), step_width - Inches(0.2), Inches(0.3),
                       f"Step {i+1}", size=BODY_SMALL, bold=True, color=self.theme["text_light"],
                       align=PP_ALIGN.CENTER, font_name=FONT_EN_BODY)
            add_textbox(slide, x + Inches(0.1), y + Inches(0.5), step_width - Inches(0.2), Inches(1.1),
                       steps[i], size=Pt(11), color=self.theme["text_light"],
                       align=PP_ALIGN.CENTER)
            
            # Arrow between steps
            if i < num - 1:
                arrow_x = x + step_width
                add_textbox(slide, arrow_x - Inches(0.05), y + Inches(0.7), Inches(0.3), Inches(0.3),
                           "→", size=Pt(18), bold=True, color=self.theme.get("text_muted", (120,120,120)),
                           align=PP_ALIGN.CENTER, font_name=FONT_EN_TITLE)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── Matrix & Framework ─────────────────────────────────────────

    def matrix_2x2(self, title, quadrants, page_num=None):
        """
        2x2 matrix slide (e.g., BCG matrix, priority matrix).
        quadrants: [(top_left, top_right, bottom_left, bottom_right)]
        Each quadrant is (label, description)
        x_axis and y_axis labels optional
        """
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        colors = self.theme.get("chart_colors", [self.theme["primary"]] * 4)
        q_colors = [colors[0], colors[1], colors[2], colors[3]]
        
        cell_w = (self.content_w - Inches(0.1)) / 2  # 2x2 grid
        cell_h = Inches(1.8)
        gap = Inches(0.1)
        start_x = (self.cw - cell_w * 2 - gap) / 2
        start_y = Inches(1.5)
        
        positions = [
            (start_x, start_y),                          # top-left
            (start_x + cell_w + gap, start_y),          # top-right
            (start_x, start_y + cell_h + gap),          # bottom-left
            (start_x + cell_w + gap, start_y + cell_h + gap),  # bottom-right
        ]
        
        for i, (qx, qy) in enumerate(positions):
            q = quadrants[i] if i < len(quadrants) else ("", "")
            add_rounded_rect(slide, qx, qy, cell_w, cell_h, q_colors[i], alpha=0.85)
            add_textbox(slide, qx + Inches(0.2), qy + Inches(0.15), cell_w - Inches(0.4), Inches(0.3),
                       q[0], size=Pt(14), bold=True, color=self.theme["text_light"])
            if len(q) > 1 and q[1]:
                add_textbox(slide, qx + Inches(0.2), qy + Inches(0.5), cell_w - Inches(0.4), Inches(1.1),
                           q[1], size=Pt(10), color=self.theme["text_light"])
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── Quote & Section ────────────────────────────────────────────

    def quote(self, title, quote_text, attribution="", page_num=None):
        """Large quote slide."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg_dark"])
        
        add_textbox(slide, Inches(1), Inches(1.5), Inches(1), Inches(0.5),
                   "❝", size=Pt(48), color=self.theme["accent"])
        add_textbox(slide, self.ml, Inches(2.0), self.content_w, Inches(1.5),
                   quote_text, size=Pt(20), color=self.theme["text_light"],
                   align=PP_ALIGN.LEFT)
        if attribution:
            add_textbox(slide, self.ml, Inches(3.8), self.content_w, Inches(0.4),
                       f"— {attribution}", size=BODY_SIZE, color=self.theme["secondary"])
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def section_divider(self, section_title, section_number=None, subtitle=""):
        """Section divider slide (full dark background)."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg_dark"])
        
        if section_number:
            # Large section number: left column, vertically centered
            num_text = f"0{section_number}" if section_number < 10 else str(section_number)
            add_textbox(slide, self.ml, Inches(1.8), Inches(1.2), Inches(1),
                       num_text, size=Pt(72), bold=True, color=self.theme["accent"],
                       font_name=FONT_EN_TITLE, align=PP_ALIGN.LEFT)
            # Accent line: starts after number column, spans title area
            add_line(slide, self.ml + Inches(1.5), Inches(1.95), self.content_w - Inches(1.5),
                    self.theme["accent"], 0.04)
            # Section title: right of number, same baseline as number
            add_textbox(slide, self.ml + Inches(1.5), Inches(2.05), self.content_w - Inches(1.5), Inches(0.9),
                       section_title, size=Pt(32), bold=True, color=self.theme["text_light"])
        else:
            add_line(slide, self.ml, Inches(1.95), self.content_w, self.theme["accent"], 0.04)
            add_textbox(slide, self.ml, Inches(2.05), self.content_w, Inches(0.9),
                       section_title, size=Pt(32), bold=True, color=self.theme["text_light"])
        
        if subtitle:
            y_sub = Inches(3.0) if section_number else Inches(3.0)
            add_textbox(slide, self.ml, y_sub, self.content_w, Inches(0.5),
                       subtitle, size=BODY_SIZE, color=self.theme["secondary"])
        return slide

    # ── Image Slides ───────────────────────────────────────────────

    def image_full(self, title, image_path, caption="", page_num=None):
        """Full-width image slide with caption."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        add_image_safe(slide, image_path, self.ml, self.content_top, self.content_w, Inches(3.5))
        if caption:
            add_textbox(slide, self.ml, Inches(4.8), self.content_w, Inches(0.3),
                       caption, size=CAPTION_SIZE, color=self.theme.get("text_muted", (120,120,120)),
                       align=PP_ALIGN.CENTER)
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    def image_split(self, title, image_path, bullets, image_side="right", page_num=None):
        """Split slide: text on one side, image on the other."""
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        half_w = (self.content_w - Inches(0.3)) / 2
        img_w = half_w
        text_w = half_w
        
        if image_side == "right":
            text_x = self.ml
            img_x = self.ml + text_w + Inches(0.3)
        else:
            text_x = self.ml + text_w + Inches(0.3)
            img_x = self.ml
        
        add_multiline_textbox(slide, text_x, self.content_top, text_w, Inches(3.8),
                             bullets, size=BODY_SIZE, color=self.theme["text"], bullet=True)
        add_image_safe(slide, image_path, img_x, self.content_top, img_w, Inches(3.8))
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── Summary ────────────────────────────────────────────────────

    def summary(self, title: str, key_points: List[str],
               conclusion: str = "", page_num: Optional[int] = None):
        """Summary slide with key takeaways — dark bg, high contrast cards.

        Args:
            title: Summary title (required).
            key_points: List of takeaway strings (1–10 items).
            conclusion: Final call-to-action text (optional).
            page_num: Page number to display (optional).
        """
        title = _validate_text(title, "title", allow_none=False, max_len=200)
        key_points = _validate_list(key_points, "key_points", min_items=1, max_items=10)
        for i, p in enumerate(key_points):
            if not isinstance(p, str):
                raise TypeError(f"key_points[{i}] must be str, got {type(p).__name__}: {p!r}")
        conclusion = _validate_text(conclusion, "conclusion") or ""
        if page_num is not None:
            page_num = _validate_int(page_num, "page_num", min_val=1)

        slide = self._new_slide()
        set_bg(slide, self.theme["bg_dark"])
        
        add_line(slide, self.ml, Inches(1.2), Inches(1.5), self.theme["accent"], 0.045)
        add_textbox(slide, self.ml, Inches(1.35), self.content_w, Inches(0.5),
                   title, size=Pt(24), bold=True, color=self.theme["text_light"])
        
        y = Inches(2.1)
        for point in key_points:
            # Dark card background so white text has high contrast
            card_color = self.theme.get("card_bg", (50, 50, 50))
            # If card_bg is light, darken it
            r, g, b = card_color
            if r + g + b > 400:  # too light for white text
                card_color = (40, 55, 85)  # forced dark blue-gray
            add_rounded_rect(slide, self.ml, y, self.content_w, Inches(0.5), card_color)
            add_textbox(slide, self.ml + Inches(0.25), y + Inches(0.08), self.content_w - Inches(0.5), Inches(0.4),
                       point, size=Pt(13), color=self.theme["text_light"])
            y += Inches(0.6)
        
        if conclusion:
            add_textbox(slide, self.ml, y + Inches(0.3), self.content_w, Inches(0.4),
                       conclusion, size=BODY_SIZE, bold=True, color=self.theme["accent"])
        
        # Side accent
        add_rect(slide, 0, 0, Inches(0.12), self.ch, self.theme["accent"])
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── Team / Stakeholders ────────────────────────────────────────

    def team_grid(self, title, members, page_num=None):
        """
        Team member grid.
        members: list of (name, role) tuples. Max 6.
        """
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        num = min(len(members), 6)
        cols = 3 if num > 3 else num
        rows_count = (num + cols - 1) // cols
        
        card_w = Inches(2.5)
        card_h = Inches(1.8)
        gap = Inches(0.3)
        total_w = card_w * cols + gap * (cols - 1)
        start_x = (self.cw - total_w) / 2
        start_y = Inches(1.5)
        
        for i, (name, role) in enumerate(members[:num]):
            row = i // cols
            col = i % cols
            x = start_x + (card_w + gap) * col
            y = start_y + (card_h + gap) * row
            
            add_rounded_rect(slide, x, y, card_w, card_h, self.theme["card_bg"])
            # Avatar circle
            add_circle(slide, x + card_w / 2, y + Inches(0.55), Inches(0.3), self.theme["primary"])
            add_textbox(slide, x + card_w / 2 - Inches(0.25), y + Inches(0.38), Inches(0.5), Inches(0.3),
                       name[0].upper() if name else "?", size=Pt(12), bold=True,
                       color=self.theme["text_light"], align=PP_ALIGN.CENTER, font_name=FONT_EN_TITLE)
            add_textbox(slide, x + Inches(0.1), y + Inches(1.0), card_w - Inches(0.2), Inches(0.3),
                       name, size=Pt(12), bold=True, color=self.theme["primary"],
                       align=PP_ALIGN.CENTER)
            add_textbox(slide, x + Inches(0.1), y + Inches(1.3), card_w - Inches(0.2), Inches(0.3),
                       role, size=Pt(9), color=self.theme.get("text_muted", (120,120,120)),
                       align=PP_ALIGN.CENTER)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── KPI Dashboard ──────────────────────────────────────────────

    def kpi_dashboard(self, title, kpis, page_num=None):
        """
        KPI dashboard with stat cards + small charts.
        kpis: list of (label, value, target, unit) tuples
        """
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        num = min(len(kpis), 4)
        card_w = Inches(2.1)
        gap = Inches(0.2)
        total_w = card_w * num + gap * (num - 1)
        start_x = (self.cw - total_w) / 2
        y = Inches(1.5)
        
        for i, (label, value, target, unit) in enumerate(kpis[:num]):
            x = start_x + (card_w + gap) * i
            add_rounded_rect(slide, x, y, card_w, Inches(3.5), self.theme["card_bg"])
            
            # Label
            add_textbox(slide, x + Inches(0.1), y + Inches(0.15), card_w - Inches(0.2), Inches(0.3),
                       label, size=Pt(10), color=self.theme.get("text_muted", (120,120,120)),
                       align=PP_ALIGN.CENTER)
            
            # Value
            pct = min(value / max(target, 1), 1.0)
            color = self.theme["accent"] if pct >= 0.9 else self.theme["primary"]
            add_textbox(slide, x + Inches(0.1), y + Inches(0.5), card_w - Inches(0.2), Inches(0.6),
                       f"{value}{unit}", size=Pt(24), bold=True, color=color,
                       align=PP_ALIGN.CENTER, font_name=FONT_EN_TITLE)
            
            # Progress bar
            bar_y = y + Inches(1.3)
            bar_w = card_w - Inches(0.4)
            add_rect(slide, x + Inches(0.2), bar_y, bar_w, Inches(0.1), (220, 220, 220))
            add_rect(slide, x + Inches(0.2), bar_y, int(bar_w * pct), Inches(0.1), color)
            
            # Target
            add_textbox(slide, x + Inches(0.1), bar_y + Inches(0.2), card_w - Inches(0.2), Inches(0.3),
                       f"Target: {target}{unit}", size=Pt(9),
                       color=self.theme.get("text_muted", (120,120,120)),
                       align=PP_ALIGN.CENTER)
            
            # Percentage
            add_textbox(slide, x + Inches(0.1), bar_y + Inches(0.5), card_w - Inches(0.2), Inches(0.3),
                       f"{pct*100:.0f}%", size=Pt(18), bold=True, color=color,
                       align=PP_ALIGN.CENTER, font_name=FONT_EN_TITLE)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── Agenda / List ──────────────────────────────────────────────

    def checklist(self, title, items, checked=None, page_num=None):
        """
        Checklist slide.
        items: list of strings
        checked: list of booleans (True = checked)
        """
        slide = self._new_slide()
        set_bg(slide, self.theme["bg"])
        self._add_top_bar(slide)
        self._add_title_area(slide, title)
        
        if checked is None:
            checked = [False] * len(items)
        
        y = self.content_top
        for i, item in enumerate(items):
            is_checked = checked[i] if i < len(checked) else False
            icon = "✅" if is_checked else "⬜"
            text_color = self.theme.get("text_muted", (120,120,120)) if is_checked else self.theme["text"]
            add_textbox(slide, self.ml, y, self.content_w, Inches(0.35),
                       f"{icon}  {item}", size=BODY_SIZE, color=text_color)
            y += Inches(0.45)
        
        if page_num:
            self._add_page_number(slide, page_num)
        return slide

    # ── Save ───────────────────────────────────────────────────────

    def save(self, path: str) -> str:
        """Save presentation to file.

        Args:
            path: Output .pptx file path. Parent directories are auto-created.

        Returns:
            The path argument (for chaining).

        Raises:
            TypeError: If path is not a string.
            OSError: If the file cannot be written.
        """
        if not isinstance(path, str):
            raise TypeError(f"path must be str, got {type(path).__name__}: {path!r}")
        if not path.lower().endswith(".pptx"):
            import warnings
            warnings.warn(f"path should end with .pptx, got: {path!r}", UserWarning, stacklevel=2)
        os.makedirs(os.path.dirname(path) if os.path.dirname(path) else ".", exist_ok=True)
        self.prs.save(path)
        return path

    @property
    def slide_count(self) -> int:
        """Number of slides created so far."""
        return self._slide_count
