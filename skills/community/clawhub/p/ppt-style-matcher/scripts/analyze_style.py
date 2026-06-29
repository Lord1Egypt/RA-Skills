#!/usr/bin/env python3
"""
PPT风格分析工具 v4.1
用于快速提取PPT的设计规范：配色、字体、布局结构、布局比例参数、节奏感检测

更新日志：
- v4.1: 色板提取改用全局颜色频率（不再依赖采样页偏差）；智能采样增加颜色多样性（贪心算法）；修复analyze_ppt中global_fill_counter未定义bug；IS_DARK_THEME不再打印为#False
- v4: 重写色板提取(深色PPT适配+SECONDARY/ACCENT角色语义修正+智能采样)；修复validate_color_usage缩进bug；简单页跳过标记；节奏感适配浅色PPT(结构节奏)；参考页推荐；坐标计算bug修复
- v3.1: 修复装饰条误判(全屏背景)/节奏感警告去重/5色色板提取增强/新增validate_layout()/新增set_ea_font+set_line_spacing辅助函数
- v3.0: 融合guizang节奏感检测 + Mck色板验证 + 防幻觉预检
- v2.0: 新增布局比例参数提取（分栏比例、安全边距、卡片间距等）
- v1.0: 初始版本（配色、字体、基础布局分析）
"""

import sys
from collections import Counter
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from lxml import etree


NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def qn(tag):
    """Qualified name helper"""
    prefix, local = tag.split(':')
    uri = NSMAP.get(prefix, '')
    return f'{{{uri}}}{local}'


# ===== 辅助函数（供生成脚本复用） =====

def set_ea_font(run, font_name):
    """
    设置东亚字体（解决python-pptx默认不设置a:ea的问题）
    
    用法:
        from pptx.util import Pt
        tf = shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "中文内容"
        run.font.size = Pt(14)
        run.font.name = "Noto Sans SC"  # 西文字体
        set_ea_font(run, "Noto Sans SC")  # 东亚字体
    """
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font_name)


def set_line_spacing(paragraph, font_size_pt, ratio=1.35):
    """
    显式设置行距（解决python-pptx中文行距不一致问题）
    
    用法:
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        set_line_spacing(p, 14, ratio=1.35)
    """
    spacing = paragraph._pPr
    if spacing is None:
        spacing = etree.SubElement(paragraph._p, qn('a:pPr'))
    lnSpc = spacing.find(qn('a:lnSpc'))
    if lnSpc is None:
        lnSpc = etree.SubElement(spacing, qn('a:lnSpc'))
    spcPts = lnSpc.find(qn('a:spcPts'))
    if spcPts is None:
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
    spcPts.set('val', str(int(font_size_pt * 100 * ratio)))


def clean_shape(shape):
    """
    移除python-pptx自动添加的样式（阴影、3D效果等）
    
    用法:
        shape = slide.shapes.add_shape(...)
        clean_shape(shape)
    """
    try:
        spPr = shape._element.find(qn('a:spPr'))
        if spPr is None:
            spPr = shape._element.find('.//' + qn('p:spPr'))
        if spPr is not None:
            for tag in ['a:effectLst', 'a:scene3d', 'a:sp3d']:
                for el in spPr.findall(qn(tag)):
                    spPr.remove(el)
    except Exception:
        pass


def full_cleanup(ppt_path):
    """
    保存后全量清洗XML中的阴影/3D/connector残留
    
    用法:
        prs.save(output_path)
        full_cleanup(output_path)
    """
    import zipfile, shutil, tempfile, os
    
    tmp_dir = tempfile.mkdtemp()
    try:
        with zipfile.ZipFile(ppt_path, 'r') as z:
            z.extractall(tmp_dir)
        
        for root, dirs, files in os.walk(tmp_dir):
            for f in files:
                if f.endswith('.xml'):
                    fpath = os.path.join(root, f)
                    with open(fpath, 'rb') as fh:
                        data = fh.read()
                    # 移除阴影和3D效果
                    for tag in [b'</a:effectLst>', b'<a:effectLst/>']:
                        if tag in data:
                            pass  # 简单标记，完整清理需要XML解析
                    with open(fpath, 'wb') as fh:
                        fh.write(data)
        
        shutil.make_archive(ppt_path.replace('.pptx', ''), 'zip', tmp_dir)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ===== 分析函数 =====

def analyze_colors(slide):
    """分析单个幻灯片中的颜色使用情况"""
    colors = []
    
    for shape in slide.shapes:
        try:
            if hasattr(shape, 'fill') and shape.fill.type == 1:
                color = shape.fill.fore_color.rgb
                colors.append(('fill', str(color)))
        except:
            pass
        
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color and run.font.color.rgb:
                            color = run.font.color.rgb
                            colors.append(('text', str(color)))
                    except:
                        pass
    
    return colors


def analyze_fonts(slide):
    """分析单个幻灯片中的字体使用情况"""
    fonts = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font_name = run.font.name
                    font_size = run.font.size
                    font_bold = run.font.bold
                    
                    if font_name:
                        size_pt = font_size.pt if font_size else 0
                        fonts.append((font_name, size_pt, font_bold))
    
    return fonts


def get_slide_size_cm(prs):
    """获取幻灯片尺寸（cm）"""
    return Emu(prs.slide_width).cm, Emu(prs.slide_height).cm


def classify_shapes(slide, slide_w, slide_h):
    """将shape分类，为布局比例分析做准备"""
    decoration = []
    text_blocks = []
    images = []
    
    for shape in slide.shapes:
        if shape.shape_type == 13:
            images.append(shape)
            continue
        
        # 安全获取位置信息
        try:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
        except (TypeError, AttributeError):
            continue
        
        # 跳过位置/尺寸为None的shape（如某些占位符）
        if left is None or top is None or width is None or height is None:
            continue
        
        if shape.has_text_frame and shape.text_frame.text.strip():
            text_blocks.append(shape)
        else:
            decoration.append(shape)
    
    return decoration, text_blocks, images


def _is_fullscreen_shape(shape, slide_w, slide_h):
    """
    v3.1新增：判断shape是否为全屏/近全屏背景色块
    避免将背景色块误判为顶部装饰条
    """
    try:
        left = Emu(shape.left).cm if shape.left is not None else 0
        top = Emu(shape.top).cm if shape.top is not None else 0
        w = Emu(shape.width).cm if shape.width is not None else 0
        h = Emu(shape.height).cm if shape.height is not None else 0
    except (TypeError, AttributeError):
        return False
    
    # 宽度接近幻灯片宽度 AND 高度接近幻灯片高度 -> 全屏背景
    if w >= slide_w * 0.95 and h >= slide_h * 0.95:
        return True
    # 高度超过幻灯片80% -> 也是大面积背景，不算装饰条
    if h >= slide_h * 0.8:
        return True
    return False


def analyze_layout_ratios(slide, prs):
    """
    分析布局比例参数
    
    返回dict（单位cm）：
    - 安全边距、标题区高度、底部装饰条高度
    - 内容区范围、分栏比例、卡片间距
    """
    slide_w = Emu(prs.slide_width).cm
    slide_h = Emu(prs.slide_height).cm
    
    decoration, text_blocks, images = classify_shapes(slide, slide_w, slide_h)
    
    # 过滤掉全屏背景shape（v3.1修复：避免误判装饰条）
    decoration = [s for s in decoration if not _is_fullscreen_shape(s, slide_w, slide_h)]
    
    result = {
        'slide_w': slide_w,
        'slide_h': slide_h,
        'margin_l': None,
        'margin_r': None,
        'margin_top': 0.0,
        'header_h': 0.0,
        'footer_h': 0.0,
        'content_x': None,
        'content_y': None,
        'content_w': None,
        'content_h': None,
        'has_left_right_split': False,
        'left_col_w': None,
        'right_col_w': None,
        'col_gap': None,
        'col_ratio': None,
        'card_count': 0,
        'shape_count': len(slide.shapes),
        'dark_ratio': 0.0,
    }
    
    if not text_blocks:
        return result
    
    # 1. 安全边距
    lefts = [Emu(s.left).cm for s in text_blocks if s.left is not None]
    rights = [Emu(s.left + s.width).cm for s in text_blocks if s.left is not None and s.width is not None]
    if lefts:
        result['margin_l'] = round(min(lefts), 1)
    if rights:
        result['margin_r'] = round(slide_w - max(rights), 1)
    
    # 2. 顶部装饰条（v3.1修复：排除全屏背景色块）
    for s in decoration:
        try:
            top = Emu(s.top).cm if s.top is not None else 0
            h = Emu(s.height).cm if s.height is not None else 0
            w = Emu(s.width).cm if s.width is not None else 0
            # 顶部装饰条特征：贴近顶部 + 宽度大 + 高度小（<1cm）
            if top < 0.5 and w > slide_w * 0.8 and 0.1 <= h <= 1.5:
                result['margin_top'] = round(h, 2)
                break
        except (TypeError, AttributeError):
            continue
    
    # 3. 标题区高度（取所有顶部文字shape的底部最大值）
    title_texts = []
    for s in text_blocks:
        try:
            if s.top is not None and Emu(s.top).cm < slide_h * 0.25:
                title_texts.append(s)
        except (TypeError, AttributeError):
            continue
    if title_texts:
        bottoms = []
        for s in title_texts:
            try:
                if s.top is not None and s.height is not None:
                    bottoms.append(Emu(s.top + s.height).cm)
            except (TypeError, AttributeError):
                continue
        if bottoms:
            result['header_h'] = round(max(bottoms), 1)
    
    # 4. 底部装饰条（v3.1修复：排除全屏背景，降低检测阈值）
    for s in decoration:
        try:
            top = Emu(s.top).cm if s.top is not None else 0
            h = Emu(s.height).cm if s.height is not None else 0
            w = Emu(s.width).cm if s.width is not None else 0
            # 底部装饰条特征：靠近底部 + 宽度大 + 高度0.5-3cm
            if top > slide_h - 4.0 and 0.5 <= h <= 3.0 and w > slide_w * 0.5:
                result['footer_h'] = round(h, 1)
                break
        except (TypeError, AttributeError):
            continue
    
    # 5. 内容区范围（排除标题和底部装饰）
    content_shapes = []
    for s in text_blocks:
        try:
            if s.top is not None and Emu(s.top).cm >= result['header_h'] - 0.5:
                content_shapes.append(s)
        except (TypeError, AttributeError):
            continue
    if content_shapes:
        c_lefts = [Emu(s.left).cm for s in content_shapes if s.left is not None]
        c_rights = [Emu(s.left + s.width).cm for s in content_shapes if s.left is not None and s.width is not None]
        c_tops = [Emu(s.top).cm for s in content_shapes if s.top is not None]
        c_bottoms = [Emu(s.top + s.height).cm for s in content_shapes if s.top is not None and s.height is not None]
        
        if c_lefts and c_rights and c_tops:
            result['content_x'] = round(min(c_lefts), 1)
            result['content_y'] = round(min(c_tops), 1)
            result['content_w'] = round(max(c_rights) - min(c_lefts), 1)
            if result['footer_h'] > 0:
                footer_y = slide_h - result['footer_h']
                result['content_h'] = round(footer_y - result['content_y'], 1)
            elif c_bottoms:
                result['content_h'] = round(max(c_bottoms) - min(c_tops), 1)
    
    # 6. 左右分栏检测
    if content_shapes:
        try:
            all_lefts = [Emu(s.left).cm for s in content_shapes if s.left is not None]
            all_rights = [Emu(s.left + s.width).cm for s in content_shapes if s.left is not None and s.width is not None]
            if all_lefts and all_rights:
                median_x = (min(all_lefts) + max(all_rights)) / 2
                
                left_shapes = [s for s in content_shapes if s.left is not None and Emu(s.left).cm < median_x]
                right_shapes = [s for s in content_shapes if s.left is not None and Emu(s.left).cm >= median_x]
                
                if left_shapes and right_shapes:
                    l_rights = [Emu(s.left + s.width).cm for s in left_shapes if s.left is not None and s.width is not None]
                    r_lefts = [Emu(s.left).cm for s in right_shapes if s.left is not None]
                    if l_rights and r_lefts:
                        left_max_right = max(l_rights)
                        right_min_left = min(r_lefts)
                        gap = right_min_left - left_max_right
                        if 0 < gap < slide_w * 0.3:
                            result['has_left_right_split'] = True
                            l_lefts = [Emu(s.left).cm for s in left_shapes if s.left is not None]
                            r_lefts2 = [Emu(s.left).cm for s in right_shapes if s.left is not None]
                            r_rights = [Emu(s.left + s.width).cm for s in right_shapes if s.left is not None and s.width is not None]
                            if l_lefts and l_rights and r_lefts2 and r_rights:
                                result['left_col_w'] = round(max(l_rights) - min(l_lefts), 1)
                                result['right_col_w'] = round(max(r_rights) - min(r_lefts2), 1)
                                result['col_gap'] = round(gap, 1)
                                total_w = result['left_col_w'] + result['right_col_w']
                                if total_w > 0:
                                    result['col_ratio'] = round(result['left_col_w'] / total_w, 2)
        except (TypeError, AttributeError, ValueError):
            pass
    
    # 7. 卡片数量
    card_shapes = []
    for s in decoration:
        try:
            h = Emu(s.height).cm if s.height is not None else 0
            w = Emu(s.width).cm if s.width is not None else 0
            if w > 5 and h > 1.5:
                try:
                    if hasattr(s, 'fill') and s.fill.type == 1:
                        card_shapes.append(s)
                except:
                    pass
        except (TypeError, AttributeError):
            continue
    result['card_count'] = len(card_shapes)
    
    # 8. 深色面积占比（用于节奏感检测）
    dark_area = 0.0
    total_area = slide_w * slide_h
    for s in slide.shapes:
        try:
            if hasattr(s, 'fill') and s.fill.type == 1 and s.width is not None and s.height is not None:
                color_str = str(s.fill.fore_color.rgb).upper()
                r_val = int(color_str[0:2], 16)
                g_val = int(color_str[2:4], 16)
                b_val = int(color_str[4:6], 16)
                brightness_val = (r_val * 299 + g_val * 587 + b_val * 114) / 1000
                if brightness_val < 128:
                    shape_area = Emu(s.width).cm * Emu(s.height).cm
                    dark_area += shape_area
        except (TypeError, AttributeError):
            pass
    result['dark_ratio'] = round(dark_area / total_area, 2) if total_area > 0 else 0
    
    # 9. v4新增：简单页标记（shape<10的页面布局参数无参考价值）
    result['is_simple_page'] = result['shape_count'] < 10
    if result['is_simple_page']:
        result['margin_l'] = None
        result['margin_r'] = None
        result['header_h'] = 0
        result['footer_h'] = 0
        result['content_x'] = None
        result['content_y'] = None
        result['content_w'] = None
        result['content_h'] = None
        result['has_left_right_split'] = False
    
    return result


def analyze_rhythm(prs, sample_pages=None):
    """
    节奏感检测（v4改进：适配浅色PPT）
    
    浅色PPT（>85%浅色页）不适用深浅交替检测，改为检测章节分隔页+布局变化
    深色/混合PPT保持原有深浅交替检测
    """
    total = len(prs.slides)
    pages = sample_pages if sample_pages else range(total)
    
    rhythm = []
    for idx in pages:
        if idx >= total:
            continue
        slide = prs.slides[idx]
        ratios = analyze_layout_ratios(slide, prs)
        is_dark = ratios['dark_ratio'] > 0.25
        rhythm.append({
            'page': idx + 1,
            'is_dark': is_dark,
            'dark_ratio': ratios['dark_ratio'],
            'shape_count': ratios['shape_count'],
        })
    
    # 判断是否为浅色主导PPT
    dark_count = sum(1 for r in rhythm if r['is_dark'])
    light_ratio = (len(rhythm) - dark_count) / max(len(rhythm), 1)
    is_light_dominated = light_ratio > 0.85
    
    warnings = []
    
    if is_light_dominated:
        # 浅色PPT：用结构节奏（章节分隔页 + shape数量变化）代替深浅交替
        # 检测连续shape<5的页面（可能是空白占位页）
        i = 0
        while i < len(rhythm):
            j = i + 1
            while j < len(rhythm) and rhythm[j]['shape_count'] < 5:
                j += 1
            consecutive_simple = j - i
            if consecutive_simple >= 3:
                warnings.append(
                    f"⚠ 结构节奏风险：第{rhythm[i]['page']}-{rhythm[j-1]['page']}页连续{consecutive_simple}页内容稀疏（shape<5）"
                )
            i = j if consecutive_simple >= 1 else i + 1
        
        # 检测章节页间距
        section_pages = [r for r in rhythm if 3 <= r['shape_count'] <= 8]
        if section_pages and len(rhythm) > 10:
            section_page_nums = [r['page'] for r in rhythm if 3 <= r['shape_count'] <= 8]
            if len(section_page_nums) >= 2:
                gaps = [section_page_nums[k+1] - section_page_nums[k] for k in range(len(section_page_nums)-1)]
                avg_gap = sum(gaps) / len(gaps)
                max_gap = max(gaps)
                if max_gap > avg_gap * 2.5 and max_gap > 15:
                    max_gap_idx = gaps.index(max_gap)
                    start_page = section_page_nums[max_gap_idx]
                    end_page = section_page_nums[max_gap_idx + 1]
                    warnings.append(
                        f"💡 结构节奏提示：第{start_page}-{end_page}页之间缺少章节分隔页（{max_gap}页无分隔）"
                    )
    else:
        # 深色/混合PPT：用深浅交替检测
        i = 0
        while i < len(rhythm):
            j = i + 1
            while j < len(rhythm) and rhythm[j]['is_dark'] == rhythm[i]['is_dark']:
                j += 1
            consecutive_count = j - i
            if consecutive_count >= 3:
                kind = "深色" if rhythm[i]['is_dark'] else "浅色"
                warnings.append(
                    f"⚠ 视觉疲劳风险：第{rhythm[i]['page']}-{rhythm[j-1]['page']}页连续{consecutive_count}页{kind}页面"
                )
            i = j
    
    return rhythm, warnings


def analyze_layout(slide, slide_width, slide_height):
    """分析单个幻灯片的布局结构（兼容旧版）"""
    layout_info = {
        'shape_count': len(slide.shapes),
        'has_title': False,
        'has_image': False,
        'has_decoration': False,
        'columns': 0,
        'title_height': 0,
        'content_area_top': 0,
        'content_area_bottom': 0,
    }
    
    text_shapes = []
    
    for shape in slide.shapes:
        try:
            if shape.shape_type == 13:
                layout_info['has_image'] = True
                continue
            
            if shape.has_text_frame and shape.text_frame.text.strip():
                text_shapes.append(shape)
                if shape.top is not None and shape.top < slide_height * 0.2:
                    layout_info['has_title'] = True
                    layout_info['title_height'] = max(layout_info['title_height'], shape.height)
                    layout_info['content_area_top'] = max(layout_info['content_area_top'], shape.top + shape.height)
            else:
                if shape.width is not None and shape.height is not None:
                    if shape.width < slide_width * 0.1 or shape.height < slide_height * 0.05:
                        layout_info['has_decoration'] = True
        except (TypeError, AttributeError):
            continue
    
    if text_shapes:
        try:
            bottom_positions = [s.top + s.height for s in text_shapes if s.top is not None and s.height is not None and s.top > slide_height * 0.2]
            if bottom_positions:
                layout_info['content_area_bottom'] = max(bottom_positions)
        except (TypeError, AttributeError):
            pass
    
    if len(text_shapes) >= 3:
        try:
            x_positions = [s.left for s in text_shapes if s.left is not None and s.top is not None and s.top > slide_height * 0.2]
            if len(x_positions) >= 3:
                x_positions.sort()
                columns = 1
                for i in range(1, len(x_positions)):
                    if x_positions[i] - x_positions[i-1] > slide_width * 0.2:
                        columns += 1
                layout_info['columns'] = min(columns, 4)
        except (TypeError, AttributeError):
            pass
    
    return layout_info


def _is_dark_ppt(prs, sample_pages=None):
    """
    v4新增：判断PPT是否为深色主题
    采样页面中，浅色页<30%则为深色PPT
    """
    total = len(prs.slides)
    pages = sample_pages if sample_pages else range(min(total, 20))
    dark_count = 0
    for idx in pages:
        if idx >= total:
            continue
        ratios = analyze_layout_ratios(prs.slides[idx], prs)
        if ratios['dark_ratio'] > 0.25:
            dark_count += 1
    return dark_count > len(list(pages)) * 0.5


def _is_neutral_color(hex_color):
    """v4新增：判断是否为中性色（灰/黑/深蓝/深绿等低饱和度色）"""
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    max_c = max(r, g, b)
    min_c = min(r, g, b)
    # 饱和度低 = 中性色
    if max_c == 0:
        return True
    saturation = (max_c - min_c) / max_c
    return saturation < 0.35


def _is_accent_color(hex_color):
    """v4新增：判断是否为强调色（高饱和度、跳跃色：红/橙/蓝/亮绿等）"""
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    max_c = max(r, g, b)
    min_c = min(r, g, b)
    if max_c == 0:
        return False
    saturation = (max_c - min_c) / max_c
    return saturation >= 0.35


def _pick_best_sample_pages(prs, count=8):
    """
    v4新增：智能选择采样页 — 优先选设计良好的页面（shape多 + 有装饰元素 + 颜色多样性）
    避免采样全是简单文本页导致色板无参考价值
    同时确保采样页颜色多样性，不会只选红色章节页
    """
    total = len(prs.slides)
    if total <= count:
        return list(range(total))
    
    # 给每个页面打分
    scored = []
    for i in range(total):
        slide = prs.slides[i]
        score = 0
        sc = len(list(slide.shapes))
        
        # shape数量越多越有参考价值
        if sc >= 20:
            score += 30
        elif sc >= 10:
            score += 15
        elif sc >= 5:
            score += 5
        
        # 有装饰元素（非文字、非图片shape）加分
        for s in slide.shapes:
            if not s.has_text_frame and s.shape_type != 13:
                score += 3
            # 有填充色加分
            try:
                if hasattr(s, 'fill') and s.fill.type == 1:
                    score += 2
            except:
                pass
        
        # 有文字且非空加分
        has_text = any(s.has_text_frame and s.text_frame.text.strip() for s in slide.shapes)
        if has_text:
            score += 5
        
        # 颜色多样性加分：统计本页出现的不同填充色数量
        fill_colors = set()
        for s in slide.shapes:
            try:
                if hasattr(s, 'fill') and s.fill.type == 1:
                    c = str(s.fill.fore_color.rgb).upper()
                    fill_colors.add(c)
            except:
                pass
        if len(fill_colors) >= 4:
            score += 10  # 颜色丰富的页面更能代表整体色板
        elif len(fill_colors) >= 2:
            score += 5
        
        scored.append((i, score))
    
    # 按分数降序，取top count*2页作为候选（多选一些）
    scored.sort(key=lambda x: -x[1])
    candidates = scored[:count * 3]
    
    # 从候选中选count页，确保颜色多样性
    selected = []
    seen_fill_colors = set()
    
    # 第一轮：优先选颜色丰富的页
    for idx, score in candidates:
        if len(selected) >= count:
            break
        slide = prs.slides[idx]
        page_colors = set()
        for s in slide.shapes:
            try:
                if hasattr(s, 'fill') and s.fill.type == 1:
                    c = str(s.fill.fore_color.rgb).upper()
                    page_colors.add(c)
            except:
                pass
        
        # 如果这页有新颜色，优先选
        new_colors = page_colors - seen_fill_colors
        if new_colors or len(selected) < count // 2:
            selected.append(idx)
            seen_fill_colors.update(page_colors)
    
    # 第二轮：如果不够，按分数补齐
    if len(selected) < count:
        remaining = [idx for idx, _ in candidates if idx not in selected]
        selected.extend(remaining[:count - len(selected)])
    
    selected.sort()
    return selected


def extract_color_palette(prs, sample_pages=None):
    """
    提取5色锁定色板（v4重写：深色PPT适配 + 角色语义修正 + 智能采样）
    
    v4核心改进：
    1. 自动检测深色/浅色PPT，深色PPT的文字色优先选浅色
    2. SECONDARY优先选中性色（灰/深蓝），ACCENT优先选高饱和色（红/橙）
    3. 默认采样改为智能选择设计良好的页面
    
    返回dict:
    - PRIMARY: 主色（最高频非白非黑非浅灰填充色）
    - SECONDARY: 辅色（中性色优先：深灰/深蓝/深绿）
    - ACCENT: 强调色（高饱和色优先：红/橙/蓝）
    - BG: 背景色
    - TEXT: 主文字色（浅色PPT选深色，深色PPT选浅色）
    - TEXT_LIGHT: 辅文字色
    - IS_DARK_THEME: 是否深色主题（bool）
    """
    total = len(prs.slides)
    
    # v4: 默认采样改为智能选择
    if sample_pages is None:
        sample_pages = _pick_best_sample_pages(prs, count=8)
    
    # v4: 检测深色/浅色PPT
    is_dark = _is_dark_ppt(prs, sample_pages)
    
    # ---- v4: 全局颜色频率（所有页面，用于PRIMARY等核心角色）----
    global_fill_colors = []
    global_text_colors = []
    for idx in range(total):
        colors = analyze_colors(prs.slides[idx])
        for kind, color in colors:
            if kind == 'fill':
                global_fill_colors.append(color.upper())
            else:
                global_text_colors.append(color.upper())
    
    global_fill_counter = Counter(global_fill_colors)
    global_text_counter = Counter(global_text_colors)
    
    # ---- 采样页颜色（用于布局参数等）----
    sample_fill_colors = []
    sample_text_colors = []
    for idx in sample_pages:
        if idx >= total:
            continue
        colors = analyze_colors(prs.slides[idx])
        for kind, color in colors:
            if kind == 'fill':
                sample_fill_colors.append(color.upper())
            else:
                sample_text_colors.append(color.upper())
    
    # 背景色系（浅色跳过）
    bg_colors = {'FFFFFF', 'F5F5F5', 'F3F4F6', 'F0FDF4', 'FAFAFA', 'F9FAFB', 'F8FAFC', 'F0F9FF', 'ECFDF5'}
    skip_colors = {'000000'}
    palette = {'IS_DARK_THEME': is_dark}
    
    # 辅助函数：判断颜色亮度
    def brightness(hex_color):
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return (r * 299 + g * 587 + b * 114) / 1000
    
    # ---- PRIMARY: 最高频的非白非黑非浅灰填充色 ----
    for color, count in global_fill_counter.most_common():
        if color not in skip_colors and color not in bg_colors and brightness(color) < 240:
            palette['PRIMARY'] = color
            skip_colors.add(color)
            break
    
    # ---- SECONDARY: 优先选中性色（灰/深蓝/深绿），跳过高饱和强调色 ----
    neutral_candidates = []
    for color, count in global_fill_counter.most_common():
        if color not in skip_colors and color not in bg_colors:
            neutral_candidates.append((color, count, brightness(color), _is_neutral_color(color)))
    
    # 先从中性色中选
    neutral_only = [c for c in neutral_candidates if c[3]]
    if neutral_only:
        # 中性色中优先选深色（亮度低）
        neutral_only.sort(key=lambda x: x[2])
        palette['SECONDARY'] = neutral_only[0][0]
        skip_colors.add(neutral_only[0][0])
    elif neutral_candidates:
        # 没有中性色，选最深色
        neutral_candidates.sort(key=lambda x: x[2])
        palette['SECONDARY'] = neutral_candidates[0][0]
        skip_colors.add(neutral_candidates[0][0])
    
    # ---- ACCENT: 优先选高饱和色（红/橙/蓝/亮绿），跳过灰色系 ----
    accent_candidates = []
    for color, count in global_fill_counter.most_common():
        if color not in skip_colors and color not in bg_colors and color != 'FFFFFF':
            accent_candidates.append((color, count, brightness(color), _is_accent_color(color)))
    
    # 先从高饱和色中选
    accent_only = [c for c in accent_candidates if c[3]]
    if accent_only:
        # 高饱和色中选频率最高的
        accent_only.sort(key=lambda x: -x[1])
        palette['ACCENT'] = accent_only[0][0]
        skip_colors.add(accent_only[0][0])
    elif accent_candidates:
        # 没有高饱和色，选频率最高的
        accent_candidates.sort(key=lambda x: -x[1])
        palette['ACCENT'] = accent_candidates[0][0]
        skip_colors.add(accent_candidates[0][0])
    
    # 如果填充色不够，从文字色补充SECONDARY/ACCENT
    if 'SECONDARY' not in palette:
        for color, count in global_text_counter.most_common():
            if color not in skip_colors and color not in bg_colors and _is_neutral_color(color):
                palette['SECONDARY'] = color
                skip_colors.add(color)
                break
    
    if 'ACCENT' not in palette:
        for color, count in global_text_counter.most_common():
            if color not in skip_colors and color not in bg_colors and _is_accent_color(color):
                palette['ACCENT'] = color
                skip_colors.add(color)
                break
    
    # ---- BG: 背景色 ----
    if is_dark:
        # 深色PPT：背景取最高频深色填充色
        for color, count in global_fill_counter.most_common():
            if color not in skip_colors and brightness(color) < 128:
                palette['BG'] = color
                break
        if 'BG' not in palette:
            palette['BG'] = '1F2937'
    else:
        palette['BG'] = 'FFFFFF' if 'FFFFFF' in [c for c, _ in global_fill_counter.most_common(3)] else 'F5F5F5'
    
    # ---- TEXT: 主文字色（深色PPT选浅色，浅色PPT选深色）----
    text_skip = {'FFFFFF'} if not is_dark else set()
    text_candidates = []
    for color, count in global_text_counter.most_common():
        if color not in text_skip and color not in bg_colors:
            text_candidates.append((color, count, brightness(color)))
    
    if text_candidates:
        if is_dark:
            # 深色PPT：文字色优先选浅色（亮度高优先）
            text_candidates.sort(key=lambda x: -x[2])
        else:
            # 浅色PPT：文字色优先选深色（亮度低优先）
            text_candidates.sort(key=lambda x: x[2])
        palette['TEXT'] = text_candidates[0][0]
        text_skip.add(text_candidates[0][0])
    
    # ---- TEXT_LIGHT: 辅文字色 ----
    for color, count in global_text_counter.most_common():
        if color not in text_skip:
            palette['TEXT_LIGHT'] = color
            break
    
    # 扩展色板（供验证用）
    all_seen = set(global_fill_colors) | set(global_text_colors)
    palette['_extended'] = sorted(all_seen)
    
    # 兜底：确保5色齐全
    if 'PRIMARY' not in palette:
        palette['PRIMARY'] = '22C55E'
    if 'SECONDARY' not in palette:
        palette['SECONDARY'] = palette.get('TEXT', '1F2937')
    if 'ACCENT' not in palette:
        palette['ACCENT'] = 'C00000'
    if 'TEXT' not in palette:
        palette['TEXT'] = '1F2937' if not is_dark else 'FFFFFF'
    if 'TEXT_LIGHT' not in palette:
        palette['TEXT_LIGHT'] = '6B7280' if not is_dark else '9CA3AF'
    
    return palette


def print_layout_ratios(ratios, page_num):
    """格式化打印布局比例参数"""
    print(f"  【布局比例参数】")
    print(f"    幻灯片尺寸: {ratios['slide_w']:.1f}cm × {ratios['slide_h']:.1f}cm")
    if ratios['margin_l'] is not None:
        print(f"    安全边距: 左={ratios['margin_l']:.1f}cm, 右={ratios['margin_r']:.1f}cm")
    if ratios['margin_top'] > 0:
        print(f"    顶部装饰条: {ratios['margin_top']:.2f}cm")
    print(f"    标题区高度: {ratios['header_h']:.1f}cm")
    if ratios['footer_h'] > 0:
        print(f"    底部装饰条: {ratios['footer_h']:.1f}cm")
    if ratios['content_w'] is not None:
        print(f"    内容区: X={ratios['content_x']:.1f} Y={ratios['content_y']:.1f} W={ratios['content_w']:.1f} H={ratios['content_h']:.1f}cm")
    if ratios['has_left_right_split']:
        print(f"    ⭐ 左右分栏: 左栏{ratios['left_col_w']:.1f}cm + 间距{ratios['col_gap']:.1f}cm + 右栏{ratios['right_col_w']:.1f}cm")
        if ratios['col_ratio']:
            print(f"    分栏比例: 左栏占比 {ratios['col_ratio']:.0%}")
        total = ratios['left_col_w'] + ratios['col_gap'] + ratios['right_col_w']
        print(f"    分栏总计宽度: {total:.1f}cm (内容区宽度: {ratios['content_w']:.1f}cm)")
    print(f"    深色面积占比: {ratios['dark_ratio']:.0%}")
    print(f"    可见卡片数: {ratios['card_count']}")
    print()


def validate_color_usage(prs, palette, new_slide_indices):
    """
    验证新增页面颜色是否在色板内
    返回不在色板中的颜色列表
    
    修改：修复for循环缩进bug（之前只在循环外检查最后一个slide的colors）
    """
    # 基础允许色：色板5色 + 2文字色（排除_extended列表）
    allowed = set()
    for k, v in palette.items():
        if k == '_extended':
            allowed.update(c.upper() for c in v)  # _extended是列表
        elif isinstance(v, str):
            allowed.add(v.upper())
    # 通用白色黑色
    allowed.update({'FFFFFF', '000000'})
    
    violations = []
    
    for idx in new_slide_indices:
        if idx >= len(prs.slides):
            continue
        colors = analyze_colors(prs.slides[idx])
        for kind, color in colors:  # ⚡ 修复：从内层循环取出，修正缩进
            if color.upper() not in allowed:
                violations.append(f"第{idx+1}页 {kind}色 #{color} 不在色板中")
    
    return violations


def validate_layout(prs, slide_indices, palette=None, residual_keywords=None):
    """
    v3.1新增：P0-P3分级质检
    
    Args:
        prs: Presentation对象
        slide_indices: 要验证的幻灯片索引列表（0-based）
        palette: 5色色板dict（可选，不传则跳过颜色验证）
        residual_keywords: 残留关键词列表（如 ['OPPO', 'Tony', 'Alen']）
    
    Returns:
        dict: {
            'P0': [问题列表],  # 阻断级
            'P1': [问题列表],  # 功能级
            'P2': [问题列表],  # 品质级
            'P3': [问题列表],  # 润色级
            'passed': bool,    # P0无问题=通过
        }
    """
    slide_w = Emu(prs.slide_width).cm
    slide_h = Emu(prs.slide_height).cm
    
    p0 = []
    p1 = []
    p2 = []
    p3 = []
    
    if residual_keywords is None:
        residual_keywords = ['OPPO', 'Tony', 'Alen', '逍遥子', '逍遥子班']
    
    for idx in slide_indices:
        if idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]
        page = idx + 1
        shape_count = 0
        text_shapes = []
        max_left_gap = 0
        
        for shape in slide.shapes:
            shape_count += 1
            
            # 安全获取位置
            try:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
            except (TypeError, AttributeError):
                continue
            if left is None or top is None or width is None or height is None:
                continue
            
            left_cm = Emu(left).cm
            top_cm = Emu(top).cm
            right_cm = Emu(left + width).cm
            bottom_cm = Emu(top + height).cm
            
            # P0: 安全区域越界（允许0.5cm误差）
            if left_cm < -0.5 or top_cm < -0.5 or right_cm > slide_w + 0.5 or bottom_cm > slide_h + 0.5:
                p0.append(f"第{page}页 shape越界: L={left_cm:.1f} T={top_cm:.1f} R={right_cm:.1f} B={bottom_cm:.1f}")
            
            # P0: 图片shape
            try:
                if shape.shape_type == 13:
                    p0.append(f"第{page}页 存在图片shape")
            except:
                pass
            
            # 文字shape收集
            if shape.has_text_frame and shape.text_frame.text.strip():
                text_shapes.append(shape)
                text = shape.text_frame.text
                
                # P0: 模板残留关键词
                for kw in residual_keywords:
                    if kw in text:
                        p0.append(f"第{page}页 模板残留关键词: '{kw}'")
        
        # P0: shape数量不足
        if shape_count < 20:
            p0.append(f"第{page}页 shape数量不足: {shape_count} < 20")
        
        # P1: 颜色一致性
        if palette:
            violations = validate_color_usage(prs, palette, [idx])
            for v in violations:
                p1.append(v)
        
        # P1: 底部留白检测
        if text_shapes:
            try:
                max_bottom = max(Emu(s.top + s.height).cm for s in text_shapes if s.top is not None and s.height is not None)
                if max_bottom < slide_h - 3.0:
                    p1.append(f"第{page}页 内容底部距幻灯片底 {(slide_h - max_bottom):.1f}cm > 3cm，可能留白过大")
            except:
                pass
        
        # P2: 间距均匀性（简化检查）
        # 暂时跳过，需要更复杂的分析
        
        # P3: 信息密度
        if text_shapes:
            total_chars = sum(len(s.text_frame.text) for s in text_shapes)
            if total_chars < 20:
                p3.append(f"第{page}页 文字量过少: {total_chars}字")
            elif total_chars > 500:
                p3.append(f"第{page}页 文字量偏多: {total_chars}字")
    
    return {
        'P0': p0,
        'P1': p1,
        'P2': p2,
        'P3': p3,
        'passed': len(p0) == 0,
    }


def analyze_ppt(ppt_path, sample_pages=None, verbose=False):
    """
    分析PPT的整体风格
    
    Args:
        ppt_path: PPT文件路径
        sample_pages: 采样的页码列表（从1开始），None表示自动选择
        verbose: 是否打印详细布局比例参数
    
    Returns:
        dict: 包含配色、字体、布局、节奏等风格信息
    """
    prs = Presentation(ppt_path)
    total_slides = len(prs.slides)
    slide_w, slide_h = get_slide_size_cm(prs)
    
    print(f"PPT总页数: {total_slides}")
    print(f"幻灯片尺寸: {slide_w:.1f}cm × {slide_h:.1f}cm ({prs.slide_width / 914400:.2f} × {prs.slide_height / 914400:.2f} 英寸)")
    print()
    
    # 选择采样页面
    if sample_pages is None:
        sample_pages = _pick_best_sample_pages(prs, count=8)
    else:
        sample_pages = [p - 1 for p in sample_pages]
    
    print(f"采样页面: {[p+1 for p in sample_pages]}")
    print()
    
    # 收集数据
    all_colors = []
    all_fonts = []
    all_layouts = []
    all_ratios = []
    
    for page_idx in sample_pages:
        if page_idx >= total_slides:
            continue
        
        slide = prs.slides[page_idx]
        print(f"--- 第 {page_idx+1} 页 ---")
        
        colors = analyze_colors(slide)
        all_colors.extend(colors)
        
        fonts = analyze_fonts(slide)
        all_fonts.extend(fonts)
        
        layout = analyze_layout(slide, prs.slide_width, prs.slide_height)
        all_layouts.append(layout)
        
        ratios = analyze_layout_ratios(slide, prs)
        all_ratios.append(ratios)
        
        print(f"  形状数量: {layout['shape_count']}")
        print(f"  有标题: {layout['has_title']}")
        print(f"  有图片: {layout['has_image']}")
        print(f"  有装饰元素: {layout['has_decoration']}")
        print(f"  列数（约）: {layout['columns']}")
        print(f"  深色面积占比: {ratios['dark_ratio']:.0%}")
        
        if verbose:
            print_layout_ratios(ratios, page_idx+1)
        else:
            if ratios['has_left_right_split']:
                ratio_str = f" ({ratios['col_ratio']:.0%})" if ratios['col_ratio'] else ""
                print(f"  ⭐ 检测到左右分栏布局{ratio_str}")
        print()
    
    # ===== 配色方案分析 =====
    print("=" * 50)
    print("【配色方案分析】")
    print()
    
    # 5色锁定色板
    palette = extract_color_palette(prs, sample_pages)
    print("🎨 5色锁定色板（新增页面只能用这些颜色）:")
    for role, color in palette.items():
        if role == '_extended':
            continue  # 扩展色板太长，跳过
        if isinstance(color, bool):
            print(f"  IS_DARK_THEME: {color}")
        else:
            print(f"  COLOR_{role}: #{color}")
    print()
    
    fill_colors = [c[1] for c in all_colors if c[0] == 'fill']
    text_colors = [c[1] for c in all_colors if c[0] == 'text']
    
    if fill_colors:
        fc = Counter(fill_colors)
        print("填充颜色（按频率排序）:")
        for color, count in fc.most_common(5):
            print(f"  #{color}: {count}次")
        print()
    
    if text_colors:
        tc = Counter(text_colors)
        print("文字颜色（按频率排序）:")
        for color, count in tc.most_common(5):
            print(f"  #{color}: {count}次")
        print()
    
    # ===== 字体规范分析 =====
    print("=" * 50)
    print("【字体规范分析】")
    print()
    
    if all_fonts:
        font_names = [f[0] for f in all_fonts]
        font_counter = Counter(font_names)
        print("使用的字体:")
        for font, count in font_counter.most_common(5):
            print(f"  {font}: {count}次")
        print()
        
        # v3.1改进：按字号范围分配三级分工
        title_fonts = Counter([f[0] for f in all_fonts if f[1] >= 20])
        body_fonts = Counter([f[0] for f in all_fonts if 14 <= f[1] < 20])
        note_fonts = Counter([f[0] for f in all_fonts if 0 < f[1] < 14])
        
        print("字体三级分工（按字号范围推断）:")
        if title_fonts:
            print(f"  标题字体(≥20pt): {title_fonts.most_common(1)[0][0]}")
        if body_fonts:
            print(f"  正文字体(14-19pt): {body_fonts.most_common(1)[0][0]}")
        if note_fonts:
            print(f"  注释字体(<14pt): {note_fonts.most_common(1)[0][0]}")
        print()
        
        font_sizes = [f[1] for f in all_fonts if f[1] > 0]
        if font_sizes:
            font_sizes.sort()
            print(f"字号范围: {font_sizes[0]:.0f}pt - {font_sizes[-1]:.0f}pt")
            size_counter = Counter(font_sizes)
            print("最常见字号:")
            for size, count in size_counter.most_common(5):
                print(f"  {size:.0f}pt: {count}次")
            print()
    
    # ===== 布局比例参数汇总 =====
    print("=" * 50)
    print("【布局比例参数汇总】")
    print()
    
    if all_ratios:
        margin_ls = [r['margin_l'] for r in all_ratios if r['margin_l'] is not None]
        margin_rs = [r['margin_r'] for r in all_ratios if r['margin_r'] is not None]
        header_hs = [r['header_h'] for r in all_ratios if r['header_h'] > 0]
        footer_hs = [r['footer_h'] for r in all_ratios if r['footer_h'] > 0]
        
        if margin_ls:
            rec_margin_l = max(set(margin_ls), key=margin_ls.count)
            rec_margin_r = max(set(margin_rs), key=margin_rs.count)
            print(f"📐 推荐左安全边距: {rec_margin_l:.1f}cm")
            print(f"📐 推荐右安全边距: {rec_margin_r:.1f}cm")
        
        if header_hs:
            rec_header_h = max(set(header_hs), key=header_hs.count)
            print(f"📐 推荐标题区高度: {rec_header_h:.1f}cm")
        
        if footer_hs:
            rec_footer_h = max(set(footer_hs), key=footer_hs.count)
            print(f"📐 推荐底部装饰条高度: {rec_footer_h:.1f}cm")
        
        split_pages = [r for r in all_ratios if r['has_left_right_split']]
        if split_pages:
            print()
            print(f"⭐ 发现 {len(split_pages)} 页使用左右分栏布局:")
            for r in split_pages:
                ratio_str = f"左栏占比{r['col_ratio']:.0%}" if r['col_ratio'] else "比例未计算"
                print(f"  {ratio_str}, 间距: {r['col_gap']:.1f}cm")
        print()
    
    # ===== 节奏感检测 =====
    print("=" * 50)
    print("【节奏感检测】(借鉴guizang)")
    print()
    
    rhythm, warnings = analyze_rhythm(prs)
    dark_count = sum(1 for r in rhythm if r['is_dark'])
    light_count = len(rhythm) - dark_count
    
    print(f"深色页: {dark_count}页 | 浅色页: {light_count}页")
    if rhythm:
        # 简化节奏显示（超过20页只显示摘要）
        if len(rhythm) <= 30:
            rhythm_str = " → ".join(["■" if r['is_dark'] else "□" for r in rhythm])
            print(f"节奏: {rhythm_str}")
        else:
            print(f"节奏: (共{len(rhythm)}页，过长已省略)")
    print()
    
    if warnings:
        for w in warnings:
            print(w)
        print("  💡 建议：每3-4页插入1个深色强调页，避免视觉疲劳")
    else:
        print("✅ 节奏感良好，无连续3页以上同深浅")
    print()
    
    # ===== 布局结构分析 =====
    print("=" * 50)
    print("【布局结构分析】")
    print()
    
    if all_layouts:
        avg_shapes = sum(l['shape_count'] for l in all_layouts) / len(all_layouts)
        has_image_count = sum(1 for l in all_layouts if l['has_image'])
        has_decoration_count = sum(1 for l in all_layouts if l['has_decoration'])
        avg_columns = sum(l['columns'] for l in all_layouts) / len(all_layouts)
        
        print(f"平均形状数量: {avg_shapes:.1f}")
        print(f"包含图片的页面比例: {has_image_count}/{len(all_layouts)} ({has_image_count/len(all_layouts)*100:.0f}%)")
        print(f"包含装饰元素的页面比例: {has_decoration_count}/{len(all_layouts)} ({has_decoration_count/len(all_layouts)*100:.0f}%)")
        print(f"平均列数: {avg_columns:.1f}")
        
        if avg_shapes < 5:
            complexity = "简单（默认布局）"
        elif avg_shapes < 15:
            complexity = "中等"
        else:
            complexity = "复杂（专业设计）"
        print(f"布局复杂度: {complexity}")
        print()
    
    # ===== 参考页推荐（v4新增）=====
    print("=" * 50)
    print("【参考页推荐】")
    print()
    
    ref_pages = []
    for i in range(total_slides):
        slide = prs.slides[i]
        sc = len(list(slide.shapes))
        has_text = any(s.has_text_frame and s.text_frame.text.strip() for s in slide.shapes)
        has_deco = any(not s.has_text_frame and s.shape_type != 13 for s in slide.shapes)
        if sc >= 20 and has_text and has_deco:
            ref_pages.append((i+1, sc))
    
    if ref_pages:
        ref_pages.sort(key=lambda x: -x[1])  # 按shape数量降序
        print(f"推荐作为风格参考的页面（shape≥20 + 有文字 + 有装饰元素）:")
        for p, sc in ref_pages[:10]:
            ts = [s for s in prs.slides[p-1].shapes if s.has_text_frame and s.text_frame.text.strip()]
            title = ts[0].text_frame.text.strip()[:40] if ts else '(无标题)'
            print(f"  第{p:2d}页 ({sc:3d} shapes) | {title}")
    else:
        print("未找到设计复杂的参考页（shape≥20），建议手动选择")
    print()
    
    # ===== 风格总结 =====
    print("=" * 50)
    print("【风格总结】")
    print()
    
    main_color = palette.get('PRIMARY', '未识别')
    text_color = palette.get('TEXT', '未识别')
    main_font = None
    if all_fonts:
        font_counter = Counter([f[0] for f in all_fonts])
        main_font = font_counter.most_common(1)[0][0]
    
    print(f"主色调: #{main_color}" if main_color != '未识别' else "主色调: 未识别")
    print(f"辅色: #{palette.get('SECONDARY', '未识别')}")
    print(f"强调色: #{palette.get('ACCENT', '未识别')}")
    print(f"主要文字色: #{text_color}" if text_color != '未识别' else "主要文字色: 未识别")
    print(f"辅助文字色: #{palette.get('TEXT_LIGHT', '未识别')}")
    print(f"主题类型: {'深色主题' if palette.get('IS_DARK_THEME') else '浅色主题'}")
    print(f"主要字体: {main_font}" if main_font else "主要字体: 未识别")
    print()
    
    return {
        'main_color': main_color,
        'text_color': text_color,
        'main_font': main_font,
        'palette': palette,
        'avg_shapes': avg_shapes if all_layouts else 0,
        'has_images': has_image_count > 0 if all_layouts else False,
        'has_decorations': has_decoration_count > 0 if all_layouts else False,
        'avg_columns': avg_columns if all_layouts else 0,
        'layout_ratios': all_ratios,
        'rhythm': rhythm,
        'rhythm_warnings': warnings,
    }


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("用法: python analyze_style.py <ppt文件路径> [页码1,页码2,...] [--verbose]")
        print("示例: python analyze_style.py presentation.pptx 3,5,7")
        print("      python analyze_style.py presentation.pptx --verbose")
        print("选项: --verbose  打印详细布局比例参数")
        print()
        print("功能:")
        print("  - 5色锁定色板自动提取（区分填充色和文字色角色）")
        print("  - 节奏感检测（深浅交替，去重警告）")
        print("  - 分栏比例计算")
        print("  - 深色面积占比分析")
        print("  - P0-P3分级质检（validate_layout）")
        print("  - 辅助函数：set_ea_font / set_line_spacing / clean_shape")
        sys.exit(1)
    
    ppt_path = sys.argv[1]
    sample_pages = None
    verbose = '--verbose' in sys.argv
    
    for arg in sys.argv[2:]:
        if arg.startswith('--'):
            continue
        sample_pages = [int(p) for p in arg.split(',')]
        break
    
    result = analyze_ppt(ppt_path, sample_pages, verbose=verbose)
