#!/usr/bin/env python3
"""
Generate PowerPoint slide deck from analysis results.

Creates a presentation with:
- Title slide
- Executive summary
- Key findings with optional charts
- Recommendations
- Appendix placeholder

Usage:
    python generate_pptx_summary.py <config.json> [--output <filename.pptx>]

Config JSON format:
{
    "title": "Analysis Title",
    "subtitle": "Optional subtitle",
    "date": "2024-01-15",
    "author": "Analyst Name",
    "executive_summary": "One sentence summary of key finding",
    "key_findings": [
        {
            "headline": "Claim-based headline",
            "bullets": ["Point 1", "Point 2"],
            "chart_path": "optional/path/to/chart.png"
        }
    ],
    "recommendations": [
        "Recommendation 1",
        "Recommendation 2"
    ],
    "caveats": [
        "Key limitation 1",
        "Key limitation 2"
    ]
}
"""

import sys
import json
import argparse
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# Slide dimensions (standard 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
COLORS = {
    'primary': RgbColor(46, 134, 171),    # #2E86AB - Blue
    'secondary': RgbColor(162, 59, 114),  # #A23B72 - Purple
    'dark': RgbColor(51, 51, 51),         # #333333 - Dark gray
    'light': RgbColor(108, 117, 125),     # #6C757D - Light gray
}


def create_title_slide(prs, config):
    """Create the title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = config.get('title', 'Analysis Report')
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    if config.get('subtitle'):
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(12.333), Inches(0.75)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = config['subtitle']
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = COLORS['light']
        subtitle_para.alignment = PP_ALIGN.CENTER

    # Date and author
    meta_text = []
    if config.get('date'):
        meta_text.append(config['date'])
    if config.get('author'):
        meta_text.append(config['author'])

    if meta_text:
        meta_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5)
        )
        meta_frame = meta_box.text_frame
        meta_para = meta_frame.paragraphs[0]
        meta_para.text = " | ".join(meta_text)
        meta_para.font.size = Pt(14)
        meta_para.font.color.rgb = COLORS['light']
        meta_para.alignment = PP_ALIGN.CENTER

    return slide


def create_executive_summary_slide(prs, config):
    """Create executive summary slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Headline
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Executive Summary"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # Key finding box
    finding_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.75), Inches(12.333), Inches(2)
    )
    finding_box.fill.solid()
    finding_box.fill.fore_color.rgb = RgbColor(240, 248, 255)  # Light blue
    finding_box.line.color.rgb = COLORS['primary']

    # Key finding text
    finding_text = slide.shapes.add_textbox(
        Inches(0.75), Inches(2), Inches(11.833), Inches(1.5)
    )
    finding_frame = finding_text.text_frame
    finding_frame.word_wrap = True

    key_para = finding_frame.paragraphs[0]
    key_para.text = "Key Finding"
    key_para.font.size = Pt(14)
    key_para.font.bold = True
    key_para.font.color.rgb = COLORS['primary']

    summary_para = finding_frame.add_paragraph()
    summary_para.text = config.get('executive_summary', 'Summary not provided')
    summary_para.font.size = Pt(20)
    summary_para.font.color.rgb = COLORS['dark']
    summary_para.space_before = Pt(12)

    # Recommendations preview
    if config.get('recommendations'):
        rec_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.25), Inches(12.333), Inches(2.5)
        )
        rec_frame = rec_box.text_frame
        rec_frame.word_wrap = True

        rec_title = rec_frame.paragraphs[0]
        rec_title.text = "Recommendations"
        rec_title.font.size = Pt(18)
        rec_title.font.bold = True
        rec_title.font.color.rgb = COLORS['dark']

        for i, rec in enumerate(config['recommendations'][:3], 1):
            rec_para = rec_frame.add_paragraph()
            rec_para.text = f"{i}. {rec}"
            rec_para.font.size = Pt(16)
            rec_para.font.color.rgb = COLORS['dark']
            rec_para.space_before = Pt(8)

    return slide


def create_finding_slide(prs, finding, slide_num):
    """Create a slide for a key finding."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Headline (should be a claim, not just a topic)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.25)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = finding.get('headline', f'Finding {slide_num}')
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # Check if there's a chart
    has_chart = finding.get('chart_path') and Path(finding['chart_path']).exists()

    if has_chart:
        # Layout with chart on right
        # Bullets on left
        bullets_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.75), Inches(5.5), Inches(5)
        )
        bullets_frame = bullets_box.text_frame
        bullets_frame.word_wrap = True

        for bullet in finding.get('bullets', []):
            para = bullets_frame.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(16)
            para.font.color.rgb = COLORS['dark']
            para.space_before = Pt(12)

        # Chart on right
        slide.shapes.add_picture(
            finding['chart_path'],
            Inches(6.5), Inches(1.75),
            width=Inches(6.333), height=Inches(5)
        )
    else:
        # Full-width bullets
        bullets_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.75), Inches(12.333), Inches(5)
        )
        bullets_frame = bullets_box.text_frame
        bullets_frame.word_wrap = True

        for bullet in finding.get('bullets', []):
            para = bullets_frame.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(18)
            para.font.color.rgb = COLORS['dark']
            para.space_before = Pt(16)

    return slide


def create_recommendations_slide(prs, config):
    """Create recommendations and next steps slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Recommendations & Next Steps"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # Recommendations
    if config.get('recommendations'):
        rec_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(6), Inches(4.5)
        )
        rec_frame = rec_box.text_frame
        rec_frame.word_wrap = True

        rec_title = rec_frame.paragraphs[0]
        rec_title.text = "Recommendations"
        rec_title.font.size = Pt(20)
        rec_title.font.bold = True
        rec_title.font.color.rgb = COLORS['primary']

        for i, rec in enumerate(config['recommendations'], 1):
            rec_para = rec_frame.add_paragraph()
            rec_para.text = f"{i}. {rec}"
            rec_para.font.size = Pt(16)
            rec_para.font.color.rgb = COLORS['dark']
            rec_para.space_before = Pt(12)

    # Caveats
    if config.get('caveats'):
        caveat_box = slide.shapes.add_textbox(
            Inches(6.833), Inches(1.5), Inches(6), Inches(4.5)
        )
        caveat_frame = caveat_box.text_frame
        caveat_frame.word_wrap = True

        caveat_title = caveat_frame.paragraphs[0]
        caveat_title.text = "Key Caveats"
        caveat_title.font.size = Pt(20)
        caveat_title.font.bold = True
        caveat_title.font.color.rgb = COLORS['secondary']

        for caveat in config['caveats']:
            caveat_para = caveat_frame.add_paragraph()
            caveat_para.text = f"• {caveat}"
            caveat_para.font.size = Pt(14)
            caveat_para.font.color.rgb = COLORS['light']
            caveat_para.space_before = Pt(10)

    return slide


def create_appendix_slide(prs):
    """Create appendix placeholder slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3), Inches(12.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Appendix"
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['light']
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_para = title_frame.add_paragraph()
    subtitle_para.text = "Detailed methodology, data sources, and supporting analysis"
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = COLORS['light']
    subtitle_para.alignment = PP_ALIGN.CENTER

    return slide


def generate_presentation(config: dict, output_path: Path) -> Path:
    """
    Generate a PowerPoint presentation from config.

    Args:
        config: Dictionary with presentation content
        output_path: Path for output .pptx file

    Returns:
        Path to created presentation
    """
    if not PPTX_AVAILABLE:
        raise ImportError(
            "python-pptx is required. Install with: pip install python-pptx"
        )

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Add slides
    create_title_slide(prs, config)
    create_executive_summary_slide(prs, config)

    # Add finding slides
    for i, finding in enumerate(config.get('key_findings', []), 1):
        create_finding_slide(prs, finding, i)

    # Add recommendations
    if config.get('recommendations') or config.get('caveats'):
        create_recommendations_slide(prs, config)

    # Add appendix placeholder
    create_appendix_slide(prs)

    # Save
    output_path = Path(output_path)
    prs.save(str(output_path))

    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Generate PowerPoint slide deck from analysis results"
    )
    parser.add_argument(
        "config",
        help="Path to JSON config file with presentation content"
    )
    parser.add_argument(
        "--output", "-o",
        default=None,
        help="Output filename (default: derived from config title)"
    )

    args = parser.parse_args()

    # Check for python-pptx
    if not PPTX_AVAILABLE:
        print("Error: python-pptx is required. Install with: pip install python-pptx")
        sys.exit(1)

    # Load config
    config_path = Path(args.config)
    if not config_path.exists():
        print(f"Error: Config file not found: {config_path}")
        sys.exit(1)

    with open(config_path) as f:
        config = json.load(f)

    # Determine output path
    if args.output:
        output_path = Path(args.output)
    else:
        title_slug = config.get('title', 'presentation').lower()
        title_slug = "".join(c if c.isalnum() else "_" for c in title_slug)
        date_str = datetime.now().strftime("%Y%m%d")
        output_path = Path(f"{title_slug}_{date_str}.pptx")

    try:
        result_path = generate_presentation(config, output_path)
        print(f"Created presentation: {result_path}")
        print(f"\nSlide count: {len(Presentation(str(result_path)).slides)}")
    except Exception as e:
        print(f"Error generating presentation: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
