"""Convert any runtime.js-based HTML deck to PPTX via Playwright 2x screenshots + python-pptx.

Usage: python convert_deck_to_pptx.py <deck_html_path>
"""
import sys, re, io, os
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

if len(sys.argv) < 2:
    print("Usage: python convert_deck_to_pptx.py <deck_html_path>")
    sys.exit(1)

html_file = Path(sys.argv[1]).resolve()
if not html_file.exists():
    print(f"File not found: {html_file}")
    sys.exit(1)

DIR = html_file.parent
deck_name = html_file.stem
output_pptx = DIR / f"{deck_name}.pptx"
print(f"Converting: {html_file.name} -> {output_pptx.name}")

# Read and inline CSS
html = html_file.read_text(encoding="utf-8")

def inline_css(content):
    def replace_link(match):
        href = match.group(1)
        css_path = DIR / href
        if css_path.exists():
            css = css_path.read_text(encoding="utf-8")
            return f"<style>{css}</style>"
        return match.group(0)
    return re.sub(r'<link[^>]+href="([^"]+)"[^>]*>', replace_link, content)

html = inline_css(html)
html = re.sub(r'<link[^>]*fonts\.googleapis\.com[^>]*>', '', html)
html = re.sub(r'<link[^>]*fonts\.gstatic\.com[^>]*>', '', html)

font_fix = """
<style>
* { font-family: 'Microsoft YaHei', 'SimHei', 'Noto Sans SC', sans-serif !important; }
h1, h2, h3, .ts-h1, .ts-h2 { font-family: 'Microsoft YaHei', 'SimHei', 'Noto Serif SC', serif !important; }
code, pre, .ts-codebox { font-family: 'Consolas', 'Courier New', monospace !important; }
</style>
"""
html = html.replace("</head>", f"{font_fix}</head>")

# base.css sets .slide{opacity:0;position:absolute} and .slide.is-active{opacity:1}
# We use the .is-active mechanism — it's already in base.css
# Just need to kill transitions so opacity changes are instant
kill_css = """
<style>.slide{transition:none!important}</style>
"""
html = html.replace("</head>", f"{kill_css}</head>")

W, H = 1920, 1080  # standard 16:9 — match browser rendering 1:1 for accurate proportions

with sync_playwright() as p:
    try:
        browser = p.chromium.launch(channel="msedge")
    except Exception:
        print("msedge 不可用，回退到默认 Chromium")
        browser = p.chromium.launch()
    page = browser.new_page(viewport={"width": W, "height": H})

    page.set_content(html, wait_until="networkidle")
    page.wait_for_timeout(1000)

    slides = page.query_selector_all(".slide")
    n = len(slides)
    print(f"Found {n} slides at {W}x{H}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    for i in range(n):
        # Activate only slide i, deactivate all others
        page.evaluate(f"""() => {{
            document.querySelectorAll(".slide").forEach(function(s, idx) {{
                if (idx === {i}) {{
                    s.classList.add("is-active");
                }} else {{
                    s.classList.remove("is-active");
                }}
            }});
        }}""")
        page.wait_for_timeout(300)

        png_bytes = page.screenshot(full_page=False)
        img = Image.open(io.BytesIO(png_bytes))

        slide_obj = prs.slides.add_slide(blank_layout)
        img_stream = io.BytesIO()
        img.save(img_stream, format="PNG")
        img_stream.seek(0)

        slide_obj.shapes.add_picture(img_stream, 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
        print(f"  Slide {i+1}/{n} done")

    browser.close()

prs.save(str(output_pptx))
print(f"PPTX saved: {output_pptx} ({n} slides)")
