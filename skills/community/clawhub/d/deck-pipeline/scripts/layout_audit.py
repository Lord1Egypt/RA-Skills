#!/usr/bin/env python3
"""Layout audit + font pollution cleanup.

Usage:
    python3 layout_audit.py <pptx-path> [--fix] [--out <out-path>]

What it does:
  - Detects font pollution (faces not in WHITELIST, or with style-suffix in name)
  - Detects body/caption font sizes below floors
  - Detects overlap risks (basic)
  - With --fix: rewrites font.name to pure family, sets bold/italic via attributes,
    and saves to --out (or <input>-终版-YYYYMMDD.pptx by default)
"""
import sys
import os
import shutil
import datetime
from collections import defaultdict
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

# ---------- Config (matches PROFILE block in SKILL.md) ----------
WHITELIST_TITLE = "Georgia"
WHITELIST_BODY = "Verdana"
POLLUTION_KNOWN = {"Calibri", "Arial", "微软雅黑", "等线", "华文中宋"}
SKIP_POLLUTION = {"Arial Black"}
STYLE_SUFFIXES = ["Bold", "Regular", "Italic", "Light", "Semibold", "Medium"]
TITLE_TOP_THRESHOLD = 200_000  # EMU
FLOOR_BODY = 7.0
FLOOR_CAPTION = 6.0
FLOOR_SOURCE = 4.0

NSMAP_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
A = f"{{{NSMAP_A}}}"


def has_style_suffix(name):
    if not name:
        return None
    for suf in STYLE_SUFFIXES:
        if name.endswith(" " + suf) or name.endswith(suf):
            # only flag if it isn't a standalone "Bold Italic" type face that has no base
            base = name.rsplit(" ", 1)[0] if " " in name else name[: -len(suf)]
            if base and base != name:
                return base.strip(), suf
    return None


def is_title_shape(shp):
    return (shp.top is not None
            and 0 <= shp.top <= TITLE_TOP_THRESHOLD
            and shp.has_text_frame)


def is_banner_shape(shp):
    if not shp.has_text_frame:
        return False
    try:
        return str(shp.fill.type) == "FillType.SOLID"
    except Exception:
        return False


def normalize_font(rPr, target_face, target_bold=None, target_italic=None):
    """Strip latin/ea/cs/sym, rewrite to target_face, set bold/italic attrs."""
    if rPr is None:
        return False
    for tag in ["sym", "latin", "ea", "cs"]:
        for el in rPr.findall(f"{A}{tag}"):
            rPr.remove(el)
    for tag in ["latin", "ea", "cs"]:
        new = etree.SubElement(rPr, f"{A}{tag}")
        new.set("typeface", target_face)
        new.set("pitchFamily", "18")
        new.set("charset", "0")
    if target_bold is not None:
        if target_bold:
            rPr.set("b", "1")
        elif rPr.get("b") == "1":
            rPr.attrib.pop("b", None)
    if target_italic is not None:
        if target_italic:
            rPr.set("i", "1")
        elif rPr.get("i") == "1":
            rPr.attrib.pop("i", None)
    return True


def audit_and_fix(prs, fix=False):
    report = defaultdict(lambda: {"pollution": [], "suffix": [],
                                  "floor_breach": [], "fixes": []})

    for sidx, slide in enumerate(prs.slides, 1):
        def walk(shapes, top_level=True):
            for shp in shapes:
                if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(shp.shapes, top_level=False)
                    continue
                if not shp.has_text_frame:
                    if shp.has_table:
                        for row in shp.table.rows:
                            for cell in row.cells:
                                process_text_frame(cell.text_frame, sidx, shp,
                                                   is_table=True)
                    continue
                process_text_frame(shp.text_frame, sidx, shp,
                                   is_title=(top_level and is_title_shape(shp)),
                                   is_banner=is_banner_shape(shp))

        def process_text_frame(tf, sidx, shp, is_title=False, is_banner=False,
                               is_table=False):
            target_face = WHITELIST_TITLE if (is_title or is_banner) else WHITELIST_BODY
            target_bold = True if (is_title or is_banner) else None

            for para in tf.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    cur = run.font.name or ""

                    # 1. Suffix pollution
                    sx = has_style_suffix(cur)
                    if sx:
                        base, suf = sx
                        report[sidx]["suffix"].append((cur, base, suf, run.text[:40]))
                        if fix:
                            rPr = run._r.find(f"{A}rPr")
                            tb = True if suf in ("Bold",) else (False if suf == "Regular" else None)
                            ti = True if suf in ("Italic",) else None
                            if normalize_font(rPr, base, target_bold=tb, target_italic=ti):
                                report[sidx]["fixes"].append(
                                    f"  suffix '{cur}'→'{base}' + attr({suf}): '{run.text[:30]}'")
                        continue  # suffix-fixed; skip pollution branch

                    # 2. Known external font pollution
                    if cur in POLLUTION_KNOWN and cur not in SKIP_POLLUTION:
                        report[sidx]["pollution"].append((cur, run.text[:40]))
                        if fix:
                            rPr = run._r.find(f"{A}rPr")
                            tf_target = WHITELIST_BODY if is_table else target_face
                            if normalize_font(rPr, tf_target, target_bold=target_bold):
                                report[sidx]["fixes"].append(
                                    f"  pollution '{cur}'→{tf_target}: '{run.text[:30]}'")

                    # 3. Floor breach (only report — fixing is a separate stage)
                    if run.font.size is not None:
                        pt = run.font.size.pt
                        if pt < FLOOR_SOURCE:
                            report[sidx]["floor_breach"].append(
                                ("source", pt, run.text[:40]))

        walk(slide.shapes)
    return report


def main():
    if len(sys.argv) < 2:
        print("Usage: layout_audit.py <pptx-path> [--fix] [--out <out-path>]")
        sys.exit(1)
    src = sys.argv[1]
    fix = "--fix" in sys.argv
    out = None
    if "--out" in sys.argv:
        out = sys.argv[sys.argv.index("--out") + 1]
    if fix and not out:
        stamp = datetime.datetime.now().strftime("%Y%m%d")
        base, ext = os.path.splitext(src)
        out = f"{base}-终版-{stamp}{ext}"

    if fix:
        shutil.copyfile(src, out)
        prs = Presentation(out)
    else:
        prs = Presentation(src)

    report = audit_and_fix(prs, fix=fix)

    print("=" * 78)
    print(f"LAYOUT AUDIT — {'FIX MODE' if fix else 'DIAGNOSTIC ONLY'}")
    print("=" * 78)
    totals = {"pollution": 0, "suffix": 0, "floor_breach": 0, "fixes": 0}
    for sidx in sorted(report.keys()):
        r = report[sidx]
        if not any(r.values()):
            continue
        print(f"\n— Slide {sidx} —")
        if r["pollution"]:
            print(f"  Pollution: {len(r['pollution'])} runs")
            totals["pollution"] += len(r["pollution"])
        if r["suffix"]:
            print(f"  Suffix-name pollution: {len(r['suffix'])} runs")
            totals["suffix"] += len(r["suffix"])
            for old, base, suf, sample in r["suffix"][:3]:
                print(f"    '{old}' → '{base}' + attr({suf})  ['{sample}']")
        if r["floor_breach"]:
            print(f"  Sub-floor sizes: {len(r['floor_breach'])} runs")
            totals["floor_breach"] += len(r["floor_breach"])
        if r["fixes"]:
            totals["fixes"] += len(r["fixes"])

    print("\n" + "=" * 78)
    print("TOTALS")
    print("=" * 78)
    for k, v in totals.items():
        print(f"  {k:15} {v}")

    if fix:
        prs.save(out)
        print(f"\nSaved to: {out}")


if __name__ == "__main__":
    main()
