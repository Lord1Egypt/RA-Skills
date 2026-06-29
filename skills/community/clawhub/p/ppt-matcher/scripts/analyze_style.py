#!/usr/bin/env python3
"""
PPT风格分析工具
用于快速提取PPT的设计规范：配色、字体、布局结构等
"""

import sys
from collections import Counter
from pptx import Presentation
from pptx.util import Pt, Emu


def analyze_colors(slide):
    """分析单个幻灯片中的颜色使用情况"""
    colors = []
    
    for shape in slide.shapes:
        # 分析填充颜色
        try:
            if hasattr(shape, 'fill') and shape.fill.type == 1:  # solid fill
                color = shape.fill.fore_color.rgb
                colors.append(('fill', str(color)))
        except:
            pass
        
        # 分析文字颜色
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color and run.font.color.rgb:
                            color = run.font.color.rgb
                            colors.append(('text', str(color)))
                    except:
                        pass
    
    return colors


def analyze_fonts(slide):
    """分析单个幻灯片中的字体使用情况"""
    fonts = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font_name = run.font.name
                    font_size = run.font.size
                    font_bold = run.font.bold
                    
                    if font_name:
                        size_pt = font_size.pt if font_size else 0
                        fonts.append((font_name, size_pt, font_bold))
    
    return fonts


def analyze_layout(slide, slide_width, slide_height):
    """分析单个幻灯片的布局结构"""
    layout_info = {
        'shape_count': len(slide.shapes),
        'has_title': False,
        'has_image': False,
        'has_decoration': False,
        'columns': 0,
        'title_height': 0,
        'content_area_top': 0,
        'content_area_bottom': 0,
    }
    
    text_shapes = []
    image_shapes = []
    decoration_shapes = []
    
    for shape in slide.shapes:
        # 判断是否是图片
        if shape.shape_type == 13:  # picture
            layout_info['has_image'] = True
            image_shapes.append(shape)
            continue
        
        # 判断是否有文字
        if shape.has_text_frame and shape.text.strip():
            text_shapes.append(shape)
            # 判断是否是标题（位置靠上，字号较大）
            if shape.top < slide_height * 0.2:
                layout_info['has_title'] = True
                layout_info['title_height'] = max(layout_info['title_height'], shape.height)
                layout_info['content_area_top'] = max(layout_info['content_area_top'], shape.top + shape.height)
        else:
            # 可能是装饰元素（小形状、线条等）
            if shape.width < slide_width * 0.1 or shape.height < slide_height * 0.05:
                layout_info['has_decoration'] = True
                decoration_shapes.append(shape)
    
    # 分析内容区域
    if text_shapes:
        bottom_positions = [s.top + s.height for s in text_shapes if s.top > slide_height * 0.2]
        if bottom_positions:
            layout_info['content_area_bottom'] = max(bottom_positions)
    
    # 分析列数（简单判断：根据水平位置分布）
    if len(text_shapes) >= 3:
        x_positions = [s.left for s in text_shapes if s.top > slide_height * 0.2]
        if len(x_positions) >= 3:
            # 简单聚类判断列数
            x_positions.sort()
            columns = 1
            for i in range(1, len(x_positions)):
                if x_positions[i] - x_positions[i-1] > slide_width * 0.2:
                    columns += 1
            layout_info['columns'] = min(columns, 4)  # 最多算4列
    
    return layout_info


def analyze_ppt(ppt_path, sample_pages=None):
    """
    分析PPT的整体风格
    
    Args:
        ppt_path: PPT文件路径
        sample_pages: 采样的页码列表（从1开始），None表示自动选择风格统一的页面
    
    Returns:
        dict: 包含配色、字体、布局等风格信息
    """
    prs = Presentation(ppt_path)
    total_slides = len(prs.slides)
    
    print(f"PPT总页数: {total_slides}")
    print(f"幻灯片尺寸: {prs.slide_width} x {prs.slide_height} ({prs.slide_width / 914400:.2f} x {prs.slide_height / 914400:.2f} 英寸)")
    print()
    
    # 选择采样页面
    if sample_pages is None:
        # 自动选择：跳过前几页可能的封面，选择中间几页
        if total_slides <= 5:
            sample_pages = list(range(total_slides))
        else:
            # 选择第3、5、7、9页等中间页面
            sample_pages = [min(i, total_slides-1) for i in [2, 4, 6, 8]]
            sample_pages = list(set(sample_pages))  # 去重
    else:
        # 转换为0-based索引
        sample_pages = [p - 1 for p in sample_pages]
    
    print(f"采样页面: {[p+1 for p in sample_pages]}")
    print()
    
    # 收集所有颜色和字体
    all_colors = []
    all_fonts = []
    all_layouts = []
    
    for page_idx in sample_pages:
        if page_idx >= total_slides:
            continue
        
        slide = prs.slides[page_idx]
        print(f"--- 第 {page_idx+1} 页 ---")
        
        # 颜色分析
        colors = analyze_colors(slide)
        all_colors.extend(colors)
        
        # 字体分析
        fonts = analyze_fonts(slide)
        all_fonts.extend(fonts)
        
        # 布局分析
        layout = analyze_layout(slide, prs.slide_width, prs.slide_height)
        all_layouts.append(layout)
        
        print(f"  形状数量: {layout['shape_count']}")
        print(f"  有标题: {layout['has_title']}")
        print(f"  有图片: {layout['has_image']}")
        print(f"  有装饰元素: {layout['has_decoration']}")
        print(f"  列数（约）: {layout['columns']}")
        print()
    
    # 统计颜色频率
    print("=" * 50)
    print("【配色方案分析】")
    print()
    
    fill_colors = [c[1] for c in all_colors if c[0] == 'fill']
    text_colors = [c[1] for c in all_colors if c[0] == 'text']
    
    if fill_colors:
        fill_counter = Counter(fill_colors)
        print("填充颜色（按频率排序）:")
        for color, count in fill_counter.most_common(5):
            print(f"  #{color}: {count}次")
        print()
    
    if text_colors:
        text_counter = Counter(text_colors)
        print("文字颜色（按频率排序）:")
        for color, count in text_counter.most_common(5):
            print(f"  #{color}: {count}次")
        print()
    
    # 统计字体频率
    print("=" * 50)
    print("【字体规范分析】")
    print()
    
    if all_fonts:
        font_names = [f[0] for f in all_fonts]
        font_counter = Counter(font_names)
        print("使用的字体:")
        for font, count in font_counter.most_common():
            print(f"  {font}: {count}次")
        print()
        
        # 分析字号分布
        font_sizes = [f[1] for f in all_fonts if f[1] > 0]
        if font_sizes:
            font_sizes.sort()
            print(f"字号范围: {font_sizes[0]:.0f}pt - {font_sizes[-1]:.0f}pt")
            print(f"最常见字号:")
            size_counter = Counter(font_sizes)
            for size, count in size_counter.most_common(5):
                print(f"  {size:.0f}pt: {count}次")
            print()
    
    # 布局总结
    print("=" * 50)
    print("【布局结构分析】")
    print()
    
    if all_layouts:
        avg_shapes = sum(l['shape_count'] for l in all_layouts) / len(all_layouts)
        has_image_count = sum(1 for l in all_layouts if l['has_image'])
        has_decoration_count = sum(1 for l in all_layouts if l['has_decoration'])
        avg_columns = sum(l['columns'] for l in all_layouts) / len(all_layouts)
        
        print(f"平均形状数量: {avg_shapes:.1f}")
        print(f"包含图片的页面比例: {has_image_count}/{len(all_layouts)} ({has_image_count/len(all_layouts)*100:.0f}%)")
        print(f"包含装饰元素的页面比例: {has_decoration_count}/{len(all_layouts)} ({has_decoration_count/len(all_layouts)*100:.0f}%)")
        print(f"平均列数: {avg_columns:.1f}")
        print()
        
        # 判断布局复杂度
        if avg_shapes < 5:
            complexity = "简单（默认布局）"
        elif avg_shapes < 15:
            complexity = "中等"
        else:
            complexity = "复杂（专业设计）"
        print(f"布局复杂度: {complexity}")
    
    print()
    print("=" * 50)
    print("【风格总结】")
    print()
    
    # 提取主色调
    main_color = None
    if fill_colors:
        fill_counter = Counter(fill_colors)
        # 排除白色和黑色背景
        for color, count in fill_counter.most_common():
            if color.upper() not in ['FFFFFF', '000000', 'WHITE', 'BLACK']:
                main_color = color
                break
    
    text_color = None
    if text_colors:
        text_counter = Counter(text_colors)
        for color, count in text_counter.most_common():
            if color.upper() not in ['FFFFFF', '000000']:
                text_color = color
                break
    
    main_font = None
    if all_fonts:
        font_counter = Counter([f[0] for f in all_fonts])
        main_font = font_counter.most_common(1)[0][0]
    
    print(f"主色调: #{main_color}" if main_color else "主色调: 未识别")
    print(f"主要文字色: #{text_color}" if text_color else "主要文字色: 未识别")
    print(f"主要字体: {main_font}" if main_font else "主要字体: 未识别")
    print()
    
    return {
        'main_color': main_color,
        'text_color': text_color,
        'main_font': main_font,
        'avg_shapes': avg_shapes if all_layouts else 0,
        'has_images': has_image_count > 0,
        'has_decorations': has_decoration_count > 0,
        'avg_columns': avg_columns if all_layouts else 0,
    }


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("用法: python analyze_style.py <ppt文件路径> [页码1,页码2,...]")
        print("示例: python analyze_style.py presentation.pptx 3,5,7")
        sys.exit(1)
    
    ppt_path = sys.argv[1]
    sample_pages = None
    
    if len(sys.argv) >= 3:
        sample_pages = [int(p) for p in sys.argv[2].split(',')]
    
    result = analyze_ppt(ppt_path, sample_pages)
