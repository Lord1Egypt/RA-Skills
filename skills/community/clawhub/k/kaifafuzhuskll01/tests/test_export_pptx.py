"""Tests for PPTX export engine."""
import json
import os
import sys
import subprocess

import pytest

# Add scripts to path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "scripts"))

from export_pptx import build_pptx

# Try importing pptx (test gracefully if not installed)
try:
    from pptx import Presentation as PptxPresentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


@pytest.fixture
def sample_data():
    return {
        "title": "测试演示文稿",
        "slides": [
            {
                "page": 1,
                "layout": "cover",
                "density": "minimal",
                "title": "测试标题",
                "subtitle": "测试副标题",
                "content": [],
            },
            {
                "page": 2,
                "layout": "content",
                "density": "expanded",
                "title": "内容页",
                "content": [
                    {"type": "bullet", "text": "要点1"},
                    {"type": "bullet", "text": "要点2"},
                    {"type": "bullet", "text": "要点3"},
                ],
            },
            {
                "page": 3,
                "layout": "ending",
                "density": "minimal",
                "title": "谢谢",
                "content": [],
            },
        ],
    }


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not installed")
def test_build_minimal_pptx(sample_data, tmp_path):
    """3 slides (cover + content + ending) generate successfully."""
    output = tmp_path / "test.pptx"
    build_pptx(sample_data, str(output))
    assert output.exists()
    assert output.stat().st_size > 0

    prs = PptxPresentation(str(output))
    assert len(prs.slides) >= 3


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not installed")
def test_all_layout_types(tmp_path):
    """Every layout type produces at least one slide."""
    data = {"title": "布局测试", "slides": []}
    layouts = [
        "cover",
        "toc",
        "chapter-cover",
        "content",
        "three-cols",
        "two-cols",
        "chart-radar",
        "chart-bar",
        "ending",
    ]
    for i, layout in enumerate(layouts, 1):
        data["slides"].append(
            {
                "page": i,
                "layout": layout,
                "density": "standard",
                "title": f"测试{layout}",
                "content": [{"type": "bullet", "text": f"测试内容{i}"}],
            }
        )

    output = tmp_path / "all_layouts.pptx"
    build_pptx(data, str(output))

    prs = PptxPresentation(str(output))
    assert len(prs.slides) == len(layouts)


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not installed")
def test_empty_content(tmp_path):
    """Slide with empty content array still generates."""
    data = {
        "title": "空内容测试",
        "slides": [
            {
                "page": 1,
                "layout": "content",
                "density": "expanded",
                "title": "无内容页",
                "content": [],
            }
        ],
    }
    output = tmp_path / "empty.pptx"
    build_pptx(data, str(output))
    assert output.exists()


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not installed")
def test_unknown_layout(tmp_path):
    """Unknown layout type falls back gracefully."""
    data = {
        "title": "未知布局测试",
        "slides": [
            {
                "page": 1,
                "layout": "nonexistent",
                "density": "minimal",
                "title": "未知",
                "content": [],
            }
        ],
    }
    output = tmp_path / "unknown.pptx"
    build_pptx(data, str(output))
    assert output.exists()


def test_invalid_json_cli(tmp_path):
    """Corrupted JSON causes non-zero exit."""
    bad_json = tmp_path / "bad.json"
    bad_json.write_text("{invalid json", encoding="utf-8")
    script = os.path.join(
        os.path.dirname(__file__), "..", "scripts", "export_pptx.py"
    )
    result = subprocess.run(
        [
            sys.executable,
            script,
            "--input",
            str(bad_json),
            "--output",
            str(tmp_path / "out.pptx"),
        ],
        capture_output=True,
        text=True,
    )
    assert result.returncode != 0


def test_missing_input_cli(tmp_path):
    """Missing input file causes non-zero exit."""
    script = os.path.join(
        os.path.dirname(__file__), "..", "scripts", "export_pptx.py"
    )
    result = subprocess.run(
        [
            sys.executable,
            script,
            "--input",
            str(tmp_path / "nonexistent.json"),
            "--output",
            str(tmp_path / "out.pptx"),
        ],
        capture_output=True,
        text=True,
    )
    assert result.returncode != 0
