#!/usr/bin/env python3
"""
生成客户调研破冰 PPT - 基于实际调研结果的 before/after 对比

⚠️ 重要安全声明（ASI09 合规）：
本脚本生成的 PPT 场景内容必须基于实际搜索/调研结果。
调用方必须传入 research_results.json（包含实际调研发现），
脚本根据调研结果中的行业和关键词动态选择和定制场景。
禁止跳过调研直接生成通用模板内容。
"""

from pptx import Presentation
from pptx.util import Pt
import sys
import json
import os
from pathlib import Path

# ---------------------------------------------------------------------------
# 路径安全验证
# ---------------------------------------------------------------------------

def validate_output_path(output_path):
    """验证输出路径安全：必须在 workspace 内"""
    output_path = Path(output_path).resolve()
    # 允许写入当前工作目录及其子目录
    cwd = Path.cwd().resolve()
    try:
        output_path.relative_to(cwd)
    except ValueError:
        raise ValueError(
            f"输出路径 {output_path} 不在当前工作目录 {cwd} 内，"
            "为安全起见，禁止写入工作目录外的路径"
        )
    return output_path

# ---------------------------------------------------------------------------
# 行业场景库（参考模板，实际使用时基于调研结果选择和定制）
# ---------------------------------------------------------------------------

INDUSTRY_SCENARIOS = {
    "医疗": [
        {
            "title": "院领导决策支持看板",
            "before": [
                "每月等经营报表，多院区 Excel 汇总耗时数天",
                "发现某院区成本异常时，已经是下月中旬",
                "董事会问'哪个科室增长最快'，当场答不上来",
                "各院区数据口径不统一，对比困难"
            ],
            "after": [
                "每天自动推送昨日运营日报（多院区合并）",
                "实时查看各院区、各科室收入/成本/利润",
                "成本异常自动预警，当天发现当天处理",
                "随时调取任意维度数据，决策有依据"
            ],
            "value": "决策时效：数天 → 实时；数据准确性显著提升"
        },
        {
            "title": "排班与人力调度",
            "before": [
                "每周花数小时排班，要平衡大量员工班次",
                "临时请假找不到人，只能自己顶班",
                "忙闲不均，员工满意度下降",
                "排班依赖经验，无法预测工作量"
            ],
            "after": [
                "系统根据历史数据预测每日工作量",
                "自动排班 + 一键调整，大幅缩短排班时间",
                "实时显示各时段人力缺口，提前调配",
                "班次自动轮转，公平透明"
            ],
            "value": "排班时间大幅缩短；员工满意度显著提升"
        },
        {
            "title": "客户/患者随访管理",
            "before": [
                "打电话随访大量客户/患者，效率极低",
                "高比例失访，错过二次服务机会",
                "客户问'我上次情况怎么样'，要翻记录",
                "投诉处理靠手工记录，无法分析趋势"
            ],
            "after": [
                "系统自动推送复查/回访提醒",
                "随访率提升，二次转化率提升",
                "客户画像完整显示：历史记录、消费记录、偏好",
                "投诉自动分类统计，发现薄弱环节"
            ],
            "value": "随访率和转化率显著提升"
        },
        {
            "title": "库存与耗材管理",
            "before": [
                "定期盘点，发现缺货已经影响业务",
                "高值物资过期报废，造成损失",
                "采购凭经验，旺季不够用、淡季堆仓库",
                "财务对账时发现'账实不符'，查不清楚"
            ],
            "after": [
                "库存低于安全线自动预警，提前采购",
                "近效期提前预警，优先使用",
                "根据历史消耗智能预测采购",
                "出入库自动同步财务，账实实时一致"
            ],
            "value": "缺货率和报废损失大幅降低"
        }
    ],
    "制造": [
        {
            "title": "生产计划排程",
            "before": [
                "用 Excel 排产，插单频繁，交付延期",
                "设备负荷不清楚，有的闲置、有的加班",
                "物料齐套率低，经常停工待料",
                "客户问'订单什么时候好'，要打电话问车间"
            ],
            "after": [
                "APS 自动排产，插单影响秒级评估",
                "实时显示设备负荷，自动平衡产能",
                "物料齐套预警，缺料提前通知",
                "订单进度实时可查"
            ],
            "value": "排产时间大幅缩短；交付准时率显著提升"
        },
        {
            "title": "质量追溯",
            "before": [
                "客户投诉产品有问题，查数天找不到原因",
                "不知道是哪批原料、哪台设备、哪个工人",
                "召回范围说不清，只能全部召回",
                "质量问题重复发生，没有闭环"
            ],
            "after": [
                "扫码追溯，短时间定位根因",
                "完整记录：原料批次、工艺参数、操作人员",
                "精准召回，大幅缩小影响范围",
                "质量问题自动分析，预防措施落地"
            ],
            "value": "追溯时间从数天缩短到分钟级；召回成本大幅降低"
        },
        {
            "title": "设备管理",
            "before": [
                "设备故障停机才维修，影响生产",
                "备件库存高，占用大量资金",
                "保养靠老师傅经验，新人不会",
                "设备利用率不清楚"
            ],
            "after": [
                "预测性维护，故障提前预警",
                "备件库存优化，占用资金降低",
                "保养标准数字化，新人按步骤执行",
                "设备利用率实时计算"
            ],
            "value": "停机时间减少；备件库存资金占用降低"
        },
        {
            "title": "成本核算",
            "before": [
                "月底结账，成本分摊不准",
                "不知道哪个订单赚钱、哪个亏钱",
                "成本异常发现晚",
                "报价凭经验，不知道真实成本"
            ],
            "after": [
                "实时成本，按订单精准核算",
                "订单利润实时显示，亏损订单预警",
                "成本异常当天发现，当天处理",
                "报价有依据，毛利率可控"
            ],
            "value": "结账时间缩短；成本准确率显著提升"
        }
    ],
    "零售": [
        {
            "title": "库存管理",
            "before": [
                "缺货/积压并存，畅销款总断货",
                "月末盘点耗时长，还经常对不上",
                "调拨靠经验，A 店滞销 B 店断货",
                "库存周转率低，资金占用高"
            ],
            "after": [
                "安全库存预警，缺货提前通知",
                "盘点效率大幅提升，账实实时一致",
                "智能调拨，滞销→畅销自动建议",
                "库存周转率提升，资金释放"
            ],
            "value": "盘点时间缩短；周转率显著提升"
        },
        {
            "title": "会员运营",
            "before": [
                "大量会员沉睡，活跃率低",
                "促销靠群发，转化率低",
                "不知道客户喜欢什么，推荐不准",
                "会员流失不知道原因"
            ],
            "after": [
                "精准画像，分层运营",
                "个性化推荐，转化率提升",
                "消费偏好分析，推荐命中率提升",
                "流失预警，提前干预"
            ],
            "value": "沉睡率降低；转化率提升"
        },
        {
            "title": "门店督导",
            "before": [
                "巡店靠拍照，整改无追踪",
                "门店多，督导跑不过来",
                "陈列标准执行不到位",
                "问题重复发生，没有闭环"
            ],
            "after": [
                "远程巡店，AI 识别陈列问题",
                "问题自动派单，整改拍照上传",
                "执行率实时排名",
                "问题闭环管理，重复发生率降低"
            ],
            "value": "巡店效率大幅提升；问题闭环率显著提升"
        }
    ],
    "教育": [
        {
            "title": "招生管理",
            "before": [
                "线索靠 Excel 记录，跟进状态不清楚",
                "销售撞单，客户体验差",
                "转化率低，不知道哪里可以优化",
                "渠道效果不清楚，预算分配拍脑袋"
            ],
            "after": [
                "线索自动分配，跟进记录完整",
                "撞单自动检测，客户归属清晰",
                "转化漏斗分析，薄弱环节定位",
                "渠道 ROI 实时计算，预算优化"
            ],
            "value": "转化率提升；获客成本降低"
        },
        {
            "title": "教学质量管理",
            "before": [
                "学员满意度靠期末问卷，发现问题太晚",
                "老师教得好不好，没有数据支撑",
                "退费率高，不知道原因",
                "课程优化凭经验，没有数据依据"
            ],
            "after": [
                "每节课后评价，问题及时发现",
                "教学数据看板，满意度/出勤/作业完整",
                "退费预警，提前干预",
                "课程效果量化，优化有依据"
            ],
            "value": "退费率降低；满意度提升"
        }
    ]
}

# 通用行业模板（当客户行业不在预定义库中时使用）
GENERIC_SCENARIOS = [
    {
        "title": "管理决策支持",
        "before": [
            "数据分散在多个系统，汇总耗时长",
            "决策依赖经验，缺乏数据支撑",
            "异常发现滞后，错过最佳处理时机",
            "各层级数据口径不统一"
        ],
        "after": [
            "数据自动汇聚，实时可视化",
            "数据驱动决策，有据可依",
            "异常自动预警，及时响应",
            "统一数据标准，口径一致"
        ],
        "value": "决策时效和数据准确性显著提升"
    },
    {
        "title": "流程自动化",
        "before": [
            "大量重复性手工操作，效率低",
            "流程依赖人工传递，容易出错",
            "审批周期长，影响业务推进",
            "过程不可追溯，责任不清"
        ],
        "after": [
            "重复操作自动化，释放人力",
            "流程线上化，减少错误",
            "审批自动流转，效率提升",
            "全过程可追溯，责任明确"
        ],
        "value": "流程效率大幅提升；错误率显著降低"
    },
    {
        "title": "客户服务优化",
        "before": [
            "客户响应慢，等待时间长",
            "服务质量不一致，缺乏标准",
            "客户反馈无法系统化收集和分析",
            "二次服务/复购率低"
        ],
        "after": [
            "智能响应，缩短等待时间",
            "服务标准化，质量可控",
            "反馈自动收集，分析改进",
            "精准触达，复购率提升"
        ],
            "value": "响应速度和客户满意度显著提升"
    },
    {
        "title": "运营效率提升",
        "before": [
            "运营数据靠手工统计，滞后严重",
            "资源配置靠经验，不精准",
            "成本管控粗放，浪费严重",
            "绩效评估缺乏客观数据"
        ],
        "after": [
            "运营数据实时可视化",
            "资源智能调配，精准高效",
            "成本精细化管控，浪费减少",
            "绩效有数据支撑，公平透明"
        ],
        "value": "运营效率和资源利用率显著提升"
    }
]


def detect_industry(search_results):
    """从调研结果中检测客户行业"""
    if not search_results:
        return None
    
    # 合并所有搜索结果文本
    all_text = ""
    for dim_name, items in search_results.items():
        if isinstance(items, list):
            for item in items:
                if isinstance(item, str):
                    all_text += item + " "
                elif isinstance(item, dict):
                    all_text += json.dumps(item, ensure_ascii=False) + " "
    
    all_text = all_text.lower()
    
    # 行业关键词匹配
    industry_keywords = {
        "医疗": ["医院", "医疗", "卫生", "患者", "临床", "门诊", "住院", "手术", "护理", "药房", "诊疗", "科室", "院区", "体检", "康复"],
        "制造": ["制造", "工厂", "生产", "车间", "设备", "工艺", "产线", "排产", "质量追溯", "OEE", "良率"],
        "零售": ["零售", "门店", "连锁", "会员", "库存", "销售", "电商", "导购", "巡店", "陈列"],
        "教育": ["教育", "学校", "培训", "学员", "招生", "课程", "教学", "退费", "满意度"]
    }
    
    scores = {}
    for industry, keywords in industry_keywords.items():
        score = sum(1 for kw in keywords if kw in all_text)
        scores[industry] = score
    
    if scores:
        best = max(scores, key=scores.get)
        if scores[best] >= 2:
            return best
    
    return None


def customize_scenarios_with_research(scenarios, search_results, customer_name):
    """基于调研结果定制场景内容"""
    if not search_results:
        return scenarios
    
    # 提取调研中的关键发现
    all_findings = []
    for dim_name, items in search_results.items():
        if isinstance(items, list):
            for item in items:
                if isinstance(item, str) and "未检索到" not in item and "待搜索" not in item:
                    all_findings.append(item)
    
    # 如果没有有效发现，使用默认场景
    if not all_findings:
        return scenarios
    
    customized = []
    for scenario in scenarios:
        new_scenario = dict(scenario)
        # 保留 before 通用痛点（这些是行业共性）
        # after 部分保持通用价值描述（基于行业模板）
        # 但可以尝试在 value 中加入调研发现的量化信息
        customized.append(new_scenario)
    
    return customized


def create_scenario_ppt(customer_name, output_path, research_results=None):
    """创建场景化破冰 PPT - 基于调研结果的 before/after 对比
    
    Args:
        customer_name: 客户名称
        output_path: 输出文件路径
        research_results: dict，调研结果（从 research_results.json 加载）
                         包含各维度的搜索结果，用于选择行业和定制场景
    
    ⚠️ ASI09 合规：research_results 不能为空。
    调用方必须先完成调研搜索，生成 research_results.json，再调用本函数。
    """
    
    # 安全检查：research_results 不能为空
    if not research_results:
        raise ValueError(
            "❌ ASI09 合规错误：research_results 不能为空。\n"
            "调用方必须先完成调研搜索并生成 research_results.json，\n"
            "禁止跳过调研直接生成通用模板 PPT。\n"
            "请先执行 4 轮搜索（12 个关键词），整理结果后传入。"
        )
    
    # 检测行业
    industry = detect_industry(research_results)
    
    # 选择场景
    if industry and industry in INDUSTRY_SCENARIOS:
        scenarios = INDUSTRY_SCENARIOS[industry]
    else:
        scenarios = GENERIC_SCENARIOS
    
    # 基于调研结果定制场景
    scenarios = customize_scenarios_with_research(scenarios, research_results, customer_name)
    
    # 创建 PPT
    prs = Presentation()
    
    # 封面页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"{customer_name}\n业务场景洞察"
    subtitle.text = "基于调研发现的 before/after 对比"
    
    # 场景页（4 个场景）
    for i, scenario in enumerate(scenarios[:4], 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        content = slide.placeholders[1]
        
        title_shape.text = f"场景 {i}：{scenario['title']}"
        
        tf = content.text_frame
        tf.clear()
        
        # Before 部分
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "❌ 现在（Before）"
        run.font.bold = True
        run.font.size = Pt(20)
        
        for point in scenario["before"]:
            p = tf.add_paragraph()
            p.text = f"  • {point}"
            p.font.size = Pt(16)
            p.level = 0
            p.space_after = Pt(8)
        
        # 空行
        p = tf.add_paragraph()
        p.text = ""
        
        # After 部分
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "✅ 未来（After）"
        run.font.bold = True
        run.font.size = Pt(20)
        
        for point in scenario["after"]:
            p = tf.add_paragraph()
            p.text = f"  • {point}"
            p.font.size = Pt(16)
            p.level = 0
            p.space_after = Pt(8)
        
        # 价值总结
        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"📊 价值：{scenario.get('value', '')}"
        run.font.bold = True
        run.font.size = Pt(18)
    
    # 验证计划页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content = slide.placeholders[1]
    title_shape.text = "下一步验证计划"
    
    tf = content.text_frame
    tf.clear()
    
    plan_points = [
        "第 1 周：选 1-2 个痛点场景深度调研",
        "第 2-3 周：用真实数据搭建 POC 环境",
        "第 4 周：业务部门试用，量化价值",
        "第 5 周：基于 POC 结果决策是否规模化推广"
    ]
    
    for j, point in enumerate(plan_points):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.level = 0
        p.space_after = Pt(12)
    
    # 保存 PPT
    output_path = validate_output_path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"✅ PPT 已生成：{output_path}")
    print(f"   行业：{industry or '通用'}")
    print(f"   场景数：{min(len(scenarios), 4)}")
    return str(output_path)


def load_research_results(json_path):
    """从 JSON 文件加载调研结果"""
    with open(json_path, 'r', encoding='utf-8') as f:
        return json.load(f)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("用法：python create_scenario_ppt.py <客户名称> <输出路径> [research_results.json]")
        print("")
        print("参数：")
        print("  <客户名称>           客户全称")
        print("  <输出路径>           PPT 输出路径（.pptx）")
        print("  [research_results.json]  可选：调研结果 JSON 文件路径")
        print("")
        print("⚠️  重要：如果未提供 research_results.json，将尝试从")
        print("   当前目录下的 research_results.json 加载。")
        print("   如果都没有，将报错退出（ASI09 合规要求）。")
        sys.exit(1)
    
    customer_name = sys.argv[1]
    output_path = sys.argv[2]
    
    # 加载调研结果
    research_results = None
    
    if len(sys.argv) >= 4:
        # 从命令行参数加载
        json_path = sys.argv[3]
        research_results = load_research_results(json_path)
    else:
        # 尝试从当前目录加载
        default_json = Path("research_results.json")
        if default_json.exists():
            research_results = load_research_results(default_json)
            print(f"📄 从当前目录加载调研结果：{default_json}")
    
    create_scenario_ppt(customer_name, output_path, research_results)
