"""
DeckCraft v6 — Native Chart Engine
Generates editable native PowerPoint charts using python-pptx chart API.
All charts are fully editable in PowerPoint — double-click to modify data, colors, labels.

No matplotlib dependency. No PNG embedding. 100% native elements.

v6.0.0: +chart_funnel, chart_gantt, chart_swot, chart_porter, chart_sankey,
         chart_heatmap, chart_radar, chart_treemap, chart_waterfall
"""

from pptx.util import Inches, Pt, Emu
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from .constants import c, chart_colors, THEMES
from .core import add_rect, add_textbox, add_circle, add_rounded_rect, add_line


def _get_colors(theme_name, count):
    """Get color list for charts."""
    colors = chart_colors(theme_name)
    if count <= len(colors):
        return colors[:count]
    return [colors[i % len(colors)] for i in range(count)]


def _add_chart_frame(slide, chart_frame, title=None, theme_name="business"):
    """Apply theme styling to a chart frame: title, legend, plot area."""
    theme = THEMES.get(theme_name, THEMES["business"])
    text_color = c(theme.get('text', (51, 51, 51)))
    
    chart = chart_frame.chart
    
    # Chart title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        title_run = chart.chart_title.text_frame.paragraphs[0].runs[0] if chart.chart_title.text_frame.paragraphs[0].runs else chart.chart_title.text_frame.paragraphs[0].add_run()
        title_run.font.size = Pt(12)
        title_run.font.bold = True
        title_run.font.color.rgb = text_color
    else:
        chart.has_title = False
    
    return chart_frame


def _style_series(series, color_tuple, has_data_labels=True):
    """Apply color and data labels to a chart series."""
    color = c(color_tuple)
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    
    if has_data_labels:
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.font.size = Pt(9)
        data_labels.font.bold = True
        data_labels.number_format = '0.0'
        # Position varies by chart type — set in caller


def bar_chart(slide, left, top, width, height,
              data, labels, title="", theme_name="business",
              orientation="vertical", series_names=None):
    """
    Add a native bar/column chart to a slide.
    
    slide: target slide
    left, top, width, height: chart frame position (Inches)
    data: list of lists (each inner list = one series values)
    labels: category labels
    orientation: "vertical" (column) or "horizontal" (bar)
    series_names: legend labels
    
    Returns: chart shape (fully editable)
    """
    chart_type = (XL_CHART_TYPE.BAR_CLUSTERED if orientation == "horizontal"
                  else XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    chart_data = CategoryChartData()
    chart_data.categories = labels
    
    for i, series_data in enumerate(data):
        name = series_names[i] if series_names and i < len(series_names) else f"Series {i+1}"
        chart_data.add_series(name, series_data)
    
    chart_frame = slide.shapes.add_chart(
        chart_type, int(left), int(top), int(width), int(height), chart_data
    )
    
    _add_chart_frame(slide, chart_frame, title, theme_name)
    chart = chart_frame.chart
    
    colors = _get_colors(theme_name, len(data))
    for i, series in enumerate(chart.series):
        _style_series(series, colors[i])
        if orientation == "vertical":
            series.data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END
        else:
            series.data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END
    
    # Style axes
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    
    # Legend
    if series_names and len(series_names) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
    else:
        chart.has_legend = False
    
    return chart_frame


def pie_chart(slide, left, top, width, height,
              data, labels, title="", theme_name="business",
              donut=True):
    """
    Add a native pie or donut chart to a slide.
    
    data: list of values
    labels: list of segment labels
    donut: True for donut, False for pie
    
    Returns: chart shape (fully editable)
    """
    chart_type = XL_CHART_TYPE.DOUGHNUT if donut else XL_CHART_TYPE.PIE
    
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series('Data', data)
    
    chart_frame = slide.shapes.add_chart(
        chart_type, int(left), int(top), int(width), int(height), chart_data
    )
    
    _add_chart_frame(slide, chart_frame, title, theme_name)
    chart = chart_frame.chart
    
    colors = _get_colors(theme_name, len(data))
    plot = chart.plots[0]
    
    # Color each point
    series = plot.series[0]
    for i, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = c(colors[i] if i < len(colors) else colors[i % len(colors)])
    
    # Data labels
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(9)
    data_labels.font.bold = True
    data_labels.number_format = '0%'
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = True
    data_labels.show_series_name = False
    data_labels.separator = '\n'
    
    # Legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    
    return chart_frame


def line_chart(slide, left, top, width, height,
               data, labels, title="", theme_name="business",
               series_names=None, fill_area=False):
    """
    Add a native line chart to a slide.
    
    data: list of lists (each inner list = one series values)
    labels: x-axis labels
    fill_area: True for area chart (not yet supported natively, uses line)
    
    Returns: chart shape (fully editable)
    """
    chart_type = XL_CHART_TYPE.AREA if fill_area else XL_CHART_TYPE.LINE
    
    chart_data = CategoryChartData()
    chart_data.categories = labels
    
    for i, series_data in enumerate(data):
        name = series_names[i] if series_names and i < len(series_names) else f"Series {i+1}"
        chart_data.add_series(name, series_data)
    
    chart_frame = slide.shapes.add_chart(
        chart_type, int(left), int(top), int(width), int(height), chart_data
    )
    
    _add_chart_frame(slide, chart_frame, title, theme_name)
    chart = chart_frame.chart
    
    colors = _get_colors(theme_name, len(data))
    for i, series in enumerate(chart.series):
        color = c(colors[i])
        series.format.line.color.rgb = color
        series.format.line.width = Pt(2)
        if not fill_area:
            series.smooth = False
        # Data labels
        series.has_data_labels = True
        series.data_labels.font.size = Pt(8)
        series.data_labels.number_format = '0'
        series.data_labels.label_position = XL_LABEL_POSITION.ABOVE
    
    # Style axes
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    
    # Legend
    if series_names and len(series_names) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
    else:
        chart.has_legend = False
    
    return chart_frame


def gauge_chart(slide, left, top, width, height,
                value, max_value=100, title="", theme_name="business",
                label=""):
    """
    Add a native doughnut chart styled as a gauge/speedometer.
    
    Uses a two-segment doughnut chart to simulate a gauge.
    The filled portion shows the value, the remainder is gray.
    
    Returns: chart shape (fully editable)
    """
    pct = min(value / max_value, 1.0) if max_value > 0 else 0
    
    chart_data = CategoryChartData()
    chart_data.categories = ['Value', 'Remaining']
    chart_data.add_series('Gauge', [pct * 100, (1 - pct) * 100])
    
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, int(left), int(top), int(width), int(height), chart_data
    )
    
    _add_chart_frame(slide, chart_frame, title, theme_name)
    chart = chart_frame.chart
    
    colors = _get_colors(theme_name, 1)
    plot = chart.plots[0]
    series = plot.series[0]
    
    # Value segment — colored
    point_value = series.points[0]
    point_value.format.fill.solid()
    point_value.format.fill.fore_color.rgb = c(colors[0])
    
    # Remaining segment — light gray
    point_remaining = series.points[1]
    point_remaining.format.fill.solid()
    point_remaining.format.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    
    # No legend for gauge
    chart.has_legend = False
    
    # Hide data labels, show value as text annotation via chart title
    plot.has_data_labels = False
    if not title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = f"{value}"
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(28)
            run.font.bold = True
    
    return chart_frame


# ═══════════════════════════════════════════════════════════════
#  v6.0.0 — 9 New Chart Types (all native python-pptx shapes)
# ═══════════════════════════════════════════════════════════════


def _add_title_textbox(slide, left, top, width, title, theme_name="business"):
    """Add a chart title textbox."""
    if not title:
        return
    theme = THEMES.get(theme_name, THEMES["business"])
    add_textbox(slide, left, top, width, Inches(0.35),
               title, size=Pt(13), bold=True,
               color=theme.get('primary', (30, 60, 114)),
               align=PP_ALIGN.LEFT)


def chart_funnel(slide, title, stages, values, **kwargs):
    """Horizontal funnel chart with trapezoid-style bars.
    
    Uses progressively narrowing horizontal bars to simulate a funnel.
    Each bar is a native rectangle shape (editable in PPT).
    
    Args:
        slide: target slide
        title: chart title
        stages: list of stage names (e.g. ["访客","注册","试用","付费","续费"])
        values: list of numeric values (same length as stages)
    """
    if not stages or not values:
        raise ValueError("stages and values must be non-empty")
    if len(stages) != len(values):
        raise ValueError(f"stages ({len(stages)}) and values ({len(values)}) must have same length")
    
    theme_name = kwargs.get("theme_name", "business")
    colors = chart_colors(theme_name)
    theme = THEMES.get(theme_name, THEMES["business"])
    text_color = theme.get('text_light', (255, 255, 255))
    
    n = len(stages)
    max_val = max(values) if values else 1
    
    chart_left = kwargs.get("left", Inches(0.6))
    chart_top = kwargs.get("top", Inches(1.5))
    chart_width = kwargs.get("width", Inches(8.8))
    chart_height = kwargs.get("height", Inches(3.5))
    
    _add_title_textbox(slide, chart_left, chart_top - Inches(0.4), chart_width, title, theme_name)
    
    bar_height = min(Inches(0.45), int(chart_height / (n * 1.4)))
    gap = int(bar_height * 0.25)
    y_start = chart_top
    
    for i in range(n):
        # Width proportional to value
        bar_w = int(chart_width * 0.5 * (values[i] / max_val) + chart_width * 0.45)
        x = int(chart_left + (chart_width - bar_w) / 2)  # center
        y = y_start + (bar_height + gap) * i
        
        color = colors[i % len(colors)]
        add_rounded_rect(slide, x, y, bar_w, bar_height, color)
        
        # Stage label + value (centered on bar)
        label_text = f"{stages[i]}  {values[i]:,}"
        add_textbox(slide, x + Inches(0.15), y + bar_height * 0.1, bar_w - Inches(0.3), bar_height * 0.8,
                   label_text, size=Pt(11), bold=True, color=text_color,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def chart_gantt(slide, title, tasks, **kwargs):
    """Gantt chart with horizontal task bars on a time axis.
    
    Args:
        slide: target slide
        title: chart title
        tasks: list of dicts with keys: name, start ("YYYY-MM-DD"), end ("YYYY-MM-DD")
    """
    if not tasks:
        raise ValueError("tasks must be non-empty")
    
    theme_name = kwargs.get("theme_name", "business")
    colors = chart_colors(theme_name)
    theme = THEMES.get(theme_name, THEMES["business"])
    text_color = theme.get('text', (51, 51, 51))
    
    chart_left = kwargs.get("left", Inches(0.6))
    chart_top = kwargs.get("top", Inches(1.5))
    chart_width = kwargs.get("width", Inches(8.8))
    chart_height = kwargs.get("height", Inches(3.5))
    
    _add_title_textbox(slide, chart_left, chart_top - Inches(0.4), chart_width, title, theme_name)
    
    from datetime import datetime
    
    # Parse dates and find range
    parsed = []
    for t in tasks:
        s = datetime.strptime(t["start"], "%Y-%m-%d")
        e = datetime.strptime(t["end"], "%Y-%m-%d")
        parsed.append((t["name"], s, e))
    
    all_starts = [p[1] for p in parsed]
    all_ends = [p[2] for p in parsed]
    date_min = min(all_starts)
    date_max = max(all_ends)
    total_days = max((date_max - date_min).days, 1)
    
    n = len(tasks)
    label_w = int(chart_width * 0.2)
    bar_area_w = chart_width - label_w
    row_h = min(Inches(0.4), int(chart_height / (n + 1)))
    gap = int(row_h * 0.2)
    
    # Draw rows
    for i, (name, start, end) in enumerate(parsed):
        y = chart_top + (row_h + gap) * i
        
        # Task name label
        add_textbox(slide, chart_left, y, label_w, row_h,
                   name, size=Pt(9), color=text_color,
                   align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        
        # Task bar
        offset_days = (start - date_min).days
        duration_days = max((end - start).days, 1)
        bar_x = chart_left + label_w + int(bar_area_w * offset_days / total_days)
        bar_w = max(int(bar_area_w * duration_days / total_days), Inches(0.1))
        
        color = colors[i % len(colors)]
        add_rounded_rect(slide, bar_x, y + int(row_h * 0.15), bar_w, int(row_h * 0.7), color)
        
        # Date label on bar
        date_label = f"{start.strftime('%m/%d')} - {end.strftime('%m/%d')}"
        add_textbox(slide, bar_x, y + int(row_h * 0.15), bar_w, int(row_h * 0.7),
                   date_label, size=Pt(7), bold=True,
                   color=theme.get('text_light', (255, 255, 255)),
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    # Time axis line at bottom
    axis_y = chart_top + (row_h + gap) * n + Inches(0.1)
    add_line(slide, chart_left + label_w, axis_y, bar_area_w, theme.get('text_muted', (120, 120, 120)), 0.015)
    
    # Date labels on axis
    num_labels = min(5, total_days)
    for j in range(num_labels + 1):
        frac = j / max(num_labels, 1)
        lx = chart_left + label_w + int(bar_area_w * frac)
        d = date_min + __import__('datetime').timedelta(days=int(total_days * frac))
        add_textbox(slide, lx - Inches(0.4), axis_y, Inches(0.8), Inches(0.2),
                   d.strftime('%m/%d'), size=Pt(7),
                   color=theme.get('text_muted', (120, 120, 120)),
                   align=PP_ALIGN.CENTER)


def chart_swot(slide, title, strengths, weaknesses, opportunities, threats, **kwargs):
    """SWOT 2x2 matrix chart.
    
    Args:
        slide: target slide
        title: chart title
        strengths: list of strings for S quadrant
        weaknesses: list of strings for W quadrant
        opportunities: list of strings for O quadrant
        threats: list of strings for T quadrant
    """
    if not strengths and not weaknesses and not opportunities and not threats:
        raise ValueError("At least one SWOT quadrant must be non-empty")
    
    theme_name = kwargs.get("theme_name", "business")
    colors = chart_colors(theme_name)
    theme = THEMES.get(theme_name, THEMES["business"])
    text_light = theme.get('text_light', (255, 255, 255))
    
    chart_left = kwargs.get("left", Inches(0.6))
    chart_top = kwargs.get("top", Inches(1.5))
    chart_width = kwargs.get("width", Inches(8.8))
    chart_height = kwargs.get("height", Inches(3.5))
    
    _add_title_textbox(slide, chart_left, chart_top - Inches(0.4), chart_width, title, theme_name)
    
    gap = Inches(0.12)
    cell_w = (chart_width - gap) // 2
    cell_h = (chart_height - gap) // 2
    
    quadrants = [
        ("S — 优势", strengths, colors[0]),
        ("W — 劣势", weaknesses, colors[3] if len(colors) > 3 else colors[1]),
        ("O — 机会", colors[2] if len(colors) > 2 else colors[0], colors[2] if len(colors) > 2 else colors[0]),
        ("T — 威胁", threats, colors[4] if len(colors) > 4 else colors[1]),
    ]
    # Fix O quadrant
    quadrants[2] = ("O — 机会", opportunities, colors[2] if len(colors) > 2 else colors[0])
    
    positions = [
        (chart_left, chart_top),  # S (top-left)
        (chart_left + cell_w + gap, chart_top),  # W (top-right)
        (chart_left, chart_top + cell_h + gap),  # O (bottom-left)
        (chart_left + cell_w + gap, chart_top + cell_h + gap),  # T (bottom-right)
    ]
    
    for i, ((qx, qy), (label, items, color)) in enumerate(zip(positions, quadrants)):
        add_rounded_rect(slide, qx, qy, cell_w, cell_h, color, alpha=0.88)
        
        # Quadrant label
        add_textbox(slide, qx + Inches(0.15), qy + Inches(0.08), cell_w - Inches(0.3), Inches(0.28),
                   label, size=Pt(12), bold=True, color=text_light)
        
        # Items
        if items:
            item_text = "\n".join(f"• {item}" for item in items)
            txBox = slide.shapes.add_textbox(int(qx + Inches(0.15)), int(qy + Inches(0.4)),
                                             int(cell_w - Inches(0.3)), int(cell_h - Inches(0.5)))
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, item in enumerate(items):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = f"• {item}"
                run.font.size = Pt(9)
                run.font.color.rgb = c(text_light)
                p.space_after = Pt(3)


def chart_porter(slide, title, forces, **kwargs):
    """Porter's Five Forces diagram.
    
    Args:
        slide: target slide
        title: chart title
        forces: dict with keys: new_entrants, suppliers, buyers, substitutes, rivalry
    """
    if not forces:
        raise ValueError("forces must be non-empty")
    
    required_keys = ["new_entrants", "suppliers", "buyers", "substitutes", "rivalry"]
    for k in required_keys:
        if k not in forces:
            raise ValueError(f"forces must contain key '{k}'")
    
    theme_name = kwargs.get("theme_name", "business")
    colors = chart_colors(theme_name)
    theme = THEMES.get(theme_name, THEMES["business"])
    text_light = theme.get('text_light', (255, 255, 255))
    
    chart_left = kwargs.get("left", Inches(0.6))
    chart_top = kwargs.get("top", Inches(1.5))
    chart_width = kwargs.get("width", Inches(8.8))
    chart_height = kwargs.get("height", Inches(3.5))
    
    _add_title_textbox(slide, chart_left, chart_top - Inches(0.4), chart_width, title, theme_name)
    
    cx = chart_left + chart_width // 2  # center x
    cy = chart_top + chart_height // 2  # center y
    
    center_r = min(chart_width, chart_height) // 5
    outer_r = min(chart_width, chart_height) // 7
    
    # Position mapping: center + 4 surrounding forces
    positions = {
        "rivalry": (cx, cy),
        "new_entrants": (cx, chart_top + Inches(0.3)),       # top
        "buyers": (chart_left + chart_width - Inches(0.8), cy),  # right
        "substitutes": (cx, chart_top + chart_height - Inches(0.3)),  # bottom
        "suppliers": (chart_left + Inches(0.8), cy),         # left
    }
    
    labels = {
        "rivalry": "行业竞争",
        "new_entrants": "新进入者威胁",
        "suppliers": "供应商议价",
        "buyers": "买方议价",
        "substitutes": "替代品威胁",
    }
    
    for i, key in enumerate(required_keys):
        px, py = positions[key]
        label = labels[key]
        desc = forces[key]
        is_center = (key == "rivalry")
        r = center_r if is_center else outer_r
        color = colors[i % len(colors)]
        
        # Circle
        add_circle(slide, px, py, r, color)
        
        # Label inside circle
        add_textbox(slide, px - r, py - Inches(0.12), r * 2, Inches(0.25),
                   label, size=Pt(8) if not is_center else Pt(10), bold=True,
                   color=text_light, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        
        # Description below circle (skip for center)
        if not is_center and desc:
            desc_trunc = desc[:40] + "..." if len(desc) > 40 else desc
            add_textbox(slide, px - Inches(1.2), py + r + Inches(0.05), Inches(2.4), Inches(0.35),
                       desc_trunc, size=Pt(7), color=theme.get('text', (51, 51, 51)),
                       align=PP_ALIGN.CENTER)
        elif is_center and desc:
            desc_trunc = desc[:40] + "..." if len(desc) > 40 else desc
            add_textbox(slide, px - Inches(1.0), py + center_r + Inches(0.05), Inches(2.0), Inches(0.3),
                       desc_trunc, size=Pt(8), color=theme.get('text_muted', (120, 120, 120)),
                       align=PP_ALIGN.CENTER)
        
        # Draw connector line to center (skip for center itself)
        if not is_center:
            _draw_dashed_line(slide, px, py, cx, cy, theme.get('text_muted', (180, 180, 180)))


def _draw_dashed_line(slide, x1, y1, x2, y2, color_tuple):
    """Draw a simple thin rectangle line between two points (horizontal/vertical approximation)."""
    # Use a thin connector-style rectangle
    from pptx.enum.shapes import MSO_SHAPE
    min_x = min(x1, x2)
    min_y = min(y1, y2)
    # Approximate with a very thin rect along the midpoint
    # For simplicity, just draw a line shape
    import math
    dx = x2 - x1
    dy = y2 - y1
    length = int(math.sqrt(dx*dx + dy*dy))
    if length < Inches(0.1):
        return
    shape = slide.shapes.add_connector(
        1,  # straight connector
        int(x1), int(y1), int(x2), int(y2)
    )
    shape.line.color.rgb = c(color_tuple)
    shape.line.width = Pt(0.75)
    # Make it dashed
    ln = shape.line._ln
    prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(prstDash)


def chart_sankey(slide, title, nodes, links, **kwargs):
    """Simplified Sankey/flow diagram with rectangular nodes and trapezoid flow bands.
    
    Args:
        slide: target slide
        title: chart title
        nodes: list of node names (e.g. ["A","B","C","D"])
        links: list of dicts with keys: source, target, value
               (e.g. [{"source":"A","target":"C","value":10}])
    """
    if not nodes or not links:
        raise ValueError("nodes and links must be non-empty")
    
    theme_name = kwargs.get("theme_name", "business")
    colors = chart_colors(theme_name)
    theme = THEMES.get(theme_name, THEMES["business"])
    text_light = theme.get('text_light', (255, 255, 255))
    
    chart_left = kwargs.get("left", Inches(0.6))
    chart_top = kwargs.get("top", Inches(1.5))
    chart_width = kwargs.get("width", Inches(8.8))
    chart_height = kwargs.get("height", Inches(3.5))
    
    _add_title_textbox(slide, chart_left, chart_top - Inches(0.4), chart_width, title, theme_name)
    
    max_val = max(link["value"] for link in links) if links else 1
    
    # Split nodes into left half and right half
    half = (len(nodes) + 1) // 2
    left_nodes = nodes[:half]
    right_nodes = nodes[half:]
    
    node_w = Inches(0.6)
    node_h = Inches(0.35)
    left_x = chart_left + Inches(0.5)
    right_x = chart_left + chart_width - Inches(1.1)
    
    def _node_y(idx, total):
        spacing = chart_height / (total + 1)
        return chart_top + spacing * (idx + 1) - node_h // 2
    
    # Draw left nodes
    for i, name in enumerate(left_nodes):
        y = _node_y(i, len(left_nodes))
        color = colors[i % len(colors)]
        add_rounded_rect(slide, left_x, y, node_w, node_h, color)
        add_textbox(slide, left_x, y + node_h // 2 - Inches(0.1), node_w, Inches(0.2),
                   name, size=Pt(9), bold=True, color=text_light,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    # Draw right nodes
    for i, name in enumerate(right_nodes):
        y = _node_y(i, len(right_nodes))
        color = colors[(half + i) % len(colors)]
        add_rounded_rect(slide, right_x, y, node_w, node_h, color)
        add_textbox(slide, right_x, y + node_h // 2 - Inches(0.1), node_w, Inches(0.2),
                   name, size=Pt(9), bold=True, color=text_light,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    # Draw flow bands (trapezoid approximated as semi-transparent rectangles)
    for link in links:
        src_name = link["source"]
        tgt_name = link["target"]
        value = link["value"]
        
        # Find source and target positions
        src_idx = left_nodes.index(src_name) if src_name in left_nodes else None
        tgt_idx = right_nodes.index(tgt_name) if tgt_name in right_nodes else None
        
        if src_idx is None:
            # Source might be on right side — try all nodes
            if src_name in nodes:
                src_idx = nodes.index(src_name)
                if src_idx >= half:
                    continue  # can't draw right-to-left easily
            else:
                continue
        if tgt_idx is None:
            if tgt_name in nodes:
                tgt_idx = nodes.index(tgt_name) - half
                if tgt_idx < 0:
                    continue
            else:
                continue
        
        src_y = _node_y(src_idx, len(left_nodes)) + node_h // 2
        tgt_y = _node_y(tgt_idx, len(right_nodes)) + node_h // 2
        
        # Flow band thickness proportional to value
        thickness = max(Inches(0.05), int(Inches(0.4) * value / max_val))
        
        # Draw as semi-transparent rectangle
        band_color = colors[src_idx % len(colors)]
        flow_x = left_x + node_w
        flow_w = right_x - flow_x
        flow_y = min(src_y, tgt_y) - thickness // 2
        flow_h = abs(tgt_y - src_y) + thickness
        
        add_rect(slide, flow_x, flow_y, flow_w, thickness if abs(tgt_y - src_y) < Inches(0.05) else flow_h,
                 band_color, alpha=0.3)
        
        # Value label in middle of flow
        mid_x = flow_x + flow_w // 2 - Inches(0.3)
        mid_y = (src_y + tgt_y) // 2 - Inches(0.08)
        add_textbox(slide, mid_x, mid_y, Inches(0.6), Inches(0.2),
                   str(value), size=Pt(8), bold=True,
                   color=theme.get('text', (51, 51, 51)),
                   align=PP_ALIGN.CENTER)


def chart_heatmap(slide, title, rows, cols, values, color_scale=None, **kwargs):
    """Heatmap grid with color-coded cells.
    
    Args:
        slide: target slide
        title: chart title
        rows: list of row labels
        cols: list of column labels
        values: 2D list of numeric values [row][col]
        color_scale: optional list of (R,G,B) tuples for gradient (default: green→yellow→red)
    """
    if not rows or not cols or not values:
        raise ValueError("rows, cols, and values must be non-empty")
    if len(values) != len(rows):
        raise ValueError(f"values rows ({len(values)}) must match rows ({len(rows)})")
    
    theme_name = kwargs.get("theme_name", "business")
    theme = THEMES.get(theme_name, THEMES["business"])
    text_color = theme.get('text', (51, 51, 51))
    text_light = theme.get('text_light', (255, 255, 255))
    
    chart_left = kwargs.get("left", Inches(0.6))
    chart_top = kwargs.get("top", Inches(1.5))
    chart_width = kwargs.get("width", Inches(8.8))
    chart_height = kwargs.get("height", Inches(3.5))
    
    _add_title_textbox(slide, chart_left, chart_top - Inches(0.4), chart_width, title, theme_name)
    
    # Default color scale: green → yellow → red
    if color_scale is None:
        color_scale = [(46, 139, 87), (255, 193, 7), (178, 34, 34)]
    
    num_rows = len(rows)
    num_cols = len(cols)
    
    # Find value range
    flat = [v for row in values for v in row]
    v_min, v_max = min(flat), max(flat)
    v_range = v_max - v_min if v_max != v_min else 1
    
    label_col_w = Inches(0.9)
    header_h = Inches(0.3)
    grid_left = chart_left + label_col_w
    grid_top = chart_top + header_h
    grid_w = chart_width - label_col_w
    grid_h = chart_height - header_h
    
    cell_w = grid_w // num_cols
    cell_h = grid_h // num_rows
    
    def _interpolate_color(val):
        """Interpolate color from color_scale based on value."""
        frac = (val - v_min) / v_range
        frac = max(0.0, min(1.0, frac))
        n_colors = len(color_scale)
        if n_colors == 1:
            return color_scale[0]
        segment = frac * (n_colors - 1)
        idx = int(segment)
        if idx >= n_colors - 1:
            return color_scale[-1]
        t = segment - idx
        c1 = color_scale[idx]
        c2 = color_scale[idx + 1]
        return tuple(int(c1[j] + (c2[j] - c1[j]) * t) for j in range(3))
    
    # Column headers
    for j, col_name in enumerate(cols):
        x = grid_left + cell_w * j
        add_textbox(slide, x, chart_top, cell_w, header_h,
                   col_name, size=Pt(9), bold=True, color=text_color,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    # Grid cells
    for i, row_name in enumerate(rows):
        y = grid_top + cell_h * i
        
        # Row label
        add_textbox(slide, chart_left, y, label_col_w, cell_h,
                   row_name, size=Pt(9), bold=True, color=text_color,
                   align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        
        for j in range(num_cols):
            x = grid_left + cell_w * j
            val = values[i][j]
            cell_color = _interpolate_color(val)
            
            # Cell rectangle
            add_rect(slide, x + Inches(0.02), y + Inches(0.02),
                     cell_w - Inches(0.04), cell_h - Inches(0.04), cell_color)
            
            # Value label
            add_textbox(slide, x, y, cell_w, cell_h,
                       str(val), size=Pt(9), bold=True,
                       color=text_light,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def chart_radar(slide, title, axes, series_data, series_names=None, **kwargs):
    """Radar/spider chart using native shapes.
    
    Simplified: draws axis lines from center + polygon outlines for data series.
    Uses native shape drawing (not embedded images).
    
    Args:
        slide: target slide
        title: chart title
        axes: list of axis labels (e.g. ["性能","价格","易用","支持","生态"])
        series_data: list of lists, each inner list has values 0-10 for each axis
        series_names: optional list of series names for legend
    """
    if not axes or not series_data:
        raise ValueError("axes and series_data must be non-empty")
    n = len(axes)
    for sd in series_data:
        if len(sd) != n:
            raise ValueError(f"Each series must have {n} values (got {len(sd)})")
    
    theme_name = kwargs.get("theme_name", "business")
    colors = chart_colors(theme_name)
    theme = THEMES.get(theme_name, THEMES["business"])
    text_color = theme.get('text', (51, 51, 51))
    
    chart_left = kwargs.get("left", Inches(0.6))
    chart_top = kwargs.get("top", Inches(1.5))
    chart_width = kwargs.get("width", Inches(8.8))
    chart_height = kwargs.get("height", Inches(3.5))
    
    _add_title_textbox(slide, chart_left, chart_top - Inches(0.4), chart_width, title, theme_name)
    
    import math
    
    cx = chart_left + chart_width // 2
    cy = chart_top + chart_height // 2
    max_r = min(chart_width, chart_height) // 2 - Inches(0.4)
    
    # Draw axis lines + labels
    angle_step = 2 * math.pi / n
    axis_endpoints = []
    for i in range(n):
        angle = -math.pi / 2 + angle_step * i  # start from top
        ex = cx + int(max_r * math.cos(angle))
        ey = cy + int(max_r * math.sin(angle))
        axis_endpoints.append((ex, ey, angle))
        
        # Axis line
        _draw_dashed_line(slide, cx, cy, ex, ey, theme.get('text_muted', (180, 180, 180)))
        
        # Axis label
        lx = cx + int((max_r + Inches(0.25)) * math.cos(angle))
        ly = cy + int((max_r + Inches(0.25)) * math.sin(angle))
        add_textbox(slide, lx - Inches(0.4), ly - Inches(0.1), Inches(0.8), Inches(0.2),
                   axes[i], size=Pt(8), color=text_color,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    # Draw reference polygon (pentagon outline at 50%)
    ref_points = []
    for i in range(n):
        angle = -math.pi / 2 + angle_step * i
        rx = cx + int(max_r * 0.5 * math.cos(angle))
        ry = cy + int(max_r * 0.5 * math.sin(angle))
        ref_points.append((rx, ry))
    for i in range(n):
        x1, y1 = ref_points[i]
        x2, y2 = ref_points[(i + 1) % n]
        _draw_dashed_line(slide, x1, y1, x2, y2, (200, 200, 200))
    
    # Draw data series polygons
    for si, series in enumerate(series_data):
        color = colors[si % len(colors)]
        points = []
        for i in range(n):
            angle = -math.pi / 2 + angle_step * i
            val = max(0, min(10, series[i])) / 10.0  # normalize to 0-1
            px = cx + int(max_r * val * math.cos(angle))
            py = cy + int(max_r * val * math.sin(angle))
            points.append((px, py))
        
        # Draw polygon edges as lines
        for i in range(n):
            x1, y1 = points[i]
            x2, y2 = points[(i + 1) % n]
            connector = slide.shapes.add_connector(1, int(x1), int(y1), int(x2), int(y2))
            connector.line.color.rgb = c(color)
            connector.line.width = Pt(2)
        
        # Draw data points as small circles
        for px, py in points:
            add_circle(slide, px, py, Inches(0.06), color)
    
    # Legend
    if series_names:
        legend_y = chart_top + chart_height - Inches(0.3)
        legend_x = chart_left
        for si, name in enumerate(series_names):
            color = colors[si % len(colors)]
            add_circle(slide, legend_x + Inches(0.1), legend_y + Inches(0.08), Inches(0.06), color)
            add_textbox(slide, legend_x + Inches(0.2), legend_y, Inches(1.0), Inches(0.2),
                       name, size=Pt(8), color=text_color)
            legend_x += Inches(1.5)


def chart_treemap(slide, title, items, **kwargs):
    """Treemap chart using squarified layout (simplified single-level).
    
    Args:
        slide: target slide
        title: chart title
        items: list of (label, value) tuples (e.g. [("产品A", 50), ("产品B", 30)])
    """
    if not items:
        raise ValueError("items must be non-empty")
    
    theme_name = kwargs.get("theme_name", "business")
    colors = chart_colors(theme_name)
    theme = THEMES.get(theme_name, THEMES["business"])
    text_light = theme.get('text_light', (255, 255, 255))
    
    chart_left = kwargs.get("left", Inches(0.6))
    chart_top = kwargs.get("top", Inches(1.5))
    chart_width = kwargs.get("width", Inches(8.8))
    chart_height = kwargs.get("height", Inches(3.5))
    
    _add_title_textbox(slide, chart_left, chart_top - Inches(0.4), chart_width, title, theme_name)
    
    # Simplified squarified treemap: alternate horizontal/vertical splits
    total = sum(v for _, v in items)
    if total <= 0:
        raise ValueError("total value must be > 0")
    
    gap = Inches(0.04)
    
    def _layout(items_rect, rect_x, rect_y, rect_w, rect_h, depth=0):
        """Recursive squarified layout."""
        if not items_rect:
            return
        if len(items_rect) == 1:
            name, val = items_rect[0]
            color = colors[depth % len(colors)]
            add_rounded_rect(slide, rect_x, rect_y, rect_w, rect_h, color)
            
            # Label
            pct = val / total * 100
            label = f"{name}\n{val} ({pct:.0f}%)"
            # Determine font size based on rect size
            min_dim = min(rect_w, rect_h)
            fs = max(7, min(14, int(min_dim / Inches(0.15))))
            add_textbox(slide, rect_x + gap, rect_y + gap,
                       rect_w - gap * 2, rect_h - gap * 2,
                       label, size=Pt(fs), bold=True, color=text_light,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            return
        
        # Split: horizontal if wider, vertical if taller
        sub_total = sum(v for _, v in items_rect)
        if rect_w >= rect_h:
            # Horizontal split
            x = rect_x
            for i, (name, val) in enumerate(items_rect):
                w = int(rect_w * val / sub_total) if i < len(items_rect) - 1 else (rect_x + rect_w - x)
                if w > Inches(0.05):
                    _layout([(name, val)], x, rect_y, w, rect_h, depth + i)
                x += w
        else:
            # Vertical split
            y = rect_y
            for i, (name, val) in enumerate(items_rect):
                h = int(rect_h * val / sub_total) if i < len(items_rect) - 1 else (rect_y + rect_h - y)
                if h > Inches(0.05):
                    _layout([(name, val)], rect_x, y, rect_w, h, depth + i)
                y += h
    
    # Sort items by value descending for better layout
    sorted_items = sorted(items, key=lambda x: x[1], reverse=True)
    _layout(sorted_items, chart_left, chart_top, chart_width, chart_height)


def chart_waterfall(slide, title, categories, values, is_total=None, **kwargs):
    """Waterfall chart showing cumulative increases and decreases.
    
    Uses native bar chart as base with custom positioning for waterfall effect.
    Simplified: draws rectangles directly (not via chart API) for precise control.
    
    Args:
        slide: target slide
        title: chart title
        categories: list of category names
        values: list of numeric values (positive=up, negative=down)
        is_total: list of booleans (True=total bar starts from 0)
    """
    if not categories or not values:
        raise ValueError("categories and values must be non-empty")
    if len(categories) != len(values):
        raise ValueError(f"categories ({len(categories)}) and values ({len(values)}) must have same length")
    
    if is_total is None:
        is_total = [False] * len(values)
        if values:
            is_total[0] = True
            is_total[-1] = True
    
    theme_name = kwargs.get("theme_name", "business")
    colors = chart_colors(theme_name)
    theme = THEMES.get(theme_name, THEMES["business"])
    text_color = theme.get('text', (51, 51, 51))
    
    chart_left = kwargs.get("left", Inches(0.6))
    chart_top = kwargs.get("top", Inches(1.5))
    chart_width = kwargs.get("width", Inches(8.8))
    chart_height = kwargs.get("height", Inches(3.5))
    
    _add_title_textbox(slide, chart_left, chart_top - Inches(0.4), chart_width, title, theme_name)
    
    n = len(categories)
    
    # Calculate running totals
    running = []
    cumulative = 0
    for i in range(n):
        if is_total[i]:
            running.append(0)  # total starts from 0
            cumulative = values[i]
        else:
            running.append(cumulative)
            cumulative += values[i]
    
    # Scale to chart area
    all_vals = []
    c = 0
    for i in range(n):
        if is_total[i]:
            all_vals.append(abs(values[i]))
            c = values[i]
        else:
            all_vals.append(abs(values[i]))
            c += values[i]
    max_val = max(abs(v) for v in values) if values else 1
    # Use cumulative max for scale
    cum_vals = []
    cv = 0
    for i in range(n):
        cv += values[i] if not is_total[i] else 0
        cum_vals.append(cv)
    scale_max = max(max(abs(v) for v in cum_vals), max(abs(v) for v in values), 1) * 1.15
    
    bar_area_left = chart_left + Inches(0.3)
    bar_area_width = chart_width - Inches(0.6)
    bar_w = bar_area_width // n - Inches(0.08)
    baseline_y = chart_top + chart_height * 0.7  # 70% from top as baseline
    scale_height = chart_height * 0.6  # available height for bars
    
    total_color = colors[0] if len(colors) > 0 else (30, 60, 114)
    up_color = colors[1] if len(colors) > 1 else (46, 139, 87)
    down_color = colors[3] if len(colors) > 3 else (178, 34, 34)
    
    prev_top_y = baseline_y  # for connector line
    
    for i in range(n):
        x = bar_area_left + (bar_area_width // n) * i + Inches(0.04)
        
        if is_total[i]:
            bar_val = abs(values[i])
            bar_h = int(scale_height * bar_val / scale_max) if scale_max > 0 else 0
            bar_y = baseline_y - bar_h
            color = total_color
        elif values[i] >= 0:
            bar_h = int(scale_height * values[i] / scale_max) if scale_max > 0 else 0
            bar_y = baseline_y - int(scale_height * running[i] / scale_max) - bar_h
            color = up_color
        else:
            bar_h = int(scale_height * abs(values[i]) / scale_max) if scale_max > 0 else 0
            bar_y = baseline_y - int(scale_height * running[i] / scale_max)
            color = down_color
        
        if bar_h < Inches(0.02):
            bar_h = Inches(0.02)
        
        add_rounded_rect(slide, x, bar_y, bar_w, bar_h, color)
        
        # Value label above/below bar
        val_text = f"{values[i]:+d}" if not is_total[i] else str(values[i])
        label_y = bar_y - Inches(0.18) if values[i] >= 0 or is_total[i] else bar_y + bar_h
        add_textbox(slide, x, label_y, bar_w, Inches(0.18),
                   val_text, size=Pt(8), bold=True, color=color,
                   align=PP_ALIGN.CENTER)
        
        # Category label
        add_textbox(slide, x - Inches(0.04), baseline_y + Inches(0.05), bar_w + Inches(0.08), Inches(0.2),
                   categories[i], size=Pt(7), color=text_color,
                   align=PP_ALIGN.CENTER)
    
    # Baseline
    add_line(slide, bar_area_left, baseline_y, bar_area_width, theme.get('text_muted', (180, 180, 180)), 0.015)
