#!/usr/bin/env python3
"""
Post-process Auto-PPT outputs on Desktop:
1) merge PDFs
2) remove/merge semantically repetitive pages (text-similarity heuristic)
3) export a cleaned PPTX to Desktop

Requires:
- pypdf (`python3 -m pip install pypdf`)
- python-pptx (already importable as `pptx` in this environment)

Usage:
  python3 postprocess_ppt_outputs.py --output-name final_deck --all-desktop
  python3 postprocess_ppt_outputs.py --output-name final_deck a.pdf b.pdf c.pdf

Outputs:
  ~/Desktop/final_deck.merged.pdf
  ~/Desktop/final_deck.cleaned.pdf
  ~/Desktop/final_deck.cleaned.pptx
"""

from __future__ import annotations

import argparse
import re
import sys
from collections import Counter
from dataclasses import dataclass
from difflib import SequenceMatcher
from io import BytesIO
from pathlib import Path
from typing import Iterable, List, Sequence

try:
    from pypdf import PdfReader, PdfWriter
except Exception:
    print("ERROR: pypdf is required. Install with: python3 -m pip install pypdf", file=sys.stderr)
    raise

from pptx import Presentation
from pptx.util import Inches, Pt

DESKTOP = Path.home() / "Desktop"


@dataclass
class PageInfo:
    source_pdf: Path
    source_index: int
    global_index: int
    text: str
    norm: str
    title: str
    bullets: List[str]


def normalize_text(text: str) -> str:
    text = text or ""
    text = text.replace("\x00", " ")
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[^\w\u4e00-\u9fff ]+", " ", text.lower())
    text = re.sub(r"\b(page|slide|figure|chart|table|source|references?)\b", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def first_nontrivial_line(text: str) -> str:
    for line in (text or "").splitlines():
        line = re.sub(r"\s+", " ", line).strip(" -•·\t")
        if len(line) >= 8:
            return line[:80]
    return "Untitled"


def make_bullets(text: str, limit: int = 4) -> List[str]:
    parts = []
    for line in text.splitlines():
        line = re.sub(r"\s+", " ", line).strip(" -•·\t")
        if len(line) >= 12:
            parts.append(line)
    if not parts:
        text2 = re.sub(r"\s+", " ", text).strip()
        parts = re.split(r"(?<=[。！？.!?])\s+", text2)
        parts = [p.strip() for p in parts if len(p.strip()) >= 12]
    deduped = []
    seen = set()
    for p in parts:
        key = normalize_text(p)
        if not key or key in seen:
            continue
        seen.add(key)
        deduped.append(p[:120])
        if len(deduped) >= limit:
            break
    return deduped or ["See cleaned PDF for original layout."]


def token_set(s: str) -> set[str]:
    return set(t for t in s.split() if len(t) >= 2)


def jaccard(a: str, b: str) -> float:
    sa, sb = token_set(a), token_set(b)
    if not sa or not sb:
        return 0.0
    return len(sa & sb) / len(sa | sb)


def similarity(a: str, b: str) -> float:
    if not a or not b:
        return 0.0
    seq = SequenceMatcher(None, a[:4000], b[:4000]).ratio()
    jac = jaccard(a, b)
    return max(seq, jac)


def collect_input_pdfs(paths: Sequence[str], all_desktop: bool) -> List[Path]:
    if all_desktop:
        pdfs = sorted(DESKTOP.glob("*.pdf"), key=lambda p: p.name.lower())
    else:
        pdfs = []
        for raw in paths:
            p = Path(raw).expanduser()
            if not p.exists():
                candidate = DESKTOP / raw
                if candidate.exists():
                    p = candidate
            if not p.exists() or p.suffix.lower() != ".pdf":
                print(f"WARN: skipped missing/non-pdf: {raw}", file=sys.stderr)
                continue
            pdfs.append(p)
    pdfs = [p for p in pdfs if p.is_file()]
    if not pdfs:
        raise SystemExit("No input PDFs found.")
    return pdfs


def load_pdf_reader(path: Path) -> PdfReader:
    # Read bytes first to avoid macOS/iCloud Desktop weirdness like:
    # OSError: [Errno 11] Resource deadlock avoided
    data = path.read_bytes()
    return PdfReader(BytesIO(data))


def merge_pdfs(inputs: Sequence[Path], out_path: Path) -> None:
    writer = PdfWriter()
    for path in inputs:
        reader = load_pdf_reader(path)
        for page in reader.pages:
            writer.add_page(page)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with out_path.open("wb") as f:
        writer.write(f)


def extract_pages(inputs: Sequence[Path]) -> List[PageInfo]:
    pages: List[PageInfo] = []
    global_idx = 0
    for pdf in inputs:
        reader = load_pdf_reader(pdf)
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            norm = normalize_text(text)
            title = first_nontrivial_line(text)
            bullets = make_bullets(text)
            pages.append(PageInfo(pdf, i, global_idx, text, norm, title, bullets))
            global_idx += 1
    return pages


def should_drop(page: PageInfo, kept: Iterable[PageInfo], threshold: float) -> bool:
    if len(page.norm) < 40:
        return False
    for prev in kept:
        if len(prev.norm) < 40:
            continue
        score = similarity(page.norm, prev.norm)
        shorter = min(len(page.norm), len(prev.norm))
        longer = max(len(page.norm), len(prev.norm))
        containment = shorter / longer if longer else 0
        if score >= threshold:
            return True
        if containment >= 0.92 and (page.norm in prev.norm or prev.norm in page.norm):
            return True
    return False


def dedupe_pages(pages: Sequence[PageInfo], threshold: float) -> List[PageInfo]:
    kept: List[PageInfo] = []
    for p in pages:
        if not should_drop(p, kept, threshold):
            kept.append(p)
    return kept


def write_clean_pdf(inputs: Sequence[Path], kept: Sequence[PageInfo], out_path: Path) -> None:
    writer = PdfWriter()
    cache = {p: load_pdf_reader(p) for p in {k.source_pdf for k in kept}}
    for info in kept:
        writer.add_page(cache[info.source_pdf].pages[info.source_index])
    with out_path.open("wb") as f:
        writer.write(f)


def choose_slide_title(info: PageInfo, idx: int) -> str:
    title = info.title.strip()
    if title.lower() == "untitled":
        title = f"Slide {idx}"
    if len(title) > 70:
        title = title[:67] + "..."
    return title


def build_pptx(kept: Sequence[PageInfo], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "Cleaned Presentation"
    cover.placeholders[1].text = f"Merged + de-duplicated from {len(kept)} pages"

    for idx, info in enumerate(kept, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = choose_slide_title(info, idx)
        tf = slide.placeholders[1].text_frame
        tf.clear()
        bullets = info.bullets[:4]
        for j, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)
        tx = slide.shapes.add_textbox(Inches(0.4), Inches(6.7), Inches(12.0), Inches(0.4))
        p = tx.text_frame.paragraphs[0]
        p.text = f"Source: {info.source_pdf.name} · page {info.source_index + 1}"
        p.font.size = Pt(9)

    prs.save(str(out_path))


def print_summary(inputs: Sequence[Path], pages: Sequence[PageInfo], kept: Sequence[PageInfo], merged: Path, cleaned: Path, pptx: Path) -> None:
    removed = len(pages) - len(kept)
    source_counts = Counter(p.source_pdf.name for p in kept)
    print("=== Auto-PPT postprocess summary ===")
    print(f"Inputs: {len(inputs)} PDFs")
    for p in inputs:
        print(f" - {p.name}")
    print(f"Pages before: {len(pages)}")
    print(f"Pages kept:   {len(kept)}")
    print(f"Pages removed:{removed}")
    print(f"Merged PDF:   {merged}")
    print(f"Clean PDF:    {cleaned}")
    print(f"PPTX:         {pptx}")
    print("Page distribution after cleaning:")
    for name, count in source_counts.items():
        print(f" - {name}: {count}")


def parse_args() -> argparse.Namespace:
    ap = argparse.ArgumentParser()
    ap.add_argument("pdfs", nargs="*", help="Input PDF files (names or paths).")
    ap.add_argument("--all-desktop", action="store_true", help="Use all PDFs on Desktop.")
    ap.add_argument("--output-name", default="auto_ppt_final", help="Base output name without extension.")
    ap.add_argument("--threshold", type=float, default=0.88, help="Similarity threshold for duplicate removal.")
    return ap.parse_args()


def main() -> int:
    args = parse_args()
    inputs = collect_input_pdfs(args.pdfs, args.all_desktop)
    base = re.sub(r"[^A-Za-z0-9._-]+", "_", args.output_name).strip("_") or "auto_ppt_final"
    merged = DESKTOP / f"{base}.merged.pdf"
    cleaned = DESKTOP / f"{base}.cleaned.pdf"
    pptx = DESKTOP / f"{base}.cleaned.pptx"

    merge_pdfs(inputs, merged)
    pages = extract_pages(inputs)
    kept = dedupe_pages(pages, args.threshold)
    write_clean_pdf(inputs, kept, cleaned)
    build_pptx(kept, pptx)
    print_summary(inputs, pages, kept, merged, cleaned, pptx)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
