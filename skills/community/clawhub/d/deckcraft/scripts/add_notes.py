#!/usr/bin/env python3
"""
DeckCraft v6 — Add Speaker Notes to PPTX

Adds or replaces speaker notes on each slide from a JSON file.

Usage:
    python3 add_notes.py deck.pptx --notes-file notes.json
    python3 add_notes.py deck.pptx --notes-file notes.json -o deck_with_notes.pptx

notes.json format:
    {
        "1": "Welcome everyone. Today we'll cover...",
        "2": "This is our agenda for today.",
        "3": "Let me walk you through the key numbers..."
    }

    Keys are slide numbers (1-based). Slides without entries are left unchanged.
"""
import sys
import os
import json
import argparse
from typing import Dict

from pptx import Presentation


def add_notes(pptx_path: str, notes: Dict[str, str],
              output_path: str = None) -> int:
    """Add speaker notes to a PPTX file.

    Args:
        pptx_path: Path to input .pptx file.
        notes: Dict mapping slide numbers (as strings) to note text.
        output_path: Output path. Default: overwrite input file.

    Returns:
        Number of slides that had notes added/updated.
    """
    if not os.path.isfile(pptx_path):
        print(f"ERROR: File not found: {pptx_path}", file=sys.stderr)
        return -1

    prs = Presentation(pptx_path)
    updated = 0

    for i, slide in enumerate(prs.slides):
        slide_num = str(i + 1)
        note_text = notes.get(slide_num, "").strip()
        if not note_text:
            continue

        # Get or create notes slide
        notes_slide = slide.notes_slide
        if notes_slide is None:
            # python-pptx auto-creates notes_slide on access
            notes_slide = slide.notes_slide

        text_frame = notes_slide.notes_text_frame
        text_frame.text = note_text
        updated += 1

    if output_path is None:
        output_path = pptx_path

    prs.save(output_path)
    total = len(prs.slides)
    print(f"✓ Notes added to {updated}/{total} slides → {output_path}")
    return updated


def main():
    parser = argparse.ArgumentParser(
        description="DeckCraft v6 — Add Speaker Notes to PPTX",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
notes.json format:
    {
        "1": "Welcome everyone. Today we'll cover...",
        "2": "This is our agenda.",
        "3": "Key numbers to highlight..."
    }
    Keys are 1-based slide numbers. Unlisted slides keep existing notes.
        """,
    )
    parser.add_argument("pptx", help="Input PPTX file")
    parser.add_argument("--notes-file", required=True,
                        help="JSON file mapping slide numbers to note text")
    parser.add_argument("-o", "--output", default=None,
                        help="Output PPTX path (default: overwrite input)")

    args = parser.parse_args()

    if not os.path.isfile(args.notes_file):
        print(f"ERROR: Notes file not found: {args.notes_file}", file=sys.stderr)
        return 2

    try:
        with open(args.notes_file, "r", encoding="utf-8") as f:
            notes = json.load(f)
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid JSON in {args.notes_file}: {e}", file=sys.stderr)
        return 2

    if not isinstance(notes, dict):
        print("ERROR: Notes file must be a JSON object {\"slide_num\": \"note text\"}",
              file=sys.stderr)
        return 2

    result = add_notes(args.pptx, notes, args.output)
    if result < 0:
        return 2
    return 0


if __name__ == "__main__":
    sys.exit(main())
