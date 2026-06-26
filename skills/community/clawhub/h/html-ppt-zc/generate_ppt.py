#!/usr/bin/env python3

import argparse, os, sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    sys.stderr.write('Missing python-pptx library. Install with: pip install -U python-pptx\n')
    sys.exit(1)

def create_ppt(title, paragraphs, output_path):
    prs = Presentation()
    # --- Title slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    # subtitle optional
    if len(paragraphs) > 0:
        try:
            slide.placeholders[1].text = paragraphs[0]
        except Exception:
            pass
    # --- Content slides ---
    bullet_slide_layout = prs.slide_layouts[1]
    for para in paragraphs[1:]:
        if not para.strip():
            continue  # skip empty lines
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        for line in para.split('\n'):
            p = tf.add_paragraph()
            p.text = line.strip()
            p.level = 0
            p.font.size = Pt(18)
    prs.save(output_path)
    print(f'PPT saved to {output_path}')

def export_html(title, paragraphs, output_path):
    # simple HTML with each paragraph as a slide div
    html_parts = ["<html><head><meta charset='utf-8'><title>{}</title></head><body>".format(title)]
    for idx, para in enumerate(paragraphs):
        html_parts.append(f"<section style='margin:2em;'><h2>Slide {idx+1}</h2><pre>{para}</pre></section>")
    html_parts.append("</body></html>")
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(html_parts))
    print(f'HTML slide deck saved to {output_path}')

def main():
    parser = argparse.ArgumentParser(description='Generate PPT or HTML slide deck from plain text')
    parser.add_argument('--title', required=True, help='Presentation title')
    parser.add_argument('--text', required=True, help='Content text. Use "\\n" for line breaks or separate paragraphs with blank lines.')
    parser.add_argument('--output', required=True, help='Output file path (.pptx or .html)')
    parser.add_argument('--html', action='store_true', help='Export as simple HTML instead of PPTX')
    args = parser.parse_args()

    # Split into paragraphs using double line breaks
    paragraphs = [p.strip() for p in args.text.split('\n\n') if p.strip()]
    if args.html:
        export_html(args.title, paragraphs, args.output)
    else:
        # ensure output dir exists
        os.makedirs(os.path.dirname(args.output), exist_ok=True)
        create_ppt(args.title, paragraphs, args.output)

if __name__ == '__main__':
    main()
