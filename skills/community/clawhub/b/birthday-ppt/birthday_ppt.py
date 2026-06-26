#!/usr/bin/env python3
"""
幼儿园生日PPT生成器
使用方法: python3 birthday_ppt.py --name 小雨 --birth-date 2020-05-20 --output birthday.pptx --photos /path/to/photo1.jpg /path/to/photo2.jpg --theme pink
"""

import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from PIL import Image
import copy

# 主题配色
THEMES = {
    "pink": {
        "primary": RGBColor(0xFF, 0x69, 0xB4),  # 粉红
        "secondary": RGBColor(0xFF, 0xB6, 0xC1),  # 浅粉
        "accent": RGBColor(0xFF, 0x14, 0x93),  # 深粉
        "text": RGBColor(0x4A, 0x00, 0x4A),  # 深紫文字
        "bg": RGBColor(0xFF, 0xF0, 0xF5),  # 淡粉背景
    },
    "blue": {
        "primary": RGBColor(0x00, 0xB0, 0xF0),  # 天蓝
        "secondary": RGBColor(0x87, 0xCE, 0xEB),  # 浅蓝
        "accent": RGBColor(0x00, 0x70, 0xC0),  # 深蓝
        "text": RGBColor(0x00, 0x40, 0x80),  # 深蓝文字
        "bg": RGBColor(0xE0, 0xF7, 0xFF),  # 淡蓝背景
    },
    "purple": {
        "primary": RGBColor(0x9B, 0x59, 0xB6),  # 紫色
        "secondary": RGBColor(0xD8, 0xB4, 0xE8),  # 浅紫
        "accent": RGBColor(0x8E, 0x44, 0xAD),  # 深紫
        "text": RGBColor(0x4A, 0x23, 0x5A),  # 深紫文字
        "bg": RGBColor(0xF5, 0xEE, 0xF8),  # 淡紫背景
    },
    "yellow": {
        "primary": RGBColor(0xFF, 0xD7, 0x00),  # 金黄
        "secondary": RGBColor(0xFF, 0xE4, 0xB5),  # 浅黄
        "accent": RGBColor(0xFF, 0xA5, 0x00),  # 深橙黄
        "text": RGBColor(0x8B, 0x45, 0x00),  # 棕色文字
        "bg": RGBColor(0xFF, 0xFD, 0xD5),  # 淡黄背景
    },
}


def add_shape_with_gradient(slide, shape_type, left, top, width, height, color1, color2):
    """添加渐变色形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.background()
    return shape


def set_shape_solid_fill(shape, color):
    """设置形状纯色填充"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=32, font_color=None, bold=False, align=PP_ALIGN.CENTER):
    """添加文本框"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if font_color:
        p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = align
    return textbox


def add_balloon_shapes(slide, theme, count=5):
    """添加气球装饰"""
    import random
    colors = [theme["primary"], theme["secondary"], theme["accent"]]
    positions = [
        (Inches(0.3), Inches(0.2)),
        (Inches(8.5), Inches(0.5)),
        (Inches(0.8), Inches(4.5)),
        (Inches(9.0), Inches(3.8)),
        (Inches(4.5), Inches(0.1)),
    ]
    for i in range(min(count, len(positions))):
        x, y = positions[i]
        w, h = Inches(0.8), Inches(1.0)
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.fill.background()


def add_star_shapes(slide, theme, count=8):
    """添加星星装饰"""
    import random
    colors = [theme["primary"], theme["accent"]]
    for i in range(count):
        x = Inches(random.uniform(0.5, 9.0))
        y = Inches(random.uniform(0.5, 6.5))
        w, h = Inches(0.3), Inches(0.3)
        shape = slide.shapes.add_shape(MSO_SHAPE.STAR_4_POINT, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.fill.background()


def create_title_slide(prs, name, birth_date, theme):
    """创建封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    t = theme
    
    # 背景色块 - 顶部装饰
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = t["primary"]
    shape.line.fill.background()
    
    # 底部装饰条
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.5), Inches(10), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = t["secondary"]
    shape.line.fill.background()
    
    # 添加气球
    add_balloon_shapes(slide, t, 6)
    
    # 添加星星
    add_star_shapes(slide, t, 10)
    
    # 标题文字
    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(1.2),
                 "🎂 Happy Birthday 🎂", font_size=54, font_color=t["accent"], bold=True)
    
    # 名字
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(9), Inches(1.0),
                 f"{name} 的生日Party", font_size=44, font_color=t["text"], bold=True)
    
    # 生日日期
    add_text_box(slide, Inches(0.5), Inches(5.2), Inches(9), Inches(0.6),
                 f"🎈 {birth_date} 🎈", font_size=28, font_color=t["text"], bold=False)


def create_intro_slide(prs, name, theme):
    """创建目录页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    t = theme
    
    # 顶部装饰条
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = t["primary"]
    shape.line.fill.background()
    
    # 添加气球
    add_balloon_shapes(slide, t, 4)
    
    # 页面标题
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(0.8),
                 "📖 今日目录", font_size=36, font_color=t["accent"], bold=True)
    
    # 目录项
    items = [
        ("👶 第一章：成长足迹", "从襒褒中到活泫乱跳的小天使"),
        ("📷 第二章：精彩瞬间", "那些难忘的笑容和游戏时光"),
        ("🎁 第三章：祝福与期待", "爸爸妈妈的温馨祝福"),
    ]
    
    y = Inches(3.0)
    for title, desc in items:
        add_text_box(slide, Inches(1.5), y, Inches(7.5), Inches(0.6),
                     title, font_size=24, font_color=t["text"], bold=True)
        add_text_box(slide, Inches(1.8), y + Inches(0.5), Inches(7), Inches(0.5),
                     desc, font_size=18, font_color=t["primary"], bold=False)
        y += Inches(1.3)


def create_photo_slide(prs, photos, page_num, total, theme):
    """创建照片页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    t = theme
    
    # 顶部装饰条
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = t["secondary"]
    shape.line.fill.background()
    
    # 页码标题
    add_text_box(slide, Inches(0.3), Inches(0.3), Inches(9), Inches(0.7),
                 f"📷 成长足迹 ({page_num}/{total})", font_size=28, font_color=t["text"], bold=True)
    
    # 添加星星装饰
    add_star_shapes(slide, t, 5)
    
    # 根据照片数量决定布局
    photo_count = len(photos)
    if photo_count == 1:
        # 单张照片大图
        add_photo_to_slide(slide, photos[0], Inches(2), Inches(1.8), Inches(6), Inches(4.5))
    elif photo_count == 2:
        # 两张照片并排
        add_photo_to_slide(slide, photos[0], Inches(0.8), Inches(2.0), Inches(4), Inches(3.5))
        add_photo_to_slide(slide, photos[1], Inches(5.2), Inches(2.0), Inches(4), Inches(3.5))
    elif photo_count == 3:
        # 三张照片三角布局
        add_photo_to_slide(slide, photos[0], Inches(3.0), Inches(1.5), Inches(4), Inches(3.0))
        add_photo_to_slide(slide, photos[1], Inches(0.8), Inches(4.0), Inches(3.5), Inches(2.5))
        add_photo_to_slide(slide, photos[2], Inches(5.5), Inches(4.0), Inches(3.5), Inches(2.5))
    else:
        # 四张照片网格布局
        positions = [
            (Inches(0.8), Inches(1.8)),
            (Inches(5.2), Inches(1.8)),
            (Inches(0.8), Inches(4.0)),
            (Inches(5.2), Inches(4.0)),
        ]
        for i, photo in enumerate(photos[:4]):
            x, y = positions[i]
            add_photo_to_slide(slide, photo, x, y, Inches(4), Inches(2.2))


def add_photo_to_slide(slide, photo_path, left, top, width, height):
    """添加照片到幻灯片"""
    from pathlib import Path
    
    p = Path(photo_path)
    if not p.exists():
        # 创建占位符
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shape.line.color.rgb = RGBColor(0xB0, 0xB0, 0xB0)
        shape.line.width = Pt(2)
        
        # 添加占位符文字
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"📷\n{photo_path}"
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
    else:
        try:
            # 添加图片
            slide.shapes.add_picture(str(p), left, top, width=width, height=height)
        except Exception as e:
            # 如果添加失败，创建占位符
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            shape.line.color.rgb = RGBColor(0xB0, 0xB0, 0xB0)
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = f"📷\n{p.name}"
            p.font.size = Pt(12)


def create_blessing_slide(prs, name, theme):
    """创建祝福页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    t = theme
    
    # 全屏渐变背景效果 - 用两个矩形模拟
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = t["bg"]
    shape1.line.fill.background()
    
    # 中央装饰框
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape2.line.color.rgb = t["primary"]
    shape2.line.width = Pt(3)
    
    # 标题
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8),
                 "💝 爸爸妈妈的祝福 💝", font_size=32, font_color=t["accent"], bold=True)
    
    # 祝福语
    blessing_text = f"""
亲爱的{name}：

在你生日的美好日子里，爸爸妈妈想对你说：
谢谢你来到我们的生活里，带给我们无限的快乐和幸福！

{name}是一个善良、勇敢、充满好奇心的小宝贝。
每天看到你的笑容，就是我们最大的满足。

愿你健康快乐地成长，平安幸福每一天！
永远爱你的爸爸妈妈 💕
"""
    textbox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(7), Inches(3.5))
    tf = textbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(blessing_text.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        if line.startswith("亲爱的"):
            p.font.size = Pt(24)
            p.font.color.rgb = t["text"]
            p.font.bold = True
        elif line.startswith("在你") or line.startswith("愿你") or line.startswith("永远"):
            p.font.size = Pt(18)
            p.font.color.rgb = t["text"]
            p.font.italic = True
        else:
            p.font.size = Pt(20)
            p.font.color.rgb = t["primary"]
    
    # 添加气球装饰
    add_balloon_shapes(slide, t, 4)


def create_ending_slide(prs, name, theme):
    """创建结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    t = theme
    
    # 全屏背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = t["primary"]
    shape.line.fill.background()
    
    # 添加大量装饰
    add_balloon_shapes(slide, t, 10)
    add_star_shapes(slide, t, 15)
    
    # 主标题
    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(9), Inches(1.2),
                 "🎉 谢谢观看 🎉", font_size=60, font_color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    
    # 副标题
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(9), Inches(0.8),
                 f"祝 {name} 生日快乐！", font_size=36, font_color=RGBColor(0xFF, 0xF0, 0xF5), bold=True)
    
    # 底部装饰文字
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.6),
                 "🎂 愿每一天都如生日般快乐 🎂", font_size=24, font_color=RGBColor(0xFF, 0xE4, 0xB5), bold=False)


def create_birthday_ppt(name, birth_date, photos, output_path, theme_name="pink"):
    """创建生日PPT"""
    theme = THEMES.get(theme_name, THEMES["pink"])
    
    # 创建演示文稿（16:9 宽屏）
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. 封面
    create_title_slide(prs, name, birth_date, theme)
    
    # 2. 目录
    create_intro_slide(prs, name, theme)
    
    # 3. 照片页（每页最多4张照片，分组展示）
    photos_per_page = 4
    total_pages = (len(photos) + photos_per_page - 1) // photos_per_page if photos else 0
    
    for page in range(total_pages):
        start = page * photos_per_page
        end = start + photos_per_page
        page_photos = photos[start:end]
        create_photo_slide(prs, page_photos, page + 1, total_pages, theme)
    
    # 如果没有照片，添加一个提示页
    if total_pages == 0:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        add_text_box(slide, Inches(2), Inches(3), Inches(6), Inches(1),
                     "📷 等待添加成长照片...", font_size=28, font_color=theme["text"], bold=True)
    
    # 4. 祝福页
    create_blessing_slide(prs, name, theme)
    
    # 5. 结束页
    create_ending_slide(prs, name, theme)
    
    # 保存文件
    prs.save(output_path)
    print(f"✅ 生日PPT已生成: {output_path}")
    return output_path


def main():
    parser = argparse.ArgumentParser(description="幼儿园生日PPT生成器")
    parser.add_argument("--name", "-n", required=True, help="小朋友名字")
    parser.add_argument("--birth-date", "-b", required=True, help="出生日期 (YYYY-MM-DD)")
    parser.add_argument("--photos", "-p", nargs="+", default=[], help="照片路径列表")
    parser.add_argument("--output", "-o", default="birthday.pptx", help="输出文件路径")
    parser.add_argument("--theme", "-t", default="pink", choices=["pink", "blue", "purple", "yellow"],
                       help="主题色 (default: pink)")
    
    args = parser.parse_args()
    
    create_birthday_ppt(
        name=args.name,
        birth_date=args.birth_date,
        photos=args.photos,
        output_path=args.output,
        theme_name=args.theme
    )


if __name__ == "__main__":
    main()