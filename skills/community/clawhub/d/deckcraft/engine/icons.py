"""
DeckCraft v6 — Icon Library

Draws vector icons on slides using python-pptx freeform shapes.
No SVG embedding — all shapes are native PowerPoint objects (fully editable).

Usage:
    from engine.icons import icon, ICON_NAMES
    icon("rocket", slide, left=Inches(1), top=Inches(2), size=Inches(0.6))
    icon("chart-bar", slide, Inches(3), Inches(2), Inches(0.5), color=(0, 102, 204))
"""

from typing import Tuple, Optional, List

from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Type alias ───────────────────────────────────────────────────────
ColorTuple = Tuple[int, int, int]


# ── SVG-path-style icon definitions ─────────────────────────────────
# Each icon is a list of drawing commands:
#   ("M", x, y)          — move to (x, y) in [0,1]×[0,1] unit square
#   ("L", x, y)          — line to
#   ("C", x1, y1, x2, y2, x, y) — cubic bezier
#   ("Q", x1, y1, x, y) — quadratic bezier
#   ("Z",)               — close path
#   ("A", rx, ry, rot, large_arc, sweep, x, y) — arc (simplified as line-to)
#   ("CIRCLE", cx, cy, r) — circle (special, rendered as oval shape)
#   ("RECT", x, y, w, h, [rx]) — rounded rect
#
# All coordinates normalized to [0, 1]. Rotated/scaled at render time.

_ICONS = {
    # ── Simple geometric icons ──────────────────────────────────
    "check": [
        ("M", 0.15, 0.52), ("L", 0.38, 0.75), ("L", 0.85, 0.25),
    ],
    "x": [
        ("M", 0.2, 0.2), ("L", 0.8, 0.8), ("M", 0.8, 0.2), ("L", 0.2, 0.8),
    ],
    "plus": [
        ("M", 0.5, 0.15), ("L", 0.5, 0.85), ("M", 0.15, 0.5), ("L", 0.85, 0.5),
    ],
    "minus": [
        ("M", 0.15, 0.5), ("L", 0.85, 0.5),
    ],

    # ── Arrows ──────────────────────────────────────────────────
    "arrow-up": [
        ("M", 0.5, 0.15), ("L", 0.8, 0.5), ("M", 0.5, 0.15), ("L", 0.2, 0.5),
        ("M", 0.5, 0.15), ("L", 0.5, 0.85),
    ],
    "arrow-right": [
        ("M", 0.85, 0.5), ("L", 0.5, 0.2), ("M", 0.85, 0.5), ("L", 0.5, 0.8),
        ("M", 0.15, 0.5), ("L", 0.85, 0.5),
    ],
    "arrow-up-right": [
        ("M", 0.25, 0.75), ("L", 0.75, 0.25),
        ("M", 0.45, 0.25), ("L", 0.75, 0.25), ("L", 0.75, 0.55),
    ],

    # ── Star ────────────────────────────────────────────────────
    "star": [
        ("M", 0.5, 0.05), ("L", 0.62, 0.35), ("L", 0.97, 0.38),
        ("L", 0.71, 0.6), ("L", 0.79, 0.95), ("L", 0.5, 0.77),
        ("L", 0.21, 0.95), ("L", 0.29, 0.6), ("L", 0.03, 0.38),
        ("L", 0.38, 0.35), ("Z",),
    ],

    # ── Heart ───────────────────────────────────────────────────
    "heart": [
        ("M", 0.5, 0.9), ("C", 0.1, 0.6, 0.0, 0.3, 0.2, 0.15),
        ("C", 0.35, 0.0, 0.5, 0.1, 0.5, 0.3),
        ("C", 0.5, 0.1, 0.65, 0.0, 0.8, 0.15),
        ("C", 1.0, 0.3, 0.9, 0.6, 0.5, 0.9),
    ],

    # ── Bell ────────────────────────────────────────────────────
    "bell": [
        ("M", 0.25, 0.55), ("L", 0.25, 0.4),
        ("C", 0.25, 0.15, 0.75, 0.15, 0.75, 0.4),
        ("L", 0.75, 0.55), ("L", 0.82, 0.65), ("L", 0.18, 0.65), ("Z",),
        ("M", 0.4, 0.65), ("C", 0.4, 0.78, 0.6, 0.78, 0.6, 0.65),
        ("M", 0.5, 0.1), ("L", 0.5, 0.18),
    ],

    # ── Calendar ────────────────────────────────────────────────
    "calendar": [
        ("RECT", 0.1, 0.2, 0.8, 0.7, 0.05),
        ("M", 0.1, 0.38), ("L", 0.9, 0.38),
        ("M", 0.35, 0.1), ("L", 0.35, 0.3),
        ("M", 0.65, 0.1), ("L", 0.65, 0.3),
        # Day dots
        ("M", 0.28, 0.5), ("L", 0.38, 0.5), ("M", 0.46, 0.5), ("L", 0.56, 0.5),
        ("M", 0.64, 0.5), ("L", 0.74, 0.5),
        ("M", 0.28, 0.62), ("L", 0.38, 0.62), ("M", 0.46, 0.62), ("L", 0.56, 0.62),
        ("M", 0.64, 0.62), ("L", 0.74, 0.62),
    ],

    # ── Rocket ──────────────────────────────────────────────────
    "rocket": [
        ("M", 0.5, 0.05), ("C", 0.5, 0.05, 0.7, 0.25, 0.7, 0.5),
        ("L", 0.7, 0.7), ("L", 0.6, 0.8), ("L", 0.5, 0.7),
        ("L", 0.4, 0.8), ("L", 0.3, 0.7), ("L", 0.3, 0.5),
        ("C", 0.3, 0.25, 0.5, 0.05, 0.5, 0.05), ("Z",),
        ("CIRCLE", 0.5, 0.38, 0.07),
    ],

    # ── Chart-bar ───────────────────────────────────────────────
    "chart-bar": [
        ("RECT", 0.1, 0.55, 0.18, 0.35, 0.02),
        ("RECT", 0.35, 0.35, 0.18, 0.55, 0.02),
        ("RECT", 0.6, 0.2, 0.18, 0.7, 0.02),
        ("M", 0.05, 0.92), ("L", 0.95, 0.92),
    ],

    # ── Chart-line (trending-up) ────────────────────────────────
    "trending-up": [
        ("M", 0.1, 0.8), ("L", 0.35, 0.5), ("L", 0.55, 0.6), ("L", 0.9, 0.15),
        ("M", 0.6, 0.15), ("L", 0.9, 0.15), ("L", 0.9, 0.45),
    ],

    # ── Users ───────────────────────────────────────────────────
    "users": [
        ("CIRCLE", 0.5, 0.25, 0.13),
        ("M", 0.25, 0.85), ("C", 0.25, 0.55, 0.75, 0.55, 0.75, 0.85), ("Z",),
        ("CIRCLE", 0.25, 0.3, 0.1),
        ("M", 0.05, 0.8), ("C", 0.05, 0.58, 0.45, 0.58, 0.45, 0.8),
    ],

    # ── User (single) ───────────────────────────────────────────
    "user": [
        ("CIRCLE", 0.5, 0.28, 0.15),
        ("M", 0.2, 0.9), ("C", 0.2, 0.55, 0.8, 0.55, 0.8, 0.9), ("Z",),
    ],

    # ── Cog / Settings ──────────────────────────────────────────
    "cog": [
        ("CIRCLE", 0.5, 0.5, 0.15),
        # Gear teeth (8 lines radiating out)
        ("M", 0.5, 0.05), ("L", 0.5, 0.18),
        ("M", 0.5, 0.82), ("L", 0.5, 0.95),
        ("M", 0.05, 0.5), ("L", 0.18, 0.5),
        ("M", 0.82, 0.5), ("L", 0.95, 0.5),
        ("M", 0.18, 0.18), ("L", 0.27, 0.27),
        ("M", 0.73, 0.73), ("L", 0.82, 0.82),
        ("M", 0.82, 0.18), ("L", 0.73, 0.27),
        ("M", 0.27, 0.73), ("L", 0.18, 0.82),
    ],
    "settings": [  # alias
        ("CIRCLE", 0.5, 0.5, 0.15),
        ("M", 0.5, 0.05), ("L", 0.5, 0.18),
        ("M", 0.5, 0.82), ("L", 0.5, 0.95),
        ("M", 0.05, 0.5), ("L", 0.18, 0.5),
        ("M", 0.82, 0.5), ("L", 0.95, 0.5),
        ("M", 0.18, 0.18), ("L", 0.27, 0.27),
        ("M", 0.73, 0.73), ("L", 0.82, 0.82),
        ("M", 0.82, 0.18), ("L", 0.73, 0.27),
        ("M", 0.27, 0.73), ("L", 0.18, 0.82),
    ],

    # ── Circle-checkmark ────────────────────────────────────────
    "circle-checkmark": [
        ("CIRCLE", 0.5, 0.5, 0.42),
        ("M", 0.28, 0.52), ("L", 0.43, 0.68), ("L", 0.72, 0.32),
    ],

    # ── Target ──────────────────────────────────────────────────
    "target": [
        ("CIRCLE", 0.5, 0.5, 0.45),
        ("CIRCLE", 0.5, 0.5, 0.3),
        ("CIRCLE", 0.5, 0.5, 0.12),
    ],

    # ── Clock ───────────────────────────────────────────────────
    "clock": [
        ("CIRCLE", 0.5, 0.5, 0.42),
        ("M", 0.5, 0.5), ("L", 0.5, 0.22),
        ("M", 0.5, 0.5), ("L", 0.72, 0.5),
    ],

    # ── File ────────────────────────────────────────────────────
    "file": [
        ("M", 0.2, 0.05), ("L", 0.2, 0.95), ("L", 0.8, 0.95), ("L", 0.8, 0.25),
        ("L", 0.6, 0.05), ("Z",),
        ("M", 0.6, 0.05), ("L", 0.6, 0.25), ("L", 0.8, 0.25),
        ("M", 0.32, 0.5), ("L", 0.68, 0.5),
        ("M", 0.32, 0.62), ("L", 0.68, 0.62),
        ("M", 0.32, 0.74), ("L", 0.55, 0.74),
    ],

    # ── Mail ────────────────────────────────────────────────────
    "mail": [
        ("RECT", 0.05, 0.2, 0.9, 0.6, 0.04),
        ("M", 0.05, 0.2), ("L", 0.5, 0.55), ("L", 0.95, 0.2),
    ],

    # ── Download ────────────────────────────────────────────────
    "download": [
        ("M", 0.5, 0.1), ("L", 0.5, 0.65),
        ("M", 0.3, 0.5), ("L", 0.5, 0.7), ("L", 0.7, 0.5),
        ("M", 0.15, 0.85), ("L", 0.85, 0.85),
    ],

    # ── Edit / Pencil ───────────────────────────────────────────
    "edit": [
        ("M", 0.75, 0.1), ("L", 0.9, 0.25), ("L", 0.4, 0.75),
        ("L", 0.2, 0.8), ("L", 0.25, 0.6), ("Z",),
        ("M", 0.15, 0.9), ("L", 0.85, 0.9),
    ],

    # ── Filter / Funnel ────────────────────────────────────────
    "filter": [
        ("M", 0.1, 0.15), ("L", 0.9, 0.15), ("L", 0.6, 0.55),
        ("L", 0.6, 0.85), ("L", 0.4, 0.75), ("L", 0.4, 0.55), ("Z",),
    ],

    # ── Map-pin ─────────────────────────────────────────────────
    "map-pin": [
        ("M", 0.5, 0.95), ("C", 0.5, 0.95, 0.15, 0.55, 0.15, 0.35),
        ("C", 0.15, 0.12, 0.5, 0.0, 0.5, 0.0),
        ("C", 0.5, 0.0, 0.85, 0.12, 0.85, 0.35),
        ("C", 0.85, 0.55, 0.5, 0.95, 0.5, 0.95),
        ("CIRCLE", 0.5, 0.35, 0.1),
    ],

    # ── Phone ───────────────────────────────────────────────────
    "phone": [
        ("M", 0.25, 0.1), ("C", 0.25, 0.1, 0.15, 0.15, 0.15, 0.25),
        ("C", 0.15, 0.5, 0.5, 0.85, 0.75, 0.85),
        ("C", 0.85, 0.85, 0.9, 0.75, 0.9, 0.75),
        ("L", 0.75, 0.65), ("L", 0.65, 0.75),
        ("C", 0.5, 0.65, 0.35, 0.5, 0.25, 0.35),
        ("L", 0.35, 0.25), ("Z",),
    ],

    # ── Search / Magnifying glass ──────────────────────────────
    "search": [
        ("CIRCLE", 0.4, 0.4, 0.28),
        ("M", 0.6, 0.6), ("L", 0.9, 0.9),
    ],

    # ── Shield ──────────────────────────────────────────────────
    "shield": [
        ("M", 0.5, 0.05), ("L", 0.9, 0.2), ("L", 0.9, 0.5),
        ("C", 0.9, 0.75, 0.5, 0.95, 0.5, 0.95),
        ("C", 0.5, 0.95, 0.1, 0.75, 0.1, 0.5),
        ("L", 0.1, 0.2), ("Z",),
    ],

    # ── Video / Play ────────────────────────────────────────────
    "video": [
        ("RECT", 0.05, 0.15, 0.6, 0.7, 0.04),
        ("M", 0.65, 0.35), ("L", 0.95, 0.15), ("L", 0.95, 0.85), ("L", 0.65, 0.65), ("Z",),
    ],

    # ── Zap / Lightning ─────────────────────────────────────────
    "zap": [
        ("M", 0.55, 0.05), ("L", 0.25, 0.5), ("L", 0.48, 0.5),
        ("L", 0.4, 0.95), ("L", 0.75, 0.42), ("L", 0.52, 0.42), ("Z",),
    ],

    # ── Bookmark ────────────────────────────────────────────────
    "bookmark": [
        ("M", 0.2, 0.05), ("L", 0.8, 0.05), ("L", 0.8, 0.95),
        ("L", 0.5, 0.72), ("L", 0.2, 0.95), ("Z",),
    ],
}

# Aliases
_ICONS["magnifying-glass"] = _ICONS["search"]
_ICONS["lightning"] = _ICONS["zap"]
_ICONS["funnel"] = _ICONS["filter"]
_ICONS["pencil"] = _ICONS["edit"]

ICON_NAMES = sorted(set(_ICONS.keys()))


def _build_freeform_paths(commands: list, size_emu: int):
    """Convert normalized icon commands to list of (is_move, x_emu, y_emu) segments.

    We separate the commands into individual strokes (sub-paths).
    Each stroke is a sequence of points starting with a move.
    Returns list of strokes, where each stroke is [(x_emu, y_emu), ...].
    """
    strokes = []
    current_stroke = None

    for cmd in commands:
        op = cmd[0]
        if op == "M":
            if current_stroke and len(current_stroke) >= 2:
                strokes.append(current_stroke)
            current_stroke = [(int(cmd[1] * size_emu), int(cmd[2] * size_emu))]
        elif op == "L":
            if current_stroke is not None:
                current_stroke.append((int(cmd[1] * size_emu), int(cmd[2] * size_emu)))
        elif op == "Z":
            if current_stroke and len(current_stroke) >= 2:
                # Close: connect back to start
                current_stroke.append(current_stroke[0])
                strokes.append(current_stroke)
                current_stroke = None
        elif op == "C":
            # Cubic bezier — approximate with line segments
            if current_stroke is None:
                current_stroke = []
            if not current_stroke:
                continue
            x0, y0 = current_stroke[-1]
            x1, y1, x2, y2, x3, y3 = cmd[1], cmd[2], cmd[3], cmd[4], cmd[5], cmd[6]
            pts = _cubic_bezier_points(
                x0 / size_emu, y0 / size_emu,
                x1, y1, x2, y2, x3, y3,
                steps=8,
            )
            for px, py in pts:
                current_stroke.append((int(px * size_emu), int(py * size_emu)))
        elif op == "Q":
            if current_stroke is None:
                current_stroke = []
            if not current_stroke:
                continue
            x0, y0 = current_stroke[-1]
            x1, y1, x2, y2 = cmd[1], cmd[2], cmd[3], cmd[4]
            pts = _quad_bezier_points(
                x0 / size_emu, y0 / size_emu, x1, y1, x2, y2, steps=6
            )
            for px, py in pts:
                current_stroke.append((int(px * size_emu), int(py * size_emu)))
        # RECT and CIRCLE handled separately

    if current_stroke and len(current_stroke) >= 2:
        strokes.append(current_stroke)

    return strokes


def _cubic_bezier_points(x0, y0, x1, y1, x2, y2, x3, y3, steps=8):
    """Generate points along a cubic bezier curve."""
    pts = []
    for i in range(1, steps + 1):
        t = i / steps
        t2 = t * t
        t3 = t2 * t
        mt = 1 - t
        mt2 = mt * mt
        mt3 = mt2 * mt
        x = mt3 * x0 + 3 * mt2 * t * x1 + 3 * mt * t2 * x2 + t3 * x3
        y = mt3 * y0 + 3 * mt2 * t * y1 + 3 * mt * t2 * y2 + t3 * y3
        pts.append((x, y))
    return pts


def _quad_bezier_points(x0, y0, x1, y1, x2, y2, steps=6):
    """Generate points along a quadratic bezier curve."""
    pts = []
    for i in range(1, steps + 1):
        t = i / steps
        mt = 1 - t
        x = mt * mt * x0 + 2 * mt * t * x1 + t * t * x2
        y = mt * mt * y0 + 2 * mt * t * y1 + t * t * y2
        pts.append((x, y))
    return pts


def _add_freeform_stroke(slide, left_emu: int, top_emu: int, points: list,
                         stroke_color: RGBColor, stroke_width_pt: float = 2.0,
                         fill_color: RGBColor = None):
    """Add a single freeform line/polygon to a slide using python-pptx."""
    if len(points) < 2:
        return

    # Build freeform builder
    x0, y0 = points[0]
    ff_builder = slide.shapes.build_freeform(
        Emu(left_emu + x0), Emu(top_emu + y0)
    )
    if len(points) >= 2:
        # Determine if closed
        is_closed_shape = len(points) >= 3 and points[0] == points[-1]
        remaining = points[1:-1] if is_closed_shape else points[1:]
        ff_builder.add_line_segments(
            [(Emu(left_emu + x), Emu(top_emu + y)) for x, y in remaining],
            close=is_closed_shape,
        )

    shape = ff_builder.convert_to_shape()

    # Style: no fill by default, stroke color
    shape.line.color.rgb = stroke_color
    shape.line.width = Pt(stroke_width_pt)

    if fill_color is None:
        shape.fill.background()  # transparent fill
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color

    return shape


def _add_circle(slide, left_emu: int, top_emu: int, size_emu: int,
                cx: float, cy: float, r: float,
                stroke_color: RGBColor, stroke_width_pt: float = 2.0,
                fill_color: RGBColor = None):
    """Add a circle shape at normalized position (cx, cy) with radius r."""
    from pptx.util import Emu as EmuType
    circle_size = int(r * 2 * size_emu)
    circle_left = left_emu + int((cx - r) * size_emu)
    circle_top = top_emu + int((cy - r) * size_emu)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(circle_left), Emu(circle_top),
        Emu(circle_size), Emu(circle_size),
    )
    shape.line.color.rgb = stroke_color
    shape.line.width = Pt(stroke_width_pt)
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def _add_rounded_rect(slide, left_emu: int, top_emu: int, size_emu: int,
                      x: float, y: float, w: float, h: float, rx: float,
                      stroke_color: RGBColor, stroke_width_pt: float = 2.0,
                      fill_color: RGBColor = None):
    """Add a rounded rectangle shape."""
    rect_left = left_emu + int(x * size_emu)
    rect_top = top_emu + int(y * size_emu)
    rect_w = int(w * size_emu)
    rect_h = int(h * size_emu)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(rect_left), Emu(rect_top),
        Emu(rect_w), Emu(rect_h),
    )
    shape.line.color.rgb = stroke_color
    shape.line.width = Pt(stroke_width_pt)
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def icon(name: str, slide, left, top, size, color: ColorTuple = None,
         fill_color: ColorTuple = None, stroke_width: float = 2.0,
         line_only: bool = False):
    """Draw an icon on a slide.

    Args:
        name: Icon name (e.g. "rocket", "chart-bar", "check").
              See ICON_NAMES for all available icons.
        slide: python-pptx Slide object.
        left: Left position (Emu, Inches, etc).
        top: Top position (Emu, Inches, etc).
        size: Icon size (width = height, square bounding box).
        color: Stroke color as (R, G, B) tuple. Default: (51, 51, 51) dark gray.
        fill_color: Fill color for closed shapes. Default: None (transparent).
                    Only applies when line_only=False.
        stroke_width: Line width in points. Default 2.0.
        line_only: If True, only draw strokes (no fills on closed shapes).

    Returns:
        List of shapes created.
    """
    if name not in _ICONS:
        raise ValueError(
            f"Unknown icon: {name!r}. Available: {', '.join(ICON_NAMES[:20])}... "
            f"({len(ICON_NAMES)} total)"
        )

    if color is None:
        color = (51, 51, 51)
    rgb = RGBColor(*color)
    fill_rgb = RGBColor(*fill_color) if fill_color else None

    # Convert all position/size to EMU
    left_emu = int(left) if not isinstance(left, int) else left
    top_emu = int(top) if not isinstance(top, int) else top
    size_emu = int(size) if not isinstance(size, int) else size

    commands = _ICONS[name]
    shapes_created = []

    # First pass: separate special commands from path commands
    path_commands = []
    for cmd in commands:
        if cmd[0] == "CIRCLE":
            _, cx, cy, r = cmd
            s = _add_circle(slide, left_emu, top_emu, size_emu,
                            cx, cy, r, rgb, stroke_width, fill_rgb)
            shapes_created.append(s)
        elif cmd[0] == "RECT":
            _, x, y, w, h = cmd[:5]
            rx = cmd[5] if len(cmd) > 5 else 0.0
            s = _add_rounded_rect(slide, left_emu, top_emu, size_emu,
                                  x, y, w, h, rx, rgb, stroke_width, fill_rgb)
            shapes_created.append(s)
        else:
            path_commands.append(cmd)

    # Build and render freeform strokes
    strokes = _build_freeform_paths(path_commands, size_emu)
    for stroke in strokes:
        # Check if this is a closed shape (first == last point)
        is_closed = len(stroke) >= 3 and stroke[0] == stroke[-1]
        s = _add_freeform_stroke(
            slide, left_emu, top_emu, stroke,
            stroke_color=rgb,
            stroke_width_pt=stroke_width,
            fill_color=fill_rgb if (is_closed and not line_only) else None,
        )
        if s is not None:
            shapes_created.append(s)

    return shapes_created
