#!/usr/bin/env python3
"""
AI PPT 生成器
逐页调用 Nano Banana Pro 生成幻灯片图片，组装成 PPTX。
支持 Ofox（无需 VPN）和 OpenRouter。

用法：
    # 1. 编辑 slides.json 定义每页的 prompt
    # 2. 运行生成
    python3 generate_deck.py

    # 只生成指定范围
    python3 generate_deck.py --start 1 --end 5

    # 重跑某一页（先删旧图）
    rm output/slide_03.jpg && python3 generate_deck.py --start 3 --end 3
"""

import json
import base64
import os
import sys
import time
import argparse
import urllib.request

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_DIR = os.path.join(SCRIPT_DIR, "output")
SLIDES_JSON = os.path.join(SCRIPT_DIR, "slides.json")


def detect_provider():
    """自动检测 API provider"""
    key = os.environ.get("OFOX_API_KEY") or os.environ.get("OPENROUTER_API_KEY")
    if not key:
        return None, None
    if key.startswith("sk-of-"):
        return "ofox", key
    else:
        return "openrouter", key


def generate_image_ofox(prompt, api_key):
    """Ofox: /v1/images/generations"""
    payload = json.dumps({
        "model": "google/gemini-3-pro-image-preview",
        "prompt": prompt,
        "size": "1792x1024",
        "response_format": "b64_json"
    }).encode("utf-8")
    req = urllib.request.Request(
        "https://api.ofox.ai/v1/images/generations",
        data=payload,
        headers={"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}
    )
    with urllib.request.urlopen(req, timeout=180) as resp:
        data = json.loads(resp.read().decode("utf-8"))
    images = data.get("data", [])
    if images and images[0].get("b64_json"):
        return base64.b64decode(images[0]["b64_json"])
    raise RuntimeError(f"No image: {json.dumps(data)[:200]}")


def generate_image_openrouter(prompt, api_key):
    """OpenRouter: /v1/chat/completions"""
    payload = json.dumps({
        "model": "google/gemini-3-pro-image-preview",
        "messages": [{"role": "user", "content": prompt}]
    }).encode("utf-8")
    req = urllib.request.Request(
        "https://openrouter.ai/api/v1/chat/completions",
        data=payload,
        headers={"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}
    )
    with urllib.request.urlopen(req, timeout=180) as resp:
        data = json.loads(resp.read().decode("utf-8"))
    msg = data.get("choices", [{}])[0].get("message", {})
    for img in msg.get("images", []):
        url = img.get("image_url", {}).get("url", "")
        if "base64," in url:
            return base64.b64decode(url.split(",", 1)[1])
    content = msg.get("content", "")
    if isinstance(content, list):
        for part in content:
            if isinstance(part, dict) and part.get("type") == "image_url":
                url = part.get("image_url", {}).get("url", "")
                if "base64," in url:
                    return base64.b64decode(url.split(",", 1)[1])
    raise RuntimeError(f"No image: {json.dumps(data)[:200]}")


def assemble_pptx(output_dir, total_slides, output_file):
    """组装 JPG 为 PPTX"""
    from pptx import Presentation
    from pptx.util import Inches, Emu
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i in range(1, total_slides + 1):
        img_path = os.path.join(output_dir, f"slide_{i:02d}.jpg")
        if not os.path.exists(img_path):
            print(f"  WARNING: slide_{i:02d}.jpg missing, skipping")
            continue
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(img_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    prs.save(output_file)
    print(f"PPTX saved: {output_file} ({len(prs.slides)} slides)")


def validate_slides_json(path):
    """校验 slides.json：合法 JSON + 发现中文内部双引号隐患"""
    with open(path, "r") as f:
        raw = f.read()
    try:
        config = json.loads(raw)
    except json.JSONDecodeError as e:
        # 自动尝试修复中文内部的 " 为 「」
        print(f"WARN: slides.json JSON 错误 at pos {e.pos}；尝试自动修复中文内引号…")
        result = []
        for i, ch in enumerate(raw):
            if ch == '"':
                prev = raw[i-1] if i > 0 else ''
                nxt = raw[i+1] if i+1 < len(raw) else ''
                # structural quote: surrounded by whitespace/delimiters
                structural = prev in '{[,:' or nxt in '},]:' or prev in ' \n\t' or nxt in ' \n\t'
                escaped = i > 0 and raw[i-1] == '\\'
                if structural or escaped:
                    result.append(ch)
                else:
                    result.append('」')
            else:
                result.append(ch)
        fixed = ''.join(result)
        try:
            config = json.loads(fixed)
            with open(path, "w") as f:
                f.write(fixed)
            print("OK: 自动修复成功，已写回 slides.json")
        except json.JSONDecodeError as e2:
            print(f"ERROR: 自动修复失败 at pos {e2.pos}")
            print(f"  Context: {fixed[max(0,e2.pos-60):e2.pos+60]}")
            sys.exit(1)
    return config


def generate_with_retry(prompt, provider, api_key, max_attempts=3):
    """带指数退避的自动重试"""
    last_err = None
    for attempt in range(1, max_attempts + 1):
        try:
            if provider == "ofox":
                return generate_image_ofox(prompt, api_key)
            else:
                return generate_image_openrouter(prompt, api_key)
        except Exception as e:
            last_err = e
            if attempt < max_attempts:
                backoff = 10 * (2 ** (attempt - 1))  # 10s, 20s, 40s
                print(f"  retry {attempt}/{max_attempts - 1} in {backoff}s: {str(e)[:80]}")
                time.sleep(backoff)
    raise last_err


def main():
    parser = argparse.ArgumentParser(description="AI PPT 生成器")
    parser.add_argument("--slides", default=SLIDES_JSON, help="Slides JSON 文件路径")
    parser.add_argument("--start", type=int, default=1, help="起始页码")
    parser.add_argument("--end", type=int, default=999, help="结束页码")
    parser.add_argument("--output-dir", default=OUTPUT_DIR, help="图片输出目录")
    parser.add_argument("--assemble", action="store_true", help="生成完后自动组装 PPTX")
    parser.add_argument("--delay", type=int, default=3, help="每页间隔秒数")
    parser.add_argument("--max-retries", type=int, default=3, help="每页最多重试次数")
    args = parser.parse_args()

    # 读取 + 校验 slides 定义
    config = validate_slides_json(args.slides)

    style = config.get("style", "")
    slides = config.get("slides", [])

    # 检测 provider
    provider, api_key = detect_provider()
    if not provider:
        print("ERROR: 未设置 API key")
        print("  export IMAGE_API_KEY=your_key")
        print("  export IMAGE_API_KEY=your_key")
        sys.exit(1)

    os.makedirs(args.output_dir, exist_ok=True)
    to_gen = [s for s in slides if args.start <= s["num"] <= args.end]

    print(f"Provider: {provider}")
    print(f"Slides: {args.start}-{args.end} ({len(to_gen)} pages)")
    print()

    success, failed = 0, 0
    for i, slide in enumerate(to_gen):
        num = slide["num"]
        output_file = os.path.join(args.output_dir, f"slide_{num:02d}.jpg")

        if os.path.exists(output_file) and os.path.getsize(output_file) > 10000:
            print(f"[{i+1}/{len(to_gen)}] Slide {num:02d} exists, skipping")
            success += 1
            continue

        prompt = style + "\n\n" + slide["prompt"]
        print(f"[{i+1}/{len(to_gen)}] Generating slide {num:02d}...", flush=True)

        try:
            img_bytes = generate_with_retry(prompt, provider, api_key, max_attempts=args.max_retries)
            with open(output_file, "wb") as f:
                f.write(img_bytes)
            print(f"  Saved ({len(img_bytes)} bytes)")
            success += 1
        except Exception as e:
            print(f"  ERROR (all retries failed): {e}")
            failed += 1

        if slide != to_gen[-1]:
            time.sleep(args.delay)

    print(f"\nDone: {success} success, {failed} failed")

    # 自动组装
    if args.assemble:
        pptx_file = os.path.join(args.output_dir, "deck.pptx")
        assemble_pptx(args.output_dir, max(s["num"] for s in slides), pptx_file)


if __name__ == "__main__":
    main()
