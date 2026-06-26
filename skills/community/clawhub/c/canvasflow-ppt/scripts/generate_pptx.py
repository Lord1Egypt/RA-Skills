#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Magazine PPT Generator - PPTX Export
Author: WuWenBin-BeiJing-ST
Version: 1.0.0

生成可编辑的 PPTX 格式演示文稿
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Dict, List, Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("❌ 缺少依赖: python-pptx")
    print("请运行: pip3 install python-pptx")
    sys.exit(1)

# 主题色配置
THEMES = {
    "ink-classic": {
        "primary": (26, 26, 46),  # RGB
        "secondary": (74, 78, 105),
        "background": (250, 250, 250),
        "text": (45, 45, 45),
    },
    "indigo-porcelain": {
        "primary": (44, 62, 80),
        "secondary": (52, 152, 219),
        "background": (255, 255, 255),
        "text": (44, 62, 80),
    },
    "forest-ink": {
        "primary": (45, 52, 54),
        "secondary": (0, 184, 148),
        "background": (248, 249, 250),
        "text": (45, 52, 54),
    },
    "kraft-paper": {
        "primary": (212, 163, 115),
        "secondary": (233, 196, 106),
        "background": (254, 250, 224),
        "text": (38, 70, 83),
    },
    "dune": {
        "primary": (38, 70, 83),
        "secondary": (231, 111, 81),
        "background": (255, 255, 255),
        "text": (38, 70, 83),
    },
}


def rgb_to_rgbcolor(rgb_tuple: tuple) -> RgbColor:
    """将 RGB 元组转换为 RgbColor"""
    return RgbColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])


def create_presentation(theme_name: str = "ink-classic") -> Presentation:
    """创建演示文稿对象"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 比例
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs: Presentation, title: str, subtitle: str = "", theme: dict = None):
    """添加封面幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_to_rgbcolor(theme.get("background", (255, 255, 255)))

    # 标题
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(12.333)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title

    # 格式化标题
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_to_rgbcolor(theme.get("primary", (0, 0, 0)))

    # 副标题
    if subtitle:
        sub_top = Inches(4.2)
        sub_height = Inches(1)

        sub_box = slide.shapes.add_textbox(left, sub_top, width, sub_height)
        sub_frame = sub_box.text_frame
        sub_frame.text = subtitle

        sub_para = sub_frame.paragraphs[0]
        sub_para.alignment = PP_ALIGN.CENTER
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = rgb_to_rgbcolor(theme.get("text", (100, 100, 100)))


def add_content_slide(
    prs: Presentation, title: str, content: List[str], theme: dict = None
):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_to_rgbcolor(theme.get("background", (255, 255, 255)))

    # 标题
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(12.333)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title

    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_to_rgbcolor(theme.get("primary", (0, 0, 0)))

    # 内容
    content_top = Inches(2)
    content_height = Inches(5)

    content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, item in enumerate(content):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()

        p.text = f"• {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = rgb_to_rgbcolor(theme.get("text", (50, 50, 50)))
        p.space_after = Pt(12)


def add_three_column_slide(
    prs: Presentation,
    title: str,
    columns: List[Dict],
    theme: dict = None,
):
    """添加三栏幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_to_rgbcolor(theme.get("background", (255, 255, 255)))

    # 标题
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(12.333)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_to_rgbcolor(theme.get("primary", (0, 0, 0)))

    # 三栏内容
    col_width = Inches(3.5)
    col_height = Inches(4)
    col_top = Inches(2.5)
    gap = Inches(0.9)

    for i, col in enumerate(columns[:3]):
        col_left = Inches(1) + (col_width + gap) * i

        # 栏标题
        col_title_box = slide.shapes.add_textbox(col_left, col_top, col_width, Inches(0.8))
        col_title_frame = col_title_box.text_frame
        col_title_frame.text = col.get("title", "")
        col_title_para = col_title_frame.paragraphs[0]
        col_title_para.alignment = PP_ALIGN.CENTER
        col_title_para.font.size = Pt(24)
        col_title_para.font.bold = True
        col_title_para.font.color.rgb = rgb_to_rgbcolor(theme.get("secondary", (100, 100, 100)))

        # 栏描述
        col_desc_box = slide.shapes.add_textbox(
            col_left, col_top + Inches(1), col_width, col_height
        )
        col_desc_frame = col_desc_box.text_frame
        col_desc_frame.word_wrap = True
        col_desc_frame.text = col.get("description", "")
        col_desc_para = col_desc_frame.paragraphs[0]
        col_desc_para.alignment = PP_ALIGN.CENTER
        col_desc_para.font.size = Pt(16)
        col_desc_para.font.color.rgb = rgb_to_rgbcolor(theme.get("text", (80, 80, 80)))


def add_quote_slide(prs: Presentation, quote: str, author: str, theme: dict = None):
    """添加引用幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_to_rgbcolor(theme.get("background", (255, 255, 255)))

    # 引号装饰
    quote_mark_left = Inches(5.5)
    quote_mark_top = Inches(1.5)
    quote_mark_box = slide.shapes.add_textbox(
        quote_mark_left, quote_mark_top, Inches(2), Inches(1.5)
    )
    quote_mark_frame = quote_mark_box.text_frame
    quote_mark_frame.text = "❝"
    quote_mark_para = quote_mark_frame.paragraphs[0]
    quote_mark_para.alignment = PP_ALIGN.CENTER
    quote_mark_para.font.size = Pt(120)
    quote_mark_para.font.color.rgb = RgbColor(200, 200, 200)

    # 引用文字
    quote_left = Inches(1)
    quote_top = Inches(2.8)
    quote_width = Inches(11.333)
    quote_height = Inches(2)

    quote_box = slide.shapes.add_textbox(quote_left, quote_top, quote_width, quote_height)
    quote_frame = quote_box.text_frame
    quote_frame.word_wrap = True
    quote_frame.text = f'"{quote}"'

    quote_para = quote_frame.paragraphs[0]
    quote_para.alignment = PP_ALIGN.CENTER
    quote_para.font.size = Pt(32)
    quote_para.font.italic = True
    quote_para.font.color.rgb = rgb_to_rgbcolor(theme.get("text", (50, 50, 50)))

    # 作者
    if author:
        author_top = Inches(5.5)
        author_box = slide.shapes.add_textbox(quote_left, author_top, quote_width, Inches(0.6))
        author_frame = author_box.text_frame
        author_frame.text = f"—— {author}"

        author_para = author_frame.paragraphs[0]
        author_para.alignment = PP_ALIGN.CENTER
        author_para.font.size = Pt(18)
        author_para.font.color.rgb = rgb_to_rgbcolor(theme.get("text", (120, 120, 120)))


def add_end_slide(prs: Presentation, theme: dict = None):
    """添加结尾幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_to_rgbcolor(theme.get("background", (255, 255, 255)))

    # 感谢文字
    left = Inches(0.5)
    top = Inches(3)
    width = Inches(12.333)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "感谢观看"

    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_to_rgbcolor(theme.get("primary", (0, 0, 0)))


def parse_outline(outline_path: str) -> Dict:
    """解析大纲 JSON 文件"""
    with open(outline_path, "r", encoding="utf-8") as f:
        return json.load(f)


def generate_pptx(outline: Dict, theme_name: str, output_path: str) -> str:
    """生成 PPTX 文件"""
    theme = THEMES.get(theme_name, THEMES["ink-classic"])
    prs = create_presentation(theme_name)

    slides = outline.get("slides", [])

    for slide_data in slides:
        layout = slide_data.get("layout", "L01")
        title = slide_data.get("title", "")
        subtitle = slide_data.get("subtitle", "")
        content = slide_data.get("content", [])
        quote = slide_data.get("quote", "")
        author = slide_data.get("author", "")

        if layout == "L01":  # 封面
            add_title_slide(prs, title, subtitle, theme)
        elif layout == "L05":  # 三栏
            columns = content if isinstance(content, list) else []
            add_three_column_slide(prs, title, columns, theme)
        elif layout == "L08":  # 引用
            add_quote_slide(prs, quote or title, author, theme)
        elif layout == "L10":  # 结尾
            add_end_slide(prs, theme)
        else:  # 默认内容页
            items = content if isinstance(content, list) else []
            add_content_slide(prs, title, items, theme)

    # 写入文件
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))

    return str(output_file)


def main():
    parser = argparse.ArgumentParser(description="生成 PPTX 格式演示文稿")
    parser.add_argument("--content", required=True, help="PPT 主题/内容描述")
    parser.add_argument("--outline", help="大纲 JSON 文件路径")
    parser.add_argument(
        "--theme", default="ink-classic", choices=THEMES.keys(), help="主题色"
    )
    parser.add_argument(
        "--output", default="./output/presentation.pptx", help="输出路径"
    )

    args = parser.parse_args()

    # 如果提供了大纲文件，使用大纲生成
    if args.outline:
        outline = parse_outline(args.outline)
    else:
        # 否则使用内容描述生成简单大纲
        outline = {
            "title": args.content,
            "slides": [
                {"layout": "L01", "title": args.content, "subtitle": "演示文稿"},
                {
                    "layout": "L05",
                    "title": "核心要点",
                    "content": [
                        {"title": "要点一", "description": "详细说明"},
                        {"title": "要点二", "description": "详细说明"},
                        {"title": "要点三", "description": "详细说明"},
                    ],
                },
                {"layout": "L01", "title": "感谢观看"},
            ],
        }

    # 生成 PPTX
    output_file = generate_pptx(outline, args.theme, args.output)
    print(f"✅ PPTX 已生成: {output_file}")

    # 自动打开（macOS）
    if sys.platform == "darwin":
        os.system(f'open "{output_file}"')


if __name__ == "__main__":
    main()
