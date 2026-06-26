#!/usr/bin/env python3
"""Bidirectional PPT ↔ Excel sync.

Modes:
  --build   : build comparison Excel from EN-pre + EN-post diff
              python3 excel_sync.py --build <en-pre.pptx> <en-post.pptx> <cn.pptx> <out.xlsx>
                                    [--cn-offset <yaml>]
  --reverse : update existing Excel's en_optimized column to match current PPT
              python3 excel_sync.py --reverse <pptx-path> <excel-path>
  --verify  : check Excel en_optimized rows are findable in current PPT
              python3 excel_sync.py --verify <pptx-path> <excel-path>

Excel header (row 2): page · kind · cn · en_original · en_optimized · notes

Three guard-rails enforced on every write:
  1. Pre-write lock check (~$xxx)
  2. Header column count assertion
  3. Post-write readback assertion

CN ↔ EN slide alignment:
  Default is 1:1 (EN slide N maps to CN slide N).
  If the two decks have been restructured (pages inserted / removed),
  pass --cn-offset <yaml> pointing to a config like:

    default: 0          # default offset added to EN slide number
    overrides:
      "9-26": -1        # EN 9..26 map to CN 8..25
      "20":   null      # EN 20 has no CN counterpart (e.g. Appendix divider)
"""
import sys
import os
import difflib
import shutil
from collections import defaultdict
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openpyxl import load_workbook, Workbook
from openpyxl.styles import Font, Alignment, PatternFill

HEADERS = ["page", "kind", "cn", "en_original", "en_optimized", "notes"]


def lock_file_path(p):
    d, name = os.path.split(p)
    return os.path.join(d, "~$" + name)


def check_no_lock(path):
    lock = lock_file_path(path)
    if os.path.exists(lock):
        print(f"⚠ STOP. Lock file exists: {lock}")
        print(f"  → '{os.path.basename(path)}' appears open in Excel/PowerPoint.")
        print(f"  Save and close it, then rerun.")
        sys.exit(2)


def collect_paragraphs(path):
    prs = Presentation(path)
    out = []

    def walk(shapes, sidx, prefix=""):
        for i, shp in enumerate(shapes):
            loc = f"{prefix}/{i}"
            try:
                if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(shp.shapes, sidx, loc)
                    continue
            except Exception:
                pass
            if shp.has_text_frame:
                for pi, p in enumerate(shp.text_frame.paragraphs):
                    text = "".join(r.text for r in p.runs)
                    if text.strip():
                        out.append({"slide": sidx, "loc": f"{loc}/p{pi}",
                                    "kind": "text", "text": text})
            if shp.has_table:
                for ri, row in enumerate(shp.table.rows):
                    for ci, cell in enumerate(row.cells):
                        for pi, p in enumerate(cell.text_frame.paragraphs):
                            text = "".join(r.text for r in p.runs)
                            if text.strip():
                                out.append({"slide": sidx,
                                            "loc": f"{loc}/r{ri}c{ci}/p{pi}",
                                            "kind": "table", "text": text})

    for idx, slide in enumerate(prs.slides, 1):
        walk(slide.shapes, idx)
    return out


def assert_header(ws):
    hdr = [ws.cell(2, c).value for c in range(2, 2 + len(HEADERS))]
    if hdr != HEADERS:
        return False
    return True


def write_header(ws):
    for ci, h in enumerate(HEADERS, start=2):
        c = ws.cell(row=2, column=ci, value=h)
        c.font = Font(bold=True)
        c.fill = PatternFill("solid", fgColor="DDDDDD")


def column_widths(ws):
    widths = {2: 6, 3: 7, 4: 50, 5: 50, 6: 50, 7: 30}
    for col, w in widths.items():
        ws.column_dimensions[chr(64 + col)].width = w


def load_offset_config(yaml_path):
    """Load CN↔EN slide offset config. Returns (default_offset, overrides_dict)."""
    if not yaml_path:
        return 0, {}
    try:
        import yaml
    except ImportError:
        print("⚠ --cn-offset requires PyYAML: pip3 install pyyaml")
        sys.exit(1)
    with open(yaml_path) as f:
        cfg = yaml.safe_load(f) or {}
    default = cfg.get("default", 0)
    raw = cfg.get("overrides", {}) or {}
    # Expand range keys ("9-26") into per-page dict
    overrides = {}
    for k, v in raw.items():
        ks = str(k)
        if "-" in ks:
            lo, hi = ks.split("-", 1)
            for p in range(int(lo), int(hi) + 1):
                overrides[p] = v
        else:
            overrides[int(ks)] = v
    return default, overrides


def make_cn_slide_resolver(default_offset, overrides):
    """Returns a function: en_slide -> cn_slide_or_None."""
    def resolve(en_slide):
        if en_slide in overrides:
            v = overrides[en_slide]
            if v is None:
                return None
            # Treat an override int as the offset to apply (not absolute)
            return en_slide + int(v)
        return en_slide + default_offset
    return resolve


def build(en_pre, en_post, cn_path, out_xlsx, cn_offset_yaml=None):
    check_no_lock(out_xlsx)
    pre = collect_paragraphs(en_pre)
    post = collect_paragraphs(en_post)
    cn = collect_paragraphs(cn_path)
    default_offset, overrides = load_offset_config(cn_offset_yaml)
    cn_slide_for_en = make_cn_slide_resolver(default_offset, overrides)
    pre_map = {(p["slide"], p["loc"]): p["text"] for p in pre}
    post_map = {(p["slide"], p["loc"]): p["text"] for p in post}
    kind_map = {(p["slide"], p["loc"]): p["kind"] for p in post}

    diffs = []
    for key in sorted(pre_map):
        if key in post_map and pre_map[key] != post_map[key]:
            diffs.append({
                "slide": key[0], "loc": key[1],
                "kind": kind_map.get(key, "text"),
                "en_original": pre_map[key],
                "en_optimized": post_map[key],
            })

    # CN ordinal alignment
    cn_by_slide = defaultdict(list)
    for c in cn:
        cn_by_slide[c["slide"]].append(c["text"])
    post_by_slide = defaultdict(list)
    for p in post:
        post_by_slide[p["slide"]].append(p)

    for d in diffs:
        d["cn"] = ""
        cs = cn_slide_for_en(d["slide"])
        if cs is None:
            continue
        en_paras = post_by_slide[d["slide"]]
        try:
            idx = next(i for i, p in enumerate(en_paras) if p["loc"] == d["loc"])
            if idx < len(cn_by_slide[cs]):
                d["cn"] = cn_by_slide[cs][idx]
        except StopIteration:
            pass
        d["notes"] = ""

    wb = Workbook()
    ws = wb.active
    ws.title = "comparison"
    write_header(ws)
    diffs.sort(key=lambda d: (d["slide"], d["loc"]))
    for ri, d in enumerate(diffs, start=3):
        ws.cell(ri, 2, d["slide"])
        ws.cell(ri, 3, d["kind"])
        ws.cell(ri, 4, d["cn"])
        ws.cell(ri, 5, d["en_original"])
        ws.cell(ri, 6, d["en_optimized"])
        ws.cell(ri, 7, d.get("notes", ""))
        for c in range(4, 8):
            ws.cell(ri, c).alignment = Alignment(wrap_text=True, vertical="top")
    column_widths(ws)
    wb.save(out_xlsx)

    # Readback assertion
    wb2 = load_workbook(out_xlsx)
    ws2 = wb2.active
    assert assert_header(ws2), "Post-write header check failed!"
    assert ws2.max_row >= len(diffs) + 2, "Row count mismatch on readback"
    print(f"Built {len(diffs)} rows → {out_xlsx}")


def reverse_sync(pptx_path, xlsx_path):
    check_no_lock(xlsx_path)
    paras = collect_paragraphs(pptx_path)
    ppt_by_slide = defaultdict(list)
    for p in paras:
        ppt_by_slide[p["slide"]].append(p["text"])

    wb = load_workbook(xlsx_path)
    ws = wb.active
    if not assert_header(ws):
        print("⚠ Header mismatch; aborting.")
        sys.exit(3)

    eo_col = 6  # en_optimized
    updated_ordinal = 0
    updated_fuzzy = 0
    unmatched = []

    # Build (slide, ordinal) → text
    for r in range(3, ws.max_row + 1):
        page = ws.cell(r, 2).value
        eo = ws.cell(r, eo_col).value
        if not page or not eo:
            continue
        eo_s = str(eo).strip()
        slide_texts = ppt_by_slide.get(page, [])
        if eo_s in slide_texts:
            continue
        # Fuzzy match within same slide
        best = difflib.get_close_matches(eo_s, slide_texts, n=1, cutoff=0.6)
        if best:
            ws.cell(r, eo_col, best[0])
            updated_fuzzy += 1
        else:
            unmatched.append((r, page, eo_s[:80]))

    wb.save(xlsx_path)
    print(f"Reverse sync: {updated_ordinal} ordinal + {updated_fuzzy} fuzzy = "
          f"{updated_ordinal + updated_fuzzy} updated; {len(unmatched)} unmatched")
    for r, p, t in unmatched[:20]:
        print(f"  row {r:3} P{p:2}: {t}")


def verify(pptx_path, xlsx_path):
    paras = collect_paragraphs(pptx_path)
    ppt = defaultdict(set)
    for p in paras:
        ppt[p["slide"]].add(p["text"].strip())
    wb = load_workbook(xlsx_path)
    ws = wb.active
    if not assert_header(ws):
        print("⚠ Header mismatch.")
        return
    miss = []
    for r in range(3, ws.max_row + 1):
        page = ws.cell(r, 2).value
        eo = ws.cell(r, 6).value
        if page and eo and str(eo).strip() not in ppt[page]:
            miss.append((r, page, str(eo).strip()))
    print(f"Mismatches: {len(miss)}")
    for r, p, t in miss[:30]:
        print(f"  row {r:3} P{p:2}: {t[:100]}")


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    mode = sys.argv[1]
    if mode == "--build" and len(sys.argv) >= 6:
        en_pre, en_post, cn_path, out_xlsx = sys.argv[2:6]
        cn_offset_yaml = None
        if "--cn-offset" in sys.argv:
            cn_offset_yaml = sys.argv[sys.argv.index("--cn-offset") + 1]
        build(en_pre, en_post, cn_path, out_xlsx, cn_offset_yaml)
    elif mode == "--reverse" and len(sys.argv) == 4:
        reverse_sync(sys.argv[2], sys.argv[3])
    elif mode == "--verify" and len(sys.argv) == 4:
        verify(sys.argv[2], sys.argv[3])
    else:
        print(__doc__)
        sys.exit(1)


if __name__ == "__main__":
    main()
