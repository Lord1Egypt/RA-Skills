#!/usr/bin/env python3
"""
学术汇报PPT导出器 - 将PPT结构导出为PowerPoint文件
"""

import json
import sys
from typing import Dict, List, Any
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


class AcademicPPTExporter:
    """学术PPT导出器"""
    
    # 配色方案 - 学术风格
    COLORS = {
        'primary': RgbColor(0x1a, 0x5f, 0x7a),      # 深蓝绿色 - 主色
        'secondary': RgbColor(0x57, 0xa0, 0xd3),    # 浅蓝色 - 辅助
        'accent': RgbColor(0xf4, 0xa2, 0x61),       # 橙色 - 强调
        'dark': RgbColor(0x2c, 0x3e, 0x50),         # 深色 - 文字
        'light': RgbColor(0xec, 0xf0, 0xf1),        # 浅色 - 背景
        'white': RgbColor(0xff, 0xff, 0xff),
    }
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 比例
        self.prs.slide_height = Inches(7.5)
    
    def _add_title_slide(self, title: str, authors: str, venue: str, year: str):
        """添加标题页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
            self.prs.slide_width, self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS['primary']
        shape.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 作者
        author_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
        )
        tf = author_box.text_frame
        p = tf.paragraphs[0]
        p.text = authors
        p.font.size = Pt(24)
        p.font.color.rgb = self.COLORS['light']
        p.alignment = PP_ALIGN.CENTER
        
        # 期刊/年份
        venue_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.2), Inches(12.333), Inches(0.6)
        )
        tf = venue_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{venue}, {year}"
        p.font.size = Pt(18)
        p.font.color.rgb = self.COLORS['secondary']
        p.alignment = PP_ALIGN.CENTER
    
    def _add_section_divider(self, section_title: str, section_num: str):
        """添加章节分隔页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 左侧色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(4), self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS['primary']
        shape.line.fill.background()
        
        # 章节编号
        num_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8), Inches(3), Inches(1.5)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_num
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        
        # 章节标题
        title_box = slide.shapes.add_textbox(
            Inches(4.5), Inches(3), Inches(8), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
    
    def _add_content_slide(self, title: str, bullets: List[str], notes: str = ""):
        """添加内容页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 顶部装饰条
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            self.prs.slide_width, Inches(0.15)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLORS['primary']
        bar.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        
        # 内容
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.8)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # 处理加粗标记 **text**
            if bullet.startswith('**') and bullet.endswith('**'):
                p.text = bullet[2:-2]
                p.font.bold = True
            elif '**' in bullet:
                # 简单处理行内加粗
                parts = bullet.split('**')
                p.text = ''.join(parts)
                p.font.bold = False
            else:
                p.text = bullet
            
            p.font.size = Pt(20) if bullet.startswith('•') or bullet.startswith(('1.', '2.', '3.', '4.', '5.', '→')) else Pt(22)
            p.font.color.rgb = self.COLORS['dark']
            p.space_before = Pt(12) if bullet == "" else Pt(6)
            p.line_spacing = 1.3
        
        # 演讲者备注
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    
    def export(self, structure: Dict[str, Any], output_path: str):
        """导出PPT"""
        # 标题页
        self._add_title_slide(
            structure['title'],
            structure['authors'],
            structure['venue'],
            structure['year']
        )
        
        slides_data = structure['slides']
        
        # 第一部分：科学问题
        if slides_data['part1_problem']:
            self._add_section_divider("科学问题与背景", "01")
            for slide_data in slides_data['part1_problem']:
                self._add_content_slide(
                    slide_data['title'],
                    slide_data['bullets'],
                    slide_data.get('notes', '')
                )
        
        # 第二部分：方法
        if slides_data['part2_method']:
            self._add_section_divider("研究方法", "02")
            for slide_data in slides_data['part2_method']:
                self._add_content_slide(
                    slide_data['title'],
                    slide_data['bullets'],
                    slide_data.get('notes', '')
                )
        
        # 第三部分：结论
        if slides_data['part3_results']:
            self._add_section_divider("结论与验证", "03")
            for slide_data in slides_data['part3_results']:
                self._add_content_slide(
                    slide_data['title'],
                    slide_data['bullets'],
                    slide_data.get('notes', '')
                )
        
        # 第四部分：局限
        if slides_data['part4_limitations']:
            self._add_section_divider("局限与展望", "04")
            limitation = slides_data['part4_limitations']
            self._add_content_slide(
                limitation['title'],
                limitation['bullets'],
                limitation.get('notes', '')
            )
        
        # 总结页
        if slides_data['summary']:
            summary = slides_data['summary']
            self._add_content_slide(
                summary['title'],
                summary['bullets'],
                summary.get('notes', '')
            )
        
        # 保存
        self.prs.save(output_path)
        print(f"PPT已保存至: {output_path}")
        print(f"总页数: {len(self.prs.slides)}")


def main():
    """主函数"""
    if len(sys.argv) < 3:
        print("Usage: python export_ppt.py <structure_json_file> <output_pptx_file>")
        sys.exit(1)
    
    structure_file = sys.argv[1]
    output_file = sys.argv[2]
    
    # 读取结构
    with open(structure_file, 'r', encoding='utf-8') as f:
        structure = json.load(f)
    
    # 导出
    exporter = AcademicPPTExporter()
    exporter.export(structure, output_file)


if __name__ == "__main__":
    main()
