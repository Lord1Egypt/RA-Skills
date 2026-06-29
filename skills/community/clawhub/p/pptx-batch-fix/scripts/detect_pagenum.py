#!/usr/bin/env python3
"""
detect_pagenum.py — 检测 PPTX 文件中的页码格式
用法: python detect_pagenum.py <pptx文件路径>
"""
import sys
import re
import os

sys.stdout.reconfigure(encoding='utf-8')

def detect_page_numbers(pptx_path):
    from pptx import Presentation

    if not os.path.exists(pptx_path):
        print(f"错误：文件不存在 — {pptx_path}")
        return

    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    total = len(slides)

    print(f"文件：{os.path.basename(pptx_path)}")
    print(f"总页数：{total}")
    print()

    # 幻灯片尺寸（EMU）
    W = prs.slide_width
    H = prs.slide_height
    print(f"幻灯片尺寸：{W} x {H} EMU（{W/914400:.1f}\" x {H/914400:.1f}\")")
    print()

    # 1. 扫描所有幻灯片中的页码候选
    print("=" * 60)
    print("扫描结果：幻灯片中的页码候选")
    print("=" * 60)

    candidates = []  # (slide_idx, shape_name, text, left, top, reason)

    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            # 格式1：X/Y 或 X / Y
            if re.match(r'^\d{1,3}\s*/\s*\d{1,4}$', text.replace('\n', ' ')):
                candidates.append((i, shape.name, text, shape.left, shape.top, '格式 X/Y'))

            # 格式2：纯数字，且在右下角区域
            elif text.isdigit() and len(text) <= 4:
                # 右下角判断：left > W*0.6 且 top > H*0.75
                if shape.left > W * 0.55 and shape.top > H * 0.70:
                    candidates.append((i, shape.name, text, shape.left, shape.top, '右下角纯数字'))

    if candidates:
        # 按页面分组
        by_slide = {}
        for c in candidates:
            s = c[0]
            if s not in by_slide:
                by_slide[s] = []
            by_slide[s].append(c)

        print(f"共发现 {len(candidates)} 个页码候选，分布在 {len(by_slide)} 页")
        for s in sorted(by_slide.keys())[:15]:
            items = by_slide[s]
            print(f"\n  第{s+1}页（{len(items)} 处）：")
            for c in items:
                print(f"    [{c[1]}] '{c[2]}' @ ({c[3]//10000},{c[4]//10000}) — {c[5]}")
        if len(by_slide) > 15:
            print(f"\n  ...（共 {len(by_slide)} 页，仅显示前 15 页）")
    else:
        print("未检测到 X/Y 格式的页码。")
        print("\n可能的原因：")
        print("  1. 页码位于母版/布局的幻灯片编号占位符中")
        print("  2. 页码格式不是 X/Y，而是其他格式")
        print("  3. 页码是页脚字段（需要通过 PowerPoint 关闭页脚显示）")

    # 2. 检查布局中的幻灯片编号占位符
    print()
    print("=" * 60)
    print("布局中的幻灯片编号占位符")
    print("=" * 60)

    has_placeholder = False
    for li, layout in enumerate(prs.slide_layouts):
        for shape in layout.shapes:
            if shape.is_placeholder and shape.placeholder_format:
                ph = shape.placeholder_format
                if ph.type == 13:  # SLIDE_NUMBER
                    has_placeholder = True
                    text = ''
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                    print(f"  布局 {li}（{layout.name}）：幻灯片编号占位符 [{shape.name}] 文本='{text}'")

    if not has_placeholder:
        print("  未找到幻灯片编号占位符。")

    # 3. 检查母版
    print()
    print("=" * 60)
    print("母版中的页码相关形状")
    print("=" * 60)

    for mi, master in enumerate(prs.slide_masters):
        found = False
        for shape in master.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and ('页' in text or '/' in text or text.isdigit()):
                    found = True
                    print(f"  母版 {mi} [{shape.name}]：'{text[:60]}'")
            if shape.is_placeholder and shape.placeholder_format:
                ph = shape.placeholder_format
                if ph.type == 13:
                    found = True
                    print(f"  母版 {mi}：幻灯片编号占位符 [{shape.name}]")
        if not found:
            print(f"  母版 {mi}：未找到页码相关形状")


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("用法：python detect_pagenum.py <pptx文件路径>")
        sys.exit(1)
    detect_page_numbers(sys.argv[1])
