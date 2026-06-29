#!/usr/bin/env python3
"""
remove_pagenum.py — 删除 PPTX 文件中的页码
用法: python remove_pagenum.py <输入pptx> <输出pptx>
安全策略：先输出到指定路径，不直接覆盖原文件。
"""
import sys
import re
import os

sys.stdout.reconfigure(encoding='utf-8')

def remove_page_numbers(input_path, output_path):
    from pptx import Presentation

    if not os.path.exists(input_path):
        print(f"错误：输入文件不存在 — {input_path}")
        return False

    prs = Presentation(input_path)
    slides = list(prs.slides)
    total = len(slides)

    print(f"处理文件：{os.path.basename(input_path)}")
    print(f"总页数：{total}")
    print()

    W = prs.slide_width
    H = prs.slide_height

    removed_count = 0
    removed_details = []

    # ——————————————————————————
    # Step 1: 删除幻灯片中的 X/Y 格式页码形状
    # ——————————————————————————
    print("【Step 1】删除幻灯片中的页码形状...")

    for i, slide in enumerate(slides):
        shapes_to_remove = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            normalized = text.replace('\n', ' ').replace('  ', ' ').strip()

            # 只删除 X/Y 格式的页码（最安全的判断）
            if re.match(r'^\d{1,3}\s*/\s*\d{1,4}$', normalized):
                shapes_to_remove.append(shape)
                removed_details.append(f"  第{i+1}页：删除 '{normalized}' [{shape.name}]")
                removed_count += 1

        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    if removed_details:
        print(f"  共删除 {removed_count} 个 X/Y 格式页码：")
        for detail in removed_details[:10]:
            print(detail)
        if len(removed_details) > 10:
            print(f"  ...（共 {len(removed_details)} 个）")
    else:
        print("  未检测到 X/Y 格式的页码形状。")

    # ——————————————————————————
    # Step 2: 删除布局中的幻灯片编号占位符
    # ——————————————————————————
    print()
    print("【Step 2】删除布局中的幻灯片编号占位符...")

    layout_removed = 0
    for li, layout in enumerate(prs.slide_layouts):
        shapes_to_remove = []
        for shape in layout.shapes:
            if shape.is_placeholder and shape.placeholder_format:
                ph = shape.placeholder_format
                if ph.type == 13:  # SLIDE_NUMBER
                    shapes_to_remove.append(shape)
                    print(f"  布局 {li}（{layout.name}）：删除占位符 [{shape.name}]")
                    layout_removed += 1

        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    if layout_removed > 0:
        print(f"  共删除 {layout_removed} 个幻灯片编号占位符。")
    else:
        print("  未找到幻灯片编号占位符。")

    # ——————————————————————————
    # Step 3: 保守处理纯数字（仅在右下角）
    # ——————————————————————————
    print()
    print("【Step 3】检查是否有残留的右下角纯数字页码...")

    extra_removed = 0
    for i, slide in enumerate(slides):
        shapes_to_remove = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # 纯数字，在右下角区域
            if text.isdigit() and len(text) <= 3:
                if shape.left > W * 0.55 and shape.top > H * 0.70:
                    shapes_to_remove.append(shape)
                    print(f"  第{i+1}页：删除右下角数字 '{text}' [{shape.name}]")
                    extra_removed += 1
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    if extra_removed > 0:
        print(f"  共删除 {extra_removed} 个右下角纯数字。")
    else:
        print("  未发现残留的右下角纯数字。")

    # ——————————————————————————
    # 保存
    # ——————————————————————————
    print()
    print("【保存】")
    try:
        prs.save(output_path)
        print(f"  ✓ 已保存到：{output_path}")
        print(f"  删除总计：{removed_count + layout_removed + extra_removed} 处")
        return True
    except PermissionError:
        print(f"  ✗ 无法保存到 {output_path}（文件被占用或被只读打开）")
        print("  请关闭 PowerPoint/WPS，或换一个输出路径。")
        return False
    except Exception as e:
        print(f"  ✗ 保存失败：{e}")
        return False


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("用法：python remove_pagenum.py <输入pptx> <输出pptx>")
        sys.exit(1)
    success = remove_page_numbers(sys.argv[1], sys.argv[2])
    sys.exit(0 if success else 1)
