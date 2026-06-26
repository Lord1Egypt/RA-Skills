#!/usr/bin/env python3
"""
DeckCraft v6 — Icon Demo

Showcases several icons from the icon library on a single slide.
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from engine.icons import icon, ICON_NAMES

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Title
from pptx.util import Pt
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "DeckCraft v6 — Icon Library Demo"
p.font.size = Pt(28)
p.font.bold = True

# Demo icons in a grid
demo_icons = [
    ("rocket", (0, 120, 215)),
    ("chart-bar", (46, 139, 87)),
    ("users", (148, 103, 189)),
    ("cog", (255, 152, 0)),
    ("star", (255, 193, 7)),
    ("zap", (233, 30, 99)),
    ("heart", (244, 67, 54)),
    ("search", (0, 150, 136)),
]

cols = 4
icon_size = Inches(0.8)
x_start = Inches(1.5)
y_start = Inches(1.8)
x_gap = Inches(2.8)
y_gap = Inches(2.2)

for idx, (icon_name, color) in enumerate(demo_icons):
    col = idx % cols
    row = idx // cols
    x = x_start + col * x_gap
    y = y_start + row * y_gap

    # Draw icon
    icon(icon_name, slide, x, y, icon_size, color=color, fill_color=color, stroke_width=2.5)

    # Label
    label_box = slide.shapes.add_textbox(
        x - Inches(0.3), y + icon_size + Inches(0.1),
        icon_size + Inches(0.6), Inches(0.4),
    )
    lp = label_box.text_frame.paragraphs[0]
    lp.text = icon_name
    lp.font.size = Pt(11)
    lp.font.color.rgb = RGBColor(100, 100, 100)
    lp.alignment = 1  # center

# Summary text
summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.6))
sp = summary_box.text_frame.paragraphs[0]
sp.text = f"Showing {len(demo_icons)} of {len(ICON_NAMES)} available icons"
sp.font.size = Pt(12)
sp.font.color.rgb = RGBColor(130, 130, 130)

output = "/tmp/deckcraft_icons_demo.pptx"
prs.save(output)
print(f"✓ Icon demo saved to {output}")
print(f"  {len(demo_icons)} icons drawn, {len(ICON_NAMES)} total available")
