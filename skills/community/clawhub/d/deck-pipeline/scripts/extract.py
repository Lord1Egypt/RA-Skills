#!/usr/bin/env python3
"""Paragraph-level text extraction from a .pptx.

Usage:
    python3 extract.py <path-to-pptx> [--json]

Default output: tab-separated (slide, loc, kind, text)
--json: emits a JSON array of {slide, loc, kind, text} objects
"""
import sys
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def collect(path):
    prs = Presentation(path)
    out = []

    def walk(shapes, sidx, prefix=""):
        for i, shp in enumerate(shapes):
            loc = f"{prefix}/{i}"
            try:
                if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(shp.shapes, sidx, loc)
                    continue
            except Exception:
                pass
            if shp.has_text_frame:
                for pi, p in enumerate(shp.text_frame.paragraphs):
                    text = "".join(r.text for r in p.runs)
                    if text.strip():
                        out.append({"slide": sidx, "loc": f"{loc}/p{pi}",
                                    "kind": "text", "text": text})
            if shp.has_table:
                for ri, row in enumerate(shp.table.rows):
                    for ci, cell in enumerate(row.cells):
                        for pi, p in enumerate(cell.text_frame.paragraphs):
                            text = "".join(r.text for r in p.runs)
                            if text.strip():
                                out.append({"slide": sidx,
                                            "loc": f"{loc}/r{ri}c{ci}/p{pi}",
                                            "kind": "table", "text": text})

    for idx, slide in enumerate(prs.slides, 1):
        walk(slide.shapes, idx)
    return out


def main():
    if len(sys.argv) < 2:
        print("Usage: extract.py <pptx-path> [--json]")
        sys.exit(1)
    rows = collect(sys.argv[1])
    if "--json" in sys.argv:
        print(json.dumps(rows, ensure_ascii=False, indent=2))
    else:
        for r in rows:
            print(f"{r['slide']}\t{r['loc']}\t{r['kind']}\t{r['text']}")


if __name__ == "__main__":
    main()
