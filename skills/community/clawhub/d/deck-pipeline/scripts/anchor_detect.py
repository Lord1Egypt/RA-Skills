#!/usr/bin/env python3
"""Structural-anchor detection — find cross-page consistent shapes.

A shape is an anchor candidate if its (position, fill, size-class) signature
recurs on ≥ MIN_PAGES slides. Anchors become the per-page compression
protection list.

Usage:
    python3 anchor_detect.py <pptx-path> [--min-pages N]

Output: JSON list of {signature, page_list, sample_text, count}.
"""
import sys
import json
from collections import defaultdict
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

MIN_PAGES_DEFAULT = 3
POS_BUCKET_EMU = 200_000  # bucket position to ~0.2"


def rgb_hex(color):
    try:
        r, g, b = color.rgb
        return f"#{r:02X}{g:02X}{b:02X}"
    except Exception:
        return None


def size_class(pt):
    if pt is None: return None
    if pt >= 18: return "H1"
    if pt >= 12: return "H2"
    if pt >= 9:  return "body"
    if pt >= 6:  return "caption"
    return "source"


def shape_signature(shp):
    if shp.left is None or shp.top is None:
        return None
    pos = (shp.left // POS_BUCKET_EMU, shp.top // POS_BUCKET_EMU)
    fill = None
    try:
        if str(shp.fill.type) == "FillType.SOLID":
            fill = rgb_hex(shp.fill.fore_color)
    except Exception:
        pass
    sz = None
    sample_text = ""
    if shp.has_text_frame:
        for p in shp.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip() and r.font.size is not None:
                    sz = size_class(r.font.size.pt)
                    break
            if sz: break
        sample_text = "".join(r.text for p in shp.text_frame.paragraphs
                              for r in p.runs).strip()[:50]
    return (pos, fill, sz), sample_text


def main():
    if len(sys.argv) < 2:
        print("Usage: anchor_detect.py <pptx-path> [--min-pages N]")
        sys.exit(1)
    path = sys.argv[1]
    min_pages = MIN_PAGES_DEFAULT
    if "--min-pages" in sys.argv:
        min_pages = int(sys.argv[sys.argv.index("--min-pages") + 1])

    prs = Presentation(path)
    sig_to_pages = defaultdict(set)
    sig_to_sample = {}

    for sidx, slide in enumerate(prs.slides, 1):
        def walk(shapes):
            for shp in shapes:
                try:
                    if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                        walk(shp.shapes)
                        continue
                except Exception:
                    pass
                result = shape_signature(shp)
                if result is None:
                    continue
                sig, sample = result
                sig_to_pages[sig].add(sidx)
                if sig not in sig_to_sample:
                    sig_to_sample[sig] = sample
        walk(slide.shapes)

    anchors = []
    for sig, pages in sig_to_pages.items():
        if len(pages) >= min_pages:
            anchors.append({
                "signature": str(sig),
                "pages": sorted(pages),
                "page_count": len(pages),
                "sample_text": sig_to_sample[sig],
            })
    anchors.sort(key=lambda a: -a["page_count"])

    print(f"=== Anchors (≥{min_pages} pages) ===")
    print(f"Found {len(anchors)} anchor signatures")
    for a in anchors:
        print(f"\n  pages={a['page_count']}  sig={a['signature']}")
        print(f"    sample: '{a['sample_text']}'")
        print(f"    on slides: {a['pages']}")

    # Also emit JSON for programmatic consumption
    print("\n--- JSON ---")
    print(json.dumps(anchors, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
