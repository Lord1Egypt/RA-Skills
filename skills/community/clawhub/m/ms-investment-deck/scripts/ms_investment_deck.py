"""Morgan Stanley 真实研究报告风格 PPTX 生成器。

设计规范摘要
------------
- 颜色：MS NAVY #0B2C5C（主色 / 标题栏 / 章节底色）
        MS GOLD #C8A951（辅助强调 / 关键数字）
        OW GREEN #1F7A3E（Overweight）
        EW ORANGE #D97706（Equal-weight）
        UW RED #B91C1C（Underweight）
- 字体：英文 Calibri Light → Calibri；中文 微软雅黑 → Source Han Sans SC
- 关键数字采用金色放大字号（36-48pt）
- 16:9 widescreen（13.333" x 7.5"）
- 全部 slide 类型独立函数：Cover / Key Takeaways / Section Divider
                             Content (左文右图) / Stat Blocks
                             Rating Table / Shovel Stocks Table
                             Financial Chart (Bar+Line + Revenue Forecast)
                             Executive Summary (4-Quadrant)
                             Scenario Comparison / WACC Breakdown
                             Valuation Bridge / Sensitivity Heatmap
                             Dual Chart Panel / 2x2 Strategic Matrix
                             Asset Allocation Dot Matrix / Donut Chart
                             Stacked Bar Chart
                             Disclosure

主入口
------
``make_deck(data, output_path, theme="classic", language="zh")``

示例
----
``sample_data()`` 会返回一份符合接口的样例字典。
直接运行本脚本：``python morgan_stanley_pptx.py`` 将在 scripts 目录生成
``_ms_sample_zh.pptx``。
"""

from __future__ import annotations

import os
import re
from typing import Any, Dict, List, Optional, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


# =============================================================================
# 1. 颜色与字体
# =============================================================================

# 颜色
NAVY = RGBColor(0x0B, 0x2C, 0x5C)        # MS NAVY 主色
NAVY_ALT = RGBColor(0x1F, 0x38, 0x64)     # 次级深蓝
GOLD = RGBColor(0xC8, 0xA9, 0x51)         # MS GOLD 强调
ORANGE = RGBColor(0xE3, 0x7C, 0x2B)       # 橙色高亮
OW_GREEN = RGBColor(0x1F, 0x7A, 0x3E)     # Overweight
EW_ORANGE = RGBColor(0xD9, 0x77, 0x06)    # Equal-weight
UW_RED = RGBColor(0xB9, 0x1C, 0x1C)       # Underweight
TEXT = RGBColor(0x1A, 0x1A, 0x1A)         # 正文
MUTED = RGBColor(0x5A, 0x5A, 0x5A)        # 辅助说明
TABLE_ALT = RGBColor(0xF5, 0xF6, 0xF8)    # 表格交替行
RULE = RGBColor(0xD4, 0xD4, 0xD4)         # 分隔线
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xFA, 0xFA, 0xFB)

# ---- MS 报告正文颜色（从 MS PDF 学习） ----
MS_BODY_TEXT = RGBColor(0x25, 0x1A, 0x1A)   # 深棕黑色正文
MS_BRAND_BLUE = RGBColor(0x00, 0x55, 0x9F)  # 摩根品牌蓝
MS_CHART_BLUE = RGBColor(0x3B, 0x81, 0xB9)   # 图表蓝
MS_CHART_RED = RGBColor(0xFB, 0x03, 0x01)    # 图表红
MS_GREEN = RGBColor(0x00, 0xAF, 0x50)        # 正面绿
MS_MUTED = RGBColor(0x66, 0x66, 0x66)        # 辅助灰

# ---- DCF 语义化颜色 ----
DEEP_NAVY = RGBColor(0x1F, 0x38, 0x64)     # 深海军蓝 - 主标题背景
MEDIUM_BLUE = RGBColor(0x2E, 0x75, 0xB6)  # 中蓝色 - 副标题/表头背景
INPUT_BG = RGBColor(0xFF, 0xF2, 0xCC)      # 浅黄色 - 输入假设背景
CALC_BG = RGBColor(0xE2, 0xEF, 0xDA)      # 浅绿色 - 关键计算行
SUMMARY_BG = RGBColor(0xD9, 0xE1, 0xF2)   # 浅蓝紫色 - 汇总行
CAPEX_BG = RGBColor(0xFC, 0xE4, 0xD6)     # 浅橙色 - CapEx/ΔNWC行
BEAR_COLOR = RGBColor(0x84, 0x3C, 0x0C)   # 棕橙色 - 熊市标识
BASE_COLOR = RGBColor(0x37, 0x56, 0x23)   # 深绿色 - 基准标识
BULL_COLOR = RGBColor(0x1F, 0x38, 0x64)   # 深海军蓝 - 牛市标识
BEAR_BG = RGBColor(0xFC, 0xE4, 0xD6)       # 浅橙色 - 熊市单元格
BASE_BG = RGBColor(0xE2, 0xEF, 0xDA)       # 浅绿色 - 基准单元格
BULL_BG = RGBColor(0xDD, 0xEE, 0xFF)      # 浅紫色 - 牛市单元格
ALT_ROW = RGBColor(0xF2, 0xF2, 0xF2)      # 浅灰色 - 交替行
HIGHLIGHT_YELLOW = RGBColor(0xFF, 0xFF, 0x00)  # 纯黄色 - 基准交叉点
NEGATIVE_RED = RGBColor(0xC0, 0x00, 0x00) # 深红色 - 负数
MUTED_GRAY = RGBColor(0x59, 0x59, 0x59)   # 灰色 - 说明文字

# 字号（pt）
SIZE_TITLE = 36          # 封面英文主标题
SIZE_TITLE_CN = 22       # 封面中文翻译标题
SIZE_SUB = 14            # 封面副标题
SIZE_KP_NUM = 44         # Key Takeaways 数字
SIZE_KP_TEXT = 18        # Key Takeaways 要点
SIZE_METRIC_NUM = 36     # Stat 数字
SIZE_METRIC_UNIT = 13    # Stat 单位
SIZE_METRIC_LABEL = 11   # Stat 说明
SIZE_BODY = 13           # 正文
SIZE_SMALL = 9           # source / 页脚
SIZE_HDR = 20            # 内容页标题栏中文

# 16:9 幻灯片
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 字体栈（英文 / 数字 / 中文）
FONT_EN_LIGHT = "Arial"
FONT_EN = "Arial"
FONT_EN_FALLBACK = "Helvetica Neue"
FONT_CN = "微软雅黑"
FONT_CN_FALLBACK = "Source Han Sans SC"


# =============================================================================
# 2. 文本 / 字体辅助函数（含 eastAsia 支持）
# =============================================================================

_CJK_RE = re.compile(r"[\u3400-\u9FFF\u3000-\u303F\uFF00-\uFFEF"
                     r"\u20000-\u2FFFF\uF900-\uFAFF]")


def _has_cjk(text: str) -> bool:
    return bool(_CJK_RE.search(text or ""))


def _apply_font(run, *, bold: bool = False, italic: bool = False,
                size: int = 11, color: RGBColor = TEXT,
                cjk: Optional[bool] = None) -> None:
    """为 run 设置英文字体 + eastAsia 中文字体，避免中文显示为默认字体。

    如果文本中包含 CJK 字符，则 typeface 使用英文字体（作为数字/符号基础），
    同时在 rPr 中加入 <a:ea typeface="微软雅黑"/> 使 PowerPoint 使用
    中文 font-family 渲染其中汉字。
    """
    f = run.font
    f.name = FONT_EN
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color

    # 判断是否需要中文字体
    need_cjk = cjk if cjk is not None else _has_cjk(run.text or "")
    if need_cjk:
        rPr = run._r.get_or_add_rPr()
        # eastAsia 字体
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", FONT_CN)
        # latin 字体（保留数字/英文为 Calibri）
        latin = rPr.find(qn("a:latin"))
        if latin is None:
            latin = rPr.makeelement(qn("a:latin"), {})
            rPr.append(latin)
        latin.set("typeface", FONT_EN)


# =============================================================================
# 3. 形状与工具函数
# =============================================================================

def _add_rect(slide, left: int, top: int, width: int, height: int,
              fill: RGBColor, line: Optional[RGBColor] = None) -> object:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if line is None else None
    if line is not None:
        shape.line.color.rgb = line
    return shape


def _add_line(slide, left: int, top: int, width: int, height: int,
              color: RGBColor) -> object:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left: int, top: int, width: int, height: int,
                 anchor: str = "t") -> object:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if anchor == "m":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "b":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    return tb


def _write_text(tb, text: str, *, align: int = PP_ALIGN.LEFT,
                size: int = 11, bold: bool = False,
                italic: bool = False, color: RGBColor = TEXT,
                cjk: Optional[bool] = None,
                paragraph_spacing_after: Optional[int] = None) -> None:
    """向 textbox 写入文本（支持中英文混排 + eastAsia 字体）。"""
    tf = tb.text_frame
    # 对于多段文本支持换行符 \n 作为分段
    lines = str(text).split("\n") if text else [""]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        if paragraph_spacing_after is not None:
            p.space_after = Pt(paragraph_spacing_after)
        r = p.add_run()
        r.text = line
        _apply_font(r, bold=bold, italic=italic, size=size, color=color, cjk=cjk)


def _add_text(slide, left: int, top: int, width: int, height: int,
              text: str, *, size: int = 11, bold: bool = False,
              italic: bool = False, color: RGBColor = TEXT,
              align: int = PP_ALIGN.LEFT, anchor: str = "t",
              cjk: Optional[bool] = None,
              paragraph_spacing_after: Optional[int] = None) -> object:
    tb = _add_textbox(slide, left, top, width, height, anchor=anchor)
    _write_text(tb, text, align=align, size=size, bold=bold, italic=italic,
                color=color, cjk=cjk,
                paragraph_spacing_after=paragraph_spacing_after)
    return tb


def _rating_color(label: str) -> RGBColor:
    s = (label or "").strip().upper()
    if "OVERWEIGHT" in s or s.startswith("OW") or "买入" in label:
        return OW_GREEN
    if "EQUAL" in s or s.startswith("EW") or "持有" in label:
        return EW_ORANGE
    if "UNDERWEIGHT" in s or s.startswith("UW") or "卖出" in label:
        return UW_RED
    return NAVY


def _rating_label_normalize(label: str) -> str:
    s = (label or "").strip()
    upper = s.upper()
    if "OVERWEIGHT" in upper or upper == "OW":
        return "Overweight"
    if "EQUAL" in upper or upper == "EW":
        return "Equal-weight"
    if "UNDERWEIGHT" in upper or upper == "UW":
        return "Underweight"
    return s


def _style_chart_axis(axis, tick_color='#4A4A4A', grid_color='#E0E0E0'):
    """统一设置坐标轴样式。

    Parameters
    ----------
    axis : pptx.chart.axis.CategoryAxis or ValueAxis
        要设置样式的坐标轴。
    tick_color : str
        刻度标签和轴线颜色，默认 #4A4A4A。
    grid_color : str
        主网格线颜色，默认 #E0E0E0。
    """
    axis.has_major_gridlines = True
    axis.major_gridlines.format.line.color.rgb = RGBColor.from_string(
        grid_color.replace('#', ''))
    axis.major_gridlines.format.line.width = Pt(0.5)
    axis.format.line.color.rgb = RGBColor.from_string(
        tick_color.replace('#', ''))
    axis.format.line.width = Pt(0.5)
    axis.tick_labels.font.size = Pt(8)
    axis.tick_labels.font.color.rgb = RGBColor.from_string(
        tick_color.replace('#', ''))


# =============================================================================
# 4. Slide 01 — Cover 封面页
# =============================================================================

def slide_cover(slide, *, title_cn: str = "", title_en: str = "",
                subtitle: str = "", rating: str = "",
                target_price: Optional[Union[str, float]] = None,
                current_price: Optional[Union[str, float]] = None,
                company_name: str = "", analyst: str = "",
                date_str: str = "", research_type: str = "Foundation",
                industry_view: str = "",
                whats_changed: Optional[List[Dict[str, str]]] = None,
                language: str = "zh") -> None:
    """封面页：顶部 NAVY 横条 + 金色大数字 / 标题 + 底部金色粗线。

    新增参数（从 MS PDF 学习）：
    - industry_view: 行业观点标签 "Attractive" / "In-Line" / "Cautious"
    - whats_changed: 变动列表 [{"item": "...", "from": "...", "to": "..."}]
    """
    # 背景
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    # 顶部 NAVY 横条 1.2"
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), NAVY)
    # 金色左竖条（装饰）
    _add_rect(slide, 0, 0, Inches(0.12), Inches(1.2), GOLD)

    # 左侧金色 M logo（用 "MS" 字样代替）
    _add_text(slide, Inches(0.6), Inches(0.32), Inches(1.6), Inches(0.7),
              "MS", size=44, bold=True, color=GOLD, anchor="m", cjk=False)

    # 右侧分类（Foundation / Global Insight 等）
    if research_type:
        _add_text(slide, SLIDE_W - Inches(5.5), Inches(0.4),
                  Inches(4.8), Inches(0.45),
                  str(research_type).upper(),
                  size=12, bold=True, color=GOLD, align=PP_ALIGN.RIGHT,
                  cjk=False)
    # 分析师 / 日期
    if analyst or date_str:
        sub_line = []
        if analyst:
            sub_line.append(str(analyst))
        if date_str:
            sub_line.append(str(date_str))
        _add_text(slide, SLIDE_W - Inches(5.5), Inches(0.8),
                  Inches(4.8), Inches(0.35),
                  " · ".join(sub_line),
                  size=10, color=WHITE, align=PP_ALIGN.RIGHT,
                  italic=True, cjk=True)

    # ---- 行业观点标签（Attractive / In-Line / Cautious） ----
    if industry_view:
        iv = str(industry_view).strip().lower()
        if iv in ("attractive", "有吸引力"):
            tag_color = MS_GREEN
            tag_text = "Attractive"
        elif iv in ("in-line", "inline", "中性"):
            tag_color = RGBColor(0xD4, 0x9B, 0x00)  # 黄色
            tag_text = "In-Line"
        elif iv in ("cautious", "谨慎"):
            tag_color = MS_CHART_RED
            tag_text = "Cautious"
        else:
            tag_color = MS_BRAND_BLUE
            tag_text = str(industry_view)
        tag_w = Inches(1.6)
        tag_h = Inches(0.35)
        tag_left = SLIDE_W - Inches(0.5) - tag_w
        tag_top = Inches(1.45)
        _add_rect(slide, tag_left, tag_top, tag_w, tag_h, tag_color)
        _add_text(slide, tag_left, tag_top, tag_w, tag_h,
                  tag_text, size=11, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor="m", cjk=False)

    # 中部英文主标题（36pt NAVY bold）
    en_title = str(title_en or title_cn or "Morgan Stanley Research")
    _add_text(slide, Inches(0.8), Inches(2.3), SLIDE_W - Inches(1.6),
              Inches(1.2), en_title,
              size=SIZE_TITLE, bold=True, color=NAVY, anchor="m",
              cjk=_has_cjk(en_title))

    # 中文翻译标题（22pt）
    if title_cn and str(title_cn).strip() != str(title_en).strip():
        _add_text(slide, Inches(0.8), Inches(3.45), SLIDE_W - Inches(1.6),
                  Inches(0.7), str(title_cn),
                  size=SIZE_TITLE_CN, color=MUTED, anchor="m",
                  cjk=True)

    # 副标题一行核心观点（italic 14pt MUTED）
    if subtitle:
        _add_text(slide, Inches(0.8), Inches(4.3), SLIDE_W - Inches(1.6),
                  Inches(0.6), str(subtitle),
                  size=SIZE_SUB, italic=True, color=MUTED,
                  align=PP_ALIGN.CENTER, anchor="m",
                  cjk=True)

    # 评级 / 目标价 —— 金色关键数字
    r_line_parts = []
    if rating:
        r_line_parts.append(_rating_label_normalize(str(rating)))
    if target_price is not None and current_price is not None:
        r_line_parts.append(f"TP {target_price}  ·  Last {current_price}")
    elif target_price is not None:
        r_line_parts.append(f"TP {target_price}")

    if r_line_parts:
        _add_text(slide, Inches(0.8), Inches(5.2), SLIDE_W - Inches(1.6),
                  Inches(1.0), "   |   ".join(r_line_parts),
                  size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
                  anchor="m", cjk=False)

    # ---- What's Changed 区域（封面底部变动摘要） ----
    if whats_changed and len(whats_changed) > 0:
        wc_top = SLIDE_H - Inches(1.55)
        # 小标题
        wc_label = "What's Changed" if language != "zh" else "本期变动"
        _add_text(slide, Inches(0.5), wc_top - Inches(0.3),
                  Inches(2.0), Inches(0.25),
                  wc_label, size=10, bold=True, color=MS_BRAND_BLUE,
                  cjk=(language == "zh"))
        # 表格区域
        wc_items = list(whats_changed)[:4]  # 最多显示4条
        wc_cols = 3
        wc_rows = len(wc_items) + 1
        wc_table_w = SLIDE_W - Inches(1.0)
        wc_row_h = Inches(0.22)
        ts_wc = slide.shapes.add_table(
            wc_rows, wc_cols,
            Inches(0.5), wc_top,
            wc_table_w, wc_row_h * wc_rows)
        tbl_wc = ts_wc.table
        tbl_wc.columns[0].width = int(wc_table_w * 0.45)
        tbl_wc.columns[1].width = int(wc_table_w * 0.25)
        tbl_wc.columns[2].width = int(wc_table_w * 0.30)
        # 表头
        hdr_labels = ["Item / 项目", "From", "To"]
        for j, hl in enumerate(hdr_labels):
            _fill_table_cell(tbl_wc.cell(0, j), hl,
                             size=8, bold=True, color=WHITE,
                             bg=MS_BRAND_BLUE, align=PP_ALIGN.CENTER, cjk=True)
        # 数据行
        for i, wc in enumerate(wc_items):
            bg = WHITE if i % 2 == 0 else TABLE_ALT
            _fill_table_cell(tbl_wc.cell(i + 1, 0),
                             str(wc.get("item", wc.get("metric", wc.get("field", "")))),
                             size=8, color=MS_BODY_TEXT, bg=bg,
                             align=PP_ALIGN.LEFT, cjk=True)
            _fill_table_cell(tbl_wc.cell(i + 1, 1),
                             str(wc.get("from", wc.get("old_value", wc.get("previous", "")))),
                             size=8, color=MS_MUTED, bg=bg,
                             align=PP_ALIGN.CENTER, cjk=True)
            # "To" 列用颜色标识变动方向
            to_val = str(wc.get("to", wc.get("new_value", wc.get("current", ""))))
            arrow = wc.get("arrow", wc.get("direction", ""))
            to_color = MS_BODY_TEXT
            if arrow == "up":
                to_color = MS_GREEN
            elif arrow == "down":
                to_color = MS_CHART_RED
            _fill_table_cell(tbl_wc.cell(i + 1, 2),
                             to_val,
                             size=8, bold=True, color=to_color, bg=bg,
                             align=PP_ALIGN.CENTER, cjk=True)

    # 底部金色粗线
    _add_rect(slide, Inches(0.5), SLIDE_H - Inches(0.65),
              SLIDE_W - Inches(1.0), Emu(28000), GOLD)
    # 公司品牌线
    brand_text = company_name or "Morgan Stanley Research"
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.45),
              SLIDE_W - Inches(1.0), Inches(0.35),
              brand_text, size=10, bold=True, color=NAVY,
              align=PP_ALIGN.LEFT, cjk=True)
    if analyst:
        _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.45),
                  SLIDE_W - Inches(1.0), Inches(0.35),
                  f"Equity Research · {analyst}",
                  size=10, color=MUTED, align=PP_ALIGN.RIGHT,
                  italic=True, cjk=True)


# =============================================================================
# 5. Slide 02 — Key Takeaways 核心要点页
# =============================================================================

def _slide_top_bar(slide, *, title_cn: str = "", title_en: str = "",
                   bar_height_in: float = 0.9,
                   show_rule: bool = True,
                   page_number: Optional[int] = None,
                   total_pages: Optional[int] = None) -> None:
    """通用：顶部 NAVY 横条 + 左侧金竖条 + 中文主标题 / 英文副标题。

    可选参数：page_number / total_pages —— 在右上角金色小字显示页码，
    形如 "6 \u00b7 15"。两者必须同时提供才会显示。
    """
    _add_rect(slide, 0, 0, SLIDE_W, Inches(bar_height_in), NAVY)
    _add_rect(slide, 0, 0, Inches(0.12), Inches(bar_height_in), GOLD)

    # 右上角页码（若提供）
    if page_number is not None and total_pages is not None:
        _add_text(slide, SLIDE_W - Inches(1.5), Inches(0.15),
                  Inches(1.2), Inches(bar_height_in * 0.5),
                  f"{int(page_number)} \u00b7 {int(total_pages)}",
                  size=10, bold=True, color=GOLD, align=PP_ALIGN.RIGHT,
                  anchor="m", cjk=False)

    title_w = (SLIDE_W - Inches(1.5)
               if (page_number is not None and total_pages is not None)
               else SLIDE_W - Inches(1.0))

    if title_cn:
        _add_text(slide, Inches(0.55), Inches(0.12),
                  title_w - Inches(0.55),
                  Inches(bar_height_in * 0.55),
                  str(title_cn),
                  size=22, bold=True, color=WHITE, anchor="m", cjk=True)
    if title_en:
        _add_text(slide, Inches(0.55), Inches(bar_height_in - 0.42),
                  title_w - Inches(0.55), Inches(0.35),
                  str(title_en),
                  size=11, italic=True, color=GOLD, anchor="m",
                  cjk=False)

    if show_rule:
        _add_rect(slide, Inches(0.5), Inches(bar_height_in + 0.1),
                  SLIDE_W - Inches(1.0), Emu(10000), RULE)


def slide_key_takeaways(slide, *, items: List[str],
                        title_cn: str = "核心观点",
                        title_en: str = "Key Takeaways",
                        language: str = "zh",
                  page_number: Optional[int] = None,
                  total_pages: Optional[int] = None) -> None:
    """Key Takeaways：每条 = 金色大号两位数字 + 右侧要点句。"""
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    items = list(items or [])
    if not items:
        items = ["（未提供要点）"]

    # 限制每页 6 条
    items = items[:6]
    n = len(items)
    top_start = Inches(1.3)
    total_h = SLIDE_H - top_start - Inches(0.6)
    per_h = int(total_h / max(n, 1))

    for i, text in enumerate(items):
        top = top_start + per_h * i
        # 左侧金色两位数字
        _add_text(slide, Inches(0.5), top, Inches(1.6), Inches(per_h),
                  f"{i + 1:02d}",
                  size=44, bold=True, color=GOLD, align=PP_ALIGN.LEFT,
                  anchor="t", cjk=False)
        # 右侧要点句
        _add_text(slide, Inches(2.3), top + Inches(0.08),
                  SLIDE_W - Inches(2.9), Inches(per_h - 0.1),
                  str(text),
                  size=18, color=TEXT, anchor="t", cjk=True)
        # 底部细线（非最后一条）
        if i < n - 1:
            _add_rect(slide, Inches(0.5), top + per_h - Emu(10000),
                      SLIDE_W - Inches(1.0), Emu(6000), RULE)

    # 底部 rule
    _add_rect(slide, Inches(0.5), SLIDE_H - Inches(0.35),
              SLIDE_W - Inches(1.0), Emu(6000), GOLD)
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.30),
              SLIDE_W - Inches(1.0), Inches(0.25),
              "Source: Morgan Stanley Research",
              size=8, italic=True, color=MUTED, align=PP_ALIGN.RIGHT,
              cjk=False)


# =============================================================================
# 6. Slide 03 — Section Divider 章节分隔页
# =============================================================================

def slide_section_divider(slide, *, title_cn: str = "章节标题",
                          title_en: str = "Section Title",
                          category: str = "Foundation",
                          language: str = "zh") -> None:
    """全屏 NAVY + 金色横条 + 中文大标题 / 英文小标题。"""
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # 左上金色 MS logo
    _add_text(slide, Inches(0.8), Inches(0.7), Inches(2.5), Inches(0.8),
              "MS", size=40, bold=True, color=GOLD, anchor="m", cjk=False)
    # 右上金色大写分类
    if category:
        _add_text(slide, SLIDE_W - Inches(5.5), Inches(0.8),
                  Inches(4.5), Inches(0.5),
                  str(category).upper(),
                  size=12, bold=True, color=GOLD, align=PP_ALIGN.RIGHT,
                  cjk=False)
    # 中部金色横线（4pt 粗，1.2" 长）
    _add_rect(slide, Inches(0.8), Inches(3.35), Inches(1.2),
              Emu(28000), GOLD)
    # 中文大标题（40pt white）
    _add_text(slide, Inches(0.8), Inches(3.75), SLIDE_W - Inches(1.6),
              Inches(1.2), str(title_cn or title_en),
              size=40, bold=True, color=WHITE, anchor="t", cjk=True)
    # 英文小标题（18pt light gold italic）
    if title_en and str(title_en) != str(title_cn):
        _add_text(slide, Inches(0.8), Inches(5.05),
                  SLIDE_W - Inches(1.6), Inches(0.6), str(title_en),
                  size=18, italic=True, color=GOLD, anchor="t",
                  cjk=False)


# =============================================================================
# 7. Slide 04 — Content Page 标准分析页（左文右图）
# =============================================================================

def _write_bullet_body(slide, left, top, width, height,
                         body: List[str]) -> None:
    """向 slide 写入双层 bullet 正文；以 "- " 开头的视为第二层；空行产生间距。"""
    tb = _add_textbox(slide, left, top, width, height, anchor="t")
    tf = tb.text_frame
    tf.word_wrap = True

    lines = [str(x) for x in (body or [])]
    first = True
    for raw in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        stripped = raw.strip()
        is_blank = not stripped
        is_level2 = stripped.startswith(("- ", "\u2013 ", "-- "))
        if is_blank:
            r = p.add_run()
            r.text = ""
            _apply_font(r, size=SIZE_BODY, color=TEXT, cjk=True)
            p.space_after = Pt(4)
            continue
        if is_level2:
            p.space_before = Pt(2)
            p.level = 1
            text_body = stripped[2:].strip() if len(stripped) > 2 else ""
            r = p.add_run()
            r.text = "\u2013  " + text_body
            _apply_font(r, size=11, color=MUTED, italic=True, cjk=True)
        else:
            r1 = p.add_run()
            r1.text = "\u2013  "
            _apply_font(r1, size=SIZE_BODY, bold=True, color=NAVY, cjk=False)
            r2 = p.add_run()
            r2.text = stripped
            _apply_font(r2, size=SIZE_BODY, color=TEXT, cjk=True)
        p.space_after = Pt(6)


def _add_default_pie_chart(slide, left, top, width, height,
                           sectors=None):
    """在 slide 上插入一个原生饼图（行业配置）。"""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import (XL_CHART_TYPE, XL_LEGEND_POSITION,
                                  XL_LABEL_POSITION)

    _add_rect(slide, left, top, width, Inches(0.4), TABLE_ALT)
    _add_text(slide, left + Inches(0.1), top, width - Inches(0.2),
              Inches(0.4),
              "Exhibit: Sector Allocation (%)",
              size=10, bold=True, color=NAVY, anchor="m", cjk=False)

    chart_left2 = left
    chart_top2 = top + Inches(0.45)
    chart_h = height - Inches(0.45)

    default_sectors = {
        "Semiconductors": 32.0, "Power & Utilities": 22.0,
        "Optics / Networking": 15.0, "AI Software": 13.0,
        "Memory / Storage": 10.0, "Other": 8.0,
    }
    sectors = sectors or default_sectors
    # 兼容两种格式：{"label": value, ...} 或 {"labels": [...], "values": [...]}
    if isinstance(sectors, dict) and "labels" in sectors and "values" in sectors:
        names = [str(x) for x in sectors["labels"]]
        values = [float(v) for v in sectors["values"]]
    else:
        names = list(sectors.keys())
        values = [float(v) for v in sectors.values()]
    total = sum(values) or 1.0
    values = [v * 100.0 / total for v in values]

    chart_data = CategoryChartData()
    chart_data.categories = names
    chart_data.add_series("Sector", values)

    gf = slide.shapes.add_chart(XL_CHART_TYPE.PIE,
                                 chart_left2, chart_top2,
                                 width, chart_h, chart_data)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.horz_offset = 0
    chart.legend.font.size = Pt(8)
    chart.legend.font.color.rgb = MUTED

    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = '0.0"%"'
    dl.show_percentage = True
    dl.show_category_name = True
    dl.show_legend_key = False
    dl.font.size = Pt(9)
    dl.font.bold = True
    dl.font.color.rgb = TEXT
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.separator = "\n"

    # MS 8色专业调色板
    palette = [
        RGBColor(0x1F, 0x38, 0x64),  # 深海军蓝
        RGBColor(0x2E, 0x75, 0xB6),  # 中蓝
        RGBColor(0xC8, 0xA9, 0x51),  # MS GOLD
        RGBColor(0x1F, 0x7A, 0x3E),  # OW GREEN
        RGBColor(0xE3, 0x7C, 0x2B),  # ORANGE
        RGBColor(0xB9, 0x1C, 0x1C),  # UW RED
        RGBColor(0x5A, 0x5A, 0x5A),  # MUTED
        RGBColor(0x4A, 0x6F, 0xA8),  # 浅蓝
    ]
    series = plot.series[0]
    n_points = len(series.points)
    for i in range(n_points):
        pt = series.points[i]
        fill = pt.format.fill
        fill.solid()
        fill.fore_color.rgb = palette[i % len(palette)]
        # 白色分隔线，linewidth=1.5pt
        line = pt.format.line
        line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        line.width = Pt(1.5)

    # 第一扇区（最大）起始角度优化：从顶部开始（-90度）
    try:
        plot.first_slice_angle = 0  # 0 = 3点钟方向, 90=12点钟方向
    except Exception:
        pass


def slide_content(slide, *, title_cn: str = "", title_en: str = "",
                  body: Optional[List[str]] = None,
                  chart_img: Optional[str] = None,
                  data: Optional[Dict[str, Any]] = None,
                  source: str = "Source: Morgan Stanley Research",
                  language: str = "zh",
                  page_number: Optional[int] = None,
                  total_pages: Optional[int] = None) -> None:
    """左文右图标准分析页：顶部 NAVY 标题栏 + 左侧双层 bullet + 右侧图表。

    body 格式：
      - 普通字符串 —— 视为第一层 bullet
      - 以 "- " 开头的字符串 —— 视为第二层（缩进 + 斜体 muted）
    右侧图表：优先使用 chart_img 图片；若无图则尝试从 ``data.sectors_allocation``
    生成原生饼图；否则退回占位卡片。
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   bar_height_in=0.7,
                   page_number=page_number, total_pages=total_pages)

    left_x = Inches(0.5)
    top_y = Inches(1.0)
    body_w = Inches(5.6)
    body_h = Inches(5.8)

    body = list(body or [])
    if body:
        _write_bullet_body(slide, left_x, top_y, body_w, body_h, body)

    chart_left = Inches(6.2)
    chart_top = Inches(1.0)
    chart_w = SLIDE_W - chart_left - Inches(0.3)
    chart_h = Inches(5.4)

    if chart_img and os.path.isfile(str(chart_img)):
        slide.shapes.add_picture(str(chart_img), chart_left, chart_top,
                                 width=chart_w)
    else:
        sectors = None
        if isinstance(data, dict):
            sectors = data.get("sectors_allocation")
        if isinstance(sectors, dict) and sectors:
            _add_default_pie_chart(slide, chart_left, chart_top,
                                   chart_w, chart_h, sectors)
        else:
            _add_rect(slide, chart_left, chart_top, chart_w, chart_h,
                      TABLE_ALT)
            _add_rect(slide, chart_left, chart_top,
                      Inches(0.06), chart_h, GOLD)
            _add_text(slide, chart_left + Inches(0.25),
                      chart_top + Inches(0.2), chart_w - Inches(0.5),
                      Inches(0.4), "Exhibit",
                      size=10, bold=True, color=GOLD, cjk=False)
            _add_text(slide, chart_left, chart_top + Inches(2.2), chart_w,
                      Inches(0.5), "[ Chart / Exhibit placeholder ]",
                      size=14, italic=True, color=MUTED,
                      align=PP_ALIGN.CENTER, cjk=False)
            _add_text(slide, chart_left, chart_top + Inches(2.7), chart_w,
                      Inches(0.5),
                      "请将 chart_img 指向 PNG/JPG 文件，或使用原生 chart。",
                      size=10, color=MUTED, align=PP_ALIGN.CENTER,
                      cjk=True)

    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.4),
              SLIDE_W - Inches(1.0), Inches(0.3),
              str(source), size=SIZE_SMALL, italic=True, color=MUTED,
              align=PP_ALIGN.LEFT, cjk=True)


# =============================================================================
# 8. Slide 05 — Stat / Metric Blocks 关键数字页
# =============================================================================

def slide_metric_blocks(slide, *, metrics: List[Dict[str, Any]],
                        title_cn: str = "关键指标",
                        title_en: str = "Key Metrics",
                        language: str = "zh",
                  page_number: Optional[int] = None,
                  total_pages: Optional[int] = None) -> None:
    """2xN 信息块栅格：左侧金色竖线 + 白色卡片 + 金色数字 + 深蓝副标题。"""
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    metrics = list(metrics or []) or [
        {"value": "—", "unit": "Metric", "label": "（未提供）"},
    ]
    # 限制最多 6 块（2x3）
    metrics = metrics[:6]
    n = len(metrics)

    # 决定布局：<=3 时 1x3；>3 且 <=6 时 2x3
    if n <= 3:
        rows, cols = 1, n
    elif n <= 4:
        rows, cols = 2, 2
    else:
        rows, cols = 2, 3

    area_top = Inches(1.45)
    area_left = Inches(0.5)
    area_w = SLIDE_W - Inches(1.0)
    area_h = SLIDE_H - area_top - Inches(0.7)
    gap_x = Inches(0.35)
    gap_y = Inches(0.35)
    card_w = int((area_w - gap_x * (cols - 1)) / cols)
    card_h = int((area_h - gap_y * (rows - 1)) / rows)

    for idx, m in enumerate(metrics):
        r = idx // cols
        c = idx % cols
        left = area_left + (card_w + gap_x) * c
        top = area_top + (card_h + gap_y) * r

        # 白色卡片（1px 灰边）
        _add_rect(slide, left, top, card_w, card_h, WHITE, line=RULE)
        # 左侧 5pt 金色竖线
        _add_rect(slide, left, top, Inches(0.06), card_h, GOLD)

        # 大号金色数字
        value = str(m.get("value", m.get("num", "—")))
        _add_text(slide, left + Inches(0.4), top + Inches(0.35),
                  card_w - Inches(0.6), Inches(1.3), value,
                  size=SIZE_METRIC_NUM, bold=True, color=GOLD,
                  anchor="t", cjk=False)
        # 深蓝色副标题单位
        unit = str(m.get("unit", ""))
        if unit:
            _add_text(slide, left + Inches(0.4), top + card_h - Inches(1.3),
                      card_w - Inches(0.6), Inches(0.4), unit,
                      size=SIZE_METRIC_UNIT, bold=True, color=NAVY,
                      anchor="t", cjk=True)
        # 灰黑色说明
        label = str(m.get("label", ""))
        if label:
            _add_text(slide, left + Inches(0.4),
                      top + card_h - Inches(0.85),
                      card_w - Inches(0.6), Inches(0.65), label,
                      size=SIZE_METRIC_LABEL, color=MUTED,
                      italic=True, anchor="t", cjk=True)


# =============================================================================
# 9. Slide 06 — Rating Table 评级总表页
# =============================================================================

def _fill_table_cell(cell, text: str, *, size: int = 11,
                     bold: bool = False, color: RGBColor = TEXT,
                     bg: Optional[RGBColor] = None,
                     align: int = PP_ALIGN.LEFT,
                     cjk: Optional[bool] = None) -> None:
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    cell.margin_left = Inches(0.1)
    cell.margin_right = Inches(0.1)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    _apply_font(r, size=size, bold=bold, color=color, cjk=cjk)


def slide_rating_table(slide, *, rows: List[List[Any]],
                       title_cn: str = "投资评级总览",
                       title_en: str = "Ratings & Targets",
                       headers: Optional[List[str]] = None,
                       language: str = "zh",
                  page_number: Optional[int] = None,
                  total_pages: Optional[int] = None) -> None:
    """5 列表格：Ticker | Rating | Last Close | Company | Reason。"""
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    if headers is None:
        headers = ["Ticker", "Rating", "Last Close", "Company", "Reason / 观点"]

    rows = [list(r) for r in (rows or [])]
    # 填充到 5 列
    for r in rows:
        while len(r) < 5:
            r.append("")

    table_left = Inches(0.5)
    table_top = Inches(1.25)
    table_w = SLIDE_W - Inches(1.0)
    n_rows = len(rows) + 1
    row_h = Inches(0.45)
    table_h = row_h * n_rows

    table_shape = slide.shapes.add_table(n_rows, 5, table_left, table_top,
                                         table_w, table_h)
    table = table_shape.table
    # 列宽：Ticker / Rating / Last Close / Company / Reason
    col_w_pct = [0.16, 0.16, 0.16, 0.26, 0.26]
    for i, pct in enumerate(col_w_pct):
        table.columns[i].width = int(table_w * pct)

    # 表头
    for j, h in enumerate(headers):
        _fill_table_cell(table.cell(0, j), str(h),
                         size=11, bold=True, color=WHITE, bg=NAVY,
                         align=PP_ALIGN.CENTER if j in (1, 2) else PP_ALIGN.LEFT,
                         cjk=True)

    # 数据行
    for i, row in enumerate(rows):
        bg = WHITE if i % 2 == 0 else TABLE_ALT
        ticker, rating, last, company, reason = row[0], row[1], row[2], row[3], row[4]
        _fill_table_cell(table.cell(i + 1, 0), str(ticker),
                         size=11, bold=True, color=NAVY, bg=bg,
                         align=PP_ALIGN.LEFT, cjk=False)
        rating_text = _rating_label_normalize(str(rating))
        rc = _rating_color(str(rating))
        _fill_table_cell(table.cell(i + 1, 1), rating_text,
                         size=11, bold=True, color=rc, bg=bg,
                         align=PP_ALIGN.CENTER, cjk=True)
        _fill_table_cell(table.cell(i + 1, 2), str(last),
                         size=11, color=TEXT, bg=bg,
                         align=PP_ALIGN.CENTER, cjk=False)
        _fill_table_cell(table.cell(i + 1, 3), str(company),
                         size=11, color=TEXT, bg=bg,
                         align=PP_ALIGN.LEFT, cjk=True)
        _fill_table_cell(table.cell(i + 1, 4), str(reason),
                         size=10, color=MUTED, bg=bg,
                         align=PP_ALIGN.LEFT, cjk=True)

    # 底部 source
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.4),
              SLIDE_W - Inches(1.0), Inches(0.3),
              "Source: Morgan Stanley Research · Prices as of most recent close",
              size=SIZE_SMALL, italic=True, color=MUTED,
              align=PP_ALIGN.LEFT, cjk=False)


# =============================================================================
# 10. Slide NEW A — Value Chain 产业链分析页
# =============================================================================

def slide_value_chain(slide, *, data: Dict[str, Any],
                     title_cn: str = "产业链分析",
                     title_en: str = "Value Chain Analysis",
                     language: str = "zh",
                     page_number: Optional[int] = None,
                     total_pages: Optional[int] = None) -> None:
    """产业链/价值链分析：按价值链环节分行展示，每个环节一个色块内含公司列表。

    数据格式：
        data["value_chain"] = [
            {"stage_en": "Brain / AI Algorithms", "stage_cn": "大脑 / AI算法",
             "companies": ["OpenAI", "Google DeepMind", ...]},
            {"stage_en": "Sensors", "stage_cn": "传感器",
             "companies": ["Cognex", "Keyence", ...]},
            ...
        ]
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    chain = list(data.get("value_chain") or [])
    if not chain:
        _add_text(slide, Inches(0.5), Inches(2.5),
                  SLIDE_W - Inches(1.0), Inches(1.0),
                  "（未提供价值链数据）" if language == "zh"
                  else "（No value chain data provided）",
                  size=14, italic=True, color=MUTED,
                  align=PP_ALIGN.CENTER, cjk=True)
        return

    # 阶段颜色调色板（循环使用）
    stage_colors = [
        MS_BRAND_BLUE, MS_CHART_BLUE, MS_GREEN,
        RGBColor(0x6B, 0x4C, 0x9A), RGBColor(0xD4, 0x9B, 0x00),
        RGBColor(0xE3, 0x7C, 0x2B), NAVY, MS_CHART_RED,
    ]

    n_stages = len(chain)
    # 布局：水平排列色块，每个色块宽度自适应
    area_left = Inches(0.5)
    area_top = Inches(1.5)
    area_w = SLIDE_W - Inches(1.0)
    area_h = SLIDE_H - area_top - Inches(1.0)
    gap = Inches(0.15)
    total_gap = gap * (n_stages - 1)
    card_w = int((area_w - total_gap) / n_stages)
    card_w = min(card_w, Inches(2.8))  # 限制最大宽度

    # 如果卡片太多放不下，分两行
    max_per_row = int(area_w / (Inches(1.5) + gap))
    rows_of_stages = []
    for i in range(0, n_stages, max_per_row):
        rows_of_stages.append(chain[i:i + max_per_row])

    row_gap = Inches(0.2)
    total_rows = len(rows_of_stages)
    card_h = int((area_h - row_gap * (total_rows - 1)) / total_rows)

    for row_idx, row_stages in enumerate(rows_of_stages):
        n_in_row = len(row_stages)
        row_total_gap = gap * (n_in_row - 1)
        row_card_w = int((area_w - row_total_gap) / n_in_row)
        row_card_w = min(row_card_w, Inches(2.8))
        row_total_w = row_card_w * n_in_row + row_total_gap
        row_offset = int((area_w - row_total_w) / 2)  # 居中

        for col_idx, stage in enumerate(row_stages):
            left = area_left + row_offset + (row_card_w + gap) * col_idx
            top = area_top + (card_h + row_gap) * row_idx
            color = stage_colors[(row_idx * max_per_row + col_idx) % len(stage_colors)]

            # 色块背景
            _add_rect(slide, left, top, row_card_w, card_h, color)

            # 阶段名称（英文 + 中文）
            stage_en = str(stage.get("stage_en", ""))
            stage_cn = str(stage.get("stage_cn", ""))
            header_h = Inches(0.55)
            _add_text(slide, left + Inches(0.1), top + Inches(0.05),
                      row_card_w - Inches(0.2), header_h,
                      stage_en,
                      size=11, bold=True, color=WHITE, anchor="m",
                      cjk=False)
            if stage_cn and stage_cn != stage_en:
                _add_text(slide, left + Inches(0.1), top + header_h,
                          row_card_w - Inches(0.2), Inches(0.3),
                          stage_cn,
                          size=9, color=RGBColor(0xE0, 0xE0, 0xE0),
                          cjk=True)

            # 公司列表
            companies = list(stage.get("companies") or [])
            list_top = top + header_h + Inches(0.3)
            list_h = card_h - header_h - Inches(0.4)
            if companies:
                # 兼容字符串列表和字典列表两种格式
                formatted = []
                for c in companies[:8]:
                    if isinstance(c, dict):
                        name = c.get("name", c.get("company", ""))
                        ticker = c.get("ticker", "")
                        note = c.get("note", c.get("description", ""))
                        if ticker:
                            formatted.append(f"  {name} ({ticker})")
                        elif note:
                            formatted.append(f"  {name} - {note}")
                        else:
                            formatted.append(f"  {name}")
                    else:
                        formatted.append(f"  {c}")
                company_text = "\n".join(formatted)
                _add_text(slide, left + Inches(0.1), list_top,
                          row_card_w - Inches(0.2), list_h,
                          company_text,
                          size=9, color=WHITE, cjk=True)

            # 箭头连接（同一行内，两色块之间）
            if col_idx < n_in_row - 1:
                arrow_left = left + row_card_w
                arrow_top = top + int(card_h / 2) - Inches(0.1)
                _add_text(slide, arrow_left - Inches(0.05), arrow_top,
                          Inches(0.25), Inches(0.2),
                          ">", size=14, bold=True, color=MUTED,
                          align=PP_ALIGN.CENTER, anchor="m", cjk=False)

    # 底部 source
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.4),
              SLIDE_W - Inches(0.8), Inches(0.3),
              "Source: Morgan Stanley Research",
              size=SIZE_SMALL, italic=True, color=MUTED,
              align=PP_ALIGN.LEFT, cjk=False)


# =============================================================================
# 10. Slide NEW B — What's Changed 变动摘要页
# =============================================================================

def slide_whats_changed(slide, *, data: Dict[str, Any],
                        title_cn: str = "本期变动",
                        title_en: str = "What's Changed",
                        language: str = "zh",
                        page_number: Optional[int] = None,
                        total_pages: Optional[int] = None) -> None:
    """变动摘要页：表格展示项目变动，用箭头/颜色标识方向。

    数据格式：
        data["whats_changed"] = [
            {"item": "Rating", "from": "Equal-weight", "to": "Overweight",
             "arrow": "up"},
            {"item": "TP", "from": "$85", "to": "$125", "arrow": "up"},
            ...
        ]
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    changes = list(data.get("whats_changed") or [])
    if not changes:
        _add_text(slide, Inches(0.5), Inches(2.5),
                  SLIDE_W - Inches(1.0), Inches(1.0),
                  "（未提供变动数据）" if language == "zh"
                  else "（No changes data provided）",
                  size=14, italic=True, color=MUTED,
                  align=PP_ALIGN.CENTER, cjk=True)
        return

    changes = changes[:12]  # 最多12条

    n_cols = 4  # Item | From | Arrow | To
    n_rows = len(changes) + 1

    table_left = Inches(0.5)
    table_top = Inches(1.5)
    table_w = SLIDE_W - Inches(1.0)
    row_h = Inches(0.45)
    table_h = row_h * n_rows

    ts = slide.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                table_w, table_h)
    table = ts.table

    col_pcts = [0.35, 0.25, 0.10, 0.30]
    for i, p in enumerate(col_pcts):
        table.columns[i].width = int(table_w * p)

    # 表头
    headers = ["Item / 项目", "From / 原值", "", "To / 新值"]
    for j, h in enumerate(headers):
        _fill_table_cell(table.cell(0, j), h,
                         size=10, bold=True, color=WHITE, bg=MS_BRAND_BLUE,
                         align=PP_ALIGN.CENTER, cjk=True)

    # 数据行
    for i, ch in enumerate(changes):
        bg = WHITE if i % 2 == 0 else TABLE_ALT
        item_text = str(ch.get("item", ch.get("metric", ch.get("field", ""))))
        from_text = str(ch.get("from", ch.get("from_val", ch.get("old_value", ch.get("previous", "")))))
        to_text = str(ch.get("to", ch.get("to_val", ch.get("new_value", ch.get("current", "")))))
        arrow = str(ch.get("arrow", ch.get("direction", ""))).strip().lower()

        # 箭头符号和颜色
        if arrow == "up":
            arrow_sym = "\u2191"  # ↑
            arrow_color = MS_GREEN
            to_color = MS_GREEN
        elif arrow == "down":
            arrow_sym = "\u2193"  # ↓
            arrow_color = MS_CHART_RED
            to_color = MS_CHART_RED
        else:
            arrow_sym = "\u2192"  # →
            arrow_color = MS_MUTED
            to_color = MS_BODY_TEXT

        _fill_table_cell(table.cell(i + 1, 0), item_text,
                         size=11, bold=True, color=MS_BODY_TEXT, bg=bg,
                         align=PP_ALIGN.LEFT, cjk=True)
        _fill_table_cell(table.cell(i + 1, 1), from_text,
                         size=11, color=MS_MUTED, bg=bg,
                         align=PP_ALIGN.CENTER, cjk=True)
        _fill_table_cell(table.cell(i + 1, 2), arrow_sym,
                         size=14, bold=True, color=arrow_color, bg=bg,
                         align=PP_ALIGN.CENTER, cjk=False)
        _fill_table_cell(table.cell(i + 1, 3), to_text,
                         size=11, bold=True, color=to_color, bg=bg,
                         align=PP_ALIGN.CENTER, cjk=True)

    # 底部 source
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.4),
              SLIDE_W - Inches(0.8), Inches(0.3),
              "Source: Morgan Stanley Research",
              size=SIZE_SMALL, italic=True, color=MUTED,
              align=PP_ALIGN.LEFT, cjk=False)


# =============================================================================
# 10. Slide NEW C — Thesis in Charts 核心论点图示页
# =============================================================================

def slide_thesis_in_charts(slide, *, data: Dict[str, Any],
                           title_cn: str = "核心论点",
                           title_en: str = "Our Thesis in Charts",
                           language: str = "zh",
                           page_number: Optional[int] = None,
                           total_pages: Optional[int] = None) -> None:
    """N张图展示核心论点：2x2或3x2网格布局，每个格子一个关键数据点。

    数据格式：
        data["thesis_charts"] = [
            {"title": "AI CapEx Surge", "value": "$600B",
             "subtitle": "2026E global AI capex", "color": "blue"},
            {"title": "GPU Shipments", "value": "1.1M",
             "subtitle": "Units shipped 2026E", "color": "green"},
            ...
        ]
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    charts = list(data.get("thesis_charts") or [])
    if not charts:
        _add_text(slide, Inches(0.5), Inches(2.5),
                  SLIDE_W - Inches(1.0), Inches(1.0),
                  "（未提供论点图表数据）" if language == "zh"
                  else "（No thesis chart data provided）",
                  size=14, italic=True, color=MUTED,
                  align=PP_ALIGN.CENTER, cjk=True)
        return

    charts = charts[:6]  # 最多6个（3x2）

    # 确定网格布局
    n = len(charts)
    if n <= 2:
        cols, rows = 2, 1
    elif n <= 4:
        cols, rows = 2, 2
    else:
        cols, rows = 3, 2

    area_left = Inches(0.5)
    area_top = Inches(1.5)
    area_w = SLIDE_W - Inches(1.0)
    area_h = SLIDE_H - area_top - Inches(1.0)
    gap = Inches(0.2)
    card_w = int((area_w - gap * (cols - 1)) / cols)
    card_h = int((area_h - gap * (rows - 1)) / rows)

    # 颜色映射
    color_map = {
        "blue": MS_CHART_BLUE,
        "red": MS_CHART_RED,
        "green": MS_GREEN,
        "navy": NAVY,
        "orange": ORANGE,
        "gold": GOLD,
        "brand": MS_BRAND_BLUE,
    }

    for idx, chart in enumerate(charts):
        r = idx // cols
        c = idx % cols
        left = area_left + (card_w + gap) * c
        top = area_top + (card_h + gap) * r

        # 卡片背景
        _add_rect(slide, left, top, card_w, card_h, LIGHT_BG)
        # 左侧色条
        bar_color_name = str(chart.get("color", "brand")).strip().lower()
        bar_color = color_map.get(bar_color_name, MS_BRAND_BLUE)
        _add_rect(slide, left, top, Inches(0.08), card_h, bar_color)

        # 标题
        chart_title = str(chart.get("title", ""))
        _add_text(slide, left + Inches(0.2), top + Inches(0.15),
                  card_w - Inches(0.4), Inches(0.3),
                  chart_title,
                  size=11, bold=True, color=MS_BODY_TEXT, cjk=True)

        # 核心数值（大字号）
        chart_value = str(chart.get("value", ""))
        _add_text(slide, left + Inches(0.2), top + Inches(0.5),
                  card_w - Inches(0.4), Inches(0.6),
                  chart_value,
                  size=32, bold=True, color=bar_color,
                  anchor="m", cjk=False)

        # 副标题/说明
        chart_sub = str(chart.get("subtitle", ""))
        _add_text(slide, left + Inches(0.2), top + card_h - Inches(0.5),
                  card_w - Inches(0.4), Inches(0.4),
                  chart_sub,
                  size=9, italic=True, color=MS_MUTED, cjk=True)

    # 底部 source
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.4),
              SLIDE_W - Inches(0.8), Inches(0.3),
              "Source: Morgan Stanley Research",
              size=SIZE_SMALL, italic=True, color=MUTED,
              align=PP_ALIGN.LEFT, cjk=False)


# =============================================================================
# 10. Slide 07 — Shovel Stocks 铲子股表格页
# =============================================================================

def slide_shovel_stocks(slide, *, stocks: List[Dict[str, Any]],
                        title_cn: str = "60 铲子股 · 行业机会清单",
                        title_en: str = "Shovel Stocks · Thematic Picks",
                        language: str = "zh",
                  page_number: Optional[int] = None,
                  total_pages: Optional[int] = None) -> None:
    """8 列：序号 | 公司 | Ticker | 评级 | 价格 | 分析师 | 市值($mn) | 1年回报(%)。

    升级：新增评级列（O/E/U），颜色编码 O=绿 E=黄 U=红。
    如果 stock 数据中包含 rating 字段则显示评级列，否则保持7列。
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    stocks = list(stocks or [])
    if not stocks:
        _add_text(slide, Inches(0.5), Inches(2.5),
                  SLIDE_W - Inches(1.0), Inches(1.0),
                  "（未提供铲子股数据）",
                  size=14, italic=True, color=MUTED,
                  align=PP_ALIGN.CENTER, cjk=True)
        return

    # 取前若干行，保证一页容纳
    max_rows = 16
    stocks = stocks[:max_rows]

    # 检测是否有评级数据
    has_rating = any(s.get("rating") for s in stocks)

    if has_rating:
        headers = ["#", "Company / 公司", "Ticker", "Rating / 评级",
                   "Price / 价格", "Analyst / 分析师", "Mkt Cap ($mn)", "1Y Return (%)"]
        n_cols = 8
        pcts = [0.04, 0.17, 0.09, 0.09, 0.09, 0.13, 0.12, 0.12]
        # 保留备注列
        has_notes = any(s.get("notes") for s in stocks)
        if has_notes:
            headers.append("Notes / 备注")
            n_cols = 9
            pcts = [0.04, 0.15, 0.08, 0.08, 0.08, 0.12, 0.11, 0.11, 0.14]
    else:
        headers = ["#", "Company / 公司", "Ticker", "Product / 核心产品",
                   "Analyst / 分析师", "Mkt Cap ($mn)", "1Y Return (%)"]
        n_cols = 7
        pcts = [0.05, 0.20, 0.11, 0.24, 0.14, 0.13, 0.13]

    n_rows = len(stocks) + 1

    table_left = Inches(0.4)
    table_top = Inches(1.2)
    table_w = SLIDE_W - Inches(0.8)
    row_h = Inches((SLIDE_H - table_top - Inches(0.8)) / n_rows)
    row_h = min(row_h, Inches(0.48))
    table_h = row_h * n_rows

    ts = slide.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                table_w, table_h)
    table = ts.table

    for i, p in enumerate(pcts):
        table.columns[i].width = int(table_w * p)

    # 表头：深蓝底白字
    for j, h in enumerate(headers):
        _fill_table_cell(table.cell(0, j), str(h),
                         size=10, bold=True, color=WHITE, bg=NAVY,
                         align=PP_ALIGN.CENTER if j in (0, 2, 5, 6, 7) else PP_ALIGN.LEFT,
                         cjk=True)

    # 数据行：白 / 浅灰交替
    for i, s in enumerate(stocks):
        bg = WHITE if i % 2 == 0 else TABLE_ALT
        rank = s.get("rank", i + 1)
        company = str(s.get("company", s.get("name", "")))
        ticker = str(s.get("ticker", ""))
        analyst = str(s.get("analyst", ""))
        mcap = s.get("market_cap_mn", s.get("mcap_mn", s.get("mkt_cap", s.get("market_cap", ""))))
        perf = s.get("perf_1y_pct", s.get("perf_1y", s.get("return_1y_pct", "")))

        # 数字格式化
        try:
            mcap_v = float(mcap)
            mcap_s = f"{mcap_v:,.0f}"
        except (TypeError, ValueError):
            mcap_s = str(mcap) if mcap else "—"

        try:
            perf_v = float(perf)
            perf_s = f"{perf_v:+.1f}%"
        except (TypeError, ValueError):
            perf_s = str(perf) if perf else "—"

        # 判断高亮：>100% 橙色加粗
        highlight_orange = False
        try:
            highlight_orange = float(perf) > 100
        except (TypeError, ValueError):
            highlight_orange = False

        if has_rating:
            # 评级颜色编码
            rating_raw = str(s.get("rating", "")).strip().upper()
            if rating_raw in ("O", "OW", "OVERWEIGHT", "买入"):
                rating_display = "O"
                rating_color = MS_GREEN
            elif rating_raw in ("E", "EW", "EQUAL-WEIGHT", "EQUALWEIGHT", "持有"):
                rating_display = "E"
                rating_color = RGBColor(0xD4, 0x9B, 0x00)
            elif rating_raw in ("U", "UW", "UNDERWEIGHT", "卖出"):
                rating_display = "U"
                rating_color = MS_CHART_RED
            else:
                rating_display = rating_raw or "—"
                rating_color = MS_MUTED

            price = str(s.get("price", s.get("target_price", "")))
            notes = str(s.get("notes", ""))

            cells = [
                (str(rank), False, NAVY_ALT, PP_ALIGN.CENTER),
                (company, True, TEXT, PP_ALIGN.LEFT),
                (ticker, True, NAVY, PP_ALIGN.CENTER),
                (rating_display, True, rating_color, PP_ALIGN.CENTER),
                (price, False, TEXT, PP_ALIGN.CENTER),
                (analyst, False, MUTED, PP_ALIGN.LEFT),
                (mcap_s, False, TEXT, PP_ALIGN.CENTER),
                (perf_s, highlight_orange,
                 (ORANGE if highlight_orange else TEXT),
                 PP_ALIGN.CENTER),
            ]
            if has_notes:
                cells.append((notes, False, MS_MUTED, PP_ALIGN.LEFT))
        else:
            product = str(s.get("product", s.get("core_product", "")))
            cells = [(str(rank), False, NAVY_ALT, PP_ALIGN.CENTER),
                     (company, True, TEXT, PP_ALIGN.LEFT),
                     (ticker, True, NAVY, PP_ALIGN.CENTER),
                     (product, False, TEXT, PP_ALIGN.LEFT),
                     (analyst, False, MUTED, PP_ALIGN.LEFT),
                     (mcap_s, False, TEXT, PP_ALIGN.CENTER),
                     (perf_s, highlight_orange,
                      (ORANGE if highlight_orange else TEXT),
                      PP_ALIGN.CENTER)]

        for j, (txt, bold, color, align) in enumerate(cells):
            _fill_table_cell(table.cell(i + 1, j), txt,
                             size=10, bold=bold, color=color, bg=bg,
                             align=align, cjk=True)

    # 底部 source
    source_note = "Source: Morgan Stanley Research"
    if has_rating:
        source_note += " · O=Overweight(green) E=Equal-weight(yellow) U=Underweight(red)"
    else:
        source_note += " · Orange shading denotes 1Y return > 100%"
    _add_text(slide, Inches(0.4), SLIDE_H - Inches(0.4),
              SLIDE_W - Inches(0.8), Inches(0.3),
              source_note,
              size=SIZE_SMALL, italic=True, color=MUTED,
              align=PP_ALIGN.LEFT, cjk=True)


# =============================================================================
# 11. Slide 08 — Market Monitor 市场监测页（可选）
# =============================================================================

def slide_market_monitor(slide, *, blocks: Optional[List[Dict[str, Any]]] = None,
                         title_cn: str = "市场监测",
                         title_en: str = "Market Monitor",
                         language: str = "zh",
                  page_number: Optional[int] = None,
                  total_pages: Optional[int] = None) -> None:
    """四色分栏色块：Bonds / Equities / Alternatives / Transition。"""
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    default_blocks = [
        {"name_en": "Bonds", "name_cn": "固收",
         "metric": "3.85%", "sub": "10Y UST yield",
         "color": RGBColor(0x0B, 0x2C, 0x5C)},
        {"name_en": "Equities", "name_cn": "股票",
         "metric": "21,420", "sub": "S&P 500 level",
         "color": RGBColor(0xC8, 0xA9, 0x51)},
        {"name_en": "Alternatives", "name_cn": "另类",
         "metric": "$2.4T", "sub": "PE dry powder",
         "color": RGBColor(0xE3, 0x7C, 0x2B)},
        {"name_en": "Transition", "name_cn": "能源转型",
         "metric": "+18%", "sub": "Clean-tech capex YoY",
         "color": RGBColor(0x1F, 0x7A, 0x3E)},
    ]
    blocks = list(blocks or default_blocks)
    blocks = blocks[:4]

    # 颜色归一化：hex 字符串 → RGBColor
    def _ensure_rgb(c):
        if isinstance(c, RGBColor):
            return c
        if isinstance(c, str) and c.startswith("#") and len(c) == 7:
            return RGBColor(int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16))
        return default_blocks[0]["color"]

    area_top = Inches(1.4)
    area_h = SLIDE_H - area_top - Inches(1.0)
    gap = Inches(0.3)
    total_gap = gap * (len(blocks) - 1)
    card_w = int((SLIDE_W - Inches(1.0) - total_gap) / len(blocks))

    for idx, b in enumerate(blocks):
        left = Inches(0.5) + (card_w + gap) * idx
        color = _ensure_rgb(b.get("color")) or default_blocks[idx % 4]["color"]
        # 主色块
        _add_rect(slide, left, area_top, card_w, area_h, color)
        # 金色左上角小块
        _add_rect(slide, left, area_top, Inches(0.35), Inches(0.35), GOLD)

        # 分类名称
        _add_text(slide, left + Inches(0.3), area_top + Inches(0.4),
                  card_w - Inches(0.5), Inches(0.45),
                  str(b.get("name_en", "")),
                  size=14, bold=True, color=WHITE, cjk=False)
        if b.get("name_cn"):
            _add_text(slide, left + Inches(0.3), area_top + Inches(0.85),
                      card_w - Inches(0.5), Inches(0.35),
                      str(b.get("name_cn", "")),
                      size=11, color=GOLD, cjk=True)

        # 金色大指标
        _add_text(slide, left + Inches(0.3), area_top + Inches(1.55),
                  card_w - Inches(0.5), Inches(1.5),
                  str(b.get("metric", "—")),
                  size=36, bold=True, color=GOLD, cjk=False)

        # 副标题 + 描述
        _add_text(slide, left + Inches(0.3), area_top + Inches(3.3),
                  card_w - Inches(0.5), Inches(0.4),
                  str(b.get("sub", "")),
                  size=11, bold=True, color=WHITE, cjk=True)

        if b.get("body"):
            _add_text(slide, left + Inches(0.3), area_top + Inches(3.75),
                      card_w - Inches(0.5), area_h - Inches(4.0),
                      str(b.get("body", "")),
                      size=10, color=WHITE, cjk=True)


# =============================================================================
# 12. Slide 09 — Disclosure 免责页
# =============================================================================

def slide_disclosure(slide, *, text: Optional[str] = None,
                     title_cn: str = "重要信息披露",
                     title_en: str = "Important Disclosures",
                     language: str = "zh",
                  page_number: Optional[int] = None,
                  total_pages: Optional[int] = None) -> None:
    """免责声明页：顶部 NAVY 横条 + 灰色正文小字。"""
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    if text is None:
        text = (
            "本材料由 Morgan Stanley 研究部门准备，仅供参考，不构成任何投资建议。"
            "文中观点可能与 Morgan Stanley 其他业务部门的观点不同。"
            "过往业绩不代表未来表现。投资涉及风险，包括本金损失。"
            "接收者在做出投资决策前应自行评估相关风险，并咨询独立财务顾问。"
            "\n\nThis material is prepared by Morgan Stanley Research for informational "
            "purposes only and does not constitute investment advice. Past performance "
            "is not indicative of future results. Investments involve risks including "
            "loss of principal. Recipients should assess risks independently and consult "
            "their own financial advisors."
        )

    _add_text(slide, Inches(0.6), Inches(1.2), SLIDE_W - Inches(1.2),
              Inches(5.6), str(text),
              size=10, color=MUTED, cjk=True)


# =============================================================================
# 12b-0. 辅助函数：收入预测路径叠加图
# =============================================================================

def _add_revenue_forecast_overlay(slide, scenarios: Dict[str, Any],
                                 years: List[str],
                                 revenue: List[float]) -> None:
    """在 slide 上添加一个小的收入预测路径折线图（三情景分叉）。

    三条线（熊/基/牛）从最后历史年份分叉，展示不同增速路径。
    用虚线样式区分预测期。
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    scenario_keys = ["bear", "base", "bull"]
    labels_map = {"bear": "Bear", "base": "Base", "bull": "Bull"}
    colors_map = {"bear": UW_RED, "base": GOLD, "bull": OW_GREEN}

    # 收集预测数据
    forecast_data = {}
    for sk in scenario_keys:
        scen = scenarios.get(sk, {})
        rf = scen.get("revenue_forecast", [])
        if rf:
            forecast_data[sk] = [float(v) for v in rf]

    if not forecast_data:
        return

    # 构建图表 categories：历史年份 + 预测年份
    n_hist = len(years) if years else 0
    n_forecast = max(len(v) for v in forecast_data.values())
    if n_forecast == 0:
        return

    # 预测年份标签
    last_hist_year = years[-1] if years else "FY24"
    # 尝试从最后历史年份推导预测年份
    try:
        base_year = int(re.search(r'(\d{4})', last_hist_year).group(1))
        forecast_years = [f"FY{base_year + i + 1}" for i in range(n_forecast)]
    except (AttributeError, ValueError):
        forecast_years = [f"FY{i}" for i in range(n_forecast)]

    all_years = list(years or []) + forecast_years

    # 图表位置：右下角，叠加在折线图区域下方
    fc_left = Inches(6.85)
    fc_top = Inches(5.0)
    fc_w = Inches(6.1)
    fc_h = Inches(1.8)

    # 背景
    _add_rect(slide, fc_left - Inches(0.05), fc_top - Inches(0.05),
              fc_w + Inches(0.1), fc_h + Inches(0.1), WHITE)
    _add_rect(slide, fc_left, fc_top - Inches(0.25), fc_w, Inches(0.25), NAVY)
    _add_text(slide, fc_left + Inches(0.1), fc_top - Inches(0.25),
              fc_w - Inches(0.2), Inches(0.25),
              "Revenue Forecast Paths / 收入预测路径",
              size=8, bold=True, color=GOLD, anchor="m", cjk=True)

    chart_data = CategoryChartData()
    chart_data.categories = all_years

    # 历史收入作为基础线（前 n_hist 个点有值，后面为 None）
    hist_extended = list(revenue or []) + [None] * n_forecast
    chart_data.add_series("Historical", hist_extended)

    # 三个情景预测线
    for sk in scenario_keys:
        if sk not in forecast_data:
            continue
        rf = forecast_data[sk]
        # 前面用历史最后一个值填充，后面用预测值
        last_hist = float(revenue[-1]) if revenue else 0
        vals = [None] * n_hist + rf
        chart_data.add_series(labels_map.get(sk, sk), vals)

    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE,
                                fc_left, fc_top, fc_w, fc_h,
                                chart_data)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(7)

    # 系列样式：历史线实线+circle marker，预测线虚线+diamond/triangle marker
    series_colors = [MUTED, UW_RED, GOLD, OW_GREEN]
    # marker 样式: Historical=circle, Bear=diamond, Base=triangle, Bull=diamond
    from pptx.enum.chart import XL_MARKER_STYLE
    marker_styles = [
        XL_MARKER_STYLE.CIRCLE,
        XL_MARKER_STYLE.DIAMOND,
        XL_MARKER_STYLE.TRIANGLE,
        XL_MARKER_STYLE.DIAMOND,
    ]
    for idx, s in enumerate(chart.series):
        ln = s.format.line
        ln.color.rgb = series_colors[idx % len(series_colors)]
        if idx == 0:
            # 历史线：实线, linewidth=2.5pt, 深灰色
            ln.width = Pt(2.5)
            ln.dash_style = None  # 实线
        else:
            # 预测线：虚线, linewidth=2pt
            ln.width = Pt(2.0)
            ln.dash_style = 2  # dash
        # 添加 marker
        s.marker.style = marker_styles[idx % len(marker_styles)]
        s.marker.size = 5  # marker size 是整数(2-72), 不是 Pt
        s.marker.format.fill.solid()
        s.marker.format.fill.fore_color.rgb = series_colors[idx % len(series_colors)]
        s.marker.format.line.color.rgb = series_colors[idx % len(series_colors)]

    # 数据标签：最后一个点显示数值
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(7)
    dl.font.bold = True
    dl.font.color.rgb = TEXT
    dl.number_format = '#,##0'
    dl.show_value = True
    dl.show_category_name = False
    dl.show_series_name = False


# =============================================================================
# 12b. Slide 10 — slide_financial_chart 财务图表（柱 + 折线）
# =============================================================================

def slide_financial_chart(slide, *,
                          years: Optional[List[str]] = None,
                          revenue: Optional[List[float]] = None,
                          ebitda: Optional[List[float]] = None,
                          margins: Optional[Dict[str, List[float]]] = None,
                          takeaways: Optional[List[str]] = None,
                          scenarios: Optional[Dict[str, Any]] = None,
                          title_cn: str = "收入 / 利润率一览",
                          title_en: str = "Revenue & EBITDA · Margin Trends",
                          language: str = "zh",
                  page_number: Optional[int] = None,
                  total_pages: Optional[int] = None) -> None:
    """财务图表页：顶部标题栏 + 左侧双柱（Revenue / EBITDA）
    + 右侧利润率折线 + 底部 takeaways bullet。

    增强功能：如果提供 ``scenarios`` 参数（含 bear/base/bull 的
    revenue_forecast），在右下角增加一个"收入预测路径"小折线图，
    三条线（熊/基/牛）从最后历史年份分叉，展示不同增速路径。
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import (XL_CHART_TYPE, XL_LEGEND_POSITION,
                                  XL_LABEL_POSITION, XL_LABEL_POSITION as _LP)

    # 默认示例数据
    if years is None:
        years = ["FY22", "FY23", "FY24", "FY25E", "FY26E", "FY27E", "FY28E"]
    if revenue is None:
        revenue = [1200.0, 1450.0, 1780.0, 2120.0, 2460.0, 2810.0, 3180.0]
    if ebitda is None:
        ebitda = [280.0, 348.0, 442.0, 550.0, 665.0, 780.0, 910.0]
    if margins is None:
        margins = {
            "gross":   [55.0, 56.5, 58.2, 60.0, 61.5, 62.8, 63.9],
            "ebitda":  [23.3, 24.0, 24.8, 25.9, 27.0, 27.8, 28.6],
            "net":     [14.2, 15.0, 16.1, 17.0, 17.8, 18.5, 19.1],
        }
    if takeaways is None:
        takeaways = [
            "Revenue CAGR ~21%，AI 相关产品贡献主要增量。",
            "EBITDA margin 自 23.3% 稳步扩张至 28.6%，规模效应显现。",
            "Gross margin 受产品结构改善推动至 63.9%。",
        ]

    n = len(years)
    assert len(revenue) == n and len(ebitda) == n,         "revenue/ebitda 长度应与 years 一致"
    for mk, mv in margins.items():
        assert len(mv) == n, f"margins[{mk}] 长度应与 years 一致"

    # 左侧柱状图区域
    bar_left = Inches(0.4)
    bar_top = Inches(1.15)
    bar_w = Inches(6.2)
    bar_h = Inches(4.4)

    bar_data = CategoryChartData()
    bar_data.categories = list(years)
    bar_data.add_series("Revenue ($M)", [float(v) for v in revenue])
    bar_data.add_series("EBITDA ($M)", [float(v) for v in ebitda])

    gf1 = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                  bar_left, bar_top, bar_w, bar_h,
                                  bar_data)
    chart1 = gf1.chart
    chart1.has_title = False
    chart1.has_legend = True
    chart1.legend.include_in_layout = False
    chart1.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart1.legend.horz_offset = 0
    chart1.legend.font.size = Pt(8)

    # 系列上色：Revenue=#1F3864, EBITDA=#2E75B6, 白色边框
    s_rev = chart1.series[0]
    s_ebi = chart1.series[1]
    bar_colors = [RGBColor(0x1F, 0x38, 0x64), RGBColor(0x2E, 0x75, 0xB6)]
    for s, color in [(s_rev, bar_colors[0]), (s_ebi, bar_colors[1])]:
        fill = s.format.fill
        fill.solid()
        fill.fore_color.rgb = color
        line = s.format.line
        line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        line.width = Pt(0.5)

    plot1 = chart1.plots[0]
    plot1.has_data_labels = True
    dl1 = plot1.data_labels
    dl1.font.size = Pt(7)
    dl1.font.bold = True
    dl1.font.color.rgb = TEXT
    dl1.number_format = '#,##0'

    # 坐标轴样式
    try:
        _style_chart_axis(chart1.value_axis)
        chart1.value_axis.major_gridlines.format.line.dash_style = 2  # 虚线
        chart1.value_axis.tick_labels.number_format = '#,##0'
    except Exception:
        pass

    # 右侧折线图区域
    line_left = Inches(6.85)
    line_top = Inches(1.15)
    line_w = Inches(6.1)
    line_h = Inches(4.4)

    line_data = CategoryChartData()
    line_data.categories = list(years)
    # Gross (NAVY), EBITDA (GOLD), Net (OW_GREEN)
    line_data.add_series("Gross Margin (%)",
                         [float(v) for v in margins["gross"]])
    line_data.add_series("EBITDA Margin (%)",
                         [float(v) for v in margins["ebitda"]])
    line_data.add_series("Net Margin (%)",
                         [float(v) for v in margins["net"]])

    gf2 = slide.shapes.add_chart(XL_CHART_TYPE.LINE,
                                  line_left, line_top, line_w, line_h,
                                  line_data)
    chart2 = gf2.chart
    chart2.has_title = False
    chart2.has_legend = True
    chart2.legend.include_in_layout = False
    chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart2.legend.font.size = Pt(8)

    # 三条线颜色: Gross=#1F3864, EBITDA=#C8A951, Net=#00AF50
    line_colors = [
        RGBColor(0x1F, 0x38, 0x64),  # Gross - 深海军蓝
        RGBColor(0xC8, 0xA9, 0x51),   # EBITDA - MS GOLD
        RGBColor(0x00, 0xAF, 0x50),   # Net - 绿色
    ]
    from pptx.enum.chart import XL_MARKER_STYLE
    marker_styles = [
        XL_MARKER_STYLE.CIRCLE,    # Gross
        XL_MARKER_STYLE.SQUARE,    # EBITDA
        XL_MARKER_STYLE.TRIANGLE,   # Net
    ]
    marker_sizes = [4, 4, 4]  # marker size 是整数(2-72), 不是 Pt
    for idx, s in enumerate(chart2.series):
        ln = s.format.line
        ln.color.rgb = line_colors[idx % len(line_colors)]
        ln.width = Pt(2.5)
        # 添加 marker
        s.marker.style = marker_styles[idx % len(marker_styles)]
        s.marker.size = marker_sizes[idx % len(marker_sizes)]
        s.marker.format.fill.solid()
        s.marker.format.fill.fore_color.rgb = line_colors[idx % len(line_colors)]
        s.marker.format.line.color.rgb = line_colors[idx % len(line_colors)]

    plot2 = chart2.plots[0]
    plot2.has_data_labels = True
    dl2 = plot2.data_labels
    dl2.font.size = Pt(7)
    dl2.font.bold = True
    dl2.font.color.rgb = TEXT
    dl2.number_format = '0.0"%"'

    # 坐标轴样式
    try:
        _style_chart_axis(chart2.value_axis)
        chart2.value_axis.major_gridlines.format.line.dash_style = 2  # 虚线
    except Exception:
        pass

    # ---- 增强功能：收入预测路径小图（三情景分叉） ----
    if scenarios and isinstance(scenarios, dict):
        _add_revenue_forecast_overlay(slide, scenarios, years, revenue)

    # 底部 takeaways 行
    ta_top = SLIDE_H - Inches(1.05)
    _add_rect(slide, Inches(0.4), ta_top,
              SLIDE_W - Inches(0.8), Inches(0.04), GOLD)

    # bullet 文本
    tb = _add_textbox(slide, Inches(0.5), ta_top + Inches(0.1),
                      SLIDE_W - Inches(1.0), Inches(0.95), anchor="t")
    tf = tb.text_frame
    tf.word_wrap = True
    for i, tk in enumerate(takeaways or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = "–  "
        _apply_font(r1, size=10, bold=True, color=NAVY, cjk=False)
        r2 = p.add_run()
        r2.text = str(tk)
        _apply_font(r2, size=10, color=TEXT, cjk=True)
        p.space_after = Pt(2)


# =============================================================================
# 12c. Slide 11 — slide_executive_summary 四象限执行摘要
# =============================================================================

def slide_executive_summary(slide, *,
                            thesis: Optional[List[str]] = None,
                            tp_and_upside: Optional[Dict[str, Any]] = None,
                            key_risks: Optional[List[str]] = None,
                            catalysts: Optional[List[str]] = None,
                            title_cn: str = "执行摘要",
                            title_en: str = "Executive Summary · Thesis, TP, Risks & Catalysts",
                            language: str = "zh",
                  page_number: Optional[int] = None,
                  total_pages: Optional[int] = None) -> None:
    """四象限信息页：
    左上 = Investment Thesis
    右上 = Target Price & Upside
    左下 = Key Risks
    右下 = Catalysts
    中央与区块之间用金色竖线分隔。
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    def _add_quadrant(left, top, width, height, *,
                      header: str, body: List[str],
                      body_color: RGBColor = TEXT) -> None:
        # 区块背景（浅灰）
        _add_rect(slide, left, top, width, height, TABLE_ALT)
        # 顶部金色横条 + 标题
        _add_rect(slide, left, top, width, Inches(0.35), GOLD)
        _add_text(slide, left + Inches(0.15), top,
                  width - Inches(0.3), Inches(0.35),
                  header.upper(),
                  size=10, bold=True, color=WHITE,
                  align=PP_ALIGN.LEFT, anchor="m", cjk=False)
        # 正文：bullet
        tb = _add_textbox(slide, left + Inches(0.2), top + Inches(0.5),
                          width - Inches(0.4), height - Inches(0.6),
                          anchor="t")
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body or []):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(4)
            r1 = p.add_run()
            r1.text = "–  "
            _apply_font(r1, size=10, bold=True, color=NAVY, cjk=False)
            r2 = p.add_run()
            r2.text = str(line)
            _apply_font(r2, size=10, color=body_color, cjk=True)

    # 象限布局
    q_top = Inches(1.15)
    q_left_margin = Inches(0.5)
    gap = Inches(0.3)
    q_w = (SLIDE_W - q_left_margin * 2 - gap) / 2
    q_h = Inches(2.85)
    mid_top = q_top + q_h + Inches(0.3)

    # 左上 - Thesis
    thesis_body = thesis or [
        "结构性 AI 需求驱动收入快速增长，23-28E CAGR ~21%。",
        "供给侧约束短期仍紧俏，公司享有定价权与份额提升。",
        "运营杠杆释放推动 EBITDA margin 持续扩张 > 28%。",
    ]
    _add_quadrant(q_left_margin, q_top, q_w, q_h,
                  header="Investment Thesis · 投资论点",
                  body=thesis_body)

    # 右上 - TP & Upside
    if tp_and_upside is None:
        tp_and_upside = {
            "rating": "Overweight", "target_price": "$125",
            "last_price": "$98.5", "upside": "+27%",
            "notes": ["目标价基于 DCF (WACC 9.5%, g 2.5%) 与同业估值对比。",
                      "12 个月视野，对应 FY27E EV/EBITDA ~18x。"],
        }
    # 画一个带金色关键数字的象限
    left, top, width, height = (q_left_margin + q_w + gap, q_top, q_w, q_h)
    _add_rect(slide, left, top, width, height, TABLE_ALT)
    _add_rect(slide, left, top, width, Inches(0.35), GOLD)
    _add_text(slide, left + Inches(0.15), top, width - Inches(0.3),
              Inches(0.35),
              "TARGET PRICE & UPSIDE · 目标价与上行空间",
              size=10, bold=True, color=WHITE,
              align=PP_ALIGN.LEFT, anchor="m", cjk=False)
    # 关键数字：rating + target price + upside
    rating_txt = str(tp_and_upside.get("rating", "Overweight"))
    rc_color = OW_GREEN if "OVERWEIGHT" in rating_txt.upper() else (
        UW_RED if "UNDERWEIGHT" in rating_txt.upper() else EW_ORANGE)
    _add_text(slide, left + Inches(0.2), top + Inches(0.55),
              width - Inches(0.4), Inches(0.4),
              rating_txt, size=14, bold=True, color=rc_color,
              align=PP_ALIGN.LEFT, anchor="m", cjk=True)
    tp_txt = str(tp_and_upside.get("target_price", "—"))
    _add_text(slide, left + Inches(0.2), top + Inches(1.05),
              width - Inches(0.4), Inches(0.8),
              tp_txt, size=40, bold=True, color=NAVY,
              align=PP_ALIGN.LEFT, anchor="t", cjk=False)
    up_txt = str(tp_and_upside.get("upside", "—"))
    last_txt = str(tp_and_upside.get("last_price", "—"))
    _add_text(slide, left + Inches(0.2), top + Inches(1.85),
              width - Inches(0.4), Inches(0.4),
              f"Upside: {up_txt}   |   Last: {last_txt}",
              size=11, bold=True, color=GOLD,
              align=PP_ALIGN.LEFT, anchor="m", cjk=False)
    # Notes bullets
    tb = _add_textbox(slide, left + Inches(0.2), top + Inches(2.3),
                      width - Inches(0.4), height - Inches(2.4),
                      anchor="t")
    tf = tb.text_frame
    tf.word_wrap = True
    notes = tp_and_upside.get("notes", []) or []
    for i, line in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = "–  "
        _apply_font(r1, size=10, bold=True, color=NAVY, cjk=False)
        r2 = p.add_run()
        r2.text = str(line)
        _apply_font(r2, size=10, color=TEXT, cjk=True)

    # 左下 - Key Risks
    risks_body = key_risks or [
        "AI 资本开支节奏不及预期或大客户砍单。",
        "电力 / 冷却供给瓶颈导致数据中心建设延期。",
        "地缘政治摩擦导致出口管制或供应链风险。",
        "估值扩张依赖持续超预期，增速回归时存在压缩风险。",
    ]
    _add_quadrant(q_left_margin, mid_top, q_w, q_h,
                  header="Key Risks · 关键风险",
                  body=risks_body, body_color=UW_RED)

    # 右下 - Catalysts
    catalysts_body = catalysts or [
        "下一季指引上调 / 超预期订单披露。",
        "新一代 AI 加速器量产节奏与客户样品反馈。",
        "大型云厂商 capex 加总确认上行周期延续。",
        "行业并购与生态合作带来份额与毛利率改善。",
    ]
    _add_quadrant(q_left_margin + q_w + gap, mid_top, q_w, q_h,
                  header="Catalysts · 催化剂",
                  body=catalysts_body, body_color=OW_GREEN)

    # 中央金色竖线（上下象限之间共用的竖线）
    vline_left = q_left_margin + q_w + gap / 2 - Inches(0.03)
    _add_rect(slide, vline_left, q_top, Inches(0.06),
              q_h * 2 + Inches(0.3), GOLD)
    # 横向中央细线
    _add_rect(slide, q_left_margin, q_top + q_h + Inches(0.13),
              SLIDE_W - q_left_margin * 2, Inches(0.04), GOLD)


# =============================================================================
# 12d. Slide 12 — slide_scenario_comparison 三情景估值对比
# =============================================================================

def slide_scenario_comparison(slide, *, data: Dict[str, Any],
                              theme: str = "classic",
                              language: str = "zh",
                              title_cn: str = "情景分析",
                              title_en: str = "Scenario Analysis",
                              page_number: Optional[int] = None,
                              total_pages: Optional[int] = None) -> None:
    """三情景估值对比页：熊市 / 基准 / 牛市三列表格，基准列金色高亮。

    数据来自 ``data["scenarios"]``，包含 bear/base/bull 三个子字典。
    每个子字典需含：wacc_value, terminal_growth_rate, exit_multiple_ebitda,
    enterprise_value, equity_value, per_share_value, ev_revenue, ev_ebitda。
    底部一行显示加权平均/基准情景的目标价。
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    scenarios = data.get("scenarios", {})
    if not scenarios:
        _add_text(slide, Inches(0.5), Inches(3.0), SLIDE_W - Inches(1.0),
                  Inches(1.0), "（未提供情景数据 / No scenario data）",
                  size=14, italic=True, color=MUTED,
                  align=PP_ALIGN.CENTER, cjk=True)
        return

    # 情景顺序与标签
    scenario_keys = ["bear", "base", "bull"]
    labels_zh = {"bear": "熊市情景", "base": "基准情景", "bull": "牛市情景"}
    labels_en = {"bear": "Bear Case", "base": "Base Case", "bull": "Bull Case"}
    colors = {"bear": BEAR_COLOR, "base": BASE_COLOR, "bull": BULL_COLOR}

    # 表格行定义
    row_labels = [
        ("WACC", "WACC"),
        ("TGR", "Terminal Growth Rate"),
        ("Exit Multiple", "Exit Multiple (EV/EBITDA)"),
        ("EV", "Enterprise Value ($M)"),
        ("Equity Value", "Equity Value ($M)"),
        ("Per Share", "Per Share Value ($)"),
        ("EV/Revenue", "EV/Revenue (x)"),
        ("EV/EBITDA", "EV/EBITDA (x)"),
    ]

    n_cols = 4  # Label + 3 scenarios
    n_rows = len(row_labels) + 2  # header + rows + target price row

    table_left = Inches(0.5)
    table_top = Inches(1.2)
    table_w = SLIDE_W - Inches(1.0)
    row_h = Inches(0.52)
    table_h = row_h * n_rows

    ts = slide.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                table_w, table_h)
    table = ts.table

    # 列宽：Label 25% + 三列各 25%
    table.columns[0].width = int(table_w * 0.25)
    for c in range(1, 4):
        table.columns[c].width = int(table_w * 0.25)

    # 表头行
    _fill_table_cell(table.cell(0, 0), "Metric / 指标",
                     size=11, bold=True, color=WHITE, bg=NAVY,
                     align=PP_ALIGN.LEFT, cjk=True)
    for ci, sk in enumerate(scenario_keys):
        scen = scenarios.get(sk, {})
        lbl = (labels_zh.get(sk, sk) if language == "zh"
               else labels_en.get(sk, sk))
        _fill_table_cell(table.cell(0, ci + 1), lbl,
                         size=11, bold=True, color=WHITE,
                         bg=colors.get(sk, NAVY),
                         align=PP_ALIGN.CENTER, cjk=True)

    # 数据行
    for ri, (label_zh, label_en) in enumerate(row_labels):
        row_idx = ri + 1
        bg_label = NAVY_ALT if ri % 2 == 0 else RGBColor(0x2A, 0x4A, 0x7A)
        _fill_table_cell(table.cell(row_idx, 0),
                         f"{label_zh}" if language == "zh" else label_en,
                         size=10, bold=True, color=WHITE, bg=bg_label,
                         align=PP_ALIGN.LEFT, cjk=True)

        for ci, sk in enumerate(scenario_keys):
            scen = scenarios.get(sk, {})
            is_base = (sk == "base")
            if is_base:
                bg = BASE_BG
            elif sk == "bear":
                bg = BEAR_BG if ri % 2 == 0 else WHITE
            elif sk == "bull":
                bg = BULL_BG if ri % 2 == 0 else WHITE
            else:
                bg = ALT_ROW if ri % 2 == 1 else WHITE

            # 根据行标签取值
            val = _get_scenario_metric(scen, label_en)
            val_str = _fmt_scenario_val(val, label_en)

            _fill_table_cell(table.cell(row_idx, ci + 1), val_str,
                             size=11, bold=is_base,
                             color=NAVY if is_base else TEXT,
                             bg=bg, align=PP_ALIGN.CENTER, cjk=False)

    # 底部目标价行
    tp_row = len(row_labels) + 1
    _fill_table_cell(table.cell(tp_row, 0),
                     "Target Price / 目标价",
                     size=11, bold=True, color=WHITE, bg=NAVY,
                     align=PP_ALIGN.LEFT, cjk=True)
    base_scen = scenarios.get("base", {})
    tp_val = base_scen.get("per_share_value", "—")
    tp_str = f"${tp_val}" if isinstance(tp_val, (int, float)) else str(tp_val)
    _fill_table_cell(table.cell(tp_row, 1), "—",
                     size=11, color=MUTED, bg=TABLE_ALT,
                     align=PP_ALIGN.CENTER, cjk=False)
    _fill_table_cell(table.cell(tp_row, 2), tp_str,
                     size=14, bold=True, color=NAVY,
                     bg=BASE_BG,
                     align=PP_ALIGN.CENTER, cjk=False)
    _fill_table_cell(table.cell(tp_row, 3), "—",
                     size=11, color=MUTED, bg=TABLE_ALT,
                     align=PP_ALIGN.CENTER, cjk=False)

    # 底部 source
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.4),
              SLIDE_W - Inches(1.0), Inches(0.3),
              "Source: Morgan Stanley Research · DCF Scenario Analysis",
              size=SIZE_SMALL, italic=True, color=MUTED,
              align=PP_ALIGN.LEFT, cjk=False)


def _get_scenario_metric(scen: Dict[str, Any], label_en: str) -> Any:
    """根据行标签英文名从 scenario dict 中取值。"""
    mapping = {
        "WACC": "wacc_value",
        "Terminal Growth Rate": "terminal_growth_rate",
        "Exit Multiple (EV/EBITDA)": "exit_multiple_ebitda",
        "Enterprise Value ($M)": "enterprise_value",
        "Equity Value ($M)": "equity_value",
        "Per Share Value ($)": "per_share_value",
        "EV/Revenue (x)": "ev_revenue",
        "EV/EBITDA (x)": "ev_ebitda",
    }
    key = mapping.get(label_en, "")
    return scen.get(key, "—") if key else "—"


def _fmt_scenario_val(val: Any, label_en: str) -> str:
    """格式化情景指标值。"""
    if val is None or val == "—":
        return "—"
    try:
        v = float(val)
    except (TypeError, ValueError):
        return str(val)
    if "WACC" in label_en or "Growth" in label_en:
        return f"{v:.1%}"
    if "Multiple" in label_en or "EV/Revenue" in label_en or "EV/EBITDA" in label_en:
        return f"{v:.1f}x"
    if "Per Share" in label_en:
        return f"${v:.1f}"
    # Enterprise Value / Equity Value
    return f"${v:,.0f}"


# =============================================================================
# 12e. Slide 13 — slide_wacc_breakdown WACC 拆解可视化
# =============================================================================

def slide_wacc_breakdown(slide, *, data: Dict[str, Any],
                          theme: str = "classic",
                          language: str = "zh",
                          title_cn: str = "WACC 拆解",
                          title_en: str = "WACC Decomposition",
                          page_number: Optional[int] = None,
                          total_pages: Optional[int] = None) -> None:
    """WACC 拆解页：上方堆叠条形图展示 Ke vs Kd 贡献，下方参数表格。

    数据来自 ``data["scenarios"][scen]["wacc"]``。
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    scenarios = data.get("scenarios", {})
    if not scenarios:
        _add_text(slide, Inches(0.5), Inches(3.0), SLIDE_W - Inches(1.0),
                  Inches(1.0), "（未提供情景数据 / No scenario data）",
                  size=14, italic=True, color=MUTED,
                  align=PP_ALIGN.CENTER, cjk=True)
        return

    scenario_keys = ["bear", "base", "bull"]
    labels_zh = {"bear": "熊市", "base": "基准", "bull": "牛市"}
    labels_en = {"bear": "Bear", "base": "Base", "bull": "Bull"}

    # ---- 上方：堆叠条形图 ----
    chart_left = Inches(0.5)
    chart_top = Inches(1.15)
    chart_w = Inches(7.5)
    chart_h = Inches(3.0)

    # 计算 Ke 和 Kd*(1-t)*D/(D+E) 贡献
    ke_values = []
    kd_values = []
    cat_labels = []
    for sk in scenario_keys:
        scen = scenarios.get(sk, {})
        wacc_d = scen.get("wacc", {})
        if not wacc_d:
            ke_values.append(0)
            kd_values.append(0)
            cat_labels.append(labels_zh.get(sk, sk) if language == "zh"
                             else labels_en.get(sk, sk))
            continue

        rf = float(wacc_d.get("rf", 0))
        erp = float(wacc_d.get("erp", 0))
        beta = float(wacc_d.get("beta", 1.0))
        size_p = float(wacc_d.get("size_premium", 0))
        cr = float(wacc_d.get("country_risk", 0))
        kd = float(wacc_d.get("kd", 0))
        tax = float(wacc_d.get("tax_rate", 0.21))
        e_w = float(wacc_d.get("e_weight", 1.0))
        d_w = float(wacc_d.get("d_weight", 0.0))

        ke = rf + erp * beta + size_p + cr
        kd_after = kd * (1 - tax) * d_w

        ke_values.append(round(ke * 100, 1))
        kd_values.append(round(kd_after * 100, 1))
        cat_labels.append(labels_zh.get(sk, sk) if language == "zh"
                         else labels_en.get(sk, sk))

    chart_data = CategoryChartData()
    chart_data.categories = cat_labels
    chart_data.add_series("Ke (Cost of Equity)", ke_values)
    chart_data.add_series("Kd×(1-t)×D/(D+E)", kd_values)

    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED,
                                chart_left, chart_top, chart_w, chart_h,
                                chart_data)
    chart = gf.chart
    chart.has_title = True
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.paragraphs[0].text = (
        "WACC Breakdown" if language == "en" else "WACC 分解")
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = NAVY
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(8)

    # 系列上色：Ke=#1F3864, Kd=#2E75B6
    s_ke = chart.series[0]
    s_kd = chart.series[1]
    s_ke.format.fill.solid()
    s_ke.format.fill.fore_color.rgb = RGBColor(0x1F, 0x38, 0x64)
    s_kd.format.fill.solid()
    s_kd.format.fill.fore_color.rgb = RGBColor(0x2E, 0x75, 0xB6)
    # 系列间白色分隔线
    for s in [s_ke, s_kd]:
        s.format.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        s.format.line.width = Pt(0.75)

    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(8)
    dl.font.bold = True
    dl.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    dl.number_format = '0.0"%"'

    # 坐标轴样式
    try:
        _style_chart_axis(chart.value_axis)
        chart.value_axis.major_gridlines.format.line.dash_style = 2  # 虚线
    except Exception:
        pass

    # ---- 右侧：WACC 数值汇总 ----
    summary_left = Inches(8.3)
    summary_top = Inches(1.15)
    summary_w = Inches(4.5)
    summary_h = Inches(3.0)

    _add_rect(slide, summary_left, summary_top, summary_w, summary_h, TABLE_ALT)
    _add_rect(slide, summary_left, summary_top, summary_w, Inches(0.35), DEEP_NAVY)
    _add_text(slide, summary_left + Inches(0.15), summary_top,
              summary_w - Inches(0.3), Inches(0.35),
              "WACC SUMMARY" if language == "en" else "WACC 汇总",
              size=10, bold=True, color=WHITE,
              align=PP_ALIGN.LEFT, anchor="m", cjk=True)

    for i, sk in enumerate(scenario_keys):
        scen = scenarios.get(sk, {})
        wacc_v = scen.get("wacc_value", 0)
        try:
            wacc_str = f"{float(wacc_v):.1%}"
        except (TypeError, ValueError):
            wacc_str = "—"
        lbl = labels_zh.get(sk, sk) if language == "zh" else labels_en.get(sk, sk)
        y_pos = summary_top + Inches(0.5) + Inches(0.7) * i
        clr = BASE_COLOR if sk == "base" else NAVY
        _add_text(slide, summary_left + Inches(0.2), y_pos,
                  Inches(1.5), Inches(0.35), lbl,
                  size=11, bold=True, color=NAVY, anchor="m", cjk=True)
        _add_text(slide, summary_left + Inches(1.8), y_pos,
                  summary_w - Inches(2.0), Inches(0.35), wacc_str,
                  size=18, bold=True, color=clr, anchor="m", cjk=False)

    # ---- 下方：CAPM 参数表格 ----
    param_top = Inches(4.4)
    param_labels = [
        ("Rf", "rf"),
        ("ERP", "erp"),
        ("Beta", "beta"),
        ("Size Premium", "size_premium"),
        ("Country Risk", "country_risk"),
        ("Ke", None),  # 计算
        ("Kd", "kd"),
        ("Tax Rate", "tax_rate"),
        ("E Weight", "e_weight"),
        ("D Weight", "d_weight"),
        ("WACC", None),  # 计算
    ]

    n_param_rows = len(param_labels) + 1  # header + data
    n_param_cols = 4  # label + 3 scenarios
    pt_left = Inches(0.5)
    pt_w = SLIDE_W - Inches(1.0)
    pr_h = Inches(0.32)
    pt_h = pr_h * n_param_rows

    ts2 = slide.shapes.add_table(n_param_rows, n_param_cols,
                                 pt_left, param_top, pt_w, pt_h)
    t2 = ts2.table
    t2.columns[0].width = int(pt_w * 0.25)
    for c in range(1, 4):
        t2.columns[c].width = int(pt_w * 0.25)

    # 表头
    _fill_table_cell(t2.cell(0, 0), "CAPM Parameter",
                     size=9, bold=True, color=WHITE, bg=NAVY,
                     align=PP_ALIGN.LEFT, cjk=False)
    for ci, sk in enumerate(scenario_keys):
        lbl = labels_zh.get(sk, sk) if language == "zh" else labels_en.get(sk, sk)
        _fill_table_cell(t2.cell(0, ci + 1), lbl,
                         size=9, bold=True, color=WHITE, bg=NAVY,
                         align=PP_ALIGN.CENTER, cjk=True)

    # 数据行
    for ri, (param_name, param_key) in enumerate(param_labels):
        row_idx = ri + 1
        bg_label = NAVY_ALT if ri % 2 == 0 else RGBColor(0x2A, 0x4A, 0x7A)
        _fill_table_cell(t2.cell(row_idx, 0), param_name,
                         size=9, bold=True, color=WHITE, bg=bg_label,
                         align=PP_ALIGN.LEFT, cjk=False)

        for ci, sk in enumerate(scenario_keys):
            scen = scenarios.get(sk, {})
            wacc_d = scen.get("wacc", {})
            is_base = (sk == "base")
            if is_base:
                bg = BASE_BG
            elif sk == "bear":
                bg = BEAR_BG if ri % 2 == 0 else WHITE
            elif sk == "bull":
                bg = BULL_BG if ri % 2 == 0 else WHITE
            else:
                bg = ALT_ROW if ri % 2 == 1 else WHITE

            if param_key is not None:
                raw = wacc_d.get(param_key, "—")
            elif param_name == "Ke":
                # 计算 Ke = Rf + ERP*Beta + Size Premium + Country Risk
                try:
                    ke_v = (float(wacc_d.get("rf", 0))
                            + float(wacc_d.get("erp", 0)) * float(wacc_d.get("beta", 1))
                            + float(wacc_d.get("size_premium", 0))
                            + float(wacc_d.get("country_risk", 0)))
                    raw = ke_v
                except (TypeError, ValueError):
                    raw = "—"
            elif param_name == "WACC":
                raw = scen.get("wacc_value", "—")
            else:
                raw = "—"

            # 格式化
            try:
                v = float(raw)
                if param_name in ("Beta", "E Weight", "D Weight"):
                    val_str = f"{v:.2f}"
                elif param_name == "Tax Rate":
                    val_str = f"{v:.0%}"
                else:
                    val_str = f"{v:.1%}"
            except (TypeError, ValueError):
                val_str = str(raw) if raw else "—"

            _fill_table_cell(t2.cell(row_idx, ci + 1), val_str,
                             size=9, bold=is_base,
                             color=NAVY if is_base else TEXT,
                             bg=bg, align=PP_ALIGN.CENTER, cjk=False)

    # 底部 source
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.4),
              SLIDE_W - Inches(1.0), Inches(0.3),
              "Source: Morgan Stanley Research · CAPM & WACC Decomposition",
              size=SIZE_SMALL, italic=True, color=MUTED,
              align=PP_ALIGN.LEFT, cjk=False)


# =============================================================================
# 12f. Slide 14 — slide_valuation_bridge 估值桥/瀑布图
# =============================================================================

def slide_valuation_bridge(slide, *, data: Dict[str, Any],
                           theme: str = "classic",
                           language: str = "zh",
                           title_cn: str = "估值桥",
                           title_en: str = "Valuation Bridge",
                           page_number: Optional[int] = None,
                           total_pages: Optional[int] = None) -> None:
    """估值桥/瀑布图：PV of FCF → +Terminal Value → -Net Debt → Equity Value。

    使用原生 python-pptx 堆叠条形图模拟瀑布效果。
    三个情景并排对比。
    数据来自 ``data["scenarios"]``。
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    scenarios = data.get("scenarios", {})
    if not scenarios:
        _add_text(slide, Inches(0.5), Inches(3.0), SLIDE_W - Inches(1.0),
                  Inches(1.0), "（未提供情景数据 / No scenario data）",
                  size=14, italic=True, color=MUTED,
                  align=PP_ALIGN.CENTER, cjk=True)
        return

    scenario_keys = ["bear", "base", "bull"]
    labels_zh = {"bear": "熊市", "base": "基准", "bull": "牛市"}
    labels_en = {"bear": "Bear", "base": "Base", "bull": "Bull"}

    # 瀑布图组件
    bridge_components = [
        ("PV of FCF", "pv_fcf"),
        ("PV of GGM TV", "pv_ggm_tv"),
        ("PV of Exit TV", "pv_exit_tv"),
        ("Net Debt", "net_debt"),
    ]

    # ---- 上方：堆叠条形图模拟瀑布 ----
    chart_left = Inches(0.5)
    chart_top = Inches(1.15)
    chart_w = Inches(8.5)
    chart_h = Inches(4.5)

    cat_labels = []
    for sk in scenario_keys:
        lbl = labels_zh.get(sk, sk) if language == "zh" else labels_en.get(sk, sk)
        cat_labels.append(lbl)

    # 准备数据：PV of FCF（基础）, PV GGM TV（增量）, PV Exit TV（增量）
    # Net Debt 取绝对值作为独立 Series（红色），堆叠图不支持负值
    # 公式：PV_FCF + PV_GGM_TV + PV_Exit_TV - Net_Debt = Equity_Value
    pv_fcf_vals = []
    pv_ggm_vals = []
    pv_exit_vals = []
    net_debt_vals = []  # 取绝对值，用红色标识

    for sk in scenario_keys:
        scen = scenarios.get(sk, {})
        pv_fcf = float(scen.get("pv_fcf", 0))
        pv_ggm = float(scen.get("pv_ggm_tv", 0))
        pv_exit = float(scen.get("pv_exit_tv", 0))
        nd = float(scen.get("net_debt", 0))

        pv_fcf_vals.append(pv_fcf)
        pv_ggm_vals.append(pv_ggm)
        pv_exit_vals.append(pv_exit)
        net_debt_vals.append(abs(nd))  # 取绝对值，红色标识表示减项

    chart_data = CategoryChartData()
    chart_data.categories = cat_labels
    chart_data.add_series("PV of FCF", pv_fcf_vals)
    chart_data.add_series("PV of GGM TV", pv_ggm_vals)
    chart_data.add_series("PV of Exit TV", pv_exit_vals)
    chart_data.add_series("Net Debt (subtract)", net_debt_vals)

    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED,
                                chart_left, chart_top, chart_w, chart_h,
                                chart_data)
    chart = gf.chart
    chart.has_title = True
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.paragraphs[0].text = (
        "Valuation Bridge" if language == "en" else "估值分解")
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = NAVY
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(8)

    # 4个系列颜色: PV FCF=#2E75B6, PV GGM TV=#1F3864, PV Exit TV=#C8A951, Net Debt=#B91C1C
    series_colors = [
        RGBColor(0x2E, 0x75, 0xB6),  # PV FCF - 中蓝
        RGBColor(0x1F, 0x38, 0x64),  # PV GGM TV - 深海军蓝
        RGBColor(0xC8, 0xA9, 0x51),  # PV Exit TV - MS GOLD
        RGBColor(0xB9, 0x1C, 0x1C),  # Net Debt - 红色
    ]
    for idx, s in enumerate(chart.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = series_colors[idx % len(series_colors)]
        # 系列间白色分隔线
        s.format.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        s.format.line.width = Pt(0.75)

    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(7)
    dl.font.bold = True
    dl.font.color.rgb = TEXT
    dl.number_format = '#,##0'

    # 坐标轴样式
    try:
        _style_chart_axis(chart.value_axis)
        chart.value_axis.major_gridlines.format.line.dash_style = 2  # 虚线
        chart.value_axis.tick_labels.number_format = '#,##0'
    except Exception:
        pass

    # ---- Equity Value 汇总标注 ----
    for i, sk in enumerate(scenario_keys):
        scen = scenarios.get(sk, {})
        pv_fcf = float(scen.get("pv_fcf", 0))
        pv_ggm = float(scen.get("pv_ggm_tv", 0))
        pv_exit = float(scen.get("pv_exit_tv", 0))
        nd = float(scen.get("net_debt", 0))
        eq_val = pv_fcf + pv_ggm + pv_exit - nd
        # 在每个堆叠条末端添加 Equity Value 标注
        _add_text(slide,
                  chart_left + chart_w - Inches(0.1),
                  chart_top + Inches(0.3) + Inches(i * 1.0),
                  Inches(1.2), Inches(0.3),
                  f"Eq: {eq_val:,.0f}",
                  size=8, bold=True, color=NAVY,
                  align=PP_ALIGN.RIGHT, anchor="m", cjk=False)

    # ---- 图表下方：公式说明 ----
    formula_top = chart_top + chart_h + Inches(0.05)
    formula_text = ("PV of FCF + PV of GGM TV + PV of Exit TV - Net Debt = Equity Value"
                    if language == "en"
                    else "PV FCF + PV GGM TV + PV Exit TV - 净债务 = 股权价值")
    _add_text(slide, chart_left, formula_top,
              chart_w, Inches(0.25),
              formula_text,
              size=9, italic=True, color=MUTED,
              align=PP_ALIGN.LEFT, cjk=True)

    # ---- 右侧：Equity Value 汇总 ----
    summary_left = Inches(9.3)
    summary_top = Inches(1.15)
    summary_w = Inches(3.5)
    summary_h = Inches(4.5)

    _add_rect(slide, summary_left, summary_top, summary_w, summary_h, TABLE_ALT)
    _add_rect(slide, summary_left, summary_top, summary_w, Inches(0.4), DEEP_NAVY)
    _add_text(slide, summary_left + Inches(0.15), summary_top,
              summary_w - Inches(0.3), Inches(0.4),
              "EQUITY VALUE" if language == "en" else "股权价值",
              size=10, bold=True, color=WHITE,
              align=PP_ALIGN.LEFT, anchor="m", cjk=True)

    for i, sk in enumerate(scenario_keys):
        scen = scenarios.get(sk, {})
        ev = scen.get("equity_value", 0)
        psv = scen.get("per_share_value", 0)
        try:
            ev_str = f"${float(ev):,.0f}M"
        except (TypeError, ValueError):
            ev_str = "—"
        try:
            psv_str = f"${float(psv):.1f}"
        except (TypeError, ValueError):
            psv_str = "—"

        lbl = labels_zh.get(sk, sk) if language == "zh" else labels_en.get(sk, sk)
        y_pos = summary_top + Inches(0.6) + Inches(1.2) * i

        # 情景标签
        clr = BASE_COLOR if sk == "base" else NAVY
        _add_text(slide, summary_left + Inches(0.2), y_pos,
                  summary_w - Inches(0.4), Inches(0.3), lbl,
                  size=11, bold=True, color=clr, anchor="m", cjk=True)
        # Equity Value
        _add_text(slide, summary_left + Inches(0.2), y_pos + Inches(0.3),
                  summary_w - Inches(0.4), Inches(0.4), ev_str,
                  size=16, bold=True, color=DEEP_NAVY, anchor="m", cjk=False)
        # Per Share
        _add_text(slide, summary_left + Inches(0.2), y_pos + Inches(0.7),
                  summary_w - Inches(0.4), Inches(0.3),
                  f"Per Share: {psv_str}",
                  size=10, color=MUTED, anchor="m", cjk=False)

    # 底部 source
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.4),
              SLIDE_W - Inches(1.0), Inches(0.3),
              "Source: Morgan Stanley Research · DCF Valuation Bridge",
              size=SIZE_SMALL, italic=True, color=MUTED,
              align=PP_ALIGN.LEFT, cjk=False)


# =============================================================================
# 12g. Slide 15 — slide_sensitivity_heatmap 敏感性热力图
# =============================================================================

def slide_sensitivity_heatmap(slide, *, data: Dict[str, Any],
                             theme: str = "classic",
                             language: str = "zh",
                             title_cn: str = "敏感性分析",
                             title_en: str = "Sensitivity Analysis",
                             page_number: Optional[int] = None,
                             total_pages: Optional[int] = None) -> None:
    """敏感性热力图：WACC (行) x TGR (列) -> Per Share Value。

    用颜色深浅表示数值高低（深蓝=高值，浅蓝=低值）。
    基准情景位置用金色边框标记。
    数据来自 ``data["sensitivity"]``。
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    sens = data.get("sensitivity", {})
    if not sens or not sens.get("matrix"):
        _add_text(slide, Inches(0.5), Inches(3.0), SLIDE_W - Inches(1.0),
                  Inches(1.0), "（未提供敏感性数据 / No sensitivity data）",
                  size=14, italic=True, color=MUTED,
                  align=PP_ALIGN.CENTER, cjk=True)
        return

    wacc_range = sens.get("wacc_range", [])
    tgr_range = sens.get("tgr_range", [])
    matrix = sens.get("matrix", [])
    base_wacc_idx = sens.get("base_wacc_idx", -1)
    base_tgr_idx = sens.get("base_tgr_idx", -1)

    # 自动校验：如果 base_wacc_idx 超出 wacc_range 范围，尝试自动匹配
    if wacc_range and (base_wacc_idx < 0 or base_wacc_idx >= len(wacc_range)):
        base_wacc_val = sens.get("base_wacc", None)
        if base_wacc_val is not None:
            for _i, _w in enumerate(wacc_range):
                if abs(float(_w) - float(base_wacc_val)) < 0.001:
                    base_wacc_idx = _i
                    break
        if base_wacc_idx < 0 or base_wacc_idx >= len(wacc_range):
            base_wacc_idx = len(wacc_range) // 2  # fallback 到中间值

    if not wacc_range or not tgr_range or not matrix:
        return

    n_wacc = len(wacc_range)
    n_tgr = len(tgr_range)

    # 计算颜色范围
    flat_vals = []
    for row in matrix:
        for v in row:
            try:
                flat_vals.append(float(v))
            except (TypeError, ValueError):
                pass
    if not flat_vals:
        return
    min_val = min(flat_vals)
    max_val = max(flat_vals)
    val_range = max_val - min_val if max_val != min_val else 1.0

    # ---- 左侧：热力图表格 ----
    n_cols = n_tgr + 1  # WACC label + TGR columns
    n_rows = n_wacc + 1  # TGR header + WACC rows

    table_left = Inches(0.5)
    table_top = Inches(1.2)
    table_w = Inches(9.5)
    row_h = Inches(0.55)
    table_h = row_h * n_rows

    ts = slide.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                table_w, table_h)
    table = ts.table

    # 列宽
    table.columns[0].width = Inches(1.2)
    col_w = int((table_w - Inches(1.2)) / n_tgr)
    for c in range(1, n_cols):
        table.columns[c].width = col_w

    # 表头行：TGR 值
    _fill_table_cell(table.cell(0, 0), "WACC \\ TGR",
                     size=9, bold=True, color=WHITE, bg=NAVY,
                     align=PP_ALIGN.CENTER, cjk=False)
    for ci, tgr in enumerate(tgr_range):
        try:
            tgr_str = f"{float(tgr):.1%}"
        except (TypeError, ValueError):
            tgr_str = str(tgr)
        bg = BASE_COLOR if ci == base_tgr_idx else NAVY
        _fill_table_cell(table.cell(0, ci + 1), tgr_str,
                         size=9, bold=True, color=WHITE, bg=bg,
                         align=PP_ALIGN.CENTER, cjk=False)

    # 数据行
    for ri, wacc in enumerate(wacc_range):
        row_idx = ri + 1
        try:
            wacc_str = f"{float(wacc):.1%}"
        except (TypeError, ValueError):
            wacc_str = str(wacc)

        bg_label = BASE_COLOR if ri == base_wacc_idx else NAVY_ALT
        _fill_table_cell(table.cell(row_idx, 0), wacc_str,
                         size=9, bold=True, color=WHITE, bg=bg_label,
                         align=PP_ALIGN.CENTER, cjk=False)

        for ci in range(n_tgr):
            if ri < len(matrix) and ci < len(matrix[ri]):
                raw = matrix[ri][ci]
            else:
                raw = "—"

            try:
                v = float(raw)
                # 归一化到 0-1
                norm = (v - min_val) / val_range
                # 从 bull_bg（低值）到 deep_navy（高值）
                r_c = int(0xDD - norm * (0xDD - 0x1F))
                g_c = int(0xEE - norm * (0xEE - 0x38))
                b_c = int(0xFF - norm * (0xFF - 0x64))
                cell_bg = RGBColor(r_c, g_c, b_c)
                val_str = f"${v:.1f}"
                # 高值用白色字体，低值用深色字体
                txt_color = WHITE if norm > 0.6 else TEXT
            except (TypeError, ValueError):
                cell_bg = TABLE_ALT
                val_str = str(raw) if raw else "—"
                txt_color = TEXT

            is_base = (ri == base_wacc_idx and ci == base_tgr_idx)
            _fill_table_cell(table.cell(row_idx, ci + 1), val_str,
                             size=10, bold=is_base,
                             color=DEEP_NAVY if is_base else txt_color,
                             bg=cell_bg, align=PP_ALIGN.CENTER, cjk=False)

            # 基准情景黄色边框
            if is_base:
                cell = table.cell(row_idx, ci + 1)
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                ln = tcPr.makeelement(qn("a:ln"), {
                    "w": "28800",  # 2pt
                    "cap": "flat",
                    "cmpd": "sng",
                    "algn": "ctr",
                })
                solidFill = ln.makeelement(qn("a:solidFill"), {})
                srgbClr = solidFill.makeelement(qn("a:srgbClr"), {
                    "val": "FFFF00",  # HIGHLIGHT_YELLOW
                })
                solidFill.append(srgbClr)
                ln.append(solidFill)
                tcPr.append(ln)

    # ---- 右侧：图例说明 ----
    legend_left = Inches(10.3)
    legend_top = Inches(1.2)
    legend_w = Inches(2.5)
    legend_h = Inches(3.5)

    _add_rect(slide, legend_left, legend_top, legend_w, legend_h, TABLE_ALT)
    _add_rect(slide, legend_left, legend_top, legend_w, Inches(0.35), NAVY)
    _add_text(slide, legend_left + Inches(0.1), legend_top,
              legend_w - Inches(0.2), Inches(0.35),
              "LEGEND" if language == "en" else "图例",
              size=9, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, anchor="m", cjk=True)

    # 颜色渐变示意
    gradient_top = legend_top + Inches(0.5)
    gradient_h = Inches(1.5)
    n_steps = 8
    step_h = int(gradient_h / n_steps)
    for si in range(n_steps):
        norm = si / max(n_steps - 1, 1)
        r_c = int(0xDD - norm * (0xDD - 0x1F))
        g_c = int(0xEE - norm * (0xEE - 0x38))
        b_c = int(0xFF - norm * (0xFF - 0x64))
        _add_rect(slide, legend_left + Inches(0.3),
                  gradient_top + step_h * si,
                  Inches(0.8), step_h,
                  RGBColor(r_c, g_c, b_c))

    _add_text(slide, legend_left + Inches(0.15), gradient_top - Inches(0.15),
              Inches(1.0), Inches(0.2), "Low",
              size=8, color=MUTED, cjk=False)
    _add_text(slide, legend_left + Inches(0.15), gradient_top + gradient_h,
              Inches(1.0), Inches(0.2), "High",
              size=8, color=MUTED, cjk=False)

    # 基准标记说明
    marker_top = gradient_top + gradient_h + Inches(0.3)
    _add_rect(slide, legend_left + Inches(0.3), marker_top,
              Inches(0.5), Inches(0.3), HIGHLIGHT_YELLOW,
              line=DEEP_NAVY)
    _add_text(slide, legend_left + Inches(1.0), marker_top,
              legend_w - Inches(1.2), Inches(0.3),
              "Base Case" if language == "en" else "基准情景",
              size=9, bold=True, color=DEEP_NAVY, anchor="m", cjk=True)

    # 数值范围
    range_top = marker_top + Inches(0.5)
    _add_text(slide, legend_left + Inches(0.15), range_top,
              legend_w - Inches(0.3), Inches(0.25),
              f"Range: ${min_val:.1f} - ${max_val:.1f}",
              size=9, color=MUTED, cjk=False)

    # 底部 source
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.4),
              SLIDE_W - Inches(1.0), Inches(0.3),
              "Source: Morgan Stanley Research · Sensitivity Matrix (WACC x TGR)",
              size=SIZE_SMALL, italic=True, color=MUTED,
              align=PP_ALIGN.LEFT, cjk=False)


# =============================================================================
# 13. 主入口 make_deck
# =============================================================================

# =============================================================================
# Slide 19 — Dual-Panel Side-by-Side Charts (双面板并排图表)
# =============================================================================

def slide_dual_chart_panel(slide, *,
                           left_title: str = "",
                           left_chart_data: Optional[Dict[str, Any]] = None,
                           right_title: str = "",
                           right_chart_data: Optional[Dict[str, Any]] = None,
                           title_cn: str = "双面板分析",
                           title_en: str = "Dual-Panel Analysis",
                           language: str = "zh",
                           page_number: Optional[int] = None,
                           total_pages: Optional[int] = None) -> None:
    """双面板并排图表：MS 报告中最常见的布局，左右各 50%。

    每个面板可包含：柱状图、折线图、指标卡片、小表格或文本。
    通过 ``*_chart_data`` 字典的 ``type`` 字段区分：
      - "bar": 柱状图（需 categories + series）
      - "line": 折线图（需 categories + series）
      - "metric": 指标卡片（需 value + label）
      - "table": 小表格（需 headers + rows）
      - "text": 文本内容（需 body 列表）
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    left_data = left_chart_data or {}
    right_data = right_chart_data or {}

    # 面板布局参数
    panel_gap = Inches(0.3)
    panel_top = Inches(1.25)
    panel_w = (SLIDE_W - Inches(1.0) - panel_gap) // 2
    panel_h = SLIDE_H - panel_top - Inches(0.5)

    # 绘制左面板
    _draw_dual_panel(slide, Inches(0.5), panel_top, panel_w, panel_h,
                     left_title or "Left Panel", left_data, language)

    # 绘制右面板
    right_left = Inches(0.5) + panel_w + panel_gap
    _draw_dual_panel(slide, right_left, panel_top, panel_w, panel_h,
                     right_title or "Right Panel", right_data, language)


def _draw_dual_panel(slide, left: int, top: int, width: int, height: int,
                     panel_title: str, data: Dict[str, Any],
                     language: str) -> None:
    """绘制单个面板：白色背景 + 金色边框 + 标题 + 内容。"""
    # 白色背景 + 金色细边框
    _add_rect(slide, left, top, width, height, WHITE, line=GOLD)

    # 面板标题（NAVY 粗体）
    _add_text(slide, left + Inches(0.15), top + Inches(0.1),
              width - Inches(0.3), Inches(0.35),
              str(panel_title), size=14, bold=True, color=NAVY,
              anchor="t", cjk=_has_cjk(panel_title))

    # 金色分隔线
    _add_rect(slide, left + Inches(0.15), top + Inches(0.45),
              width - Inches(0.3), Emu(8000), GOLD)

    content_top = top + Inches(0.6)
    content_left = left + Inches(0.15)
    content_w = width - Inches(0.3)
    content_h = height - Inches(0.75)

    panel_type = str(data.get("type", "text")).strip().lower()

    if panel_type == "bar" and data.get("categories"):
        _draw_panel_bar_chart(slide, content_left, content_top,
                              content_w, content_h, data)
    elif panel_type == "line" and data.get("categories"):
        _draw_panel_line_chart(slide, content_left, content_top,
                               content_w, content_h, data)
    elif panel_type == "metric":
        _draw_panel_metric(slide, content_left, content_top,
                           content_w, content_h, data, language)
    elif panel_type == "table" and data.get("headers"):
        _draw_panel_table(slide, content_left, content_top,
                           content_w, content_h, data, language)
    else:
        # 默认文本面板
        body = data.get("body") or data.get("text") or ["(No data)"]
        if isinstance(body, str):
            body = [body]
        tb = _add_textbox(slide, content_left, content_top,
                          content_w, content_h, anchor="t")
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(line)
            _apply_font(r, size=11, color=TEXT, cjk=_has_cjk(str(line)))
            p.space_after = Pt(4)


def _draw_panel_bar_chart(slide, left: int, top: int, width: int, height: int,
                          data: Dict[str, Any]) -> None:
    """在面板内绘制柱状图。"""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    categories = list(data.get("categories") or [])
    series_list = data.get("series") or []
    if not categories or not series_list:
        return

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(str(s.get("name", "Series")),
                              [float(v) for v in s.get("values", [])])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data)
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = True if len(series_list) > 1 else False
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(7)

    # MS 配色
    ms_bar_colors = [
        RGBColor(0x1F, 0x38, 0x64),
        RGBColor(0x2E, 0x75, 0xB6),
        RGBColor(0xC8, 0xA9, 0x51),
        RGBColor(0x00, 0xAF, 0x50),
        RGBColor(0xFB, 0x03, 0x01),
    ]
    for idx, s in enumerate(chart.series):
        fill = s.format.fill
        fill.solid()
        fill.fore_color.rgb = ms_bar_colors[idx % len(ms_bar_colors)]

    try:
        _style_chart_axis(chart.value_axis)
        chart.value_axis.tick_labels.font.size = Pt(7)
    except Exception:
        pass


def _draw_panel_line_chart(slide, left: int, top: int, width: int, height: int,
                            data: Dict[str, Any]) -> None:
    """在面板内绘制折线图。"""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE

    categories = list(data.get("categories") or [])
    series_list = data.get("series") or []
    if not categories or not series_list:
        return

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(str(s.get("name", "Series")),
                              [float(v) for v in s.get("values", [])])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, left, top, width, height, chart_data)
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = True if len(series_list) > 1 else False
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(7)

    line_colors = [
        RGBColor(0x1F, 0x38, 0x64),
        RGBColor(0xC8, 0xA9, 0x51),
        RGBColor(0x00, 0xAF, 0x50),
        RGBColor(0xFB, 0x03, 0x01),
        RGBColor(0x3B, 0x81, 0xB9),
    ]
    marker_styles = [
        XL_MARKER_STYLE.CIRCLE,
        XL_MARKER_STYLE.SQUARE,
        XL_MARKER_STYLE.TRIANGLE,
        XL_MARKER_STYLE.DIAMOND,
        XL_MARKER_STYLE.STAR,
    ]
    for idx, s in enumerate(chart.series):
        ln = s.format.line
        ln.color.rgb = line_colors[idx % len(line_colors)]
        ln.width = Pt(2)
        s.marker.style = marker_styles[idx % len(marker_styles)]
        s.marker.size = 4
        s.marker.format.fill.solid()
        s.marker.format.fill.fore_color.rgb = line_colors[idx % len(line_colors)]
        s.marker.format.line.color.rgb = line_colors[idx % len(line_colors)]

    try:
        _style_chart_axis(chart.value_axis)
        chart.value_axis.tick_labels.font.size = Pt(7)
    except Exception:
        pass


def _draw_panel_metric(slide, left: int, top: int, width: int, height: int,
                        data: Dict[str, Any], language: str) -> None:
    """在面板内绘制指标卡片。"""
    metrics = data.get("metrics") or []
    if not metrics and data.get("value"):
        metrics = [{"value": data["value"],
                     "label": data.get("label", ""),
                     "unit": data.get("unit", "")}]
    if not metrics:
        return

    n = min(len(metrics), 4)
    cols = 2 if n > 1 else 1
    rows_count = (n + cols - 1) // cols
    cell_w = (width - Inches(0.1) * (cols - 1)) // cols
    cell_h = (height - Inches(0.1) * (rows_count - 1)) // rows_count

    for i, m in enumerate(metrics):
        col = i % cols
        row = i // cols
        cx = left + col * (cell_w + Inches(0.1))
        cy = top + row * (cell_h + Inches(0.1))

        # 浅灰背景卡片
        _add_rect(slide, cx, cy, cell_w, cell_h, LIGHT_BG)
        # 金色顶部细线
        _add_rect(slide, cx, cy, cell_w, Emu(6000), GOLD)

        # 数值
        val_str = str(m.get("value", ""))
        _add_text(slide, cx + Inches(0.1), cy + Inches(0.15),
                  cell_w - Inches(0.2), Inches(0.5),
                  val_str, size=28, bold=True, color=NAVY,
                  anchor="m", cjk=False)
        # 单位
        unit_str = str(m.get("unit", ""))
        if unit_str:
            _add_text(slide, cx + Inches(0.1), cy + Inches(0.55),
                      cell_w - Inches(0.2), Inches(0.25),
                      unit_str, size=10, color=MUTED, anchor="t", cjk=True)
        # 标签
        label_str = str(m.get("label", ""))
        if label_str:
            _add_text(slide, cx + Inches(0.1), cy + cell_h - Inches(0.35),
                      cell_w - Inches(0.2), Inches(0.25),
                      label_str, size=9, color=MS_MUTED, anchor="b",
                      cjk=_has_cjk(label_str))


def _draw_panel_table(slide, left: int, top: int, width: int, height: int,
                      data: Dict[str, Any], language: str) -> None:
    """在面板内绘制小表格。"""
    headers = list(data.get("headers") or [])
    rows = list(data.get("rows") or [])
    if not headers:
        return

    n_cols = len(headers)
    n_rows = min(len(rows), 8) + 1  # 最多8行数据 + 1行表头
    row_h = min(Inches(0.35), (height - Inches(0.1)) // n_rows)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, row_h * n_rows)
    tbl = table_shape.table

    # 列宽
    col_w = width // n_cols
    for j in range(n_cols):
        tbl.columns[j].width = col_w

    # 表头
    for j, h in enumerate(headers):
        _fill_table_cell(tbl.cell(0, j), str(h),
                         size=8, bold=True, color=WHITE,
                         bg=MS_BRAND_BLUE, align=PP_ALIGN.CENTER, cjk=True)

    # 数据行
    for i in range(min(len(rows), n_rows - 1)):
        bg = WHITE if i % 2 == 0 else TABLE_ALT
        for j in range(n_cols):
            val = rows[i][j] if j < len(rows[i]) else ""
            _fill_table_cell(tbl.cell(i + 1, j), str(val),
                             size=8, color=TEXT, bg=bg,
                             align=PP_ALIGN.CENTER, cjk=_has_cjk(str(val)))


# =============================================================================
# Slide 20 — 2x2 Strategic Matrix (2x2 战略矩阵)
# =============================================================================

def slide_2x2_matrix(slide, *,
                      quadrants: Optional[List[Dict[str, Any]]] = None,
                      x_axis_label: str = "",
                      y_axis_label: str = "",
                      title_cn: str = "战略矩阵分析",
                      title_en: str = "2x2 Strategic Matrix",
                      language: str = "zh",
                      page_number: Optional[int] = None,
                      total_pages: Optional[int] = None) -> None:
    """四象限战略矩阵：MS 报告中常见的框架布局。

    参数：
      quadrants: 4 个字典列表，每个包含：
        - title: 象限标题
        - items: 要点列表
        - color: 象限颜色（支持 hex 字符串或 RGBColor）
        - position: "top-left" / "top-right" / "bottom-left" / "bottom-right"
      x_axis_label: X 轴标签
      y_axis_label: Y 轴标签
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    # 默认四象限数据
    if not quadrants:
        quadrants = [
            {"title": "High Growth / High Share",
             "items": ["Market leader", "Strong competitive moat"],
             "color": RGBColor(0x1F, 0x38, 0x64),
             "position": "top-left"},
            {"title": "High Growth / Low Share",
             "items": ["Emerging opportunity", "Requires investment"],
             "color": RGBColor(0x2E, 0x75, 0xB6),
             "position": "top-right"},
            {"title": "Low Growth / High Share",
             "items": ["Cash cow", "Harvest & optimize"],
             "color": RGBColor(0xC8, 0xA9, 0x51),
             "position": "bottom-left"},
            {"title": "Low Growth / Low Share",
             "items": ["Divest / exit", "Minimize exposure"],
             "color": RGBColor(0xB9, 0x1C, 0x1C),
             "position": "bottom-right"},
        ]

    # 矩阵布局
    matrix_left = Inches(0.8)
    matrix_top = Inches(1.35)
    matrix_w = SLIDE_W - Inches(1.6)
    matrix_h = SLIDE_H - matrix_top - Inches(0.8)

    half_w = (matrix_w - Inches(0.06)) // 2  # 0.06 = 2 * gap
    half_h = (matrix_h - Inches(0.06)) // 2
    gap = Inches(0.06)

    # 位置映射
    pos_map = {
        "top-left": (matrix_left, matrix_top),
        "top-right": (matrix_left + half_w + gap, matrix_top),
        "bottom-left": (matrix_left, matrix_top + half_h + gap),
        "bottom-right": (matrix_left + half_w + gap, matrix_top + half_h + gap),
    }

    for q in quadrants:
        pos = str(q.get("position", "")).strip().lower()
        if pos not in pos_map:
            continue
        qx, qy = pos_map[pos]

        # 解析颜色
        q_color = q.get("color", NAVY)
        if isinstance(q_color, str):
            q_color = RGBColor.from_string(q_color.replace("#", ""))

        # 白色背景
        _add_rect(slide, qx, qy, half_w, half_h, WHITE, line=RULE)

        # 顶部彩色标题栏
        header_h = Inches(0.4)
        _add_rect(slide, qx, qy, half_w, header_h, q_color)

        # 象限标题
        _add_text(slide, qx + Inches(0.1), qy + Inches(0.02),
                  half_w - Inches(0.2), header_h - Inches(0.04),
                  str(q.get("title", "")),
                  size=12, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor="m",
                  cjk=_has_cjk(str(q.get("title", ""))))

        # 要点列表
        items = list(q.get("items") or [])
        item_top = qy + header_h + Inches(0.08)
        for i, item in enumerate(items[:5]):  # 最多5条
            _add_text(slide, qx + Inches(0.15), item_top + Inches(0.28) * i,
                      half_w - Inches(0.3), Inches(0.28),
                      f"•  {item}",
                      size=10, color=TEXT, anchor="t",
                      cjk=_has_cjk(str(item)))

    # 轴标签
    if x_axis_label:
        _add_text(slide, matrix_left, matrix_top + matrix_h + Inches(0.05),
                  matrix_w, Inches(0.3),
                  str(x_axis_label), size=10, bold=True, color=NAVY,
                  align=PP_ALIGN.CENTER, anchor="t",
                  cjk=_has_cjk(str(x_axis_label)))
    if y_axis_label:
        # Y 轴标签（竖排，放在左侧）
        tb = _add_textbox(slide, matrix_left - Inches(0.6), matrix_top,
                          Inches(0.5), matrix_h, anchor="m")
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(y_axis_label)
        _apply_font(r, size=10, bold=True, color=NAVY,
                    cjk=_has_cjk(str(y_axis_label)))


# =============================================================================
# Slide 21 — Asset Allocation Dot Matrix (资产配置点阵矩阵)
# =============================================================================

def slide_asset_allocation(slide, *,
                           matrix_data: Optional[List[Dict[str, Any]]] = None,
                           columns: Optional[List[str]] = None,
                           title_cn: str = "资产配置立场",
                           title_en: str = "Asset Allocation Stance",
                           language: str = "zh",
                           page_number: Optional[int] = None,
                           total_pages: Optional[int] = None) -> None:
    """MS 风格资产配置立场矩阵（源自 "The BEAT" 报告）。

    行 = 资产类别，列 = 时间跨度或区域。
    单元格使用点符号表示立场：
      -- : 强烈低配 (Strong Underweight)
      -  : 低配 (Underweight)
      =  : 中性 (Neutral)
      +  : 超配 (Overweight)
      ++ : 强烈超配 (Strong Overweight)

    颜色编码：红色=负面，灰色=中性，蓝/绿色=正面。
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    # 默认列（时间跨度）
    if not columns:
        columns = ["1M", "3M", "6M", "12M"]

    # 默认矩阵数据
    if not matrix_data:
        matrix_data = [
            {"asset_class": "Equities", "asset_cn": "股票",
             "stances": ["+", "+", "++", "++"]},
            {"asset_class": "Credit", "asset_cn": "信用债",
             "stances": ["=", "+", "+", "="]},
            {"asset_class": "Rates", "asset_cn": "利率",
             "stances": ["-", "-", "=", "="]},
            {"asset_class": "FX (USD)", "asset_cn": "外汇 (美元)",
             "stances": ["=", "=", "+", "+"]},
            {"asset_class": "Commodities", "asset_cn": "大宗商品",
             "stances": ["+", "=", "=", "+"]},
            {"asset_class": "Alternatives", "asset_cn": "另类资产",
             "stances": ["=", "=", "=", "+"],},
        ]

    # 布局参数
    table_left = Inches(0.8)
    table_top = Inches(1.4)
    n_cols = len(columns)
    n_rows = len(matrix_data)

    # 资产类别列宽 + 立场列宽
    asset_col_w = Inches(2.2)
    stance_col_w = Inches(1.8)
    total_w = asset_col_w + stance_col_w * n_cols
    row_h = Inches(0.55)
    header_h = Inches(0.45)

    # 列标题行
    _add_rect(slide, table_left, table_top, total_w, header_h, NAVY)
    # 资产类别表头
    _add_text(slide, table_left + Inches(0.1), table_top,
              asset_col_w - Inches(0.2), header_h,
              "Asset Class" if language != "zh" else "资产类别",
              size=11, bold=True, color=WHITE,
              align=PP_ALIGN.LEFT, anchor="m", cjk=(language == "zh"))
    # 时间跨度表头
    for j, col_name in enumerate(columns):
        cx = table_left + asset_col_w + stance_col_w * j
        _add_text(slide, cx, table_top, stance_col_w, header_h,
                  str(col_name), size=11, bold=True, color=GOLD,
                  align=PP_ALIGN.CENTER, anchor="m", cjk=False)

    # 数据行
    for i, row_data in enumerate(matrix_data):
        ry = table_top + header_h + row_h * i

        # 交替行背景
        row_bg = WHITE if i % 2 == 0 else TABLE_ALT
        _add_rect(slide, table_left, ry, total_w, row_h, row_bg, line=RULE)

        # 资产类别名称
        asset_label = str(row_data.get("asset_class", ""))
        asset_cn = str(row_data.get("asset_cn", ""))
        label_text = f"{asset_label}" if language != "zh" else f"{asset_cn}"
        _add_text(slide, table_left + Inches(0.1), ry,
                  asset_col_w - Inches(0.2), row_h,
                  label_text, size=11, bold=True, color=NAVY,
                  align=PP_ALIGN.LEFT, anchor="m",
                  cjk=_has_cjk(label_text))

        # 立场单元格
        stances = list(row_data.get("stances") or [])
        for j in range(min(len(stances), n_cols)):
            cx = table_left + asset_col_w + stance_col_w * j
            stance = str(stances[j]).strip()

            # 解析立场
            dot_text, dot_color = _parse_stance(stance)

            _add_text(slide, cx, ry, stance_col_w, row_h,
                      dot_text, size=16, bold=True, color=dot_color,
                      align=PP_ALIGN.CENTER, anchor="m", cjk=False)

    # 底部图例
    legend_top = table_top + header_h + row_h * n_rows + Inches(0.2)
    legend_items = [
        ("--", "Strong UW", UW_RED),
        ("-", "Underweight", RGBColor(0xE0, 0x60, 0x60)),
        ("=", "Neutral", MUTED),
        ("+", "Overweight", RGBColor(0x2E, 0x75, 0xB6)),
        ("++", "Strong OW", RGBColor(0x00, 0xAF, 0x50)),
    ]
    legend_x = table_left
    for symbol, label, color in legend_items:
        _add_text(slide, legend_x, legend_top, Inches(0.5), Inches(0.25),
                  symbol, size=11, bold=True, color=color,
                  align=PP_ALIGN.CENTER, anchor="m", cjk=False)
        _add_text(slide, legend_x + Inches(0.45), legend_top, Inches(1.2), Inches(0.25),
                  label, size=9, color=MUTED, anchor="m", cjk=False)
        legend_x += Inches(1.8)


def _parse_stance(stance: str):
    """解析立场符号，返回 (显示文本, 颜色)。"""
    s = stance.strip().lower()
    if s in ("--", "strong underweight", "strong uw", "suw"):
        return "--", UW_RED
    if s in ("-", "underweight", "uw"):
        return "-", RGBColor(0xE0, 0x60, 0x60)
    if s in ("+", "overweight", "ow"):
        return "+", RGBColor(0x2E, 0x75, 0xB6)
    if s in ("++", "strong overweight", "strong ow", "sow"):
        return "++", RGBColor(0x00, 0xAF, 0x50)
    # 默认中性
    return "=", MUTED


# =============================================================================
# Slide 22 — Donut Chart with Center Metric (环形图 + 中心指标)
# =============================================================================

def slide_donut_chart(slide, *,
                      segments: Optional[List[Dict[str, Any]]] = None,
                      center_metric: str = "",
                      center_label: str = "",
                      title_cn: str = "构成分析",
                      title_en: str = "Composition Analysis",
                      language: str = "zh",
                      page_number: Optional[int] = None,
                      total_pages: Optional[int] = None) -> None:
    """环形图（Donut Chart）+ 中心大数字指标。

    参数：
      segments: 列表，每个元素含 {name, value, color}
      center_metric: 中心显示的大数字
      center_label: 中心数字下方的说明文字
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    # 默认数据
    if not segments:
        segments = [
            {"name": "Semiconductors", "value": 32.0,
             "color": RGBColor(0x1F, 0x38, 0x64)},
            {"name": "Power & Utilities", "value": 22.0,
             "color": RGBColor(0x2E, 0x75, 0xB6)},
            {"name": "Optics / Networking", "value": 15.0,
             "color": RGBColor(0xC8, 0xA9, 0x51)},
            {"name": "AI Software", "value": 13.0,
             "color": RGBColor(0x00, 0xAF, 0x50)},
            {"name": "Memory / Storage", "value": 10.0,
             "color": RGBColor(0xFB, 0x03, 0x01)},
            {"name": "Other", "value": 8.0,
             "color": RGBColor(0x66, 0x66, 0x66)},
        ]
    if not center_metric:
        total_val = sum(s.get("value", 0) for s in segments)
        center_metric = f"{total_val:.0f}%"
    if not center_label:
        center_label = "Total Allocation" if language != "zh" else "总配置"

    # MS 8 色调色板（用于未指定颜色的 segment）
    ms_8_palette = [
        RGBColor(0x1F, 0x38, 0x64),  # 深海军蓝
        RGBColor(0x2E, 0x75, 0xB6),  # 中蓝
        RGBColor(0xC8, 0xA9, 0x51),  # MS GOLD
        RGBColor(0x00, 0xAF, 0x50),  # 绿色
        RGBColor(0xFB, 0x03, 0x01),  # 红色
        RGBColor(0x3B, 0x81, 0xB9),  # 图表蓝
        RGBColor(0xE3, 0x7C, 0x2B),  # 橙色
        RGBColor(0x66, 0x66, 0x66),  # 灰色
    ]

    # 构建图表数据
    chart_data = CategoryChartData()
    chart_data.categories = [str(s.get("name", "")) for s in segments]
    chart_data.add_series("Value", [float(s.get("value", 0)) for s in segments])

    # 环形图位置（偏左，留出右侧图例空间）
    donut_left = Inches(0.5)
    donut_top = Inches(1.3)
    donut_w = Inches(6.5)
    donut_h = Inches(5.5)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, donut_left, donut_top, donut_w, donut_h,
        chart_data)
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = TEXT

    # 为每个 segment 上色
    plot = chart.plots[0]
    for idx, s in enumerate(chart.series[0].points):
        seg = segments[idx] if idx < len(segments) else {}
        seg_color = seg.get("color")
        if seg_color is None:
            seg_color = ms_8_palette[idx % len(ms_8_palette)]
        elif isinstance(seg_color, str):
            seg_color = RGBColor.from_string(seg_color.replace("#", ""))

        point = s.format.fill
        point.solid()
        point.fore_color.rgb = seg_color

    # 数据标签
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(9)
    dl.font.bold = True
    dl.font.color.rgb = TEXT
    dl.number_format = '0.0"%"'
    dl.show_percentage = False
    dl.show_value = True
    dl.show_category_name = False

    # 中心指标（叠加在环形图中央）
    center_x = donut_left + donut_w // 2 - Inches(1.5)
    center_y = donut_top + donut_h // 2 - Inches(0.5)
    _add_text(slide, center_x, center_y, Inches(3.0), Inches(0.7),
              str(center_metric), size=36, bold=True, color=NAVY,
              align=PP_ALIGN.CENTER, anchor="m", cjk=False)
    _add_text(slide, center_x, center_y + Inches(0.65), Inches(3.0), Inches(0.35),
              str(center_label), size=11, color=MUTED,
              align=PP_ALIGN.CENTER, anchor="t",
              cjk=_has_cjk(str(center_label)))


# =============================================================================
# Slide 23 — Stacked Bar Chart (堆叠柱状图)
# =============================================================================

def slide_stacked_bar(slide, *,
                       categories: Optional[List[str]] = None,
                       series_data: Optional[List[Dict[str, Any]]] = None,
                       title: str = "",
                       title_cn: str = "堆叠柱状图",
                       title_en: str = "Stacked Bar Chart",
                       horizontal: bool = False,
                       language: str = "zh",
                       page_number: Optional[int] = None,
                       total_pages: Optional[int] = None) -> None:
    """堆叠柱状图（水平或垂直）。

    参数：
      categories: 类别标签列表（如年份、地区等）
      series_data: 系列数据列表，每个元素含 {name, values, color}
      title: 图表标题（可选，覆盖 title_cn/title_en）
      horizontal: True 为水平堆叠条形图，False 为垂直堆叠柱状图
    """
    _slide_top_bar(slide, title_cn=title_cn, title_en=title_en,
                   page_number=page_number, total_pages=total_pages)

    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    # 默认数据
    if not categories:
        categories = ["FY22", "FY23", "FY24", "FY25E", "FY26E"]
    if not series_data:
        series_data = [
            {"name": "Hardware", "values": [400, 480, 560, 650, 720],
             "color": RGBColor(0x1F, 0x38, 0x64)},
            {"name": "Software", "values": [300, 380, 470, 580, 680],
             "color": RGBColor(0x2E, 0x75, 0xB6)},
            {"name": "Services", "values": [200, 250, 310, 380, 450],
             "color": RGBColor(0xC8, 0xA9, 0x51)},
            {"name": "Other", "values": [100, 120, 150, 180, 210],
             "color": RGBColor(0x66, 0x66, 0x66)},
        ]

    # MS 配色（用于未指定颜色的 series）
    ms_stack_colors = [
        RGBColor(0x1F, 0x38, 0x64),
        RGBColor(0x2E, 0x75, 0xB6),
        RGBColor(0xC8, 0xA9, 0x51),
        RGBColor(0x00, 0xAF, 0x50),
        RGBColor(0xFB, 0x03, 0x01),
        RGBColor(0x3B, 0x81, 0xB9),
        RGBColor(0xE3, 0x7C, 0x2B),
        RGBColor(0x66, 0x66, 0x66),
    ]

    # 构建图表数据
    chart_data = CategoryChartData()
    chart_data.categories = list(categories)
    for s in series_data:
        chart_data.add_series(str(s.get("name", "Series")),
                              [float(v) for v in s.get("values", [])])

    # 图表位置
    chart_left = Inches(0.5)
    chart_top = Inches(1.3)
    chart_w = SLIDE_W - Inches(1.0)
    chart_h = SLIDE_H - chart_top - Inches(0.6)

    # 选择图表类型
    if horizontal:
        chart_type = XL_CHART_TYPE.BAR_STACKED
    else:
        chart_type = XL_CHART_TYPE.COLUMN_STACKED

    chart_frame = slide.shapes.add_chart(
        chart_type, chart_left, chart_top, chart_w, chart_h, chart_data)
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(9)
    chart.legend.font.color.rgb = TEXT

    # 为每个 series 上色
    for idx, s in enumerate(chart.series):
        s_color = series_data[idx].get("color") if idx < len(series_data) else None
        if s_color is None:
            s_color = ms_stack_colors[idx % len(ms_stack_colors)]
        elif isinstance(s_color, str):
            s_color = RGBColor.from_string(s_color.replace("#", ""))

        fill = s.format.fill
        fill.solid()
        fill.fore_color.rgb = s_color

    # 数据标签
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(8)
    dl.font.bold = True
    dl.font.color.rgb = TEXT
    dl.number_format = '#,##0'

    # 坐标轴样式
    try:
        _style_chart_axis(chart.value_axis)
        chart.value_axis.tick_labels.font.size = Pt(8)
        if not horizontal:
            chart.value_axis.major_gridlines.format.line.dash_style = 2
    except Exception:
        pass

    # 自定义标题（如果提供）
    if title:
        _add_text(slide, chart_left, chart_top - Inches(0.02),
                  chart_w, Inches(0.3),
                  str(title), size=13, bold=True, color=NAVY,
                  align=PP_ALIGN.LEFT, anchor="b",
                  cjk=_has_cjk(str(title)))

    # 底部金色分隔线
    _add_rect(slide, Inches(0.5), SLIDE_H - Inches(0.5),
              SLIDE_W - Inches(1.0), Emu(8000), GOLD)
    # 来源标注
    _add_text(slide, Inches(0.5), SLIDE_H - Inches(0.4),
              SLIDE_W - Inches(1.0), Inches(0.25),
              "Source: Morgan Stanley Research",
              size=8, color=MUTED, align=PP_ALIGN.LEFT, cjk=False)


# =============================================================================
# 14. 页码与语言辅助
# =============================================================================

def _normalize_language(language: Optional[str]) -> str:
    v = (language or "zh").strip().lower()
    if v in ("zh", "cn", "chinese", "zh-cn", "zh_cn"):
        return "zh"
    if v in ("en", "english", "en-us", "en_us"):
        return "en"
    return "zh"


def make_deck(data: Dict[str, Any], output_path: str,
              theme: str = "classic", language: str = "zh") -> str:
    """构建并保存一份 Morgan Stanley 风格的研究演示文稿。

    参数
    ----
    data: 结构化数据字典。字段全部可选，缺失则跳过对应页：
        - company_name, title_cn/title_en, subtitle
        - rating / target_price / current_price
        - date_str / analyst / research_type
        - industry_view: 行业观点 "Attractive" / "In-Line" / "Cautious"（可选）
        - key_takeaways: List[str]
        - metrics: List[dict]
        - content_pages: List[dict] 每个含 title_cn/title_en/body/chart_img/source/data
        - rating_table: List[List]
        - shovel_stocks: List[dict]（支持 rating 字段触发评级列）
        - shovel_stocks_list: List[dict]（别名，同 shovel_stocks）
        - market_monitor: List[dict]（可选）
        - financial_chart: dict {years, revenue, ebitda, margins, takeaways}
        - executive_summary: dict {thesis, tp_and_upside, key_risks, catalysts}
        - scenarios: dict {bear/base/bull} — DCF 情景数据（可选，触发新 slides）
        - sensitivity: dict {wacc_range, tgr_range, matrix, ...} — 敏感性矩阵（可选）
        - whats_changed: List[dict] — 变动摘要（可选，触发 What's Changed slide）
        - thesis_charts: List[dict] — 核心论点图表数据（可选，触发 Thesis in Charts slide）
        - value_chain: List[dict] — 产业链数据（可选，触发 Value Chain slide）
        - disclosure_text: str（可选）
    output_path: 输出 .pptx 文件路径。
    theme: "classic"。
    language: "zh" / "en"。

    返回
    ----
    output_path 保存成功的文件路径。
    """
    language = _normalize_language(language)
    d = data or {}

    # ---------- 先预计算总页数 ----------
    n_pages = 1  # cover 始终存在
    if d.get("key_takeaways"):
        n_pages += 1
    # content_pages 及其前面插入的 section divider（1 页）
    pages = list(d.get("content_pages") or [])
    if pages:
        n_pages += 1 + len(pages)  # 1 section divider + N content pages
    if d.get("metrics"):
        n_pages += 1  # 1 section divider
        n_pages += 1  # metric blocks
    # financial_chart 页
    if d.get("financial_chart") is not None:
        n_pages += 1
    if d.get("rating_table"):
        n_pages += 1  # section divider
        n_pages += 1  # rating table
    # shovel_stocks（兼容 shovel_stocks_list 别名）
    shovel_data = d.get("shovel_stocks") or d.get("shovel_stocks_list")
    if shovel_data:
        n_pages += 1  # section divider
        n_pages += 1  # shovel stocks
    if d.get("market_monitor") is not None:
        n_pages += 1
    # executive_summary 页
    if d.get("executive_summary") is not None:
        n_pages += 1
    # 新增：情景相关 slides（需要 scenarios 数据）
    has_scenarios = bool(d.get("scenarios"))
    has_sensitivity = bool(d.get("sensitivity"))
    if has_scenarios:
        n_pages += 1  # scenario comparison
        n_pages += 1  # wacc breakdown
        n_pages += 1  # valuation bridge
    if has_sensitivity:
        n_pages += 1  # sensitivity heatmap
    # 新增：可选 slides
    has_whats_changed = bool(d.get("whats_changed"))
    has_thesis_charts = bool(d.get("thesis_charts"))
    has_value_chain = bool(d.get("value_chain"))
    if has_whats_changed:
        n_pages += 1
    if has_thesis_charts:
        n_pages += 1
    if has_value_chain:
        n_pages += 1
    # 新增 v2：5 种新 slide 类型
    has_dual_panel = bool(d.get("dual_panel"))
    has_matrix_2x2 = bool(d.get("matrix_2x2"))
    has_asset_allocation = bool(d.get("asset_allocation"))
    has_donut_data = bool(d.get("donut_data"))
    has_stacked_bar_data = bool(d.get("stacked_bar_data"))
    if has_dual_panel:
        n_pages += 1
    if has_matrix_2x2:
        n_pages += 1
    if has_asset_allocation:
        n_pages += 1
    if has_donut_data:
        n_pages += 1
    if has_stacked_bar_data:
        n_pages += 1
    n_pages += 1  # disclosure

    total_pages = n_pages
    current_page = 0

    def _next_pn() -> int:
        nonlocal current_page
        current_page += 1
        return current_page

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # -------- 封面 --------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_cover(
        slide,
        title_cn=d.get("title_cn", ""),
        title_en=d.get("title_en", ""),
        subtitle=d.get("subtitle", ""),
        rating=d.get("rating", ""),
        target_price=d.get("target_price"),
        current_price=d.get("current_price"),
        company_name=d.get("company_name", ""),
        analyst=d.get("analyst", ""),
        date_str=d.get("date_str", ""),
        research_type=d.get("research_type", "Foundation"),
        industry_view=d.get("industry_view", ""),
        whats_changed=d.get("whats_changed"),
        language=language,
    )
    _next_pn()  # cover = 1

    # -------- What's Changed（封面后，如果有数据） --------
    if has_whats_changed:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_whats_changed(slide, data=d,
                             language=language,
                             page_number=_next_pn(),
                             total_pages=total_pages)

    # -------- Key Takeaways --------
    if d.get("key_takeaways"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_key_takeaways(slide, items=list(d["key_takeaways"]),
                            language=language,
                            page_number=_next_pn(),
                            total_pages=total_pages)

    # -------- 章节分隔 + 标准分析页 --------
    if pages:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_section_divider(
            slide,
            title_cn=pages[0].get("section_cn") or "分析内容",
            title_en=pages[0].get("section_en") or "Research Analysis",
            category=d.get("research_type", "Foundation"),
            language=language,
        )
        _next_pn()
        for page in pages:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide_content(
                slide,
                title_cn=page.get("title_cn", ""),
                title_en=page.get("title_en", ""),
                body=list(page.get("body") or []),
                chart_img=page.get("chart_img"),
                data=(page.get("data") if page.get("data") else {"sectors_allocation": d.get("sectors_allocation")}),
                source=page.get("source") or "Source: Morgan Stanley Research",
                language=language,
                page_number=_next_pn(),
                total_pages=total_pages,
            )

    # -------- Stat / Metric Blocks --------
    if d.get("metrics"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_section_divider(slide, title_cn="关键指标",
                              title_en="Key Metrics",
                              category=d.get("research_type", "Foundation"),
                              language=language)
        _next_pn()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_metric_blocks(slide, metrics=list(d["metrics"]),
                            language=language,
                            page_number=_next_pn(),
                            total_pages=total_pages)

    # -------- Financial Chart --------
    if d.get("financial_chart") is not None:
        fc = d["financial_chart"]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_financial_chart(
            slide,
            years=fc.get("years"),
            revenue=fc.get("revenue"),
            ebitda=fc.get("ebitda"),
            margins=fc.get("margins"),
            takeaways=fc.get("takeaways"),
            scenarios=d.get("scenarios"),  # 传入 scenarios 数据以启用预测路径
            title_cn=fc.get("title_cn") or "收入 / 利润率一览",
            title_en=fc.get("title_en") or "Revenue & EBITDA \u00b7 Margin Trends",
            language=language,
            page_number=_next_pn(),
            total_pages=total_pages,
        )

    # -------- Dual Chart Panel（双面板并排图表，financial_chart 后） --------
    if has_dual_panel:
        dp = d["dual_panel"]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_dual_chart_panel(
            slide,
            left_title=dp.get("left_title", ""),
            left_chart_data=dp.get("left_chart_data"),
            right_title=dp.get("right_title", ""),
            right_chart_data=dp.get("right_chart_data"),
            title_cn=dp.get("title_cn") or "双面板分析",
            title_en=dp.get("title_en") or "Dual-Panel Analysis",
            language=language,
            page_number=_next_pn(),
            total_pages=total_pages,
        )

    # -------- Stacked Bar Chart（堆叠柱状图，financial section） --------
    if has_stacked_bar_data:
        sb = d["stacked_bar_data"]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_stacked_bar(
            slide,
            categories=sb.get("categories"),
            series_data=sb.get("series_data"),
            title=sb.get("title", ""),
            title_cn=sb.get("title_cn") or "堆叠柱状图",
            title_en=sb.get("title_en") or "Stacked Bar Chart",
            horizontal=sb.get("horizontal", False),
            language=language,
            page_number=_next_pn(),
            total_pages=total_pages,
        )

    # -------- WACC Breakdown（需要 scenarios 数据） --------
    if has_scenarios:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_wacc_breakdown(
            slide,
            data=d,
            theme=theme,
            language=language,
            page_number=_next_pn(),
            total_pages=total_pages,
        )

    # -------- Valuation Bridge（需要 scenarios 数据） --------
    if has_scenarios:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_valuation_bridge(
            slide,
            data=d,
            theme=theme,
            language=language,
            page_number=_next_pn(),
            total_pages=total_pages,
        )

    # -------- Sensitivity Heatmap（需要 sensitivity 数据） --------
    if has_sensitivity:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_sensitivity_heatmap(
            slide,
            data=d,
            theme=theme,
            language=language,
            page_number=_next_pn(),
            total_pages=total_pages,
        )

    # -------- Rating Table --------
    if d.get("rating_table"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_section_divider(slide, title_cn="投资评级总览",
                              title_en="Ratings & Target Prices",
                              category=d.get("research_type", "Foundation"),
                              language=language)
        _next_pn()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 兼容两种格式：
        #   List[List]: [["NVDA.O", "Overweight", "$126.5", "NVIDIA", "reason"]]
        #   List[Dict]: [{"metric": "Revenue Growth", "score": "★★★★☆", "comment": "..."}]
        raw_rt = d["rating_table"]
        if raw_rt and isinstance(raw_rt[0], dict):
            # Dict 格式 -> 转为 List[List]，映射到 Ticker/Rating/Last/Company/Reason
            rt_rows = []
            for item in raw_rt:
                rt_rows.append([
                    str(item.get("ticker", item.get("metric", ""))),
                    str(item.get("score", item.get("rating", ""))),
                    str(item.get("last_close", item.get("target_price", ""))),
                    str(item.get("company", item.get("name", ""))),
                    str(item.get("comment", item.get("reason", item.get("note", "")))),
                ])
        else:
            rt_rows = list(raw_rt)
        slide_rating_table(slide, rows=rt_rows,
                           language=language,
                           page_number=_next_pn(),
                           total_pages=total_pages)

    # -------- Shovel Stocks Table --------
    if shovel_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_section_divider(slide, title_cn="铲子股清单",
                              title_en="Shovel Stocks \u2014 Thematic Opportunity",
                              category=d.get("research_type", "Foundation"),
                              language=language)
        _next_pn()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_shovel_stocks(slide, stocks=list(shovel_data),
                            language=language,
                            page_number=_next_pn(),
                            total_pages=total_pages)

    # -------- Market Monitor（可选） --------
    if d.get("market_monitor") is not None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_market_monitor(slide, blocks=list(d["market_monitor"]),
                             language=language,
                             page_number=_next_pn(),
                             total_pages=total_pages)

    # -------- Asset Allocation（资产配置点阵矩阵，market_monitor 后） --------
    if has_asset_allocation:
        aa = d["asset_allocation"]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_asset_allocation(
            slide,
            matrix_data=aa.get("matrix_data"),
            columns=aa.get("columns"),
            title_cn=aa.get("title_cn") or "资产配置立场",
            title_en=aa.get("title_en") or "Asset Allocation Stance",
            language=language,
            page_number=_next_pn(),
            total_pages=total_pages,
        )

    # -------- Executive Summary（可选） --------
    if d.get("executive_summary") is not None:
        es = d["executive_summary"]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_executive_summary(
            slide,
            thesis=es.get("thesis"),
            tp_and_upside=es.get("tp_and_upside"),
            key_risks=es.get("key_risks"),
            catalysts=es.get("catalysts"),
            title_cn=es.get("title_cn") or "执行摘要",
            title_en=es.get("title_en") or "Executive Summary \u00b7 Thesis, TP, Risks & Catalysts",
            language=language,
            page_number=_next_pn(),
            total_pages=total_pages,
        )

    # -------- Thesis in Charts（执行摘要后，如果有数据） --------
    if has_thesis_charts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_thesis_in_charts(slide, data=d,
                               language=language,
                               page_number=_next_pn(),
                               total_pages=total_pages)

    # -------- Donut Chart（环形图，thesis section） --------
    if has_donut_data:
        dd = d["donut_data"]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_donut_chart(
            slide,
            segments=dd.get("segments"),
            center_metric=dd.get("center_metric", ""),
            center_label=dd.get("center_label", ""),
            title_cn=dd.get("title_cn") or "构成分析",
            title_en=dd.get("title_en") or "Composition Analysis",
            language=language,
            page_number=_next_pn(),
            total_pages=total_pages,
        )

    # -------- Scenario Comparison（需要 scenarios 数据） --------
    if has_scenarios:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_scenario_comparison(
            slide,
            data=d,
            theme=theme,
            language=language,
            page_number=_next_pn(),
            total_pages=total_pages,
        )

    # -------- 2x2 Strategic Matrix（scenario_comparison 后） --------
    if has_matrix_2x2:
        m2 = d["matrix_2x2"]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_2x2_matrix(
            slide,
            quadrants=m2.get("quadrants"),
            x_axis_label=m2.get("x_axis_label", ""),
            y_axis_label=m2.get("y_axis_label", ""),
            title_cn=m2.get("title_cn") or "战略矩阵分析",
            title_en=m2.get("title_en") or "2x2 Strategic Matrix",
            language=language,
            page_number=_next_pn(),
            total_pages=total_pages,
        )

    # -------- Value Chain（情景分析后，如果有数据） --------
    if has_value_chain:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_value_chain(slide, data=d,
                          language=language,
                          page_number=_next_pn(),
                          total_pages=total_pages)

    # -------- Disclosure --------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_disclosure(slide, text=d.get("disclosure_text"),
                     language=language,
                     page_number=_next_pn(),
                     total_pages=total_pages)

    # 保存
    out_dir = os.path.dirname(os.path.abspath(output_path))
    if out_dir and not os.path.isdir(out_dir):
        os.makedirs(out_dir, exist_ok=True)
    prs.save(output_path)
    return output_path


# =============================================================================
# 15. 示例数据 sample_data_financial —— 单独返回财务序列示例
# =============================================================================

def sample_data_financial() -> Dict[str, Any]:
    """返回一份符合 ``financial_chart`` 字段要求的样例财务数据。"""
    return {
        "years": ["FY22", "FY23", "FY24", "FY25E", "FY26E", "FY27E", "FY28E"],
        "revenue": [1200.0, 1450.0, 1780.0, 2120.0, 2460.0, 2810.0, 3180.0],
        "ebitda":  [280.0,  348.0,  442.0,  550.0,  665.0,  780.0,  910.0],
        "margins": {
            "gross":  [55.0, 56.5, 58.2, 60.0, 61.5, 62.8, 63.9],
            "ebitda": [23.3, 24.0, 24.8, 25.9, 27.0, 27.8, 28.6],
            "net":    [14.2, 15.0, 16.1, 17.0, 17.8, 18.5, 19.1],
        },
        "takeaways": [
            "Revenue CAGR ~21%，AI 相关产品贡献主要增量。",
            "EBITDA margin 自 23.3% 稳步扩张至 28.6%，规模效应显现。",
            "Gross margin 受益于产品结构优化，从 55.0% 提升至 63.9%。",
        ],
    }



def sample_data() -> Dict[str, Any]:
    """返回一份符合 ``make_deck`` 接口的样例数据字典。"""
    return {
        "company_name": "极光资本控股（示例）",
        "title_cn": "为下一个 10 倍科技周期——资本是瓶颈吗？",
        "title_en": "Financing the Next 10x Tech Cycle",
        "subtitle": "资本开支浪潮已成共识，关注供给约束下的结构性机会",
        "rating": "Overweight",
        "target_price": "$125",
        "current_price": "$98.5",
        "date_str": "2026-06-12",
        "analyst": "陈嘉怡 · CFA",
        "research_type": "Foundation",

        "key_takeaways": [
            "全球 AI 相关资本开支预计在 2026 年达到 ~$600B，同比增长 ~38%。",
            "电力与供应链约束日益明显，电力采购协议（PPA）长期价格进入新稳态。",
            "We estimate 芯片供给在 2H26 前仍紧俏，领先厂商享受定价权。",
            "对下游企业而言，单位算力成本下降 ~30% 前，商业化模型仍需打磨。",
            "关注铲子股：电力、散热、光模块、HBM、低轨卫星构成上游机会组合。",
        ],

        "metrics": [
            {"value": "~$600B", "unit": "AI CapEx 2026E",
             "label": "全球 AI 相关资本开支"},
            {"value": "+38%", "unit": "YoY Growth",
             "label": "AI CapEx 同比增速"},
            {"value": "32.5–60mn", "unit": "MW under PPA",
             "label": "长期电力采购协议需求"},
            {"value": "~1,128,000", "unit": "GPUs shipped 2026E",
             "label": "全球 GPU 出货量（估算）"},
            {"value": "(2,840,890)", "unit": "Net income impact",
             "label": "（负值示例）财务影响"},
            {"value": "+18%", "unit": "Clean-tech Capex",
             "label": "能源转型资本开支 YoY"},
        ],

        "content_pages": [
            {
                "title_cn": "AI 智能体崛起——全球影响",
                "title_en": "The Rise of AI Agents — Global Implications",
                "body": [
                    "我们认为，AI agent 的规模化部署将在未来 12-18 个月内",
                    "显著改变企业软件支出结构。早期采用者集中在金融服务、",
                    "医疗与零售业，2026E agent-related software 支出预计",
                    "~$45B，占企业软件总支出约 5.2%。",
                    "",
                    "Crucially，供给侧压力仍在：HBM 供给紧张、电力与冷却",
                    "限制成为跨区域的共性约束。我们预期领先的数据中心运营商",
                    "将通过长期 PPA 与一体化散热方案维持领先。",
                    "",
                    "Key investment implications: prefer power/semicon equipment",
                    "vendors with visible 2026E order book；选择性增持软件",
                    "公司中具备 agent 产品化路径的龙头。",
                ],
                "chart_img": "",  # 留空则使用占位卡片
            },
            {
                "title_cn": "能源冲击 / 改革路线图 / 科技创新",
                "title_en": "Energy Shock · Reform Roadmap · Tech Innovation",
                "body": [
                    "从电力结构看，AI 数据中心的单位电力需求约为传统园区的",
                    "8-10 倍。我们预计 2026E 全球数据中心电力消费",
                    "c.720 TWh，相当于法国全国年度用电量的 ~1.2x。",
                    "",
                    "这一结构性变化驱动三条投资主线：",
                    "1) 低碳稳定基荷（核电、天然气与储能），",
                    "2) 先进散热与电力电子设备，",
                    "3) 工业能效与电网数字化升级。",
                ],
                "chart_img": "",
            },
        ],

        "rating_table": [
            ["NVDA.O", "Overweight", "$126.5", "NVIDIA", "AI compute 领导者"],
            ["TSM.N", "Overweight", "$178.2", "Taiwan Semiconductor",
             "先进制程独家地位"],
            ["AVGO.O", "Overweight", "$1,180", "Broadcom",
             "AI 定制 ASIC + 网络芯片"],
            ["AMZN.O", "Equal-weight", "$188.4", "Amazon",
             "零售 / 云双轮驱动，估值已较充分"],
            ["GOOGL.O", "Equal-weight", "$165.0", "Alphabet",
             "搜索广告稳健，agent 路线仍待验证"],
            ["META.O", "Underweight", "$505.0", "Meta Platforms",
             "Reels 投放效率与资本开支节奏"],
        ],

        "shovel_stocks": [
            {"rank": 1, "company": "Aurora Power Co.", "ticker": "AUR.N",
             "product": "高压直流输电与数据中心电力模块",
             "analyst": "S. Lee", "market_cap_mn": 3200.5,
             "perf_1y_pct": 142.3},
            {"rank": 2, "company": "ThermoCore Labs", "ticker": "THC.O",
             "product": "液冷散热解决方案（液冷板 + CDU）",
             "analyst": "K. Zhang", "market_cap_mn": 1820.0,
             "perf_1y_pct": 88.6},
            {"rank": 3, "company": "Lumen Optics", "ticker": "LO.O",
             "product": "800G / 1.6T 光模块",
             "analyst": "M. Sato", "market_cap_mn": 4150.2,
             "perf_1y_pct": 118.4},
            {"rank": 4, "company": "HBM Memory Inc.", "ticker": "HBM.N",
             "product": "HBM3e / HBM4 高带宽内存",
             "analyst": "J. Park", "market_cap_mn": 5620.8,
             "perf_1y_pct": 205.1},
            {"rank": 5, "company": "Starlink Comm.", "ticker": "STL.N",
             "product": "低轨卫星通信与地面站设备",
             "analyst": "A. Rossi", "market_cap_mn": 2750.0,
             "perf_1y_pct": 62.4},
            {"rank": 6, "company": "Nimbus Silicon", "ticker": "NBS.O",
             "product": "AI 推理芯片（Chiplet 架构）",
             "analyst": "陈嘉怡", "market_cap_mn": 6380.0,
             "perf_1y_pct": 95.7},
            {"rank": 7, "company": "GridEdge AI", "ticker": "GEI.N",
             "product": "数据中心智能电网调度软件",
             "analyst": "L. Müller", "market_cap_mn": 980.5,
             "perf_1y_pct": 42.1},
            {"rank": 8, "company": "Helium Cooling", "ticker": "HCL.O",
             "product": "浸入式冷却与工质服务",
             "analyst": "S. Lee", "market_cap_mn": 1520.3,
             "perf_1y_pct": 128.9},
        ],

        "market_monitor": [
            {"name_en": "Bonds", "name_cn": "固收",
             "metric": "3.85%", "sub": "10Y UST yield",
             "body": "实际利率上行限制估值扩张，增持短久期信用债。"},
            {"name_en": "Equities", "name_cn": "股票",
             "metric": "21,420", "sub": "S&P 500 level",
             "body": "青睐高质量科技 / 能源设备龙头。"},
            {"name_en": "Alternatives", "name_cn": "另类",
             "metric": "$2.4T", "sub": "PE dry powder",
             "body": "私募资本等待公开市场折价机会。"},
            {"name_en": "Transition", "name_cn": "能源转型",
             "metric": "+18%", "sub": "Clean-tech capex",
             "body": "核电、电网数字化与工业节能为结构性主线。"},
        ],

        # 新增：财务图表页数据（可覆盖 sample_data_financial() 返回值）
        "financial_chart": sample_data_financial(),
        # 新增：执行摘要页数据（四象限）
        "executive_summary": {
            "thesis": [
                "结构性 AI 需求驱动收入快速增长，23-28E CAGR ~21%。",
                "供给侧约束短期仍紧俏，公司享有定价权与份额提升。",
                "运营杠杆释放推动 EBITDA margin 持续扩张至 28%+。",
            ],
            "tp_and_upside": {
                "rating": "Overweight", "target_price": "$125",
                "last_price": "$98.5", "upside": "+27%",
                "notes": ["目标价基于 DCF (WACC 9.5%, g 2.5%) 与同业估值对比。",
                          "12 个月视野，对应 FY27E EV/EBITDA ~18x。"],
            },
            "key_risks": [
                "AI 资本开支节奏不及预期或大客户砍单。",
                "电力 / 冷却供给瓶颈导致数据中心建设延期。",
                "地缘政治摩擦导致出口管制或供应链风险。",
                "估值扩张依赖持续超预期，增速回归时存在压缩风险。",
            ],
            "catalysts": [
                "下一季指引上调 / 超预期订单披露。",
                "新一代 AI 加速器量产节奏与客户样品反馈。",
                "大型云厂商 capex 加总确认上行周期延续。",
                "行业并购与生态合作带来份额与毛利率改善。",
            ],
        },
        # 新增：内容页饼图默认行业配置（slide_content 右侧图表使用）
        "sectors_allocation": {
            "Semiconductors": 32.0, "Power & Utilities": 22.0,
            "Optics / Networking": 15.0, "AI Software": 13.0,
            "Memory / Storage": 10.0, "Other": 8.0,
        },

        # 新增：DCF 情景数据（触发情景对比 / WACC 拆解 / 估值桥 slides）
        "scenarios": {
            "bear": {
                "label_zh": "熊市情景", "label_en": "Bear Case",
                "wacc": {
                    "rf": 0.044, "erp": 0.055, "beta": 1.8,
                    "size_premium": 0.005, "country_risk": 0.0,
                    "kd": 0.065, "tax_rate": 0.21,
                    "e_weight": 0.90, "d_weight": 0.10,
                },
                "wacc_value": 0.148,
                "terminal_growth_rate": 0.015,
                "exit_multiple_ebitda": 8.0,
                "pv_fcf": 120,
                "pv_ggm_tv": 80,
                "pv_exit_tv": 150,
                "enterprise_value": 350,
                "net_debt": 50,
                "equity_value": 300,
                "per_share_value": 30.0,
                "ev_revenue": 2.3,
                "ev_ebitda": 7.5,
                "revenue_forecast": [205, 210, 215, 218, 220],
            },
            "base": {
                "label_zh": "基准情景", "label_en": "Base Case",
                "wacc": {
                    "rf": 0.044, "erp": 0.055, "beta": 1.5,
                    "size_premium": 0.003, "country_risk": 0.0,
                    "kd": 0.060, "tax_rate": 0.21,
                    "e_weight": 0.92, "d_weight": 0.08,
                },
                "wacc_value": 0.125,
                "terminal_growth_rate": 0.025,
                "exit_multiple_ebitda": 12.0,
                "pv_fcf": 150,
                "pv_ggm_tv": 200,
                "pv_exit_tv": 250,
                "enterprise_value": 600,
                "net_debt": 50,
                "equity_value": 550,
                "per_share_value": 55.0,
                "ev_revenue": 3.5,
                "ev_ebitda": 12.5,
                "revenue_forecast": [216, 229, 243, 258, 275],
            },
            "bull": {
                "label_zh": "牛市情景", "label_en": "Bull Case",
                "wacc": {
                    "rf": 0.044, "erp": 0.055, "beta": 1.2,
                    "size_premium": 0.001, "country_risk": 0.0,
                    "kd": 0.055, "tax_rate": 0.21,
                    "e_weight": 0.95, "d_weight": 0.05,
                },
                "wacc_value": 0.108,
                "terminal_growth_rate": 0.030,
                "exit_multiple_ebitda": 16.0,
                "pv_fcf": 180,
                "pv_ggm_tv": 350,
                "pv_exit_tv": 400,
                "enterprise_value": 930,
                "net_debt": 50,
                "equity_value": 880,
                "per_share_value": 88.0,
                "ev_revenue": 5.0,
                "ev_ebitda": 18.0,
                "revenue_forecast": [220, 240, 262, 288, 318],
            },
        },

        # 新增：敏感性矩阵数据（触发敏感性热力图 slide）
        "sensitivity": {
            "wacc_range": [0.09, 0.095, 0.10, 0.105, 0.11, 0.115, 0.12, 0.125, 0.13],
            "tgr_range": [0.01, 0.015, 0.02, 0.025, 0.03],
            "matrix": [
                [42.0, 46.5, 52.0, 59.0, 68.5],
                [38.5, 42.0, 47.0, 53.0, 61.0],
                [35.0, 38.5, 42.5, 48.0, 55.0],
                [32.0, 35.5, 39.0, 43.5, 49.5],
                [29.5, 32.5, 36.0, 40.0, 45.0],
                [27.0, 30.0, 33.0, 37.0, 41.5],
                [25.0, 27.5, 30.5, 34.0, 38.0],
                [23.0, 25.5, 28.0, 31.0, 35.0],
                [21.5, 23.5, 26.0, 29.0, 32.5],
            ],
            "base_wacc_idx": 7,  # 12.5% 在 wacc_range 中的索引
            "base_tgr_idx": 2,    # 2.0% 在 tgr_range 中的索引
        },

        "disclosure_text": None,  # 留空使用默认英文+中文免责文本

        # ---- 新增：MS 报告风格扩展字段 ----
        "industry_view": "Attractive",  # Attractive / In-Line / Cautious

        "whats_changed": [
            {"item": "Rating", "from": "Equal-weight", "to": "Overweight", "arrow": "up"},
            {"item": "Target Price", "from": "$85", "to": "$125", "arrow": "up"},
            {"item": "AI CapEx Estimate", "from": "$520B", "to": "$600B", "arrow": "up"},
            {"item": "GPU Supply View", "from": "Balanced", "to": "Tight", "arrow": "down"},
        ],

        "thesis_charts": [
            {"title": "AI CapEx Surge", "value": "$600B",
             "subtitle": "2026E global AI capex (+38% YoY)", "color": "blue"},
            {"title": "GPU Shipments", "value": "1.1M",
             "subtitle": "Units shipped 2026E", "color": "green"},
            {"title": "Data Center Power", "value": "720 TWh",
             "subtitle": "Estimated 2026E global consumption", "color": "orange"},
            {"title": "HBM TAM", "value": "$45B",
             "subtitle": "Total addressable market 2026E", "color": "brand"},
            {"title": "EV/EBITDA Multiple", "value": "12.5x",
             "subtitle": "Current vs 18x peak cycle", "color": "navy"},
            {"title": "Upside to TP", "value": "+27%",
             "subtitle": "Based on DCF WACC 9.5%, g 2.5%", "color": "gold"},
        ],

        "value_chain": [
            {"stage_en": "AI Algorithms / Brain", "stage_cn": "AI算法 / 大脑",
             "companies": ["OpenAI", "Google DeepMind", "Anthropic", "Meta FAIR"]},
            {"stage_en": "AI Chips / Compute", "stage_cn": "AI芯片 / 算力",
             "companies": ["NVIDIA", "AMD", "Intel Gaudi", "Broadcom ASIC"]},
            {"stage_en": "Memory / HBM", "stage_cn": "存储 / 高带宽内存",
             "companies": ["SK Hynix", "Samsung", "Micron"]},
            {"stage_en": "Networking / Optics", "stage_cn": "网络 / 光模块",
             "companies": ["Coherent", "Lumentum", "Infinera", "Zhongji Innolight"]},
            {"stage_en": "Power / Cooling", "stage_cn": "电力 / 散热",
             "companies": ["Vertiv", "Eaton", "Schneider", "Aurora Power"]},
            {"stage_en": "System Integration", "stage_cn": "系统集成商",
             "companies": ["Dell", "HPE", "Super Micro", "Lenovo"]},
        ],

        # ---- 新增 v2：5 种新 slide 类型示例数据 ----

        # Slide 19: 双面板并排图表
        "dual_panel": {
            "title_cn": "收入增长 vs 利润率趋势",
            "title_en": "Revenue Growth vs Margin Trends",
            "left_title": "Revenue by Segment ($M)",
            "left_chart_data": {
                "type": "bar",
                "categories": ["FY22", "FY23", "FY24", "FY25E", "FY26E"],
                "series": [
                    {"name": "Hardware", "values": [400, 480, 560, 650, 720]},
                    {"name": "Software", "values": [300, 380, 470, 580, 680]},
                ],
            },
            "right_title": "Gross Margin Trend (%)",
            "right_chart_data": {
                "type": "line",
                "categories": ["FY22", "FY23", "FY24", "FY25E", "FY26E"],
                "series": [
                    {"name": "Gross Margin", "values": [55.0, 56.5, 58.2, 60.0, 61.5]},
                    {"name": "EBITDA Margin", "values": [23.3, 24.0, 24.8, 25.9, 27.0]},
                ],
            },
        },

        # Slide 20: 2x2 战略矩阵
        "matrix_2x2": {
            "title_cn": "AI 投资框架：风险-收益矩阵",
            "title_en": "AI Investment Framework: Risk-Reward Matrix",
            "x_axis_label": "Market Maturity →",
            "y_axis_label": "↑ Technology Readiness",
            "quadrants": [
                {"title": "Core AI Infrastructure\n核心 AI 基础设施",
                 "items": ["GPU / TPU 供应商", "HBM 内存", "数据中心电力"],
                 "color": RGBColor(0x1F, 0x38, 0x64),
                 "position": "top-left"},
                {"title": "AI Applications\nAI 应用层",
                 "items": ["企业级 AI agent", "AI 编程助手", "AI 药物研发"],
                 "color": RGBColor(0x2E, 0x75, 0xB6),
                 "position": "top-right"},
                {"title": "Enabling Technologies\n使能技术",
                 "items": ["光模块 / 800G", "液冷散热", "智能电网"],
                 "color": RGBColor(0xC8, 0xA9, 0x51),
                 "position": "bottom-left"},
                {"title": "Early Stage / Speculative\n早期 / 投机性",
                 "items": ["量子计算", "具身智能", "脑机接口"],
                 "color": RGBColor(0xB9, 0x1C, 0x1C),
                 "position": "bottom-right"},
            ],
        },

        # Slide 21: 资产配置点阵矩阵
        "asset_allocation": {
            "title_cn": "全球资产配置立场",
            "title_en": "Global Asset Allocation Stance",
            "columns": ["1M", "3M", "6M", "12M"],
            "matrix_data": [
                {"asset_class": "US Equities", "asset_cn": "美股",
                 "stances": ["+", "+", "++", "++"]},
                {"asset_class": "EM Equities", "asset_cn": "新兴市场股票",
                 "stances": ["=", "=", "+", "+"]},
                {"asset_class": "IG Credit", "asset_cn": "投资级信用债",
                 "stances": ["=", "+", "+", "="]},
                {"asset_class": "HY Credit", "asset_cn": "高收益信用债",
                 "stances": ["-", "=", "=", "="]},
                {"asset_class": "Govt Bonds", "asset_cn": "政府债券",
                 "stances": ["--", "-", "=", "="]},
                {"asset_class": "USD", "asset_cn": "美元",
                 "stances": ["=", "=", "+", "+"]},
                {"asset_class": "Commodities", "asset_cn": "大宗商品",
                 "stances": ["+", "=", "=", "+"]},
                {"asset_class": "Gold", "asset_cn": "黄金",
                 "stances": ["+", "+", "+", "++"]},
            ],
        },

        # Slide 22: 环形图
        "donut_data": {
            "title_cn": "AI 产业链价值分布",
            "title_en": "AI Value Chain Distribution",
            "segments": [
                {"name": "Semiconductors", "value": 32, "color": "#1F3864"},
                {"name": "Power & Cooling", "value": 18, "color": "#2E75B6"},
                {"name": "Networking / Optics", "value": 14, "color": "#C8A951"},
                {"name": "AI Software", "value": 16, "color": "#00AF50"},
                {"name": "Memory / Storage", "value": 12, "color": "#FB0301"},
                {"name": "Services", "value": 8, "color": "#666666"},
            ],
            "center_metric": "$600B",
            "center_label": "Total AI TAM 2026E",
        },

        # Slide 23: 堆叠柱状图
        "stacked_bar_data": {
            "title_cn": "全球 AI 资本开支构成",
            "title_en": "Global AI Capex Breakdown",
            "categories": ["FY22", "FY23", "FY24", "FY25E", "FY26E", "FY27E"],
            "series_data": [
                {"name": "Compute (GPU/TPU)", "values": [180, 260, 380, 480, 560, 620],
                 "color": "#1F3864"},
                {"name": "Networking", "values": [40, 55, 75, 95, 115, 130],
                 "color": "#2E75B6"},
                {"name": "Power & Cooling", "values": [25, 35, 50, 70, 95, 120],
                 "color": "#C8A951"},
                {"name": "Data Center RE", "values": [30, 40, 55, 70, 85, 100],
                 "color": "#00AF50"},
                {"name": "Other", "values": [15, 20, 30, 40, 50, 60],
                 "color": "#666666"},
            ],
        },
    }


# =============================================================================
# 15. 自测：生成样例 PPTX
# =============================================================================

def _self_test() -> None:
    base_dir = os.path.dirname(os.path.abspath(__file__))
    out = os.path.join(base_dir, "_ms_sample_zh.pptx")
    make_deck(sample_data(), out, theme="classic", language="zh")
    size = os.path.getsize(out)
    from pptx import Presentation as _P
    total = len(_P(out).slides)
    status = "OK" if size > 0 and total >= 6 else "FAIL"
    print(f"[{status}] {out}  size={size} bytes, slides={total}")


if __name__ == "__main__":
    # 默认样例生成
    _self_test()

    # 生成测试文件并打印 TEST_OK slides=<总页数>
    _test_out = "/tmp/ms_test.pptx"
    make_deck(sample_data(), _test_out, theme="classic", language="zh")
    from pptx import Presentation as _PTest
    _total = len(_PTest(_test_out).slides)
    print(f"TEST_OK slides={_total}")
