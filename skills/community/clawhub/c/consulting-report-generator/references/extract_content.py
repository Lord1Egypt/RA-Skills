#!/usr/bin/env python3
"""
咨询报告资料提取脚本 v1
基于 markitdown + python-pptx 的多格式文档文字提取
支持：PDF / Word / Excel / 图片(OCR) / HTML / PPTX

用法：
    python extract_content.py <文件路径>
    
输出：
    extracted_content.md — 按来源分类的结构化文本摘要
    extracted_data.json — 识别的关键数据点（数字/指标/KPI）
"""

import os, sys, json


# ═══════════════════════════════════════════════════════════
# 提取方法矩阵
# ═══════════════════════════════════════════════════════════

# 方法A: markitdown（推荐，支持PDF/Word/Excel/图片OCR/HTML）
# 安装：pip install 'markitdown[all]'
# 用法：markitdown 文件路径 -o 输出.md

# 方法B: Read工具（WorkBuddy内置，支持PDF/图片/文本读取）
# 用法：直接调用 Read 工具

# 方法C: python-pptx（PPTX文字+图片提取）
def extract_pptx_text(pptx_path, output_dir="extracted"):
    """从PPTX提取文字和图片"""
    from pptx import Presentation
    os.makedirs(output_dir, exist_ok=True)
    
    prs = Presentation(pptx_path)
    slides_text = []
    
    for si, slide in enumerate(prs.slides):
        slide_text = []
        img_idx = 0
        for shape in slide.shapes:
            # 提取文字
            if hasattr(shape, 'text') and shape.text.strip():
                slide_text.append(shape.text)
            
            # 提取图片
            if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # 图片类型
                image = shape.image
                img_bytes = image.blob
                ext = image.content_type.split('/')[-1]
                if ext == 'jpeg': ext = 'jpg'
                img_idx += 1
                img_name = f"slide{si+1:02d}_img{img_idx:02d}.{ext}"
                img_path = os.path.join(output_dir, img_name)
                with open(img_path, 'wb') as f:
                    f.write(img_bytes)
                print(f"  📷 提取图片: {img_name}")
        
        if slide_text:
            slides_text.append(f"--- 第{si+1}页 ---\n" + "\n".join(slide_text))
    
    return "\n\n".join(slides_text)


def extract_images_from_pptx(source_path, output_dir="extracted_images"):
    """从原PPTX提取所有图片，按幻灯片编号保存"""
    from pptx import Presentation
    os.makedirs(output_dir, exist_ok=True)
    
    prs = Presentation(source_path)
    img_count = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # 图片类型
                image = shape.image
                image_bytes = image.blob
                ext = image.content_type.split('/')[-1]
                if ext == 'jpeg': ext = 'jpg'
                filename = f"slide{slide_idx+1:02d}_img{shape_idx+1:02d}.{ext}"
                filepath = os.path.join(output_dir, filename)
                with open(filepath, 'wb') as f:
                    f.write(image_bytes)
                img_count += 1
                print(f"  提取: {filename} ({len(image_bytes)} bytes)")
    
    return img_count


# 方法D: 关键数据点识别
def extract_key_data_points(text):
    """从文本中识别关键数据点"""
    import re
    
    data_points = {
        "percentages": [],      # 百分比：85%、30% 等
        "numbers_with_units": [], # 带单位的数字：1.9天、500件等
        "years": [],             # 年份：2021、2022等
        "currencies": [],        # 金额：¥8.5亿、500万等
        "kpis": [],              # KPI关键指标
    }
    
    # 百分比
    for m in re.finditer(r'(\d+\.?\d*)\s*%', text):
        data_points["percentages"].append(m.group(0))
    
    # 带单位的数字
    for m in re.finditer(r'(\d+\.?\d*)\s*(天|月|年|件|个|家|次|人|台|套|条|亿|万)', text):
        data_points["numbers_with_units"].append(m.group(0))
    
    # 年份
    for m in re.finditer(r'(19|20)\d{2}\s*年', text):
        data_points["years"].append(m.group(0))
    
    # 金额
    for m in re.finditer(r'[¥￥]?\d+\.?\d*\s*(亿|万|元)', text):
        data_points["currencies"].append(m.group(0))
    
    # KPI常见指标
    kpi_keywords = ['OEE', '交付率', '良率', '周转率', '周期', '产能', '效率',
                    '成本', '库存', 'WIP', '在制', '节拍', 'Takt', '换型', 'SMED']
    for kw in kpi_keywords:
        if kw.lower() in text.lower():
            # 找到包含该关键词的句子
            for line in text.split('\n'):
                if kw.lower() in line.lower():
                    data_points["kpis"].append(line.strip()[:80])
                    break
    
    return data_points


# ═══════════════════════════════════════════════════════════
# 主流程：判断文件类型→执行提取→输出结果
# ═══════════════════════════════════════════════════════════

def main(file_paths):
    """对多个文件执行资料提取"""
    output_dir = "extracted_content"
    os.makedirs(output_dir, exist_ok=True)
    
    all_text = []
    all_data_points = {}
    
    for file_path in file_paths:
        if not os.path.exists(file_path):
            print(f"⚠️ 文件不存在: {file_path}")
            continue
        
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        print(f"\n{'='*50}")
        print(f"📄 处理文件: {filename}")
        print(f"{'='*50}")
        
        text = ""
        
        # 方法1: markitdown (推荐)
        try:
            import subprocess
            result = subprocess.run(
                ["markitdown", file_path],
                capture_output=True, text=True, timeout=120  # 大文件可能需要较长时间
            )
            if result.returncode == 0 and result.stdout.strip():
                text = result.stdout
                print(f"  ✅ markitdown 提取成功")
        except (FileNotFoundError, subprocess.TimeoutExpired, Exception) as e:
            print(f"  ⚠️ markitdown 提取失败: {e}")
        
        # 方法2: 根据文件类型降级
        if not text:
            if ext in ['.pptx']:
                text = extract_pptx_text(file_path, output_dir)
                img_count = extract_images_from_pptx(file_path, output_dir)
                print(f"  ✅ python-pptx 提取成功（含 {img_count} 张图片）")
            elif ext in ['.txt', '.md']:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                print(f"  ✅ 直接读取成功")
            elif ext in ['.xlsx', '.xls']:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path, data_only=True)
                    lines = []
                    for ws in wb.worksheets:
                        for row in ws.iter_rows(values_only=True):
                            vals = [str(v) for v in row if v is not None]
                            if vals:
                                lines.append(" | ".join(vals))
                    text = "\n".join(lines)
                    print(f"  ✅ openpyxl 读取成功")
                except Exception as e:
                    print(f"  ⚠️ openpyxl 读取失败: {e}")
            else:
                print(f"  ❌ 无法提取: {ext} 格式")
                continue
        
        # 保存提取的文本
        output_file = os.path.join(output_dir, f"{filename}.md")
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(f"# 提取来源: {filename}\n\n{text}")
        
        all_text.append(f"## 来源: {filename}\n\n{text}")
        
        # 识别关键数据点
        data_points = extract_key_data_points(text)
        all_data_points[filename] = data_points
        
        print(f"  📊 识别数据点: ", end="")
        for k, v in data_points.items():
            if v:
                print(f"{k}={len(v)} ", end="")
        print()
    
    # 输出汇总文件
    summary_path = os.path.join(output_dir, "extracted_summary.md")
    with open(summary_path, 'w', encoding='utf-8') as f:
        f.write("# 资料提取汇总\n\n")
        f.write("\n\n".join(all_text))
    
    # 输出关键数据点
    data_path = os.path.join(output_dir, "extracted_data.json")
    with open(data_path, 'w', encoding='utf-8') as f:
        json.dump(all_data_points, f, ensure_ascii=False, indent=2)
    
    print(f"\n{'='*50}")
    print(f"✅ 资料提取完成")
    print(f"  📝 提取文本: {summary_path}")
    print(f"  📊 数据点: {data_path}")
    print(f"{'='*50}")


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("用法: python extract_content.py <文件路径1> [文件路径2 ...]")
        print("示例: python extract_content.py report.pdf data.xlsx photo.png")
        sys.exit(1)
    main(sys.argv[1:])
