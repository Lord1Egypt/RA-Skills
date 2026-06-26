from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def natural_key(path: Path):
    # Sort like slide2.png < slide10.png
    parts = re.split(r"(\d+)", path.name)
    key = []
    for part in parts:
        if part.isdigit():
            key.append(int(part))
        else:
            key.append(part.lower())
    return key


def build_pptx(images: list[Path], out: Path, width_in: float, height_in: float) -> None:
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)

    blank = prs.slide_layouts[6]

    for img in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(img),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))


def main() -> None:
    parser = argparse.ArgumentParser(description="Assemble an images-only PPTX from slide images.")
    parser.add_argument("--images-dir", required=True, help="Directory containing slide images.")
    parser.add_argument("--glob", default="slide*.png", help="Glob pattern (default: slide*.png).")
    parser.add_argument("--out", required=True, help="Output .pptx path.")
    parser.add_argument("--width-in", type=float, default=13.333, help="Slide width in inches (default: 13.333).")
    parser.add_argument("--height-in", type=float, default=7.5, help="Slide height in inches (default: 7.5).")
    args = parser.parse_args()

    images_dir = Path(args.images_dir).expanduser().resolve()
    out = Path(args.out).expanduser().resolve()

    images = sorted(images_dir.glob(args.glob), key=natural_key)
    if not images:
        raise SystemExit(f"No images found: {images_dir} ({args.glob})")

    build_pptx(images, out, args.width_in, args.height_in)
    print(f"Wrote {out} ({len(images)} slides)")


if __name__ == "__main__":
    main()

