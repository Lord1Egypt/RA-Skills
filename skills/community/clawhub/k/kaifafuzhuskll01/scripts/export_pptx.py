#!/usr/bin/env python3
"""
PPT Generator — PPTX Export Engine

Converts structured content data (JSON) into an editable .pptx file.
Uses python-pptx with the design system defined in references/design-system.md.

Usage:
    python3 scripts/export_pptx.py --input content_data.json --output output.pptx
"""

from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("❌ python-pptx 未安装。请运行: pip install python-pptx")
    sys.exit(2)


# ---------------------------------------------------------------------------
# Design constants (from design-system.md § PPTX 排版参数)
# ---------------------------------------------------------------------------

DESIGN = {
    "slide_width": Inches(13.333),
    "slide_height": Inches(7.5),
    "margin": Inches(0.8),
    "title_size": Pt(32),
    "subtitle_size": Pt(18),
    "body_size": Pt(14),
    "bullet_size": Pt(12),
    "font_name": "Noto Sans SC",
    "fallback_font": "Microsoft YaHei",
    "colors": {
        "primary": RGBColor(0, 77, 140),
        "primary_light": RGBColor(26, 109, 181),
        "bg": RGBColor(10, 22, 40),
        "card": RGBColor(17, 34, 64),
        "text": RGBColor(232, 237, 245),
        "text_secondary": RGBColor(136, 153, 170),
        "accent": RGBColor(0, 180, 216),
        "white": RGBColor(255, 255, 255),
    },
}

# Default content data for graceful fallback
DEFAULT_CONTENT_DATA: dict[str, Any] = {
    "title": "未命名演示文稿",
    "slides": [],
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _add_text_box(slide, left, top, width, height, text="", font_size=None,
                  color=None, bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    """Add a single text box to a slide and return the paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if font_size:
        p.font.size = font_size
    if color:
        p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name or DESIGN["font_name"]
    p.alignment = alignment
    return tf


def _add_bullet_list(tf, items, font_size=None, color=None):
    """Append bullet points to an existing text frame."""
    if not items:
        return
    fs = font_size or DESIGN["bullet_size"]
    c = color or DESIGN["colors"]["text"]
    for i, item in enumerate(items):
        if isinstance(item, str):
            text = item
        elif isinstance(item, dict):
            text = item.get("text", "")
        else:
            text = str(item)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {text}"
        p.font.size = fs
        p.font.color.rgb = c
        p.font.name = DESIGN["font_name"]
        p.space_after = Pt(6)


def _set_slide_bg(slide, color):
    """Fill slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# Layout builders — one function per layout type
# ---------------------------------------------------------------------------

def _make_cover_slide(prs, sd):
    """Full-bleed dark background, large white centered title + subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _set_slide_bg(slide, DESIGN["colors"]["bg"])

    m = DESIGN["margin"]
    sw = DESIGN["slide_width"]
    sh = DESIGN["slide_height"]

    title = sd.get("title", "")
    subtitle = sd.get("subtitle", "")

    _add_text_box(slide, m, sh * 0.35, sw - 2 * m, Inches(1.5),
                  title, DESIGN["title_size"], DESIGN["colors"]["white"],
                  bold=True, alignment=PP_ALIGN.CENTER)

    if subtitle:
        _add_text_box(slide, m, sh * 0.55, sw - 2 * m, Inches(0.8),
                      subtitle, DESIGN["subtitle_size"], DESIGN["colors"]["text_secondary"],
                      alignment=PP_ALIGN.CENTER)


def _make_chapter_cover_slide(prs, sd):
    """Left accent bar + chapter number + title + optional description."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DESIGN["colors"]["bg"])

    m = DESIGN["margin"]
    sw = DESIGN["slide_width"]
    sh = DESIGN["slide_height"]

    title = sd.get("title", "")
    subtitle = sd.get("subtitle", "")

    # Left accent bar
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        m, sh * 0.25, Inches(0.08), sh * 0.5,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DESIGN["colors"]["accent"]
    shape.line.fill.background()

    _add_text_box(slide, m + Inches(0.4), sh * 0.25, sw - 2 * m, Inches(1.2),
                  title, DESIGN["title_size"], DESIGN["colors"]["white"],
                  bold=True)

    if subtitle:
        _add_text_box(slide, m + Inches(0.4), sh * 0.45, sw - 2 * m, Inches(0.8),
                      subtitle, DESIGN["subtitle_size"], DESIGN["colors"]["text_secondary"])


def _make_toc_slide(prs, sd):
    """2-column grid of chapter entries with numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DESIGN["colors"]["bg"])

    m = DESIGN["margin"]
    sw = DESIGN["slide_width"]
    sh = DESIGN["slide_height"]

    title = sd.get("title", "目录")
    _add_text_box(slide, m, m, sw - 2 * m, Inches(0.8),
                  title, DESIGN["title_size"], DESIGN["colors"]["white"], bold=True)

    items = sd.get("content", [])
    if not items:
        return

    col_width = (sw - 2 * m) / 2 - Inches(0.3)
    for i, item in enumerate(items):
        text = item.get("text", str(item)) if isinstance(item, dict) else str(item)
        col = i % 2
        row = i // 2
        left = m + col * (col_width + Inches(0.6))
        top = m + Inches(1.2) + row * Inches(0.6)

        num = f"{i + 1:02d}"
        _add_text_box(slide, left, top, Inches(0.6), Inches(0.5),
                      num, DESIGN["subtitle_size"], DESIGN["colors"]["accent"], bold=True)
        _add_text_box(slide, left + Inches(0.7), top, col_width - Inches(0.7), Inches(0.5),
                      text, DESIGN["body_size"], DESIGN["colors"]["text"])


def _make_content_slide(prs, sd):
    """Title + bullet point list + optional bottom data highlight."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DESIGN["colors"]["bg"])

    m = DESIGN["margin"]
    sw = DESIGN["slide_width"]

    title = sd.get("title", "")
    subtitle = sd.get("subtitle", "")

    _add_text_box(slide, m, m, sw - 2 * m, Inches(0.7),
                  title, DESIGN["title_size"], DESIGN["colors"]["white"], bold=True)

    top_offset = m + Inches(1.0)
    if subtitle:
        _add_text_box(slide, m, top_offset, sw - 2 * m, Inches(0.5),
                      subtitle, DESIGN["subtitle_size"], DESIGN["colors"]["text_secondary"])
        top_offset += Inches(0.7)

    content = sd.get("content", [])
    if content:
        tf = _add_text_box(slide, m, top_offset, sw - 2 * m, Inches(4.5),
                          "", DESIGN["bullet_size"], DESIGN["colors"]["text"])
        texts = [c.get("text", str(c)) if isinstance(c, dict) else str(c) for c in content]
        _add_bullet_list(tf, texts)


def _make_three_cols_slide(prs, sd):
    """Three card columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DESIGN["colors"]["bg"])

    m = DESIGN["margin"]
    sw = DESIGN["slide_width"]

    title = sd.get("title", "")
    _add_text_box(slide, m, m, sw - 2 * m, Inches(0.7),
                  title, DESIGN["title_size"], DESIGN["colors"]["white"], bold=True)

    cols_data = sd.get("content", [])
    if not cols_data:
        return

    col_w = (sw - 2 * m - Inches(0.6)) / 3
    for i, col in enumerate(cols_data[:3]):
        left = m + i * (col_w + Inches(0.3))
        top = m + Inches(1.2)

        # Card background
        card = slide.shapes.add_shape(1, left, top, col_w, Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = DESIGN["colors"]["card"]
        card.line.color.rgb = DESIGN["colors"]["primary_light"]
        card.line.width = Pt(1)

        col_title = col.get("title", "") if isinstance(col, dict) else str(col)
        _add_text_box(slide, left + Inches(0.2), top + Inches(0.2),
                      col_w - Inches(0.4), Inches(0.5),
                      col_title, DESIGN["subtitle_size"], DESIGN["colors"]["accent"], bold=True)

        if isinstance(col, dict):
            items = col.get("items", col.get("content", []))
            texts = [it.get("text", str(it)) if isinstance(it, dict) else str(it) for it in items]
            tf = _add_text_box(slide, left + Inches(0.2), top + Inches(1.0),
                              col_w - Inches(0.4), Inches(3.0),
                              "", DESIGN["bullet_size"], DESIGN["colors"]["text"])
            _add_bullet_list(tf, texts)


def _make_two_cols_slide(prs, sd):
    """Two columns: left text, right description/data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DESIGN["colors"]["bg"])

    m = DESIGN["margin"]
    sw = DESIGN["slide_width"]

    title = sd.get("title", "")
    _add_text_box(slide, m, m, sw - 2 * m, Inches(0.7),
                  title, DESIGN["title_size"], DESIGN["colors"]["white"], bold=True)

    content = sd.get("content", [])
    col_w = (sw - 2 * m - Inches(0.5)) / 2

    for i, col in enumerate(content[:2]):
        left = m + i * (col_w + Inches(0.5))
        top = m + Inches(1.2)

        card = slide.shapes.add_shape(1, left, top, col_w, Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = DESIGN["colors"]["card"]
        card.line.color.rgb = DESIGN["colors"]["primary_light"]
        card.line.width = Pt(1)

        col_title = col.get("title", "") if isinstance(col, dict) else str(col)
        _add_text_box(slide, left + Inches(0.2), top + Inches(0.2),
                      col_w - Inches(0.4), Inches(0.5),
                      col_title, DESIGN["subtitle_size"], DESIGN["colors"]["accent"], bold=True)

        if isinstance(col, dict):
            items = col.get("items", col.get("content", []))
            texts = [it.get("text", str(it)) if isinstance(it, dict) else str(it) for it in items]
            tf = _add_text_box(slide, left + Inches(0.2), top + Inches(1.0),
                              col_w - Inches(0.4), Inches(3.0),
                              "", DESIGN["bullet_size"], DESIGN["colors"]["text"])
            _add_bullet_list(tf, texts)


def _make_chart_slide(prs, sd):
    """Chart summary slide — v1 uses table fallback."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DESIGN["colors"]["bg"])

    m = DESIGN["margin"]
    sw = DESIGN["slide_width"]

    title = sd.get("title", "")
    subtitle = sd.get("subtitle", "")

    _add_text_box(slide, m, m, sw - 2 * m, Inches(0.7),
                  title, DESIGN["title_size"], DESIGN["colors"]["white"], bold=True)

    top_offset = m + Inches(1.0)
    if subtitle:
        _add_text_box(slide, m, top_offset, sw - 2 * m, Inches(0.5),
                      subtitle, DESIGN["subtitle_size"], DESIGN["colors"]["text_secondary"])
        top_offset += Inches(0.7)

    # TODO v2: Use native python-pptx charts for chart-radar and chart-bar
    # v1 fallback: list items as bullet-like entries
    content = sd.get("content", [])
    if content:
        tf = _add_text_box(slide, m, top_offset, sw - 2 * m, Inches(4.0),
                          "", DESIGN["bullet_size"], DESIGN["colors"]["text"])
        for i, item in enumerate(content):
            text = item.get("text", str(item)) if isinstance(item, dict) else str(item)
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▸ {text}"
            p.font.size = DESIGN["bullet_size"]
            p.font.color.rgb = DESIGN["colors"]["text"]
            p.font.name = DESIGN["font_name"]
            p.space_after = Pt(4)


def _make_ending_slide(prs, sd):
    """Centered thank you message on dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DESIGN["colors"]["bg"])

    m = DESIGN["margin"]
    sw = DESIGN["slide_width"]
    sh = DESIGN["slide_height"]

    title = sd.get("title", "谢谢！")
    subtitle = sd.get("subtitle", "")

    _add_text_box(slide, m, sh * 0.35, sw - 2 * m, Inches(1.5),
                  title, DESIGN["title_size"], DESIGN["colors"]["white"],
                  bold=True, alignment=PP_ALIGN.CENTER)

    if subtitle:
        _add_text_box(slide, m, sh * 0.55, sw - 2 * m, Inches(0.8),
                      subtitle, DESIGN["subtitle_size"], DESIGN["colors"]["text_secondary"],
                      alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Layout dispatch
# ---------------------------------------------------------------------------

LAYOUT_BUILDERS = {
    "cover": _make_cover_slide,
    "chapter-cover": _make_chapter_cover_slide,
    "toc": _make_toc_slide,
    "content": _make_content_slide,
    "three-cols": _make_three_cols_slide,
    "two-cols": _make_two_cols_slide,
    "chart-radar": _make_chart_slide,
    "chart-bar": _make_chart_slide,
    "ending": _make_ending_slide,
}


def _make_fallback_slide(prs, sd):
    """Fallback: blank slide with layout name as title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DESIGN["colors"]["bg"])
    layout_name = sd.get("layout", "unknown")
    _add_text_box(slide, DESIGN["margin"], DESIGN["margin"],
                  DESIGN["slide_width"] - 2 * DESIGN["margin"], Inches(0.8),
                  f"[未知布局: {layout_name}]",
                  DESIGN["title_size"], DESIGN["colors"]["text_secondary"],
                  alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def build_pptx(content_data: dict, output_path: str) -> str:
    """Build a .pptx file from structured content data.

    Args:
        content_data: Dict with ``title`` and ``slides`` fields.
            Each slide has: page, layout, density, title, subtitle?, content[]
        output_path: Where to write the .pptx file.

    Returns:
        The output_path on success.

    Raises:
        SystemExit(1) on JSON parse errors.
        SystemExit(2) if python-pptx is not installed.
    """
    prs = PptxPresentation()
    prs.slide_width = DESIGN["slide_width"]
    prs.slide_height = DESIGN["slide_height"]

    slides = content_data.get("slides", [])

    for sd in slides:
        layout = sd.get("layout", "content")
        builder = LAYOUT_BUILDERS.get(layout, _make_fallback_slide)
        try:
            builder(prs, sd)
        except Exception:
            _make_fallback_slide(prs, sd)

    # Remove the default empty slide that Presentation() creates
    if len(prs.slides) > len(slides):
        rId = prs.slides._sldIdLst[0].get("r:id")
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def main():
    import argparse

    parser = argparse.ArgumentParser(description="PPT Generator — PPTX Export")
    parser.add_argument("--input", required=True, help="Path to content_data.json")
    parser.add_argument("--output", required=True, help="Path to output .pptx")
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"❌ 输入文件不存在: {input_path}")
        sys.exit(1)

    try:
        with open(input_path, "r", encoding="utf-8") as f:
            data = json.load(f)
    except json.JSONDecodeError as e:
        print(f"❌ JSON 解析失败: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"❌ 读取文件失败: {e}")
        sys.exit(1)

    output_path = args.output
    try:
        result = build_pptx(data, output_path)
        print(f"✅ PPTX 已生成: {result}")
    except Exception as e:
        print(f"❌ 生成失败: {e}")
        sys.exit(3)


if __name__ == "__main__":
    main()
