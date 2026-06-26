#!/usr/bin/env python3
"""
DeckCraft v6 — CRAP Design Optimizer

Analyzes a PPTX file and produces a Markdown report diagnosing
Contrast, Repetition, Alignment, and Proximity (CRAP) issues.

Usage:
    python3 optimize_crap.py deck.pptx
    python3 optimize_crap.py deck.pptx -o report.md
"""
import sys
import os
import json
import argparse
from typing import List, Dict, Any, Tuple, Optional
from collections import defaultdict

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN


# ── Constants ─────────────────────────────────────────────────────
MIN_TITLE_BODY_RATIO = 1.4    # Title should be ≥1.4× body text
IDEAL_TITLE_BODY_RATIO = 1.8
MAX_FONT_SIZES_PER_ROLE = 3   # Shouldn't have >3 distinct font sizes for same role
MIN_MARGIN_EMU = Emu(Inches(0.3))  # Minimum margin from slide edge
ALIGNMENT_TOLERANCE_EMU = Emu(Inches(0.08))  # Tolerance for "same position"


def _pt_from_size(size) -> Optional[float]:
    """Extract font size in pt from a run or paragraph, or None."""
    if size is None:
        return None
    try:
        return size.pt
    except (AttributeError, TypeError):
        return None


def _analyze_slide_shapes(slide, slide_idx: int) -> Dict[str, Any]:
    """Analyze all shapes on a single slide."""
    result = {
        "slide": slide_idx + 1,
        "texts": [],       # [{left, top, width, height, font_size, bold, text, align}]
        "rects": [],       # [{left, top, width, height}]
    }

    for shape in slide.shapes:
        shape_info = {
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
        }

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if not para.text.strip():
                    continue
                font_sizes = []
                bolds = []
                for run in para.runs:
                    if run.font.size:
                        font_sizes.append(_pt_from_size(run.font.size))
                    bolds.append(run.font.bold or False)

                avg_size = sum(s for s in font_sizes if s) / len(font_sizes) if font_sizes else None
                any_bold = any(bolds)

                alignment = para.alignment
                align_str = {
                    PP_ALIGN.LEFT: "left",
                    PP_ALIGN.CENTER: "center",
                    PP_ALIGN.RIGHT: "right",
                    PP_ALIGN.JUSTIFY: "justify",
                }.get(alignment, "left") if alignment else "default"

                result["texts"].append({
                    **shape_info,
                    "font_size": avg_size,
                    "bold": any_bold,
                    "text": para.text[:60],
                    "align": align_str,
                })

        # Track non-text shapes for proximity analysis
        if not shape.has_text_frame and shape.shape_type is not None:
            result["rects"].append(shape_info)

    return result


def _check_contrast(slides_data: List[Dict]) -> List[str]:
    """Check Contrast: title vs body font size ratio."""
    issues = []
    for data in slides_data:
        slide_num = data["slide"]
        texts = data["texts"]
        if not texts:
            continue

        sizes = [(t["font_size"], t["bold"], t["text"]) for t in texts if t["font_size"]]
        if len(sizes) < 2:
            continue

        # Assume largest bold text is title, rest is body
        title_candidates = [(s, b, t) for s, b, t in sizes if b and s >= 16]
        body_candidates = [(s, b, t) for s, b, t in sizes if not b or s < 16]

        if not title_candidates:
            # Fallback: largest = title
            sorted_sizes = sorted(sizes, key=lambda x: -x[0])
            title_candidates = [sorted_sizes[0]]
            body_candidates = sorted_sizes[1:]

        if title_candidates and body_candidates:
            title_size = max(s for s, _, _ in title_candidates)
            body_size = min(s for s, _, _ in body_candidates)
            if body_size > 0:
                ratio = title_size / body_size
                if ratio < MIN_TITLE_BODY_RATIO:
                    issues.append(
                        f"- 第{slide_num}页: 标题字号 {title_size:.0f}pt, "
                        f"正文字号 {body_size:.0f}pt, 差异 {ratio:.1f}x "
                        f"⚠️ 偏小（建议 ≥{MIN_TITLE_BODY_RATIO}x）"
                    )
                elif ratio >= IDEAL_TITLE_BODY_RATIO:
                    issues.append(
                        f"- 第{slide_num}页: 标题字号 {title_size:.0f}pt, "
                        f"正文字号 {body_size:.0f}pt, 差异 {ratio:.1f}x ✅ 达标"
                    )
                else:
                    issues.append(
                        f"- 第{slide_num}页: 标题字号 {title_size:.0f}pt, "
                        f"正文字号 {body_size:.0f}pt, 差异 {ratio:.1f}x 🆗 可接受"
                    )

    return issues


def _check_repetition(slides_data: List[Dict]) -> List[str]:
    """Check Repetition: consistent font sizes and styles across slides."""
    issues = []

    # Collect all font sizes
    all_title_sizes = []
    all_body_sizes = []
    slide_body_aligns = defaultdict(list)

    for data in slides_data:
        texts = data["texts"]
        slide_num = data["slide"]
        for t in texts:
            if not t["font_size"]:
                continue
            if t["bold"] and t["font_size"] >= 16:
                all_title_sizes.append((slide_num, t["font_size"]))
            else:
                all_body_sizes.append((slide_num, t["font_size"]))
                slide_body_aligns[data["slide"]].append(t["align"])

    # Title font consistency
    if all_title_sizes:
        sizes = [s for _, s in all_title_sizes]
        unique_sizes = sorted(set(sizes))
        if len(unique_sizes) > MAX_FONT_SIZES_PER_ROLE:
            issues.append(
                f"- 标题使用了 {len(unique_sizes)} 种字号 "
                f"({', '.join(f'{s:.0f}pt' for s in unique_sizes)}) "
                f"⚠️ 建议统一为 1-2 种"
            )
        else:
            issues.append(
                f"- 标题字号: {', '.join(f'{s:.0f}pt' for s in unique_sizes)} ✅ 一致"
            )

    # Body font consistency
    if all_body_sizes:
        sizes = [s for _, s in all_body_sizes]
        unique_sizes = sorted(set(sizes))
        if len(unique_sizes) > MAX_FONT_SIZES_PER_ROLE + 1:
            issues.append(
                f"- 正文使用了 {len(unique_sizes)} 种字号 "
                f"({', '.join(f'{s:.0f}pt' for s in unique_sizes[:6])}...) "
                f"⚠️ 建议统一为 2-3 种"
            )
        else:
            issues.append(
                f"- 正文字号: {', '.join(f'{s:.0f}pt' for s in unique_sizes)} ✅ 基本一致"
            )

    # Alignment consistency
    all_aligns = []
    for aligns in slide_body_aligns.values():
        all_aligns.extend(aligns)
    if all_aligns:
        from collections import Counter
        align_counts = Counter(all_aligns)
        dominant = align_counts.most_common(1)[0]
        if dominant[1] / len(all_aligns) >= 0.7:
            issues.append(
                f"- 对齐方式: 以「{dominant[0]}」为主 "
                f"({dominant[1]}/{len(all_aligns)}) ✅ 一致"
            )
        else:
            issues.append(
                f"- 对齐方式: 混合使用 "
                f"({', '.join(f'{k}×{v}' for k, v in align_counts.most_common())}) "
                f"⚠️ 建议统一"
            )

    return issues


def _check_alignment(slides_data: List[Dict], slide_width: int, slide_height: int) -> List[str]:
    """Check Alignment: consistent left margins, centering."""
    issues = []
    tol = ALIGNMENT_TOLERANCE_EMU

    for data in slides_data:
        slide_num = data["slide"]
        texts = data["texts"]
        if len(texts) < 2:
            continue

        # Check if text elements share left edges
        lefts = sorted(set(t["left"] for t in texts))
        # Group lefts within tolerance
        groups = []
        for l in lefts:
            placed = False
            for g in groups:
                if abs(l - g[0]) < tol:
                    g.append(l)
                    placed = True
                    break
            if not placed:
                groups.append([l])

        multi_groups = [g for g in groups if len(g) > 1]
        if len(groups) > 3 and len(texts) > 3:
            issues.append(
                f"- 第{slide_num}页: 存在 {len(groups)} 种不同左对齐位置 "
                f"⚠️ 建议减少到 2-3 种"
            )
        elif multi_groups:
            issues.append(
                f"- 第{slide_num}页: 对齐情况良好 "
                f"({len(multi_groups)} 组对齐元素) ✅"
            )

    return issues


def _check_proximity(slides_data: List[Dict]) -> List[str]:
    """Check Proximity: related items should be close together."""
    issues = []

    for data in slides_data:
        slide_num = data["slide"]
        texts = data["texts"]
        if len(texts) < 2:
            continue

        # Check vertical gaps between consecutive text elements
        sorted_texts = sorted(texts, key=lambda t: t["top"])
        gaps = []
        for i in range(len(sorted_texts) - 1):
            t1 = sorted_texts[i]
            t2 = sorted_texts[i + 1]
            gap = t2["top"] - (t1["top"] + t1["height"])
            if gap > 0:
                gaps.append(gap)

        if not gaps:
            continue

        avg_gap = sum(gaps) / len(gaps)
        large_gaps = [g for g in gaps if g > avg_gap * 3 and g > Emu(Inches(0.5))]

        if large_gaps:
            issues.append(
                f"- 第{slide_num}页: 存在 {len(large_gaps)} 处过大间距 "
                f"(最大 {Inches(large_gaps[0]).inches:.1f}\") "
                f"⚠️ 相关内容应靠近摆放"
            )
        else:
            issues.append(
                f"- 第{slide_num}页: 内容间距均匀 ✅"
            )

    return issues


def analyze_crap(pptx_path: str) -> str:
    """Analyze a PPTX file and return CRAP report as Markdown.

    Args:
        pptx_path: Path to .pptx file.

    Returns:
        Markdown report string.
    """
    if not os.path.isfile(pptx_path):
        return f"# CRAP 优化报告\n\n❌ 文件不存在: {pptx_path}"

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return f"# CRAP 优化报告\n\n❌ 无法打开文件: {e}"

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    num_slides = len(prs.slides)

    if num_slides == 0:
        return "# CRAP 优化报告\n\n⚠️ 文件中没有幻灯片"

    # Collect data
    slides_data = []
    for i, slide in enumerate(prs.slides):
        slides_data.append(_analyze_slide_shapes(slide, i))

    # Run checks
    contrast = _check_contrast(slides_data)
    repetition = _check_repetition(slides_data)
    alignment = _check_alignment(slides_data, slide_w, slide_h)
    proximity = _check_proximity(slides_data)

    # Build report
    lines = [
        "# CRAP 优化报告",
        "",
        f"**文件**: `{os.path.basename(pptx_path)}`",
        f"**幻灯片数**: {num_slides}",
        f"**画布**: {slide_w / 914400:.1f}\" × {slide_h / 914400:.1f}\""  ,
        "",
    ]

    lines.append("## 对比 (Contrast)")
    lines.append("")
    if contrast:
        lines.extend(contrast)
    else:
        lines.append("- 不足 2 个文本元素，跳过对比检查")
    lines.append("")

    lines.append("## 重复 (Repetition)")
    lines.append("")
    lines.extend(repetition if repetition else ["- 无文本内容，跳过重复检查"])
    lines.append("")

    lines.append("## 对齐 (Alignment)")
    lines.append("")
    lines.extend(alignment if alignment else ["- 不足 2 个文本元素，跳过对齐检查"])
    lines.append("")

    lines.append("## 亲密性 (Proximity)")
    lines.append("")
    lines.extend(proximity if proximity else ["- 不足 2 个文本元素，跳过亲密性检查"])

    # Summary
    total_issues = sum(
        1 for section in [contrast, repetition, alignment, proximity]
        for line in section if "⚠️" in line
    )
    total_ok = sum(
        1 for section in [contrast, repetition, alignment, proximity]
        for line in section if "✅" in line
    )

    lines.extend([
        "",
        "---",
        "",
        f"**总结**: {total_ok} 项达标 ✅ | {total_issues} 项需优化 ⚠️",
    ])

    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(
        description="DeckCraft v6 — CRAP Design Optimizer",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("pptx", help="Input PPTX file to analyze")
    parser.add_argument("-o", "--output", default=None,
                        help="Output Markdown report path (default: stdout)")

    args = parser.parse_args()

    report = analyze_crap(args.pptx)

    if args.output:
        os.makedirs(os.path.dirname(args.output) if os.path.dirname(args.output) else ".",
                    exist_ok=True)
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(report)
        print(f"✓ CRAP report saved to {args.output}")
    else:
        print(report)

    return 0


if __name__ == "__main__":
    sys.exit(main())
