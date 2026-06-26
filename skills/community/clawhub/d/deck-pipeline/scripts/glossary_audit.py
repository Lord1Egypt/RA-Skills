#!/usr/bin/env python3
"""Late-stage glossary re-scan + wavering detection.

Usage:
    python3 glossary_audit.py <pptx-path> <glossary.yaml>

glossary.yaml format:
    locked:
      business_terms:
        "<source term>": "<canonical translation>"
      domain_specific:
        "<source term>": "<canonical translation>"
    rejected_rewrites:
      - {source: "...", proposed: "..."}

Reports:
  1. Wavering: same source term translated multiple ways in the deck
  2. Glossary breach: locked-EN term has a competing translation in deck
  3. Rejected-rewrite reappearance: deck contains a previously-vetoed phrase
"""
import sys
import re
from collections import defaultdict, Counter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    import yaml
except ImportError:
    print("Requires PyYAML: pip3 install pyyaml")
    sys.exit(1)


def collect_text(path):
    prs = Presentation(path)
    out = []  # (slide, text)
    def walk(shapes, sidx):
        for shp in shapes:
            try:
                if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(shp.shapes, sidx); continue
            except Exception: pass
            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs).strip()
                    if t: out.append((sidx, t))
            if shp.has_table:
                for row in shp.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            t = "".join(r.text for r in p.runs).strip()
                            if t: out.append((sidx, t))
    for i, s in enumerate(prs.slides, 1):
        walk(s.shapes, i)
    return out


def main():
    if len(sys.argv) < 3:
        print("Usage: glossary_audit.py <pptx-path> <glossary.yaml>")
        sys.exit(1)
    path = sys.argv[1]
    glossary_path = sys.argv[2]

    with open(glossary_path) as f:
        gl = yaml.safe_load(f)

    locked = {}
    for cat, mapping in (gl.get("locked") or {}).items():
        if isinstance(mapping, dict):
            locked.update(mapping)
    rejected = gl.get("rejected_rewrites") or []

    texts = collect_text(path)

    # 1. Wavering: for each locked source term, find all distinct EN strings
    #    that appear in the same context as that source token.
    #    Simpler proxy: search for known canonical translations AND known
    #    sibling variants in the deck text.
    print("=" * 70)
    print("Locked glossary breach check")
    print("=" * 70)
    breaches = []
    for src, canonical in locked.items():
        # Look for canonical's word stem in deck
        for slide, text in texts:
            if canonical in text:
                continue
        # Look for common siblings (heuristic: same starting word, different ending)
        canon_first = canonical.split(" ")[0] if " " in canonical else canonical
        for slide, text in texts:
            if canon_first in text and canonical not in text:
                breaches.append((slide, src, canonical, text[:80]))
    if not breaches:
        print("  (none)")
    for s, src, ca, t in breaches[:30]:
        print(f"  P{s:2}  '{src}' should map to '{ca}', but text reads: '{t}'")

    # 2. Wavering: token frequency for known-translated source terms
    print("\n" + "=" * 70)
    print("Wavering detection (same canonical term, multiple variants)")
    print("=" * 70)
    # For each canonical EN value, check whether its near-variants also appear
    # e.g. "Acme Corporation" vs "Acme Corp" vs "ACME"
    canonical_variants = defaultdict(set)
    for src, canonical in locked.items():
        # tokenize canonical to words
        words = canonical.split()
        if len(words) >= 2:
            short = "".join(w[0].upper() for w in words if w[0].isalpha())
            for slide, text in texts:
                if short and short in text and canonical not in text:
                    canonical_variants[canonical].add((slide, "abbrev:" + short))
                if canonical.lower() in text.lower() and canonical not in text:
                    canonical_variants[canonical].add((slide, "case-variant"))
    if not canonical_variants:
        print("  (none)")
    for ca, vs in canonical_variants.items():
        print(f"  canonical '{ca}' appears with variants: {sorted(vs)[:5]}")

    # 3. Rejected-rewrite reappearance
    print("\n" + "=" * 70)
    print("Rejected-rewrite reappearance check")
    print("=" * 70)
    found = []
    for entry in rejected:
        proposed = entry.get("proposed", "")
        if not proposed: continue
        for slide, text in texts:
            if proposed in text:
                found.append((slide, proposed, text[:80]))
    if not found:
        print("  (none)")
    for s, p, t in found:
        print(f"  P{s:2}  rejected '{p}' found in: '{t}'")


if __name__ == "__main__":
    main()
