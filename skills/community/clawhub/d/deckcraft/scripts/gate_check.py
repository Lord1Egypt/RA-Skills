#!/usr/bin/env python3
"""
DeckCraft v5 — S4 QA Gate Check
Renders gate_result.json with machine-readable pass/fail.
AI must read the JSON output; verbal declaration is not accepted.

Usage: python3 gate_check.py <pptx_path> <project_dir>
"""
import sys, os, json
from pptx import Presentation
from pptx.util import Pt

def check_pptx(pptx_path, project_dir):
    prs = Presentation(pptx_path)
    SW = prs.slide_width
    SH = prs.slide_height
    sw_in = SW / 914400
    sh_in = SH / 914400

    user_code_errors = []
    warnings = []
    engine_warnings = []

    for si, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        positions = []

        for s in shapes:
            name = s.name[:50]
            l = s.left
            t = s.top
            r = l + s.width
            b = t + s.height

            # 1. Overflow — elements outside slide boundaries
            if l < -914400 * 0.05:
                user_code_errors.append({
                    "slide": si + 1, "category": "overflow_left",
                    "detail": f"{name} left={l/914400:.2f}\" (outside slide)"
                })
            if r > SW + 914400 * 0.05:
                user_code_errors.append({
                    "slide": si + 1, "category": "overflow_right",
                    "detail": f"{name} right={r/914400:.2f}\" > slide {sw_in:.2f}\""
                })
            if b > SH + 914400 * 0.05:
                user_code_errors.append({
                    "slide": si + 1, "category": "overflow_bottom",
                    "detail": f"{name} bottom={b/914400:.2f}\" > slide {sh_in:.2f}\""
                })

            # 2. Font size checks
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.size and run.font.size < Pt(7):
                            txt = run.text[:40]
                            user_code_errors.append({
                                "slide": si + 1, "category": "tiny_font",
                                "detail": f"{name} \"{txt}\" = {run.font.size/12700:.1f}pt (<7pt)"
                            })

            # 3. Track for overlap
            positions.append((l, t, r, b, name))

        # 4. Overlap detection (>80% overlap between TextBox shapes)
        for i in range(len(positions)):
            for j in range(i + 1, len(positions)):
                l1, t1, r1, b1, n1 = positions[i]
                l2, t2, r2, b2, n2 = positions[j]
                overlap_x = max(0, min(r1, r2) - max(l1, l2))
                overlap_y = max(0, min(b1, b2) - max(t1, t2))
                area1 = max(1, (r1 - l1) * (b1 - t1))
                overlap_area = overlap_x * overlap_y
                if overlap_area > area1 * 0.8 and 'TextBox' in n1 and 'TextBox' in n2:
                    warnings.append({
                        "slide": si + 1, "category": "overlap",
                        "detail": f"{n1} and {n2} overlap >80%"
                    })

    # 5. Check slide count is reasonable
    if len(prs.slides) == 0:
        user_code_errors.append({
            "slide": 0, "category": "empty_deck",
            "detail": "Presentation has no slides"
        })

    # Build result
    passed = len(user_code_errors) == 0
    score = max(0, 100 - len(user_code_errors) * 10 - len(warnings) * 3)

    result = {
        "passed": passed,
        "overall_score": score,
        "checklist": {
            "user_code_errors": len(user_code_errors),
            "warnings": len(warnings),
            "total_slides": len(prs.slides)
        },
        "verdict": "PASS — ready for delivery" if passed else "FAIL — fix user_code_errors before delivery",
        "user_code_errors": user_code_errors,
        "warnings": warnings
    }

    # Write JSON
    os.makedirs(project_dir, exist_ok=True)
    out_path = os.path.join(project_dir, "gate_result.json")
    with open(out_path, "w") as f:
        json.dump(result, f, indent=2, ensure_ascii=False)

    return result


def main():
    if len(sys.argv) < 3:
        print("Usage: python3 gate_check.py <pptx_path> <project_dir>")
        print("Output: <project_dir>/gate_result.json")
        sys.exit(1)

    pptx_path = sys.argv[1]
    project_dir = sys.argv[2]

    if not os.path.exists(pptx_path):
        result = {
            "passed": False,
            "overall_score": 0,
            "verdict": "FAIL — pptx file not found",
            "user_code_errors": [{"slide": 0, "category": "file_not_found", "detail": pptx_path}],
            "warnings": []
        }
        os.makedirs(project_dir, exist_ok=True)
        with open(os.path.join(project_dir, "gate_result.json"), "w") as f:
            json.dump(result, f, indent=2, ensure_ascii=False)
        print(json.dumps(result, indent=2))
        sys.exit(1)

    result = check_pptx(pptx_path, project_dir)

    # Also print to stdout for visibility
    print(json.dumps(result, indent=2, ensure_ascii=False))

    if result["passed"]:
        print(f"\n✅ GATE PASSED — score {result['overall_score']}/100")
    else:
        print(f"\n❌ GATE FAILED — {len(result['user_code_errors'])} error(s)")
        for err in result["user_code_errors"]:
            print(f"   Slide {err['slide']}: [{err['category']}] {err['detail']}")

    sys.exit(0 if result["passed"] else 1)


if __name__ == "__main__":
    main()
