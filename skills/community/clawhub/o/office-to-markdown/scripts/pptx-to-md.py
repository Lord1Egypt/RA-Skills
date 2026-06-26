"""
pptx-to-md.py  (v2 — security-hardened)
Converts a PowerPoint (.pptx) file to Markdown.
- Extracts text from all text shapes (no API calls)
- Falls back to Claude vision for image-only slides (only with --allow-vision)

Usage:
    python pptx-to-md.py <input.pptx> <output.md> [--allow-vision]

Dependencies installed to /tmp/office_md_deps/ (isolated, pinned versions).
"""

import sys, base64, json, urllib.request, subprocess
from pathlib import Path

# ── isolated dependency install ───────────────────────────────────────────────
_DEP_DIR = Path("/tmp/office_md_deps")
_DEP_DIR.mkdir(exist_ok=True)
subprocess.run(
    [
        sys.executable, "-m", "pip", "install", "--quiet",
        "--target", str(_DEP_DIR),
        "python-pptx==1.0.2",
    ],
    check=True,
)
if str(_DEP_DIR) not in sys.path:
    sys.path.insert(0, str(_DEP_DIR))

from pptx import Presentation


# ── vision helper (only called when --allow-vision is set) ────────────────────
def call_vision(b64_data: str, media_type: str = "image/png") -> str:
    payload = {
        "model": "claude-sonnet-4-20250514",
        "max_tokens": 2000,
        "messages": [{
            "role": "user",
            "content": [
                {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": b64_data}},
                {"type": "text", "text": (
                    "Extract ALL text and structured content from this slide image into clean Markdown.\n"
                    "Preserve headings, bullet lists, and any tables. "
                    "Output ONLY the extracted markdown — no preamble."
                )}
            ]
        }]
    }
    req = urllib.request.Request(
        "https://api.anthropic.com/v1/messages",
        data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json", "anthropic-version": "2023-06-01"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req) as resp:
            return json.loads(resp.read())["content"][0]["text"]
    except Exception as e:
        return f"*(vision extraction failed: {e})*"


# ── main ──────────────────────────────────────────────────────────────────────
def convert(src: str, dst: str, allow_vision: bool = False):
    prs = Presentation(src)
    md_lines = []
    vision_count = 0
    vision_skipped = 0

    for i, slide in enumerate(prs.slides, 1):
        md_lines.append(f"## Slide {i}")
        has_text = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        has_text = True
                        level = para.level
                        prefix = "#" * (level + 3) if level else ""
                        md_lines.append(f"{prefix} {text}" if prefix else text)

        if not has_text:
            extracted_any = False
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    if not allow_vision:
                        vision_skipped += 1
                        continue
                    try:
                        img_blob = shape.image.blob
                        mt = shape.image.content_type
                        b64 = base64.standard_b64encode(img_blob).decode("utf-8")
                        md_lines.append(call_vision(b64, mt))
                        extracted_any = True
                        vision_count += 1
                    except Exception:
                        pass

            if not extracted_any:
                if vision_skipped > 0:
                    md_lines.append(
                        "*(image-only slide — vision extraction skipped; "
                        "re-run with --allow-vision after user confirmation)*"
                    )
                else:
                    md_lines.append("*(image-only slide — no extractable content)*")

        md_lines.append("")

    Path(dst).write_text("\n".join(md_lines), encoding="utf-8")
    if vision_skipped:
        print(
            f"  VISION_REQUIRED: {vision_skipped} image-only slide(s) detected. "
            f"Re-run with --allow-vision after user confirmation."
        )
    print(f"  Saved → {dst}  ({len(prs.slides)} slides, {vision_count} via vision)")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: pptx-to-md.py <input.pptx> <output.md> [--allow-vision]")
        sys.exit(1)
    allow_vision = "--allow-vision" in sys.argv
    convert(sys.argv[1], sys.argv[2], allow_vision=allow_vision)
