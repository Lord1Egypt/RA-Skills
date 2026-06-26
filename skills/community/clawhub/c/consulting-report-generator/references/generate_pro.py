#!/usr/bin/env python3
"""
咨询报告专业PPT生成器 v5（豆包对标系）
基于 mck-ppt-design 专业布局框架
麦肯锡设计风格 + 微软雅黑（Mac兼容）+ 无BLOCK_ARC
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, zipfile
from lxml import etree

# ═══════════════════════════════════════════
# Ⅰ. 配色方案
# ═══════════════════════════════════════════

# 主色
NV = RGBColor(0x05, 0x1C, 0x2C)  # 深蓝
WH = RGBColor(0xFF, 0xFF, 0xFF)  # 白色
BK = RGBColor(0x00, 0x00, 0x00)  # 黑色
D3 = RGBColor(0x33, 0x33, 0x33)  # 正文
M6 = RGBColor(0x66, 0x66, 0x66)  # 次级文字
LG = RGBColor(0xCC, 0xCC, 0xCC)  # 分隔线
BG = RGBColor(0xF2, 0xF2, 0xF2)  # 背景面板

# 强调色
AB = RGBColor(0x00, 0x6B, 0xA6)  # 蓝
AG = RGBColor(0x00, 0x7A, 0x53)  # 绿
AO = RGBColor(0xD4, 0x6A, 0x00)  # 橙
AR = RGBColor(0xC6, 0x28, 0x28)  # 红

# 浅背景
LB = RGBColor(0xE3, 0xF2, 0xFD)
LG2 = RGBColor(0xE8, 0xF5, 0xE9)
LO = RGBColor(0xFF, 0xF3, 0xE0)
LR = RGBColor(0xFF, 0xEB, 0xEE)

# ═══ 深蓝递进色（商品定义框架用） ═══
LV1 = RGBColor(0x05, 0x1C, 0x2C)  # 顶层 - 最深
LV2 = RGBColor(0x0A, 0x30, 0x4A)  # 第2层
LV3 = RGBColor(0x10, 0x44, 0x62)  # 第3层
LV4 = RGBColor(0x16, 0x58, 0x7A)  # 第4层
LV5 = RGBColor(0x1C, 0x6C, 0x92)  # 第5层（底）- 最浅

# 深蓝递进子色（内部分块用）
SV1 = RGBColor(0x10, 0x60, 0x78)
SV2 = RGBColor(0x10, 0x68, 0x82)
SV3 = RGBColor(0x10, 0x70, 0x8C)
SV4 = RGBColor(0x18, 0x68, 0x88)
SV5 = RGBColor(0x18, 0x70, 0x92)
SV6 = RGBColor(0x18, 0x78, 0x9C)
SV7 = RGBColor(0x20, 0x6C, 0x90)
SV8 = RGBColor(0x20, 0x7C, 0xA0)
SV9 = RGBColor(0x20, 0x8C, 0xB0)
TAG_BG = RGBColor(0x28, 0x80, 0xA8)  # 标签色

# ═══════════════════════════════════════════
# Ⅱ. 常量
# ═══════════════════════════════════════════

FN = "Microsoft YaHei"   # 中文
EN = "Arial"              # 英文/数字
LM = Inches(0.8)          # 左边距
CW = Inches(11.733)       # 内容宽度
SW = Inches(13.333)       # 幻灯片宽度
SH = Inches(7.5)          # 幻灯片高度
TZ = Inches(1.3)          # 内容顶部起始

# ═══════════════════════════════════════════
# Ⅲ. 助手函数
# ═══════════════════════════════════════════

def cs(s):
    """清理p:style"""
    e = s._element
    st = e.find(f'{{http://schemas.openxmlformats.org/presentationml/2006/main}}style')
    if st is not None: e.remove(st)

def sf(run, is_en=False):
    """设置字体"""
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        latin = rPr.makeelement(qn('a:latin'), {})
        rPr.append(latin)
    latin.set('typeface', EN if is_en else FN)
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FN)

def TX(s, l, t, w, h, tx, sz=14, c=D3, b=False, al=PP_ALIGN.LEFT, an='t'):
    """统一文本助手"""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        for attr in ['lIns','tIns','rIns','bIns']:
            bp.set(attr, '45720')
        bp.set('anchor', an)
    lines = tx if isinstance(tx, list) else [tx]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(line)
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.bold = b
        p.alignment = al
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = 0.93 if sz >= 18 else Pt(sz * 1.35)
        for run in p.runs:
            sf(run, is_en=(sz >= 18 and all(ch.isascii() or ch.isdigit() for ch in line)))
    return tb

def R(s, l, t, w, h, c):
    """矩形"""
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = c
    r.line.fill.background()
    cs(r)
    return r

def RR(s, l, t, w, h, c):
    """圆角矩形"""
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = c
    r.line.fill.background()
    cs(r)
    return r

def HL(s, x, y, ln, c=BK, t=Pt(0.5)):
    """水平线（薄矩形）"""
    h = max(int(t), Emu(6350))
    return R(s, x, y, ln, h, c)

def OV(s, x, y, lb, sz=Inches(0.45), bg=NV, fg=WH):
    """圆形标签"""
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    c.fill.solid()
    c.fill.fore_color.rgb = bg
    c.line.fill.background()
    cs(c)
    tf = c.text_frame
    p = tf.paragraphs[0]
    p.text = str(lb)
    p.font.size = Pt(14)
    p.font.color.rgb = fg
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        sf(run)
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None: bp.set('anchor', 'ctr')
    return c

def AT(s, tx, sz=22):
    """Action Title"""
    TX(s, LM, Inches(0.15), CW, Inches(0.9), tx, sz, BK, True, PP_ALIGN.LEFT, 'b')
    HL(s, LM, Inches(1.05), CW, BK, Pt(0.5))

def SR(s, tx, y=Inches(7.05)):
    """来源标注"""
    TX(s, LM, y, CW, Inches(0.3), tx, 9, M6)

def PN(s, n, t):
    """页码"""
    TX(s, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3), f"{n}/{t}", 9, M6, False, PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════
# Ⅳ. 布局模式
# ═══════════════════════════════════════════

# ---- 布局：#1 Cover ----
def cover(s, title, subtitle="", date="", company=""):
    """封面"""
    R(s, 0, 0, SW, Inches(0.05), NV)
    lines = title.count('\n') + 1
    th = Inches(0.8 + 0.65 * max(lines - 1, 0))
    TX(s, Inches(1), Inches(2.0), Inches(11), th, title, 40, NV, True)
    sy = Inches(2.0) + th + Inches(0.3)
    if subtitle:
        TX(s, Inches(1), sy, Inches(11), Inches(0.6), subtitle, 24, D3)
        sy += Inches(1.0)
    if date or company:
        TX(s, Inches(1), sy + Inches(0.3), Inches(11), Inches(0.5), f"{company}  |  {date}", 14, M6)
    HL(s, Inches(1), Inches(6.8), Inches(4), NV, Pt(2))

# ---- 布局：#6 TOC ----
def toc(s, items):
    """目录 items=[(编号, 标题, 描述)]"""
    AT(s, "目录")
    iy = Inches(1.6)
    for num, title, desc in items:
        OV(s, LM, iy, num, Inches(0.5), NV)
        TX(s, LM + Inches(0.7), iy, Inches(4.0), Inches(0.4), title, 16, NV, True)
        TX(s, Inches(5.5), iy + Inches(0.05), Inches(6.5), Inches(0.35), desc, 14, M6)
        iy += Inches(0.7)
        HL(s, LM, iy, CW, LG, Pt(0.25))
        iy += Inches(0.3)

# ---- 布局：#5 Section Divider ----
def SEC(s, n, title, subtitle=""):
    """章节分隔页"""
    R(s, 0, 0, Inches(0.6), SH, NV)
    TX(s, Inches(1.2), Inches(2.0), Inches(10), Inches(0.6), f"PART {n}", 16, M6)
    TX(s, Inches(1.2), Inches(2.8), Inches(10), Inches(1.2), title, 28, NV, True)
    if subtitle:
        TX(s, Inches(1.2), Inches(4.2), Inches(10), Inches(0.6), subtitle, 14, D3)

# ---- 布局：#11 Data Table ----
def TB(s, l, t, rows, cols, data, col_ws=None):
    """原生数据表格"""
    tw = sum(col_ws) if col_ws else CW
    ts = s.shapes.add_table(rows, cols, l, t, tw, Inches(0.35 * rows))
    for i, rd in enumerate(data):
        for j, ct in enumerate(rd):
            c = ts.table.cell(i, j)
            c.text = str(ct)
            if col_ws:
                c.width = col_ws[j]
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(9)
                for run in p.runs:
                    sf(run)
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = WH if i == 0 else D3
                p.font.bold = (i == 0)
            if i == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = NV
            elif i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

# ---- 布局：#12 Metric Cards Row ----
def MCARDS(s, items, y=TZ):
    """指标卡片 items=[(数值, 标签, 示例值, 颜色)]"""
    n = len(items)
    cw = (CW - Inches(0.15) * (n - 1)) / n
    for i, (val, label, detail, clr) in enumerate(items):
        cx = LM + i * (cw + Inches(0.15))
        R(s, cx, y, cw, Inches(1.5), WH)
        R(s, cx, y, cw, Inches(0.06), clr)
        TX(s, cx + Inches(0.2), y + Inches(0.15), cw - Inches(0.4), Inches(0.5), val, 24, clr, True, PP_ALIGN.CENTER)
        TX(s, cx + Inches(0.2), y + Inches(0.7), cw - Inches(0.4), Inches(0.3), label, 11, M6, False, PP_ALIGN.CENTER)
        TX(s, cx + Inches(0.2), y + Inches(1.05), cw - Inches(0.4), Inches(0.3), detail, 10, D3, False, PP_ALIGN.CENTER)

# ---- 布局：#8 Big Number ----
def BN(s, items, y=TZ):
    """大数字对比 items=[(大数字, 标签, 颜色), ...]"""
    n = len(items)
    cw = (CW - Inches(0.3) * (n - 1)) / n
    for i, (num, label, clr) in enumerate(items):
        cx = LM + i * (cw + Inches(0.3))
        R(s, cx, y, cw, Inches(1.8), NV if i == 0 else BG)
        tc = WH if i == 0 else NV
        sc = WH if i == 0 else D3
        TX(s, cx + Inches(0.2), y + Inches(0.15), cw - Inches(0.4), Inches(0.8), num, 44, tc, True, PP_ALIGN.CENTER)
        TX(s, cx + Inches(0.2), y + Inches(1.1), cw - Inches(0.4), Inches(0.5), label, 14, sc, False, PP_ALIGN.CENTER)

# ---- 布局：#24 Executive Summary ----
def EXEC(s, title, points):
    """执行摘要"""
    AT(s, title)
    R(s, LM, TZ, CW, Inches(1.0), NV)
    TX(s, LM + Inches(0.3), TZ, CW - Inches(0.6), Inches(1.0), points[0], 16, WH)
    iy = TZ + Inches(1.3)
    for i, (num, t, d) in enumerate(points[1:]):
        OV(s, LM, iy, num)
        TX(s, LM + Inches(0.6), iy, Inches(3.5), Inches(0.4), t, 14, NV, True)
        TX(s, Inches(5.0), iy, Inches(7.5), Inches(0.35), d, 14, D3)
        iy += Inches(0.55)
        HL(s, LM, iy, CW, LG, Pt(0.25))
        iy += Inches(0.15)

# ---- 布局：#14 Three-Pillar Framework ----
def PILLAR(s, items, y=TZ):
    """三支柱 items=[(标题, [要点...], 颜色)]"""
    n = len(items)
    pw = (CW - Inches(0.2) * (n - 1)) / n
    for i, (title, points, clr) in enumerate(items):
        px = LM + i * (pw + Inches(0.2))
        R(s, px, y, pw, Inches(0.5), clr)
        TX(s, px + Inches(0.15), y, pw - Inches(0.3), Inches(0.5), title, 14, WH, True, PP_ALIGN.CENTER, 'ctr')
        R(s, px, y + Inches(0.5), pw, Inches(3.5), BG)
        TX(s, px + Inches(0.2), y + Inches(0.6), pw - Inches(0.4), Inches(3.0), points, 12, D3)

# ---- 布局：#20 Before/After ----
def BA(s, before, after, y=TZ):
    """Before/After 对比"""
    hw = Inches(5.0)
    # Before
    R(s, LM, y, hw, Inches(4.0), BG)
    TX(s, LM + Inches(0.3), y + Inches(0.1), hw - Inches(0.6), Inches(0.5), "X  现状 (Before)", 16, D3)
    HL(s, LM + Inches(0.3), y + Inches(0.6), hw - Inches(0.6), LG, Pt(0.5))
    TX(s, LM + Inches(0.3), y + Inches(0.8), hw - Inches(0.6), Inches(3.0), before, 12, D3)
    # Arrow
    TX(s, LM + hw + Inches(0.1), y + Inches(1.5), Inches(1.5), Inches(0.5), "→", 36, NV, True, PP_ALIGN.CENTER)
    # After
    ax = LM + hw + Inches(1.733)
    R(s, ax, y, hw, Inches(4.0), NV)
    TX(s, ax + Inches(0.3), y + Inches(0.1), hw - Inches(0.6), Inches(0.5), "V  目标 (After)", 16, WH)
    HL(s, ax + Inches(0.3), y + Inches(0.6), hw - Inches(0.6), WH, Pt(0.5))
    TX(s, ax + Inches(0.3), y + Inches(0.8), hw - Inches(0.6), Inches(3.0), after, 12, WH)

# ---- 布局：#29 Timeline ----
def TIMELINE(s, items, y=Inches(3.0)):
    """时间线 items=[(标题, 描述), ...]"""
    n = len(items)
    sp = Inches(10.7) / max(n - 1, 1)
    HL(s, LM + Inches(0.5), y, Inches(10.7), LG, Pt(2))
    for i, (title, desc) in enumerate(items):
        mx = LM + Inches(0.5) + sp * i
        OV(s, mx - Inches(0.225), y - Inches(0.225), str(i + 1))
        TX(s, mx - Inches(1.0), y - Inches(0.8), Inches(2.0), Inches(0.5), title, 16, NV, True, PP_ALIGN.CENTER)
        TX(s, mx - Inches(1.0), y + Inches(0.5), Inches(2.0), Inches(1.0), desc, 11, D3, False, PP_ALIGN.CENTER)

# ---- 布局：#16 Process Chevron ----
def FLOW(s, steps, y=Inches(2.5)):
    """流程箭头 steps=[(标题, 描述), ...]"""
    n = len(steps)
    sw = CW / n
    for i, (title, desc) in enumerate(steps):
        sx = LM + i * sw
        fill = NV if i == 0 else BG
        R(s, sx, y, sw - Inches(0.06), Inches(1.0), fill)
        tc = WH if i == 0 else NV
        TX(s, sx + Inches(0.1), y + Inches(0.05), sw - Inches(0.3), Inches(0.4), title, 13, tc, True, PP_ALIGN.CENTER, 'ctr')
        TX(s, sx + Inches(0.1), y + Inches(0.5), sw - Inches(0.3), Inches(0.4), desc, 9, tc if i == 0 else D3, False, PP_ALIGN.CENTER, 'ctr')
        if i < n - 1:
            TX(s, sx + sw - Inches(0.12), y + Inches(0.25), Inches(0.15), Inches(0.4), "→", 18, AO, True, PP_ALIGN.CENTER)

# ---- 布局：#39 Horizontal Bar ----
def HBAR(s, items, y=TZ, title=""):
    """水平柱状图 items=[(标签, 值), ...] 已排序"""
    if title: TX(s, LM, y - Inches(0.25), CW, Inches(0.2), title, 10, NV, True)
    max_v = max(v for _, v in items)
    bar_max = Inches(7.0)
    n = len(items)
    rh = Inches(0.5)
    for i, (lb, v) in enumerate(items):
        ry = y + i * rh
        bw = bar_max * v / max_v
        clr = NV if i == 0 else BG
        TX(s, LM, ry, Inches(1.8), rh, lb, 11, D3, i == 0, PP_ALIGN.LEFT, 'ctr')
        R(s, LM + Inches(2.0), ry + Inches(0.08), bar_max, Inches(0.34), RGBColor(0xF0, 0xF0, 0xF0))
        R(s, LM + Inches(2.0), ry + Inches(0.08), bw, Inches(0.34), clr)
        TX(s, LM + Inches(2.0) + bw + Inches(0.1), ry, Inches(1.0), rh, str(v), 11, D3, i == 0, PP_ALIGN.LEFT, 'ctr')

# ---- 布局：#37 Grouped Bar (简化版) ----
def GBAR(s, data, labels, colors, y=TZ, title=""):
    """分组柱状图 data=[[v1,v2,...], ...] 每行=系列"""
    if title: TX(s, LM, y - Inches(0.25), CW, Inches(0.2), title, 10, NV, True)
    n_groups = len(labels)
    n_series = len(data)
    chart_l = LM + Inches(1.0)
    chart_b = y + Inches(2.8)
    chart_h = Inches(2.5)
    chart_w = Inches(8.0)
    group_w = chart_w / n_groups
    bar_w = Inches(0.35)
    max_v = max(max(row) for row in data) or 1
    # Y轴
    for tick in range(0, 101, 25):
        ty_pos = chart_b - chart_h * tick / 100
        TX(s, LM, ty_pos - Inches(0.1), Inches(0.8), Inches(0.2), f"{tick}%", 8, M6, False, PP_ALIGN.RIGHT)
        if tick > 0: HL(s, chart_l, ty_pos, chart_w, RGBColor(0xE8, 0xE8, 0xE8), Pt(0.25))
    HL(s, chart_l, chart_b, chart_w, BK, Pt(0.5))
    # 柱子
    for gi, lb in enumerate(labels):
        gx = chart_l + group_w * gi + (group_w - bar_w * n_series) / 2
        for si in range(n_series):
            val = data[si][gi]
            bh = chart_h * val / max_v
            bx = gx + si * (bar_w + Inches(0.05))
            R(s, bx, chart_b - bh, bar_w, bh, colors[si])
            if val > 0:
                TX(s, bx, chart_b - bh - Inches(0.2), bar_w, Inches(0.18), str(val), 7, D3, False, PP_ALIGN.CENTER)
        TX(s, chart_l + group_w * gi, chart_b + Inches(0.03), group_w, Inches(0.2), lb, 8, D3, False, PP_ALIGN.CENTER)

# ---- 布局：#52 KPI Tracker ----
def KTRACK(s, items, y=TZ):
    """KPI追踪 items=[(名称, 进度0-1, 详情, 状态on/risk/off)]"""
    # 表头
    hy = y - Inches(0.3)
    TX(s, LM, hy, Inches(3.0), Inches(0.3), "指标", 11, M6, True)
    TX(s, LM + Inches(3.0), hy, Inches(6.0), Inches(0.3), "进度", 11, M6, True)
    TX(s, LM + Inches(9.0), hy, Inches(1.5), Inches(0.3), "达成率", 11, M6, True, PP_ALIGN.CENTER)
    HL(s, LM, hy + Inches(0.35), CW, BK, Pt(0.75))
    sc_map = {'on': AG, 'risk': AO, 'off': AR}
    for i, (nm, pct, det, st) in enumerate(items):
        ry = hy + Inches(0.5) + i * Inches(0.6)
        TX(s, LM, ry, Inches(2.8), Inches(0.5), nm, 12, D3, True, PP_ALIGN.LEFT, 'ctr')
        R(s, LM + Inches(3.0), ry + Inches(0.12), Inches(5.8), Inches(0.25), BG)
        R(s, LM + Inches(3.0), ry + Inches(0.12), Inches(5.8) * min(pct, 1.0), Inches(0.25), sc_map[st])
        TX(s, LM + Inches(9.0), ry, Inches(1.2), Inches(0.5), f"{int(pct*100)}%", 14, sc_map[st], True, PP_ALIGN.CENTER, 'ctr')
        R(s, LM + Inches(10.3), ry + Inches(0.17), Inches(0.15), Inches(0.15), sc_map[st])
        if i < len(items) - 1: HL(s, LM, ry + Inches(0.55), CW, LG, Pt(0.25))

# ---- 布局：#57 Dashboard ----
def DASH(s, kpis, data, labels, colors, insights):
    """仪表盘: KPI行+柱状图+洞见"""
    AT(s, "业务健康度仪表盘")
    # KPI行
    cw = CW / len(kpis)
    for i, (v, lb, clr) in enumerate(kpis):
        cx = LM + i * cw
        R(s, cx, TZ, cw - Inches(0.08), Inches(1.0), WH)
        R(s, cx, TZ, cw - Inches(0.08), Inches(0.06), clr)
        TX(s, cx + Inches(0.1), TZ + Inches(0.08), cw - Inches(0.2), Inches(0.4), v, 22, clr, True, PP_ALIGN.CENTER)
        TX(s, cx + Inches(0.1), TZ + Inches(0.55), cw - Inches(0.2), Inches(0.3), lb, 10, M6, False, PP_ALIGN.CENTER)
    # 柱状图
    cy = TZ + Inches(1.3)
    GBAR(s, data, labels, colors, cy)
    # 图例
    ly = cy + Inches(3.3)
    for i, (lb, clr) in enumerate(zip(labels, colors)):
        tx = LM + Inches(3.0) + i * Inches(2.0)
        R(s, tx, ly, Inches(0.2), Inches(0.15), clr)
        TX(s, tx + Inches(0.3), ly - Inches(0.02), Inches(1.5), Inches(0.2), lb, 9, D3)
    # 洞见
    bay = Inches(6.0)
    R(s, LM, bay, CW, Inches(0.7), BG)
    TX(s, LM + Inches(0.3), bay, CW - Inches(0.6), Inches(0.7), insights, 12, NV, True, PP_ALIGN.LEFT, 'ctr')

# ---- 布局：#35 Action Items ----
def ACTIONS(s, items, y=TZ):
    """行动项 items=[(行动, 时间, 负责人), ...]"""
    n = len(items)
    cw = (CW - Inches(0.15) * (n - 1)) / n
    for i, (action, time, owner) in enumerate(items):
        cx = LM + i * (cw + Inches(0.15))
        R(s, cx, y, cw, Inches(3.0), BG)
        R(s, cx, y, cw, Inches(0.06), NV)
        OV(s, cx + Inches(0.15), y + Inches(0.15), str(i + 1), Inches(0.4), NV)
        TX(s, cx + Inches(0.6), y + Inches(0.15), cw - Inches(0.3), Inches(0.4), action, 14, NV, True, PP_ALIGN.CENTER)
        TX(s, cx + Inches(0.2), y + Inches(0.8), cw - Inches(0.4), Inches(0.3), f"时间：{time}", 11, M6, False, PP_ALIGN.CENTER)
        TX(s, cx + Inches(0.2), y + Inches(1.2), cw - Inches(0.4), Inches(0.3), f"负责人：{owner}", 11, D3, False, PP_ALIGN.CENTER)

# ---- 布局：#36 Closing ----
def CLOSE(s, main_text, sub_text="", tagline=""):
    """结束页"""
    R(s, 0, 0, SW, Inches(0.05), NV)
    TX(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0), main_text, 28, NV, True, PP_ALIGN.CENTER)
    HL(s, Inches(5.5), Inches(3.3), Inches(2.3), NV, Pt(1.5))
    TX(s, Inches(1.5), Inches(3.8), Inches(10.3), Inches(1.5), sub_text, 18, D3, False, PP_ALIGN.CENTER)
    HL(s, LM, Inches(6.8), CW, NV, Pt(2))
    if tagline:
        TX(s, LM, Inches(6.0), CW, Inches(0.5), tagline, 20, AO, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Ⅴ. 新增布局（v4 深蓝递进系）
# ═══════════════════════════════════════════

# ---- 布局：#60 HeaderBar ----
def HDR(s, title, subtitle=""):
    """深蓝顶栏标题（参照星网智云样式）
    title: 主标题
    subtitle: 副标题（浅蓝英文/说明）
    """
    R(s, 0, 0, SW, Inches(0.85), NV)
    TX(s, LM, Inches(0.08), Inches(8), Inches(0.55), title, 24, WH, True)
    if subtitle:
        TX(s, LM, Inches(0.55), Inches(8), Inches(0.25), subtitle, 11, RGBColor(0x88, 0xBB, 0xDD))

# ---- 布局：#61 PYRAMID - 居中多层金字塔 ----
def PYRAMID(s, layers, header=""):
    """金字塔布局（参照星网智云商品定义框架）
    layers = [(层名, 层级序号如"第一层", 内容描述, 主色, 子块列表), ...]
      子块列表 = [(子标题, 子背景色), ...] 或 None
    header: 页脚说明文字
    """
    if header:
        TX(s, LM, Inches(7.1), Inches(8), Inches(0.3), header, 9, M6)
    
    n = len(layers)
    # 从下到上绘制，每层递增宽度
    base_w = Inches(10.0)
    base_h = Inches(0.9)
    gap = Inches(0.08)
    # 每层宽度递减
    step_w = (base_w - Inches(4.0)) / max(n - 1, 1)  # 顶层宽4寸
    
    for idx, (_, level_label, desc, bg_clr, sub_blocks) in enumerate(layers):
        # 从底层(0)到顶层(n-1)
        layer_w = base_w - step_w * (n - 1 - idx)
        layer_x = (SW - layer_w) / 2
        # Y坐标从底部开始
        layer_y = Inches(5.6) - idx * (base_h + gap)
        
        # 绘制整层背景
        R(s, layer_x, layer_y, layer_w, base_h, bg_clr)
        
        # 左侧标签
        TX(s, layer_x + Inches(0.3), layer_y + Inches(0.05), Inches(4.5), Inches(0.5),
           desc, 14, WH, True, PP_ALIGN.LEFT)
        
        # 如果有子块，绘制内部子块
        if sub_blocks and len(sub_blocks) > 0:
            total_sub = len(sub_blocks)
            sub_w = (layer_w - Inches(0.3) * 2 - Inches(0.08) * (total_sub - 1)) / total_sub
            sub_y = layer_y + Inches(0.08)
            sub_h = base_h - Inches(0.16)
            for si, (sub_title, sub_bg) in enumerate(sub_blocks):
                sub_x = layer_x + Inches(0.08) + si * (sub_w + Inches(0.08))
                R(s, sub_x, sub_y, sub_w, sub_h, sub_bg)
                TX(s, sub_x, sub_y, sub_w, sub_h, sub_title, 14, WH, True, PP_ALIGN.CENTER, 'ctr')
        
        # 右侧层级标签
        tag_x = layer_x + layer_w - Inches(3.8)
        tag_y = layer_y + base_h + Inches(0.02)
        TX(s, tag_x, tag_y, Inches(3.8), Inches(0.25), level_label, 11, RGBColor(0x44, 0xAA, 0xCC), True, PP_ALIGN.RIGHT)

# ---- 布局：#62 BLOCK_3COL - 三列内部分块 ----
def BLOCK_3COL(s, cx, cy, cw, ch, items, y_offset=0):
    """三列子块布局（在父块内部绘制三列子块）
    items = [(标题, 背景色), ...]  3个
    cx, cy, cw, ch: 父块位置尺寸
    y_offset: Y偏移
    """
    n = min(len(items), 3)
    sub_w = (cw - Inches(0.08) * (n - 1)) / n
    sub_h = ch - Inches(0.16)
    sub_y = cy + y_offset + Inches(0.08)
    for si, (sub_title, sub_bg) in enumerate(items):
        sub_x = cx + Inches(0.08) + si * (sub_w + Inches(0.08))
        R(s, sub_x, sub_y, sub_w, sub_h, sub_bg)
        TX(s, sub_x, sub_y, sub_w, sub_h, sub_title, 14, WH, True, PP_ALIGN.CENTER, 'ctr')

# ---- 布局：#63 TAG_CHIP - 层级标签 ----
def TAG_CHIP(s, text, x, y, bg=TAG_BG):
    """标签chip（层级说明标签）
    放置在层级右侧底部
    """
    w = Inches(3.3)
    h = Inches(0.28)
    R(s, x - w, y, w, h, bg)
    TX(s, x - w + Inches(0.08), y, w - Inches(0.16), h, text, 11, WH, True, PP_ALIGN.CENTER, 'ctr')

# ---- 布局：#64 LEFT_LABEL_BAR - 左标签右内容条 ----
def LEFT_LABEL_BAR(s, x, y, w, h, label, label_w, content, label_bg=NV, content_bg=NV):
    """左标签+右内容色条
    label: 左侧标签文字
    content: 右侧内容文字
    label_w: 标签宽度
    """
    gap = Inches(0.06)
    # 左侧标签
    R(s, x, y, label_w, h, label_bg)
    TX(s, x + Inches(0.08), y, label_w - Inches(0.16), h, label, 13, WH, True, PP_ALIGN.CENTER, 'ctr')
    # 右侧内容
    content_x = x + label_w + gap
    content_w = w - label_w - gap
    R(s, content_x, y, content_w, h, content_bg)
    TX(s, content_x + Inches(0.15), y, content_w - Inches(0.3), h, content, 12, WH, False, PP_ALIGN.LEFT, 'ctr')

# ═══════════════════════════════════════════
# Ⅵ. 配色主题系统（v5 新增 - 对标豆包多模板）
# ═══════════════════════════════════════════
# 支持一键切换整体配色风格

def apply_theme(theme="deepblue"):
    """应用配色主题，返回颜色字典
    可选主题: deepblue(深蓝经典), tech(科技蓝), business(商务灰), vibrant(活力橙)
    """
    themes = {
        "deepblue": {  # 默认深蓝 - 咨询专业风
            "primary": RGBColor(0x05,0x1C,0x2C),   # 主色
            "accent1": RGBColor(0x00,0x6B,0xA6),   # 强调1-蓝
            "accent2": RGBColor(0x00,0x7A,0x53),   # 强调2-绿
            "accent3": RGBColor(0xD4,0x6A,0x00),   # 强调3-橙
            "accent4": RGBColor(0xC6,0x28,0x28),   # 强调4-红
            "light1": RGBColor(0xE3,0xF2,0xFD),    # 浅色1-蓝
            "light2": RGBColor(0xE8,0xF5,0xE9),    # 浅色2-绿
            "light3": RGBColor(0xFF,0xF3,0xE0),    # 浅色3-橙
            "light4": RGBColor(0xFF,0xEB,0xEE),    # 浅色4-红
            "header_bg": RGBColor(0x05,0x1C,0x2C), # 顶栏
            "header_sub": RGBColor(0x88,0xBB,0xDD),# 顶栏副标题
        },
        "tech": {  # 科技蓝 - 偏亮蓝/白
            "primary": RGBColor(0x00,0x5A,0x96),
            "accent1": RGBColor(0x00,0x96,0xD6),
            "accent2": RGBColor(0x00,0xB0,0x6A),
            "accent3": RGBColor(0xFF,0x8C,0x00),
            "accent4": RGBColor(0xE0,0x3E,0x3E),
            "light1": RGBColor(0xE0,0xF2,0xFE),
            "light2": RGBColor(0xE6,0xF7,0xEE),
            "light3": RGBColor(0xFF,0xF4,0xE0),
            "light4": RGBColor(0xFF,0xEB,0xEE),
            "header_bg": RGBColor(0x00,0x5A,0x96),
            "header_sub": RGBColor(0x88,0xCC,0xEE),
        },
        "business": {  # 商务灰 - 沉稳专业
            "primary": RGBColor(0x2C,0x3E,0x50),
            "accent1": RGBColor(0x34,0x95,0xDB),
            "accent2": RGBColor(0x27,0xAE,0x60),
            "accent3": RGBColor(0xE6,0x7E,0x22),
            "accent4": RGBColor(0xE7,0x4C,0x3C),
            "light1": RGBColor(0xEB,0xF5,0xFB),
            "light2": RGBColor(0xE8,0xF8,0xF0),
            "light3": RGBColor(0xFE,0xF9,0xE7),
            "light4": RGBColor(0xFD,0xED,0xEC),
            "header_bg": RGBColor(0x2C,0x3E,0x50),
            "header_sub": RGBColor(0xAA,0xBB,0xCC),
        },
        "vibrant": {  # 活力橙 - 热情鲜明
            "primary": RGBColor(0xD4,0x6A,0x00),
            "accent1": RGBColor(0x00,0x96,0xD6),
            "accent2": RGBColor(0x00,0x7A,0x53),
            "accent3": RGBColor(0x6C,0x5C,0xE7),
            "accent4": RGBColor(0xE0,0x3E,0x3E),
            "light1": RGBColor(0xFF,0xF3,0xE0),
            "light2": RGBColor(0xE8,0xF5,0xE9),
            "light3": RGBColor(0xF3,0xE8,0xFF),
            "light4": RGBColor(0xFF,0xEB,0xEE),
            "header_bg": RGBColor(0xD4,0x6A,0x00),
            "header_sub": RGBColor(0xFF,0xDD,0xAA),
        }
    }
    return themes.get(theme, themes["deepblue"])

# ═══════════════════════════════════════════
# Ⅶ. 豆包对标布局（v5 新增）
# ═══════════════════════════════════════════

# ---- 布局：#70 SWOT - 战略分析矩阵 ----
def SWOT(s, strengths, weaknesses, opportunities, threats, y=TZ):
    """SWOT分析矩阵（对标豆包战略分析模板）
    strengths/weaknesses/opportunities/threats = [要点列表]
    """
    half_w = (CW - Inches(0.1)) / 2
    half_h = Inches(1.7)
    gap = Inches(0.1)
    # S - 左上（绿）
    R(s, LM, y, half_w, half_h, RGBColor(0xE8,0xF5,0xE9))
    TX(s, LM+Inches(0.15), y+Inches(0.05), half_w-Inches(0.3), Inches(0.35), "S 优势 (Strengths)", 14, RGBColor(0x00,0x7A,0x53), True)
    for si, pt in enumerate(strengths):
        TX(s, LM+Inches(0.15), y+Inches(0.4)+si*Inches(0.3), half_w-Inches(0.3), Inches(0.28), pt, 10, D3)
    # W - 右上（红）
    wx = LM + half_w + gap
    R(s, wx, y, half_w, half_h, RGBColor(0xFF,0xEB,0xEE))
    TX(s, wx+Inches(0.15), y+Inches(0.05), half_w-Inches(0.3), Inches(0.35), "W 劣势 (Weaknesses)", 14, RGBColor(0xC6,0x28,0x28), True)
    for wi, pt in enumerate(weaknesses):
        TX(s, wx+Inches(0.15), y+Inches(0.4)+wi*Inches(0.3), half_w-Inches(0.3), Inches(0.28), pt, 10, D3)
    # O - 左下（蓝）
    oy = y + half_h + gap
    R(s, LM, oy, half_w, half_h, RGBColor(0xE3,0xF2,0xFD))
    TX(s, LM+Inches(0.15), oy+Inches(0.05), half_w-Inches(0.3), Inches(0.35), "O 机会 (Opportunities)", 14, RGBColor(0x00,0x6B,0xA6), True)
    for oi, pt in enumerate(opportunities):
        TX(s, LM+Inches(0.15), oy+Inches(0.4)+oi*Inches(0.3), half_w-Inches(0.3), Inches(0.28), pt, 10, D3)
    # T - 右下（橙）
    tx2 = LM + half_w + gap
    R(s, tx2, oy, half_w, half_h, RGBColor(0xFF,0xF3,0xE0))
    TX(s, tx2+Inches(0.15), oy+Inches(0.05), half_w-Inches(0.3), Inches(0.35), "T 威胁 (Threats)", 14, RGBColor(0xD4,0x6A,0x00), True)
    for ti, pt in enumerate(threats):
        TX(s, tx2+Inches(0.15), oy+Inches(0.4)+ti*Inches(0.3), half_w-Inches(0.3), Inches(0.28), pt, 10, D3)

# ---- 布局：#71 PIE_BAR - 饼图替代柱状占比图 ----
def PIE_BAR(s, items, y=TZ, title=""):
    """饼图替代方案：用水平占比条模拟饼图效果（避免BLOCK_ARC）
    items = [(标签, 占比0-100, 颜色), ...]  总和应≈100
    """
    if title: TX(s, LM, y-Inches(0.3), CW, Inches(0.25), title, 10, NV, True)
    total = sum(v for _,v,_ in items) or 100
    bar_w = Inches(8.0)
    bar_y = y
    # 占比条（用多个矩形拼接模拟饼图）
    cx = LM + Inches(1.0)
    for i, (lb, val, clr) in enumerate(items):
        bw = bar_w * val / total
        R(s, cx, bar_y, max(bw, Inches(0.1)), Inches(0.5), clr)
        TX(s, cx + Inches(0.05), bar_y, bw - Inches(0.1), Inches(0.5), 
           f"{lb} {val}%" if bw > Inches(1.0) else "", 10, WH, True, PP_ALIGN.LEFT, 'ctr')
        cx += bw
    # 图例
    ly = bar_y + Inches(0.7)
    for i, (lb, val, clr) in enumerate(items):
        lx = LM + Inches(1.0) + i * Inches(2.5)
        R(s, lx, ly, Inches(0.25), Inches(0.2), clr)
        TX(s, lx + Inches(0.35), ly, Inches(2.0), Inches(0.2), f"{lb} {val}%", 10, D3)

# ---- 布局：#72 ICON_DECORATION - 装饰图标 ----
def ICON(s, x, y, icon_char, sz=Inches(0.4), bg=NV, fg=WH):
    """装饰图标（用带背景的圆形+字符模拟图标）
    icon_char: 单字符emoji或文字
    """
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    c.fill.solid(); c.fill.fore_color.rgb = bg; c.line.fill.background(); cs(c)
    tf = c.text_frame; p = tf.paragraphs[0]; p.text = icon_char
    p.font.size = Pt(int(sz/Inches(1.0)*18)); p.font.color.rgb = fg; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs: sf(run)
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None: bp.set('anchor', 'ctr')

def DECORATED_CARD(s, x, y, w, h, icon_char, title, desc, clr):
    """带图标的装饰卡片（对标豆包带图标卡片）
    icon_char: 图标字符(emoji)
    """
    R(s, x, y, w, h, BG)
    R(s, x, y, w, Inches(0.06), clr)
    ICON(s, x+Inches(0.15), y+Inches(0.2), icon_char, Inches(0.45), clr, WH)
    TX(s, x+Inches(0.7), y+Inches(0.15), w-Inches(0.9), Inches(0.35), title, 16, clr, True)
    TX(s, x+Inches(0.15), y+Inches(0.75), w-Inches(0.3), h-Inches(0.9), desc, 11, D3)

# ---- 布局：#73 MATRIX - 比较矩阵 ----
def MATRIX(s, headers, rows, y=TZ):
    """比较矩阵（对标豆包对比分析模板）
    headers = [列标题...]  第一列是行名
    rows = [(行名, 列值1, 列值2, ...), ...]
    第一行作为高亮行
    """
    n_cols = len(headers)
    n_rows = len(rows)
    col_w = (CW - Inches(0.1)*(n_cols-1)) / n_cols
    row_h = Inches(0.55)
    # 表头
    for ci, hdr in enumerate(headers):
        hx = LM + ci*(col_w + Inches(0.1))
        R(s, hx, y, col_w, row_h, NV)
        TX(s, hx, y, col_w, row_h, hdr, 11, WH, True, PP_ALIGN.CENTER, 'ctr')
    # 数据行
    for ri, row in enumerate(rows):
        ry = y + row_h + Inches(0.05) + ri*(row_h + Inches(0.05))
        bg_clr = NV if ri == 0 else (BG if ri % 2 == 1 else WH)
        tc = WH if ri == 0 else D3
        for ci, val in enumerate(row):
            hx = LM + ci*(col_w + Inches(0.1))
            R(s, hx, ry, col_w, row_h, bg_clr if ci == 0 else (RGBColor(0xF8,0xF8,0xF8) if ri%2==1 else WH))
            TX(s, hx, ry, col_w, row_h, str(val), 10, tc, ri==0 or ci==0, PP_ALIGN.CENTER, 'ctr')

# ═══════════════════════════════════════════
# Ⅷ. 质量审查（v5）
# ═══════════════════════════════════════════

def AUDIT(prs):
    """自审"""
    total = sum(1 for _ in prs.slides for _ in _.shapes)
    pages = len(prs.slides)
    passed, issues = [], []
    # 图形量
    if total >= 200: passed.append(f"✅ 图形量 {total}≥200")
    else: issues.append(f"❌ 图形量 {total}<200")
    # 页数
    if pages >= 18: passed.append(f"✅ 页数 {pages}≥18")
    else: issues.append(f"❌ 页数 {pages}<18")
    # 密度
    density = total // max(pages, 1)
    if density >= 10: passed.append(f"✅ 密度 {density}图/页≥10")
    else: issues.append(f"⚠️ 密度 {density}<10")
    # 字体（抽样）
    yahei = 0
    total_runs = 0
    for sl in prs.slides:
        for sh in sl.shapes:
            if hasattr(sh, 'text_frame'):
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        total_runs += 1
                        if r.font.name == FN: yahei += 1
    if total_runs > 0 and yahei >= total_runs * 0.8:
        passed.append(f"✅ 微软雅黑覆盖 {yahei}/{total_runs}")
    else:
        issues.append(f"⚠️ 字体覆盖 {yahei}/{total_runs}")
    for p in passed: print(f"  {p}")
    for i in issues: print(f"  {i}")
    print(f"总数:{pages}页 | 图形:{total} | 密度:{density}图/页")
    return passed, issues

# ═══════════════════════════════════════════
# Ⅶ. 文件清理
# ═══════════════════════════════════════════

def full_cleanup(p):
    """全XML清理（v4：支持PYRAMID布局等新增形状）"""
    tmp = p + '.tmp'
    with zipfile.ZipFile(p, 'r') as zin:
        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                d = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(d)
                    np = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    na = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for st in root.findall(f'.//{{{np}}}style'):
                        sp = st.getparent()
                        if sp is not None: sp.remove(st)
                    if 'theme' in item.filename.lower():
                        for tg in ['outerShdw','innerShdw','scene3d','sp3d']:
                            for el in root.findall(f'.//{{{na}}}{tg}'):
                                ep = el.getparent()
                                if ep is not None: ep.remove(el)
                    d = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)
                zout.writestr(item, d)
    os.replace(tmp, p)


# ═══════════════════════════════════════════
# Ⅷ. 使用示例
# ═══════════════════════════════════════════

def PYRAMID_demo():
    """演示：深蓝递进金字塔布局"""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    BL = prs.slide_layouts[6]
    s = prs.slides.add_slide(BL)
    HDR(s, "AI星工厂2.0 商品定义逻辑框架", "AI Star Factory 2.0 · 精益咨询+AI数字员工+数字化底座")
    PYRAMID(s, [
        ('L5', "第五层：技术·底座·对象", "技术&底座：ODDO·CLAW·精益工具集", LV5,
         [("ODDO部署", SV7), ("CLAW框架", SV8), ("精益工具", SV9)]),
        ('L4', "第四层：功能·模块·架构", "AI星工厂架构·实施方法论·运营体系", LV4,
         [("AI星工厂架构", SV7), ("实施方法论", SV8), ("运营体系", SV9)]),
        ('L3', "第三层：要素·策略·定义", "要素：精益AI+底座·策略：543/QCD·定义：三位一体", LV3,
         [("要素", SV4), ("策略", SV5), ("定义", SV6)]),
        ('L2', "第二层：QCD核心价值", "交付加速·成本精控·质量跃升", LV2,
         [("交付加速", SV1), ("成本精控", SV2), ("质量跃升", SV3)]),
        ('L1', "第一层：愿景使命", "让人们享受智能化带来的美好生活", LV1, None),
    ], "广西星网智云科技有限公司 | AI星工厂2.0")
    out = "/tmp/pyramid_demo.pptx"
    prs.save(out); full_cleanup(out)
    print(f"✅ 金字塔示例：{out}")

def main_demo():
    """民和电器示例（默认）"""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    BL = prs.slide_layouts[6]

    # === P1: 封面 ===
    s = prs.slides.add_slide(BL)
    cover(s, "民和电器精益咨询项目\n总结报告", "精益生产 · 计划物控 · 数字化转型", "2021—2023", "连恩信息科技（厦门）有限公司")
    PN(s, 1, 19)

    # === P2: 目录 ===
    s = prs.slides.add_slide(BL)
    toc(s, [
        ("1", "项目回顾与总体概览", "项目背景、三大旗舰项目"),
        ("2", "组织架构/职责/流程优化", "调研诊断、新架构设计"),
        ("3", "生产计划模式优化", "统一策划、联动排产"),
        ("4", "物流转运优化方案", "搬运改善、容器化"),
        ("5", "新计划模式验证成果", "500+机种验证数据"),
        ("6", "行业对标与价值分析", "库存周转、数字化率"),
        ("7", "数字化规划与MES系统", "四层架构、实施路线"),
        ("8", "项目核心价值与展望", "成效总结、三阶段路线"),
    ])
    PN(s, 2, 19)

    # === P3: SEC01 ===
    s = prs.slides.add_slide(BL); SEC(s, 1, "项目回顾与总体概览"); PN(s, 3, 19)

    # === P4: 项目背景 ===
    s = prs.slides.add_slide(BL)
    AT(s, "民和电器携手连恩，启动精益+数字化双轮驱动战略转型")
    TX(s, LM, Inches(1.3), CW, Inches(1.2), [
        "▎ ABB战略伙伴，2021年启动三大项目",
        "▎ 全球精益生产及物流体系数字化优化转型",
        "▎ 智造2025数字化车间+MES系统优化",
        "▎ 2024年电气机械业增加值同比+5.0%（国统局）",
    ], 12, D3)
    MCARDS(s, [
        ("2年+", "项目周期", "2021-2023", NV),
        ("3个", "旗舰项目", "ABB合作", AB),
        ("500+", "验证机种", "多品类覆盖", AG),
        ("1.9天", "装配周期", "优于行业3-5天", AO),
    ], Inches(3.5))
    SR(s, "来源：国家统计局·2024年"); PN(s, 4, 19)

    # === P5: 核心目标 ===
    s = prs.slides.add_slide(BL)
    AT(s, "核心目标与行业对标 — KPI总览")
    KTRACK(s, [
        ("制造周期", 0.85, "1.9天/目标2.0天", "on"),
        ("WIP降低", 0.70, "↓30-50%/目标", "risk"),
        ("OEE", 0.82, "85%/目标", "on"),
        ("交付率", 0.95, "99%/目标", "on"),
        ("排产效率", 0.75, "1-1.5h/目标", "on"),
    ])
    SR(s, "来源：TEEPTRAK·2024年OEE行业基准"); PN(s, 5, 19)

    # === P6: SEC02 ===
    s = prs.slides.add_slide(BL); SEC(s, 2, "组织架构/职责/流程优化"); PN(s, 6, 19)

    # === P7: 调研发现 ===
    s = prs.slides.add_slide(BL)
    AT(s, "深入调研发现问题 — 组织架构不匹配、职责交叉")
    BA(s,
        ["• 职能部门各自为政，缺乏流程导向",
         "• 岗位职责交叉，审批流程冗长",
         "• 信息传递效率低，跨部门协作困难"],
        ["• 新架构：从职能型转向流程型",
         "• 职责明确：每个岗位职责清晰",
         "• 流程优化：审批环节减少50%"],
    )
    SR(s, "连恩咨询诊断报告"); PN(s, 7, 19)

    # === P8: 改善方案 ===
    s = prs.slides.add_slide(BL)
    AT(s, "三大改善支柱 — 架构重组 · 流程优化 · 能力建设")
    PILLAR(s, [
        ("架构重组", ["• 新组织架构设计", "• 岗位职责明确", "• 管理层级精简"], NV),
        ("流程优化", ["• 端到端流程梳理", "• 审批环节精简", "• 信息化支撑"], AB),
        ("能力建设", ["• 精益培训体系", "• 排产方法掌握", "• 持续改善文化"], AG),
    ])
    SR(s, "连恩咨询改善方案"); PN(s, 8, 19)

    # === P9: SEC03 ===
    s = prs.slides.add_slide(BL); SEC(s, 3, "生产计划模式优化"); PN(s, 9, 19)

    # === P10: 计划流程 ===
    s = prs.slides.add_slide(BL)
    AT(s, "新计划模式 — 统一策划·联动排产·精细控制·异常快反")
    PILLAR(s, [
        ("统一策划", ["• 计划统一编制", "• 控制生产节奏"], NV),
        ("联动齐套", ["• 多车间协同", "• 减少在制停滞"], AB),
        ("精细控制", ["• 精细化排产", "• WIP显著↓"], AG),
    ], TZ)
    FLOW(s, [("需求接收", ""), ("主计划", ""), ("联动排产", ""), ("齐套检查", ""), ("工单下发", ""), ("进度跟踪", ""), ("异常闭环", "")], Inches(5.2))
    SR(s, "民和电器内部方案"); PN(s, 10, 19)

    # === P11: 实施时间线 ===
    s = prs.slides.add_slide(BL)
    AT(s, "项目实施路线图 — 分三阶段推进，30天关键里程碑")
    TIMELINE(s, [
        ("调研诊断", "5天\n需求+痛点识别"),
        ("方案设计", "5天\n方案确认"),
        ("审批确认", "3天\n管理层审批"),
        ("方案落地", "10天\n试运行"),
        ("流程优化", "12天\n持续改善"),
        ("培训推广", "8天\n全员赋能"),
        ("效果评估", "5天\n数据验证"),
    ], Inches(3.5))
    SR(s, "连恩咨询实施计划"); PN(s, 11, 19)

    # === P12: SEC04 ===
    s = prs.slides.add_slide(BL); SEC(s, 4, "物流转运优化方案"); PN(s, 12, 19)

    # === P13: 物流优化 ===
    s = prs.slides.add_slide(BL)
    AT(s, "物流转运优化 — 降低搬运活性系数，消除无效搬运与停滞")
    BA(s,
        ["• 物料散放，活性系数低",
         "• 人工搬运占比高",
         "• 转运路径长，迂回严重",
         "• 等待停滞时间占比>30%"],
        ["• 容器化/托盘化，活性系数↑",
         "• 搬运路径优化，距离↓40%",
         "• 工序间衔接顺畅",
         "• 流转效率提升50%+"],
    )
    SR(s, "搬运活性指数：物料存放状态对搬运的影响"); PN(s, 13, 19)

    # === P14: 库存分析 ===
    s = prs.slides.add_slide(BL)
    AT(s, "库存管理行业对标 — 中国制造业与国际差距1.4-2.3倍")
    TB(s, LM, TZ, 5, 6,
        [["行业", "中国", "国际", "差距", "改善方向", "数字率"],
         ["机械", "2.8次/年", "6.5次", "↓2.3倍", "精益+数字", "38%"],
         ["电子", "4.9次", "7.2次", "↓1.5倍", "计划联动", "42%"],
         ["汽车", "3.2次", "6.0次", "↓1.9倍", "供应链", "37%"],
         ["政策", "60%普及", "2026目标", "—", "国务院", "60%"]],
        [Inches(1.3), Inches(2.0), Inches(2.0), Inches(1.8), Inches(2.0), Inches(1.5)])
    TX(s, LM, Inches(4.5), CW, Inches(0.5), "▎ 广东某电子企业库存周转3.6→6.1次/年（+69%）| 来源：简道云·2026", 9, D3)
    SR(s, "简道云·2026库存周转率基准 + 国务院·数字中国规划"); PN(s, 14, 19)

    # === P15: SEC05 ===
    s = prs.slides.add_slide(BL); SEC(s, 5, "新计划模式验证成果", "L4: 量化成果"); PN(s, 15, 19)

    # === P16: 验证成果 ===
    s = prs.slides.add_slide(BL)
    AT(s, "新计划模式2个月实战验证 — 装配周期1.9天，优于行业基准3-5天")
    BN(s, [
        ("500+", "验证机种", NV),
        ("1.9天", "装配周期", AG),
        ("97.2%", "完工率", NV),
        ("1-1.5h", "排产时间", AO),
    ])
    TB(s, LM, Inches(3.5), 6, 6,
        [["车间", "批次", "有完工", "平均周期", "行业参考", "结论"],
         ["装配", "442", "432(97.7%)", "1.9天", "3-5天", "领先✅"],
         ["注塑", "98", "93(94.9%)", "2.4天", "3-6天", "优于✅"],
         ["五金", "35", "34(97.1%)", "2.8天", "4-7天", "改善✅"],
         ["跨车间", "6", "6(100%)", "9.0天", "10-15天", "优于✅"],
         ["合计", "581", "565(97.2%)", "—", "3-7天", "整体优于✅"]],
        [Inches(1.2), Inches(1.2), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)])
    SR(s, "数据：2022.2.22-3.26 | TEEPTRAK·2024"); PN(s, 16, 19)

    # === P17: 成果仪表盘 ===
    s = prs.slides.add_slide(BL)
    KPI_CARDS = [("¥8.5亿", "年营收", NV), ("97.2%", "完工率", AG), ("1.9天", "装配周期", AO), ("500+", "验证机种", AB)]
    BARS_data = [[45, 65, 82, 92, 97], [30, 50, 70, 80, 88]]
    BARS_labels = ["M1", "M2", "M3", "M4", "M5"]
    BARS_colors = [NV, AB]
    DASH(s, KPI_CARDS, BARS_data, BARS_labels, BARS_colors,
         "验证期间完工率持续提升，第5个月达到97.2%，验证通过率优于行业基准")
    PN(s, 17, 19)

    # === P18: SEC06 ===
    s = prs.slides.add_slide(BL); SEC(s, 6, "项目核心价值与展望"); PN(s, 18, 19)

    # === P19: 行动项 ===
    s = prs.slides.add_slide(BL)
    AT(s, "持续改善行动计划 — 三阶段推进精益数字化深度融合")
    ACTIONS(s, [
        ("维护BOM\n管控质量", "短期≤3月", "项目团队"),
        ("全产品线\n深化MES", "中期3-12月", "IT+运营"),
        ("智能排产\n预测计划", "长期1-3年", "全公司"),
    ])
    SR(s, "连恩咨询持续改善建议"); PN(s, 19, 19)

    # === P20: 结束页 ===
    s = prs.slides.add_slide(BL)
    CLOSE(s,
        "连恩智能 —— 帮助中国制造\"快\"起来！",
        "精益生产 · 计划物控 · 数字化转型\n一套体系 | 一支队伍 | 一个平台",
        "连恩信息科技（厦门）有限公司")

    # ── 保存 & 清理 ──
    out = "/Users/ericshao/Downloads/民和电器精益咨询项目总结报告_pro.pptx"
    prs.save(out)
    full_cleanup(out)
    print(f"\n✅ 生成完成：{out}")
    print("=" * 50)
    print("📋 PPT自审报告")
    print("=" * 50)
    AUDIT(prs)

if __name__ == '__main__':
    main_demo()
