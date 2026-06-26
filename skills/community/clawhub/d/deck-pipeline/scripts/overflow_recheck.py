#!/usr/bin/env python3
"""Overflow estimator with severity tiers.

Usage:
    python3 overflow_recheck.py <pptx-path> [--all-severity]

Default: surfaces HIGH only. `--all-severity` reports MED/LOW too.

Improved over v3.1.0:
  - Honors auto_size (SHAPE_TO_FIT_TEXT / TEXT_TO_FIT_SHAPE both skipped)
  - Reads actual margins
  - Line-height multiplier 1.15 (was 1.25)
  - Per-character width by class (narrow/wide/digit/upper/space)
  - Greedy word-wrap simulation
"""
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

NARROW = set("ijlt.,;:!|()[]{}'\"`-")
WIDE = set("MWmw@%&")
DIGITS = set("0123456789")
UPPER = set("ABCDEFGHIJKLMNOPQRSTUVWXYZ")

LINE_HEIGHT = 1.15
SEV_HIGH = 1.5
SEV_MED = 1.15
DEFER_THRESHOLD = 10


def char_width_pt(ch, size_pt):
    if ch in NARROW: return size_pt * 0.30
    if ch in WIDE: return size_pt * 0.85
    if ch in DIGITS: return size_pt * 0.60
    if ch in UPPER: return size_pt * 0.70
    if ch == " ": return size_pt * 0.30
    return size_pt * 0.55


def measure_text_width(text, size_pt):
    return sum(char_width_pt(c, size_pt) for c in text)


def wrap_lines(text, size_pt, width_pt):
    if width_pt <= 0:
        return 1
    lines = 0
    for para in text.split("\n"):
        if not para:
            lines += 1
            continue
        words = para.split(" ")
        cur_w = 0
        has = False
        for w in words:
            ww = measure_text_width(w, size_pt)
            sp = char_width_pt(" ", size_pt)
            need = ww + (sp if has else 0)
            if has and cur_w + need > width_pt:
                lines += 1
                cur_w = ww
                has = True
            else:
                cur_w += need
                has = True
        if has:
            lines += 1
    return max(1, lines)


def shape_text_and_size(shp):
    if not shp.has_text_frame:
        return "", None
    texts = []
    sizes = []
    for p in shp.text_frame.paragraphs:
        line = "".join(r.text for r in p.runs)
        texts.append(line)
        for r in p.runs:
            if r.text.strip() and r.font.size is not None:
                sizes.append(r.font.size.pt)
    text = "\n".join(texts)
    if not sizes:
        return text, None
    sizes.sort()
    return text, sizes[len(sizes) // 2]


def get_margins_pt(tf):
    try:
        l = (tf.margin_left or 91440) / 12700
        r = (tf.margin_right or 91440) / 12700
        t = (tf.margin_top or 45720) / 12700
        b = (tf.margin_bottom or 45720) / 12700
        return l, r, t, b
    except Exception:
        return 7.2, 7.2, 3.6, 3.6


def check_shape(shp, sidx, results, top_level=True):
    if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shp.shapes:
            check_shape(s, sidx, results, top_level=False)
        return
    if not shp.has_text_frame or shp.width is None or shp.height is None:
        return
    text, size_pt = shape_text_and_size(shp)
    if not text.strip() or size_pt is None:
        return
    tf = shp.text_frame
    try:
        autosize = tf.auto_size
    except Exception:
        autosize = None
    # autosize 1 = SHAPE_TO_FIT_TEXT, 2 = TEXT_TO_FIT_SHAPE (both safe)
    if autosize is not None and int(autosize) in (1, 2):
        return
    try:
        wrap = tf.word_wrap
    except Exception:
        wrap = None
    ml, mr, mt, mb = get_margins_pt(tf)
    width_pt = shp.width / 12700 - ml - mr
    height_pt = shp.height / 12700 - mt - mb
    if width_pt <= 0 or height_pt <= 0:
        return
    line_h = size_pt * LINE_HEIGHT
    lines_fit = max(1, int(height_pt / line_h))

    if wrap is False:
        w_needed = measure_text_width(text.replace("\n", " "), size_pt)
        ratio = w_needed / width_pt
        if ratio > 1.0:
            results.append({
                "slide": sidx, "ratio": ratio, "size_pt": size_pt,
                "mode": "no-wrap-width", "text": text[:80],
                "shape_name": shp.name,
            })
        return

    lines_needed = wrap_lines(text, size_pt, width_pt)
    if lines_needed > lines_fit:
        ratio = lines_needed / lines_fit
        results.append({
            "slide": sidx, "ratio": ratio, "size_pt": size_pt,
            "mode": "wrap-height",
            "lines_needed": lines_needed, "lines_fit": lines_fit,
            "text": text[:80], "shape_name": shp.name,
        })


def severity(ratio):
    if ratio > SEV_HIGH: return "HIGH"
    if ratio > SEV_MED: return "MED"
    return "LOW"


def main():
    if len(sys.argv) < 2:
        print("Usage: overflow_recheck.py <pptx-path> [--all-severity]")
        sys.exit(1)
    all_sev = "--all-severity" in sys.argv
    prs = Presentation(sys.argv[1])
    results = []
    for idx, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            check_shape(shp, idx, results)

    by_sev = {"HIGH": [], "MED": [], "LOW": []}
    for r in results:
        by_sev[severity(r["ratio"])].append(r)

    print("=" * 80)
    print("OVERFLOW RECHECK")
    print("=" * 80)
    print(f"Total: {len(results)}")
    print(f"  HIGH: {len(by_sev['HIGH'])}")
    print(f"  MED:  {len(by_sev['MED'])}")
    print(f"  LOW:  {len(by_sev['LOW'])}")

    if len(by_sev["HIGH"]) > DEFER_THRESHOLD:
        print(f"\n⚠ HIGH count > {DEFER_THRESHOLD}. Defer to user-side rendering:")
        print("  → Export deck to PDF via Keynote/PowerPoint, identify problematic")
        print("    pages visually, then ask for targeted fixes by page number.")
        return

    bands = ["HIGH"] if not all_sev else ["HIGH", "MED", "LOW"]
    for sev in bands:
        if not by_sev[sev]:
            continue
        print(f"\n=== {sev} ({len(by_sev[sev])}) ===")
        for r in by_sev[sev]:
            extra = (f"lines {r['lines_needed']}/{r['lines_fit']}"
                     if r["mode"] == "wrap-height" else "no-wrap")
            print(f"  Slide {r['slide']:>2}  ratio={r['ratio']:.2f}  "
                  f"{r['size_pt']:.1f}pt  {extra}")
            print(f"           text: '{r['text']}'")


if __name__ == "__main__":
    main()
