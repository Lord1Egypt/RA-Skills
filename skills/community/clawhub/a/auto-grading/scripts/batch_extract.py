#!/usr/bin/env python3
"""批量提取存档文档文本 v1.1 — 支持 .doc/.docx/.pdf/.pptx"""
import os, sys, json, argparse
from pathlib import Path

def extract_docx(path):
    from docx import Document
    doc = Document(path)
    text = '\n'.join([p.text for p in doc.paragraphs])
    images = sum(1 for rel in doc.part.rels.values() if "image" in rel.reltype)
    domain_codes = sum(1 for p in doc.paragraphs if any(k in p.text for k in ['TOC','HYPERLINK','PAGEREF','REF _Ref','SHAPE']))
    return {'text': text, 'images': images, 'domain_codes': domain_codes}

def extract_doc(path):
    """WPS/Word .doc 二进制提取（三级降级）"""
    # 1. textutil (macOS)
    import subprocess, tempfile
    tmp = tempfile.NamedTemporaryFile(suffix='.txt', delete=False)
    tmp.close()
    r = subprocess.run(['textutil', '-convert', 'txt', '-output', tmp.name, path],
                       capture_output=True, timeout=30)
    if r.returncode == 0:
        with open(tmp.name, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        os.unlink(tmp.name)
        if len(text) > 100:
            return {'text': text, 'method': 'textutil'}
    # 2. antiword
    r = subprocess.run(['antiword', path], capture_output=True, text=True, timeout=30)
    if r.returncode == 0 and len(r.stdout) > 100:
        return {'text': r.stdout, 'method': 'antiword'}
    # 3. olefile 二进制提取
    try:
        import olefile
        ole = olefile.OleFileIO(path)
        data = ole.openstream('WordDocument').read()
        text_parts = []
        i = 0
        while i < len(data) - 1:
            cp = data[i] | (data[i+1] << 8)
            if (0x4E00 <= cp <= 0x9FFF or 0x20 <= cp <= 0x7E or
                cp in (0x0D, 0x0A, 0xFF0C, 0x3001, 0x3002)):
                start = i
                while i < len(data) - 1:
                    cp2 = data[i] | (data[i+1] << 8)
                    if not (0x4E00 <= cp2 <= 0x9FFF or 0x20 <= cp2 <= 0x7E or
                            cp2 in (0x0D, 0x0A, 0xFF0C, 0x3001, 0x3002)):
                        break
                    i += 2
                chunk = data[start:i].decode('utf-16-le', errors='ignore')
                if len(chunk.strip()) > 3:
                    text_parts.append(chunk)
            i += 1
        ole.close()
        return {'text': '\n'.join(text_parts), 'method': 'olefile'}
    except:
        return {'text': '', 'method': 'failed'}

def extract_pdf(path):
    from PyPDF2 import PdfReader
    reader = PdfReader(open(path, 'rb'))
    text = ''
    for page in reader.pages:
        t = page.extract_text()
        if t: text += t
    return {'text': text, 'pages': len(reader.pages)}

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    imgs = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)
            if shape.shape_type == 13:
                imgs += 1
    return {'text': '\n'.join(texts), 'slides': len(prs.slides), 'images': imgs}

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('directory', help='存档目录路径')
    parser.add_argument('--output', default='/tmp/auto_grading', help='输出目录')
    args = parser.parse_args()
    
    os.makedirs(args.output, exist_ok=True)
    manifest = []
    
    for root, dirs, files in os.walk(args.directory):
        for f in sorted(files):
            if f.startswith("~$"): continue
            path = os.path.join(root, f)
            ext = Path(f).suffix.lower()
            try:
                if ext == '.docx': r = extract_docx(path)
                elif ext == '.doc': r = extract_doc(path)
                elif ext == '.pdf': r = extract_pdf(path)
                elif ext == '.pptx': r = extract_pptx(path)
                else: continue
                
                txt_name = Path(f).stem + '.txt'
                txt_path = os.path.join(args.output, txt_name)
                with open(txt_path, 'w', encoding='utf-8') as wf:
                    wf.write(r.get('text', ''))
                
                entry = {'file': f, 'path': path, 'ext': ext, 'txt_path': txt_path}
                entry.update({k:v for k,v in r.items() if k != 'text'})
                manifest.append(entry)
                method = r.get('method', 'native')
                print(f'OK: {f} ({len(r.get("text",""))} chars, {method})')
            except Exception as e:
                print(f'FAIL: {f} - {e}')
    
    manifest_path = os.path.join(args.output, '_manifest.json')
    with open(manifest_path, 'w', encoding='utf-8') as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    print(f'\nManifest: {manifest_path} ({len(manifest)} files)')

if __name__ == '__main__':
    main()
