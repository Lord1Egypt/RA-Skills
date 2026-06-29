#!/usr/bin/env python3
from __future__ import annotations

import argparse
import csv
import hashlib
import importlib.util
import json
import mimetypes
import re
import shutil
import sys
import time
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

TEXT_EXTS = {'.txt', '.md', '.html', '.htm', '.json', '.xml'}
IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.webp', '.tif', '.tiff', '.bmp'}
PDF_EXTS = {'.pdf'}
SHEET_EXTS = {'.xlsx', '.xls', '.csv'}
SUPPORTED_EXTS = TEXT_EXTS | IMAGE_EXTS | PDF_EXTS | SHEET_EXTS | {'.docx', '.doc', '.pptx', '.ppt'}


def have_module(name: str) -> bool:
    return importlib.util.find_spec(name) is not None


def have_cli(name: str) -> bool:
    return shutil.which(name) is not None


def module_version(name: str) -> Optional[str]:
    if not have_module(name):
        return None
    try:
        mod = __import__(name)
        return getattr(mod, '__version__', None)
    except Exception:
        return None


def capabilities() -> Dict[str, Any]:
    return {
        'markitdown': 'available' if have_module('markitdown') else 'missing',
        'pymupdf': 'available' if have_module('fitz') else 'missing',
        'pypdf': 'available' if have_module('pypdf') else 'missing',
        'pdfplumber': 'available' if have_module('pdfplumber') else 'missing',
        'paddleocr': 'available' if have_module('paddleocr') else 'missing',
        'python_docx': 'available' if have_module('docx') else 'missing',
        'python_pptx': 'available' if have_module('pptx') else 'missing',
        'openpyxl': 'available' if have_module('openpyxl') else 'missing',
        'pandas': 'available' if have_module('pandas') else 'missing',
        'pillow': 'available' if have_module('PIL') else 'missing',
        'libreoffice': 'available' if (have_cli('soffice') or have_cli('libreoffice')) else 'missing',
        'versions': {
            'markitdown': module_version('markitdown'),
            'pymupdf': module_version('fitz'),
            'pypdf': module_version('pypdf'),
            'pdfplumber': module_version('pdfplumber'),
            'paddleocr': module_version('paddleocr'),
            'openpyxl': module_version('openpyxl'),
            'pandas': module_version('pandas'),
        },
    }


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open('rb') as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b''):
            digest.update(chunk)
    return digest.hexdigest()


def safe_read_text(path: Path, max_chars: int = 200_000) -> str:
    for encoding in ('utf-8', 'utf-8-sig', 'gb18030', 'latin-1'):
        try:
            return path.read_text(encoding=encoding, errors='replace')[:max_chars]
        except Exception:
            continue
    return ''


def safe_name(name: str) -> str:
    return re.sub(r'[^A-Za-z0-9_\-\u4e00-\u9fff]+', '_', name).strip('_') or 'document'


def text_quality(text: str) -> Dict[str, Any]:
    if not text:
        return {'chars': 0, 'readable_ratio': 0.0, 'garbled_ratio': 1.0, 'duplicate_line_ratio': 0.0, 'warnings': ['empty_extraction']}
    chars = len(text)
    garbled = text.count('�')
    control = sum(1 for c in text if ord(c) < 32 and c not in '\n\r\t')
    readable = sum(1 for c in text if c.isalnum() or '\u4e00' <= c <= '\u9fff' or c in '，。；：、,.!?()[]{}<>《》+-=*/%:;\n\t \'')
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return {
        'chars': chars,
        'readable_ratio': round(readable / max(chars, 1), 4),
        'garbled_ratio': round(min(1.0, (garbled + control) / max(chars, 1)), 4),
        'duplicate_line_ratio': round(1.0 - len(set(lines)) / len(lines), 4) if lines else 0.0,
        'warnings': [w for w in [
            'very_low_text_volume' if chars < 50 else None,
            'low_readable_ratio' if readable / max(chars, 1) < 0.65 else None,
            'garbled_text' if min(1.0, (garbled + control) / max(chars, 1)) > 0.03 else None,
            'high_duplicate_lines' if lines and 1.0 - len(set(lines)) / len(lines) > 0.35 else None,
        ] if w],
    }


def write_json(path: Path, obj: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(obj, ensure_ascii=False, indent=2), encoding='utf-8')


def write_jsonl(path: Path, rows: Iterable[Dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open('w', encoding='utf-8') as handle:
        for row in rows:
            handle.write(json.dumps(row, ensure_ascii=False) + '\n')


def chunk_text(text: str, document_id: str, source_file: str, method: str, size: int = 900) -> List[Dict[str, Any]]:
    paragraphs = [p.strip() for p in re.split(r'\n\s*\n', text) if p.strip()]
    chunks: List[Dict[str, Any]] = []
    current = ''

    def emit(body: str) -> None:
        if not body.strip():
            return
        idx = len(chunks) + 1
        chunks.append({
            'chunk_id': f'{document_id}#c{idx:04d}',
            'document_id': document_id,
            'source_file': source_file,
            'unit_type': 'section',
            'unit_ref': None,
            'section_path': [],
            'text': body.strip(),
            'keywords': [],
            'entities': [],
            'chunk_type': 'body',
            'extraction_method': method,
            'confidence': 0.85,
        })

    for paragraph in paragraphs:
        if len(current) + len(paragraph) + 2 <= size:
            current = f'{current}\n\n{paragraph}' if current else paragraph
        else:
            emit(current)
            current = paragraph
    emit(current)
    if not chunks and text.strip():
        for start in range(0, len(text), size):
            emit(text[start:start + size])
    return chunks


def probe_pdf(path: Path) -> Dict[str, Any]:
    details: Dict[str, Any] = {'pages': None, 'sample_text_chars': 0, 'warnings': []}
    sample = ''
    if have_module('fitz'):
        try:
            import fitz  # type: ignore
            doc = fitz.open(path)
            details['pages'] = len(doc)
            texts = []
            image_blocks = 0
            for page in doc[: min(3, len(doc))]:
                texts.append(page.get_text('text') or '')
                image_blocks += len(page.get_images(full=True))
            sample = '\n'.join(texts)
            details['sample_text_chars'] = len(sample.strip())
            details['image_blocks'] = image_blocks
        except Exception as exc:
            details['warnings'].append(f'pymupdf_probe_failed: {exc}')
    elif have_module('pypdf'):
        try:
            from pypdf import PdfReader  # type: ignore
            reader = PdfReader(str(path))
            details['pages'] = len(reader.pages)
            sample = '\n'.join((p.extract_text() or '') for p in reader.pages[: min(3, len(reader.pages))])
            details['sample_text_chars'] = len(sample.strip())
        except Exception as exc:
            details['warnings'].append(f'pypdf_probe_failed: {exc}')
    details['sample_quality'] = text_quality(sample)
    return details


def classify(path: Path) -> Dict[str, Any]:
    path = path.resolve()
    ext = path.suffix.lower()
    manifest: Dict[str, Any] = {
        'document_id': sha256_file(path)[:16],
        'source_file': str(path),
        'source_hash': sha256_file(path),
        'file_type': ext.lstrip('.') or 'unknown',
        'mime_type': mimetypes.guess_type(path.name)[0],
        'size_bytes': path.stat().st_size,
        'content_class': 'unknown',
        'classification_confidence': 0.35,
        'capabilities': capabilities(),
        'chosen_routes': [],
        'missing_dependencies': [],
        'probe_details': {},
    }
    caps = manifest['capabilities']
    if ext in TEXT_EXTS:
        q = text_quality(safe_read_text(path, 20_000))
        manifest['probe_details']['sample_quality'] = q
        manifest['content_class'] = 'text_native'
        manifest['classification_confidence'] = 0.95 if q['readable_ratio'] > 0.75 else 0.7
        manifest['chosen_routes'] = ['direct_text']
    elif ext == '.csv':
        manifest['content_class'] = 'tabular'
        manifest['classification_confidence'] = 0.95
        manifest['chosen_routes'] = ['pandas_csv' if caps['pandas'] == 'available' else 'csv_builtin']
    elif ext in {'.xlsx', '.xls'}:
        manifest['content_class'] = 'tabular'
        manifest['classification_confidence'] = 0.9
        if caps['openpyxl'] == 'available' or caps['pandas'] == 'available':
            manifest['chosen_routes'] = ['openpyxl_or_pandas']
        else:
            manifest['missing_dependencies'].append({'needed': 'openpyxl or pandas', 'recommended_profile': 'minimal', 'reason': 'spreadsheet parsing'})
    elif ext in PDF_EXTS:
        probe = probe_pdf(path)
        manifest['probe_details'] = probe
        q = probe['sample_quality']
        pages = int(probe.get('pages') or 1)
        chars_per_page = q['chars'] / max(min(3, pages), 1)
        if chars_per_page >= 200 and q['readable_ratio'] >= 0.7:
            manifest['content_class'] = 'text_native'
            manifest['classification_confidence'] = 0.86
            if caps['markitdown'] == 'available':
                manifest['chosen_routes'] = ['markitdown']
            elif caps['pymupdf'] == 'available':
                manifest['chosen_routes'] = ['pymupdf']
            else:
                manifest['chosen_routes'] = ['pypdf']
        else:
            manifest['content_class'] = 'scanned_or_image'
            manifest['classification_confidence'] = 0.78
            if caps['paddleocr'] == 'available':
                manifest['chosen_routes'] = ['paddleocr']
            else:
                manifest['missing_dependencies'].append({'needed': 'PaddleOCR', 'recommended_profile': 'research_report', 'reason': 'PDF appears scanned or has weak text layer'})
    elif ext in IMAGE_EXTS:
        manifest['content_class'] = 'scanned_or_image'
        manifest['classification_confidence'] = 0.9
        if caps['paddleocr'] == 'available':
            manifest['chosen_routes'] = ['paddleocr']
        else:
            manifest['missing_dependencies'].append({'needed': 'PaddleOCR', 'recommended_profile': 'research_report', 'reason': 'image OCR'})
    elif ext == '.docx':
        manifest['content_class'] = 'text_native'
        manifest['classification_confidence'] = 0.9
        if caps['markitdown'] == 'available':
            manifest['chosen_routes'] = ['markitdown']
        elif caps['python_docx'] == 'available':
            manifest['chosen_routes'] = ['python_docx']
        else:
            manifest['missing_dependencies'].append({'needed': 'markitdown or python-docx', 'recommended_profile': 'office_full', 'reason': 'DOCX parsing'})
    elif ext == '.pptx':
        manifest['content_class'] = 'presentation'
        manifest['classification_confidence'] = 0.9
        if caps['markitdown'] == 'available':
            manifest['chosen_routes'] = ['markitdown']
        elif caps['python_pptx'] == 'available':
            manifest['chosen_routes'] = ['python_pptx']
        else:
            manifest['missing_dependencies'].append({'needed': 'markitdown or python-pptx', 'recommended_profile': 'office_full', 'reason': 'PPTX parsing'})
    elif ext in {'.doc', '.ppt', '.xls'}:
        manifest['content_class'] = 'mixed' if ext in {'.doc', '.ppt'} else 'tabular'
        manifest['classification_confidence'] = 0.65
        if caps['libreoffice'] == 'available':
            manifest['chosen_routes'] = ['libreoffice_convert']
        else:
            manifest['missing_dependencies'].append({'needed': 'LibreOffice headless', 'recommended_profile': 'office_full', 'reason': f'legacy Office format {ext}'})
    return manifest


def extract_with_markitdown(path: Path) -> Optional[str]:
    if not have_module('markitdown'):
        return None
    try:
        from markitdown import MarkItDown  # type: ignore
        result = MarkItDown().convert(str(path))
        return getattr(result, 'text_content', None) or str(result)
    except Exception:
        return None


def extract_pdf_text(path: Path) -> Optional[str]:
    if have_module('fitz'):
        try:
            import fitz  # type: ignore
            doc = fitz.open(path)
            return '\n\n'.join((page.get_text('text') or '') for page in doc)
        except Exception:
            pass
    if have_module('pypdf'):
        try:
            from pypdf import PdfReader  # type: ignore
            reader = PdfReader(str(path))
            return '\n\n'.join((page.extract_text() or '') for page in reader.pages)
        except Exception:
            pass
    return None


def extract_docx(path: Path) -> Optional[str]:
    text = extract_with_markitdown(path)
    if text:
        return text
    if have_module('docx'):
        try:
            import docx  # type: ignore
            d = docx.Document(str(path))
            parts = [p.text for p in d.paragraphs if p.text.strip()]
            for table in d.tables:
                for row in table.rows:
                    parts.append(' | '.join(cell.text.strip() for cell in row.cells))
            return '\n\n'.join(parts)
        except Exception:
            return None
    return None


def extract_pptx(path: Path) -> Optional[str]:
    text = extract_with_markitdown(path)
    if text:
        return text
    if have_module('pptx'):
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides, start=1):
                parts.append(f'# Slide {i}')
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        parts.append(shape.text.strip())
            return '\n\n'.join(parts)
        except Exception:
            return None
    return None


def extract_table(path: Path, out_dir: Path, manifest: Dict[str, Any]) -> str:
    tables_dir = out_dir / 'tables'
    tables_dir.mkdir(parents=True, exist_ok=True)
    ext = path.suffix.lower()
    if ext == '.csv':
        text = safe_read_text(path)
        copy = tables_dir / 'sheet_1.csv'
        copy.write_text(text, encoding='utf-8')
        manifest.setdefault('tables', []).append({'name': 'sheet_1', 'path': str(copy), 'format': 'csv'})
        return '# CSV Table\n\n' + text[:200_000]
    if have_module('openpyxl'):
        try:
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(str(path), data_only=False, read_only=True)
            sections = []
            for ws in wb.worksheets:
                csv_path = tables_dir / f'{safe_name(ws.title)}.csv'
                rows_written = 0
                with csv_path.open('w', newline='', encoding='utf-8-sig') as handle:
                    writer = csv.writer(handle)
                    for row in ws.iter_rows(values_only=True):
                        values = ['' if v is None else v for v in row]
                        if any(str(v).strip() for v in values):
                            writer.writerow(values)
                            rows_written += 1
                manifest.setdefault('tables', []).append({'name': ws.title, 'path': str(csv_path), 'rows': rows_written, 'format': 'csv'})
                sections.append(f'# Sheet: {ws.title}\n\nExtracted rows: {rows_written}\nTable file: {csv_path}')
            return '\n\n'.join(sections)
        except Exception as exc:
            manifest.setdefault('quality', {}).setdefault('warnings', []).append(f'openpyxl_failed: {exc}')
    return ''


def extract_image_ocr(path: Path, manifest: Dict[str, Any]) -> Optional[str]:
    if not have_module('paddleocr'):
        manifest.setdefault('missing_dependencies', []).append({'needed': 'PaddleOCR', 'recommended_profile': 'research_report', 'reason': 'image OCR'})
        return None
    try:
        from paddleocr import PaddleOCR  # type: ignore
        ocr = PaddleOCR(use_angle_cls=True, lang='ch', show_log=False)
        result = ocr.ocr(str(path), cls=True)
        lines = []
        confidences = []
        for page in result or []:
            for item in page or []:
                if len(item) >= 2 and item[1]:
                    text, conf = item[1][0], item[1][1]
                    lines.append(str(text))
                    try:
                        confidences.append(float(conf))
                    except Exception:
                        pass
        if confidences:
            manifest.setdefault('quality', {})['ocr_confidence_avg'] = round(sum(confidences) / len(confidences), 4)
        return '\n'.join(lines)
    except Exception as exc:
        manifest.setdefault('quality', {}).setdefault('warnings', []).append(f'paddleocr_failed: {exc}')
        return None


def parse_file(path: Path, out_base: Path, enable_chunks: bool = True) -> Dict[str, Any]:
    path = path.resolve()
    cls = classify(path)
    out_dir = out_base.resolve() / safe_name(path.stem)
    out_dir.mkdir(parents=True, exist_ok=True)
    manifest: Dict[str, Any] = {
        **cls,
        'metadata': {},
        'structure': {},
        'quality': {'warnings': []},
        'tables': [],
        'provenance': {'extraction_timestamp': time.strftime('%Y-%m-%dT%H:%M:%S%z'), 'tool': 'document_classifier_router.py'},
    }
    ext = path.suffix.lower()
    text: Optional[str] = None
    method = 'none'

    if ext in TEXT_EXTS:
        text = safe_read_text(path)
        method = 'direct_text'
    elif ext in {'.csv', '.xlsx', '.xls'}:
        text = extract_table(path, out_dir, manifest)
        method = 'openpyxl_or_csv'
    elif ext in PDF_EXTS:
        if cls['content_class'] == 'text_native':
            text = extract_with_markitdown(path)
            method = 'markitdown'
            if not text:
                text = extract_pdf_text(path)
                method = 'pymupdf_or_pypdf'
        else:
            text = extract_image_ocr(path, manifest)
            method = 'paddleocr'
    elif ext == '.docx':
        text = extract_docx(path)
        method = 'markitdown_or_python_docx'
    elif ext == '.pptx':
        text = extract_pptx(path)
        method = 'markitdown_or_python_pptx'
    elif ext in IMAGE_EXTS:
        text = extract_image_ocr(path, manifest)
        method = 'paddleocr'

    manifest['chosen_routes'] = [method] if method != 'none' else cls.get('chosen_routes', [])
    q = text_quality(text or '')
    manifest['quality'].update(q)
    if not text:
        manifest['status'] = 'blocked_or_failed'
        if not manifest.get('missing_dependencies'):
            manifest.setdefault('quality', {}).setdefault('warnings', []).append('no_parser_succeeded')
    else:
        manifest['status'] = 'parsed'
        md = f"---\ndocument_id: {manifest['document_id']}\nsource_file: {path.name}\ncontent_class: {manifest['content_class']}\nextraction_method: {method}\n---\n\n{text.strip()}\n"
        (out_dir / 'document.md').write_text(md, encoding='utf-8')
        if enable_chunks:
            chunks = chunk_text(text, manifest['document_id'], str(path), method)
            write_jsonl(out_dir / 'chunks.jsonl', chunks)
            manifest['structure']['chunk_count'] = len(chunks)
    write_json(out_dir / 'document.json', manifest)
    return {'out_dir': str(out_dir), 'manifest': manifest}


def iter_supported_files(input_dir: Path, recursive: bool = True) -> Iterable[Path]:
    pattern = '**/*' if recursive else '*'
    for path in input_dir.glob(pattern):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTS:
            yield path


def batch_process(input_dir: Path, output_dir: Path, copy_sources: bool = False, enable_chunks: bool = True, recursive: bool = True) -> Dict[str, Any]:
    input_dir = input_dir.resolve()
    output_dir = output_dir.resolve()
    parse_root = output_dir / 'parsed'
    sample_root = output_dir / 'samples'
    output_dir.mkdir(parents=True, exist_ok=True)
    results: List[Dict[str, Any]] = []

    for source in iter_supported_files(input_dir, recursive=recursive):
        parse_source = source
        if copy_sources:
            sample_root.mkdir(parents=True, exist_ok=True)
            target = sample_root / source.name
            if target.exists():
                target = sample_root / f"{safe_name(source.stem)}_{sha256_file(source)[:8]}{source.suffix.lower()}"
            shutil.copy2(source, target)
            parse_source = target
        try:
            parsed = parse_file(parse_source, parse_root, enable_chunks=enable_chunks)
            manifest = parsed['manifest']
            results.append({
                'source_file': str(source),
                'parsed_source_file': str(parse_source),
                'out_dir': parsed['out_dir'],
                'status': manifest.get('status'),
                'file_type': manifest.get('file_type'),
                'content_class': manifest.get('content_class'),
                'chosen_routes': manifest.get('chosen_routes', []),
                'missing_dependencies': manifest.get('missing_dependencies', []),
                'quality_warnings': manifest.get('quality', {}).get('warnings', []),
            })
        except Exception as exc:
            results.append({'source_file': str(source), 'status': 'error', 'error': str(exc)})

    summary = {
        'input_dir': str(input_dir),
        'output_dir': str(output_dir),
        'copy_sources': copy_sources,
        'recursive': recursive,
        'file_count': len(results),
        'parsed_count': sum(1 for item in results if item.get('status') == 'parsed'),
        'blocked_or_failed_count': sum(1 for item in results if item.get('status') == 'blocked_or_failed'),
        'error_count': sum(1 for item in results if item.get('status') == 'error'),
        'results': results,
    }
    write_json(output_dir / 'batch_summary.json', summary)
    return summary


def main() -> None:
    parser = argparse.ArgumentParser(description='Lightweight document classifier/router')
    sub = parser.add_subparsers(dest='cmd', required=True)
    sub.add_parser('capabilities')
    p_cls = sub.add_parser('classify')
    p_cls.add_argument('--input', required=True)
    p_parse = sub.add_parser('parse')
    p_parse.add_argument('--input', required=True)
    p_parse.add_argument('--output', required=True)
    p_parse.add_argument('--no-chunks', action='store_true')
    p_batch = sub.add_parser('batch')
    p_batch.add_argument('--input-dir', required=True)
    p_batch.add_argument('--output', required=True)
    p_batch.add_argument('--copy-sources', action='store_true')
    p_batch.add_argument('--no-chunks', action='store_true')
    p_batch.add_argument('--no-recursive', action='store_true')
    args = parser.parse_args()

    if args.cmd == 'capabilities':
        print(json.dumps(capabilities(), ensure_ascii=False, indent=2))
    elif args.cmd == 'classify':
        print(json.dumps(classify(Path(args.input)), ensure_ascii=False, indent=2))
    elif args.cmd == 'parse':
        print(json.dumps(parse_file(Path(args.input), Path(args.output), enable_chunks=not args.no_chunks), ensure_ascii=False, indent=2))
    elif args.cmd == 'batch':
        print(json.dumps(batch_process(Path(args.input_dir), Path(args.output), copy_sources=args.copy_sources, enable_chunks=not args.no_chunks, recursive=not args.no_recursive), ensure_ascii=False, indent=2))


if __name__ == '__main__':
    main()
