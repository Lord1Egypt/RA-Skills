"""PPT editing logic for text and image replacement."""

from __future__ import annotations

from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .image_service import generate_image
from .utils import build_non_overwrite_path, get_logger

logger = get_logger(__name__)


def _replace_text_in_shape(shape, text_replace: Dict[str, str]) -> List[Tuple[str, str]]:
    """Replace matching text in one shape and return applied replacements."""
    replaced_pairs: List[Tuple[str, str]] = []
    if not shape.has_text_frame:
        return replaced_pairs

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            original = run.text
            updated = original
            for old_text, new_text in text_replace.items():
                if old_text in updated:
                    updated = updated.replace(old_text, new_text)
                    replaced_pairs.append((old_text, new_text))
            if updated != original:
                run.text = updated

    return replaced_pairs


def _should_replace_picture(shape, old_keyword: str) -> bool:
    """Check if a picture should be replaced based on current lightweight rules."""
    name_text = f"{getattr(shape, 'name', '')} {getattr(shape, 'alt_text', '')}".lower()
    return old_keyword.lower() in name_text


def _replace_picture_keep_layout(slide, picture_shape, new_image_path: str) -> None:
    """Replace one picture while preserving exact position and size."""
    left = picture_shape.left
    top = picture_shape.top
    width = picture_shape.width
    height = picture_shape.height
    element = picture_shape._element
    element.getparent().remove(element)
    slide.shapes.add_picture(new_image_path, left, top, width=width, height=height)


def _resolve_image_prompt(old_keyword: str, rule_value: Any) -> str:
    """Resolve prompt text from image replacement rule value."""
    if isinstance(rule_value, dict):
        explicit_prompt = rule_value.get("prompt")
        if explicit_prompt:
            return str(explicit_prompt)
        target = str(rule_value.get("target", "")).strip()
        if target:
            return f"High quality image of {target}"
    if isinstance(rule_value, str):
        target = rule_value.strip()
        if target:
            return f"High quality image of {target}"
    # TODO: business-specific prompt strategy should be confirmed by product owner.
    return f"Image related to {old_keyword}"


def update_ppt(
    ppt_path: str | Path,
    replace_rules: Dict[str, Any],
    image_generator: Optional[Callable[[str], str]] = None,
) -> Dict[str, Any]:
    """Update PPT text/images and save as a new file with replacement logs."""
    ppt_file = Path(ppt_path)
    logger.info("Start update_ppt source=%s", ppt_file)
    prs = Presentation(str(ppt_file))
    image_generator = image_generator or generate_image

    text_rules: Dict[str, str] = replace_rules.get("text_replace", {})
    image_rules: Dict[str, Any] = replace_rules.get("image_replace", {})

    replace_log: Dict[str, Any] = {"source": str(ppt_file), "slides": []}

    for idx, slide in enumerate(prs.slides, start=1):
        slide_log: Dict[str, Any] = {"slide_index": idx, "text_changes": [], "image_changes": []}
        logger.info("Processing slide index=%d", idx)

        for shape in list(slide.shapes):
            if text_rules:
                text_changes = _replace_text_in_shape(shape, text_rules)
                for old_text, new_text in text_changes:
                    slide_log["text_changes"].append({"from": old_text, "to": new_text})

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and image_rules:
                for old_keyword, rule_value in image_rules.items():
                    if old_keyword == "__all__":
                        prompt = _resolve_image_prompt(old_keyword, rule_value)
                        new_image_path = image_generator(prompt)
                        _replace_picture_keep_layout(slide, shape, new_image_path)
                        slide_log["image_changes"].append(
                            {"match": old_keyword, "prompt": prompt, "new_image": new_image_path}
                        )
                        break

                    if _should_replace_picture(shape, old_keyword):
                        prompt = _resolve_image_prompt(old_keyword, rule_value)
                        new_image_path = image_generator(prompt)
                        _replace_picture_keep_layout(slide, shape, new_image_path)
                        slide_log["image_changes"].append(
                            {"match": old_keyword, "prompt": prompt, "new_image": new_image_path}
                        )
                        break

        if slide_log["text_changes"] or slide_log["image_changes"]:
            replace_log["slides"].append(slide_log)

    output_path = build_non_overwrite_path(ppt_file, suffix="_replaced")
    prs.save(str(output_path))
    replace_log["output"] = str(output_path)
    logger.info("update_ppt completed output=%s changed_slides=%d", output_path, len(replace_log["slides"]))
    return {"new_ppt_path": str(output_path), "replace_log": replace_log}
