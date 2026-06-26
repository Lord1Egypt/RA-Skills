"""DOCX → PPTX (段落转幻灯片)

用法:
    python docx_to_pptx.py input.docx [output.pptx]

参数:
    input.docx  - DOCX 文件路径
    output.pptx - 输出文件路径 (默认: output.pptx)

说明:
    标题段落创建新幻灯片，正文段落作为 bullet points，每张幻灯片最多 8 个要点
"""
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches
import sys
import os


def main():
    if len(sys.argv) < 2:
        print("用法: python docx_to_pptx.py input.docx [output.pptx]")
        sys.exit(1)

    docx_path = sys.argv[1]
    out_path = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"

    if not os.path.isfile(docx_path):
        print(f"错误: 文件不存在 - {docx_path}")
        sys.exit(1)

    doc = Document(docx_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    slide = None
    tf = None
    bullet_count = 0
    max_bullets = 8

    def new_slide():
        nonlocal slide, tf, bullet_count
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        bullet_count = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name.lower()
        is_heading = "heading" in style or "title" in style

        if is_heading or bullet_count >= max_bullets:
            new_slide()
            if is_heading:
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(32)
                p.font.bold = True
                bullet_count = 0
                continue

        if slide is None:
            new_slide()

        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(20)
        level = 0
        if "Heading" in para.style.name:
            level = min(para.style.name.count("Heading"), 2)
        p.level = level
        bullet_count += 1

    prs.save(out_path)
    print(f"完成: {len(prs.slides)} 张幻灯片 -> {out_path}")


if __name__ == "__main__":
    main()