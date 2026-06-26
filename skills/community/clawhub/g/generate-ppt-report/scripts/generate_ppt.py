# -*- coding: utf-8 -*-
r"""
生成 PPT 简报
按照 E:\daily\cankao\M2.pptx 模板生成每日金融市场 PPT 简报。

布局：单页 Dashboard，3个文本框（Shape 6市场/Shape 15政策/Shape 23企业）

各部分结构：
  市场表现（Shape 6）：美国 / 中国及中国香港 / 黄金（加粗20pt小标题 + 非加粗20pt正文）
  政策（Shape 15）：美国 / 欧洲 / 中国（加粗20pt小标题 + 非加粗20pt正文）
  企业（Shape 23）：美国 / 中国及中国香港（加粗20pt小标题 + 非加粗20pt正文）
  字体统一为微软雅黑
"""
import sys
sys.path.insert(0, r"C:\Users\qu669\.openclaw\workspace-yoyo")
sys.stdout.reconfigure(encoding='utf-8')
import os, json, datetime, logging, re
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree
from copy import deepcopy
from io import BytesIO
import config

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.FileHandler(config.LOG_FILE, encoding="utf-8"), logging.StreamHandler(sys.stdout)]
)
log = logging.getLogger()

FONT_NAME = "微软雅黑"
FONT_SIZE_TITLE = 20
FONT_SIZE_BODY  = 20

# ─────────────────────────────────────────────
# 辅助函数
# ─────────────────────────────────────────────

def gp(key, mdata, default=None):
    d = mdata.get(key, {})
    if isinstance(d, dict):
        v = d.get('price')
        return v if v is not None else default
    return default

def gc(key, mdata, default=None):
    d = mdata.get(key, {})
    if isinstance(d, dict):
        v = d.get('change')
        return v if v is not None else default
    return default

def fp(v):
    if v is None or v == 0: return "N/A"
    if abs(v) >= 10000: return f"{v:,.2f}"
    return f"{v:,.2f}"

def fpc_abs(v):
    if v is None or v == 0: return "N/A"
    return f"{abs(v):.2f}%"

def chn(v):
    if v is None or v == 0: return "涨跌N/A"
    return "涨" if v > 0 else "跌"

def chn_full(v):
    if v is None or v == 0: return "涨跌N/A"
    return "上涨" if v > 0 else "下跌"

def clean_text(text):
    if not text: return text
    for p in ['。', '，', '、', '；', '：']:
        text = re.sub(p + r'{2,}', p, text)
    text = re.sub(r'，。', '。', text)
    text = re.sub(r'。，', '。', text)
    return text.strip().rstrip('，。、;:')


# ─────────────────────────────────────────────
# 直接 XML 操作：重建 txBody 段落
# ─────────────────────────────────────────────

def rebuild_txBody(sp_element, paragraphs_data, title_solidFill_xml):
    """直接操作 sp 元素的 txBody XML，清空并重建段落"""
    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    txBody = sp_element.find('p:txBody', nsmap)
    if txBody is None:
        log.warning("  Shape has no p:txBody, skipping")
        return

    all_p = list(txBody.findall('a:p', nsmap))
    orig_pPr = None
    if all_p:
        first_pPr = all_p[0].find('a:pPr', nsmap)
        if first_pPr is not None:
            orig_pPr = deepcopy(first_pPr)

    ref_rPr = None
    if all_p:
        first_r = all_p[0].find('a:r', nsmap)
        if first_r is not None:
            first_rPr = first_r.find('a:rPr', nsmap)
            if first_rPr is not None:
                ref_rPr = deepcopy(first_rPr)

    for p_elem in list(txBody.findall('a:p', nsmap)):
        txBody.remove(p_elem)

    if orig_pPr is not None:
        txBody.insert(0, orig_pPr)

    for (text, bold, size_pt) in paragraphs_data:
        sz = int(size_pt * 100)
        b = "1" if bold else "0"
        sf_xml = title_solidFill_xml if bold else ""

        if ref_rPr is not None:
            new_rPr = deepcopy(ref_rPr)
            new_rPr.set('sz', str(sz))
            new_rPr.set('b', b)
            old_sf = new_rPr.find('a:solidFill', nsmap)
            if old_sf is not None:
                new_rPr.remove(old_sf)
            if sf_xml:
                sf_elem = etree.fromstring(title_solidFill_xml)
                new_rPr.insert(0, sf_elem)
            rPr_elem = new_rPr
        else:
            if sf_xml:
                rPr_elem = etree.fromstring(
                    f'<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                    f'lang="zh-CN" altLang="en-US" sz="{sz}" b="{b}" dirty="0">'
                    f'{sf_xml}'
                    f'<a:latin typeface="{FONT_NAME}" panose="020B0503020204020204" pitchFamily="34" charset="-122"/>'
                    f'<a:ea typeface="{FONT_NAME}" panose="020B0503020204020204" pitchFamily="34" charset="-122"/>'
                    f'</a:rPr>'
                )
            else:
                rPr_elem = etree.fromstring(
                    f'<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                    f'lang="zh-CN" altLang="en-US" sz="{sz}" b="{b}" dirty="0"/>'
                )

        para_elem = etree.fromstring(
            f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:pPr><a:buNone/></a:pPr>'
            f'</a:p>'
        )
        run_elem = etree.fromstring(
            f'<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )
        run_elem.append(rPr_elem)
        t_elem = etree.fromstring(f'<a:t xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">{text}</a:t>')
        run_elem.append(t_elem)
        para_elem.append(run_elem)
        txBody.append(para_elem)


def get_title_solidFill_xml(sp_element):
    """从 Shape 的 txBody 中提取第一个 run 的 solidFill XML"""
    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    txBody = sp_element.find('p:txBody', nsmap)
    if txBody is None:
        return None
    for p in txBody.findall('a:p', nsmap):
        for r in p.findall('a:r', nsmap):
            rPr = r.find('a:rPr', nsmap)
            if rPr is not None:
                sf = rPr.find('a:solidFill', nsmap)
                if sf is not None:
                    return etree.tostring(sf, encoding='unicode')
    return None


# ─────────────────────────────────────────────
# 主流程
# ─────────────────────────────────────────────

log.info("📊 生成 PPT 简报...")

if not os.path.exists(config.MARKET_DATA_FILE):
    log.error(f"❌ 数据不存在: {config.MARKET_DATA_FILE}"); sys.exit(1)

with open(config.MARKET_DATA_FILE, 'r', encoding='utf-8') as f:
    mdata = json.load(f)

# ── 从 Word 日报读取结构化内容 ──
# Word 文件路径：E:\daily\{YYYY-MM-DD}\金融市场日报_{YYYYMMDD}.docx
word_doc_path = os.path.join(config.OUTPUT_DIR, f"金融市场日报_{config.TODAY.strftime('%Y%m%d')}.docx")
word_paragraphs = []
if os.path.exists(word_doc_path):
    from docx import Document as DocxDoc
    doc_word = DocxDoc(word_doc_path)
    word_paragraphs = [p.text.strip() for p in doc_word.paragraphs]
    log.info(f"   Word 文档已加载，共 {len(word_paragraphs)} 段")
else:
    log.warning(f"   ⚠️ Word 文档未找到: {word_doc_path}，PPT 将完全依赖 market_data.json")

# 从 Word 段落列表中提取关键段落的索引（用于 Shape 6 市场表现）
# Word 结构（以 generate_word.py 为准）：
#   [7]美股内容 [9]美债内容 [11]美元指数内容 [13]就业市场内容
#   [16]A股内容 [18]港股内容 [20]黄金内容
#   [23]美国政策内容 [25]欧洲政策内容 [27]中国政策内容
def _w(idx, default=''):
    """安全读取 Word 段落"""
    return word_paragraphs[idx] if idx < len(word_paragraphs) and word_paragraphs[idx] else default

prs = Presentation(config.TEMPLATE_PPT)
log.info(f"   模板加载成功，共 {len(prs.slides)} 页幻灯片，共 {len(prs.slides[0].shapes)} 个形状")

def _get_price_change(key, section_dict):
    """从市场表现section中获取price和change"""
    d = section_dict.get(key, {})
    price = d.get('price') if isinstance(d, dict) else None
    change = d.get('change') if isinstance(d, dict) else None
    return (price, change)

# ── 提取所有市场数据（使用新的market_data结构） ──
# 市场表现类数据
market_perf = mdata.get('市场表现', {})

# 美国股市
us_stocks = market_perf.get('美国股市', {})
dji     = _get_price_change('道琼斯工业平均指数', us_stocks)
spx     = _get_price_change('标普500指数', us_stocks)
nasdaq  = _get_price_change('纳斯达克综合指数', us_stocks)
vix     = _get_price_change('VIX恐慌指数', us_stocks)

# 美国债券与外汇
us_bonds = market_perf.get('美国债券与外汇', {})
yield10 = us_bonds.get('10年期美债收益率', {}).get('price')
yield10_chg = us_bonds.get('10年期美债收益率', {}).get('change')
yield2  = us_bonds.get('2年期美债收益率', {}).get('price')
usdidx  = us_bonds.get('美元指数(DXY)', {}).get('price')
usdidx_chg = us_bonds.get('美元指数(DXY)', {}).get('change')

# 美国大宗商品
us_commodities = market_perf.get('美国大宗商品', {})
spot_gold  = us_commodities.get('现货黄金(XAUUSD)', {}).get('price')
spot_chg   = us_commodities.get('现货黄金(XAUUSD)', {}).get('change')
comex_gold = us_commodities.get('COMEX黄金期货', {}).get('price')
comex_chg  = us_commodities.get('COMEX黄金期货', {}).get('change')

# A股
cn_stocks = market_perf.get('A股', {})
sh     = _get_price_change('上证指数', cn_stocks)
sz     = _get_price_change('深证成指', cn_stocks)
cyb    = _get_price_change('创业板指', cn_stocks)
hs300  = _get_price_change('沪深300', cn_stocks)

# 港股
hk_stocks = market_perf.get('港股', {})
hsi    = _get_price_change('恒生指数', hk_stocks)
hst    = _get_price_change('恒生科技指数', hk_stocks)

# 中国外汇与贵金属
cn_fm = market_perf.get('中国外汇与贵金属', {})
usdcny = cn_fm.get('USD/CNY', {}).get('price')
eurusd = cn_fm.get('EUR/USD', {}).get('price')

# 经济数据
econ = mdata.get('经济数据', {})
cn_econ = econ.get('中国', {})
us_econ = econ.get('美国', {})

log.info(f"   道指: {fp(dji[0])} | 10Y: {yield10}% | 现货黄金: {fp(spot_gold)}")

# 直接获取 slide XML 根元素
slide_elem = prs.slides[0]._element
nsmap = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}

updated = 0

for sp in slide_elem.findall('.//p:sp', nsmap):
    cNvPr = sp.find('.//p:cNvPr', nsmap)
    if cNvPr is None: continue
    shape_id = cNvPr.get('id')
    if shape_id not in ('6', '15', '23'): continue

    title_sf_xml = get_title_solidFill_xml(sp)

    # ── Shape 6: 市场表现（格式参考模板，数据全部来自 Word 市场表现方面） ──
    if shape_id == '6':
        """
        模板格式（Shape 6）：
          美国：
          道指跌0.13%报46504.67点，标普500指数涨0.11%报6582.69点，纳指涨0.18%报21879.18点。
          10年期基准国债收益率跌1.37个基点，报4.3049%，两年期美债收益率跌0.47个基点，报3.7963%。
          美元指数涨0.46%报100.01。
          3月挑战者企业裁员人数6.062万人，前值为4.8307万人。
          上周初请失业金人数20.2万人，预期21.2万人，前值22.5万人；四周均值20.8万人。
          至当日续请失业金人数184.0万人，预期186.0万人，前值185.5万人。
          中国及中国香港：
          上证指数跌0.21%报3317.12点，深证成指跌0.32%报10835.67点，创业板指跌0.58%报2229.42点。
          港股：恒生指数跌0.84%报22487.22点，恒生科技指数跌1.32%报5159.45点。
          黄金：
          现货黄金涨0.35%报4580.48美元/盎司；COMEX黄金期货涨0.42%报4592.50美元/盎司。
        """
        mp = []

        # ── 美国 ──
        mp.append(("美国：", True, FONT_SIZE_TITLE))

        # 美股：Word[7] 第一句
        w7 = _w(7)
        if w7:
            # 取第一句（到第一个句号），截取前120字
            first_sentence = w7.split('。')[0] + '。'
            mp.append((first_sentence[:150], False, FONT_SIZE_BODY))

        # 美债：Word[9] 第一句
        w9 = _w(9)
        if w9:
            first_sentence = w9.split('。')[0] + '。'
            mp.append((first_sentence[:150], False, FONT_SIZE_BODY))

        # 美元指数：Word[11] 第一句
        w11 = _w(11)
        if w11:
            first_sentence = w11.split('。')[0] + '。'
            mp.append((first_sentence[:120], False, FONT_SIZE_BODY))

        # 就业市场：Word[13]（前150字，简写）
        w13 = _w(13)
        if w13:
            # 取前150字，句号截断
            emp_text = w13[:200].split('。')[0] + '。'
            mp.append((emp_text[:180], False, FONT_SIZE_BODY))

        # ── 中国及中国香港 ──
        mp.append(("中国及中国香港：", True, FONT_SIZE_TITLE))

        # A股：Word[16] 第一句
        w16 = _w(16)
        if w16:
            first_sentence = w16.split('。')[0] + '。'
            mp.append((first_sentence[:180], False, FONT_SIZE_BODY))

        # 港股：Word[18] 第一句
        w18 = _w(18)
        if w18:
            first_sentence = w18.split('。')[0] + '。'
            mp.append((first_sentence[:150], False, FONT_SIZE_BODY))

        # ── 黄金 ──
        mp.append(("黄金：", True, FONT_SIZE_TITLE))

        # 黄金：Word[20]（提取现货和COMEX两条数据，均取第一句）
        w20 = _w(20)
        if w20:
            sentences = w20.split('。')
            gold_lines = []
            for sent in sentences[:2]:
                sent = sent.strip()
                if sent and len(sent) > 10:
                    gold_lines.append(sent + '。')
            if gold_lines:
                mp.append(("；".join(gold_lines[:2]), False, FONT_SIZE_BODY))
            else:
                mp.append((w20[:150], False, FONT_SIZE_BODY))
        else:
            mp.append(("黄金数据待更新。", False, FONT_SIZE_BODY))

        rebuild_txBody(sp, mp, title_sf_xml)
        updated += 1
        log.info(f"   ✓ 市场表现 (Shape 6) 已更新（参考模板格式，来自 Word 市场表现方面）")

    # ── Shape 15: 政策动态 ──
    # ── Shape 15: 政策动态（数据来自 Word 三、政策方面，分美国/欧洲/中国） ──
    elif shape_id == '15':
        """
        数据来源：Word 三、政策方面
          [23] 美国政策  [25] 欧洲政策  [27] 中国政策
        每条政策提炼，最多5句话，保留核心政策内容和数据。
        """
        def extract_policy_sentences(text, max_sentences=5, max_len=180):
            """
            从政策文本中提取核心句子。
            规则：
              - 按句号拆分，最多取前 max_sentences 句
              - 每句保留关键政策内容（主体/动作/数据）
              - 过滤过长括号（>50字的解释性内容）
              - 全局不超过 max_len 字符
            """
            import re
            if not text:
                return ''
            # 去除开头媒体标注
            text = re.sub(r'^[\s\u3000]*[来源出处记者编者网站财经网][：:\s]*\S*[\s\u3000]*', '', text)
            # 去除开头「此外、值得注意的是、不过、然而」等转折引子
            text = re.sub(r'^此外[,，]\s*', '', text)
            text = re.sub(r'^值得注意的是[,，]\s*', '', text)
            # 去除过长括号内容
            text = re.sub(r'（[^）]{50,}）', '', text)
            text = re.sub(r'\([^)]{50,}\)', '', text)
            # 合并多余空格
            text = re.sub(r'\s+', ' ', text).strip()
            # 按句号拆分
            parts = text.split('。')
            sentences = []
            for p in parts:
                p = p.strip()
                if not p or len(p) < 8:
                    continue
                # 过滤纯数字/引子句
                if re.match(r'^[\d\.\%亿元万美元]+$', p):
                    continue
                sentences.append(p)
            # 取前 max_sentences 句
            selected = sentences[:max_sentences]
            if not selected:
                return text[:max_len] if text else ''
            result = '。'.join(selected)
            if not result.endswith('。'):
                result += '。'
            if len(result) > max_len:
                return result[:max_len].rstrip('，。、;：') + '…'
            return result

        # 美国政策：Word[23]
        w23 = _w(23)
        pol_us = extract_policy_sentences(w23) if w23 else "美国政策动态待更新。"

        # 欧洲政策：Word[25]
        w25 = _w(25)
        pol_eu = extract_policy_sentences(w25) if w25 else "欧洲政策动态待更新。"

        # 中国政策：Word[27]
        w27 = _w(27)
        pol_cn = extract_policy_sentences(w27) if w27 else "中国政策动态待更新。"

        pp = [
            ("美国：", True, FONT_SIZE_TITLE), (pol_us, False, FONT_SIZE_BODY),
            ("欧洲：", True, FONT_SIZE_TITLE), (pol_eu, False, FONT_SIZE_BODY),
            ("中国：", True, FONT_SIZE_TITLE), (pol_cn, False, FONT_SIZE_BODY),
        ]

        rebuild_txBody(sp, pp, title_sf_xml)
        updated += 1
        log.info(f"   ✓ 政策动态 (Shape 15) 已更新（来自 Word 三、政策方面）")

    # ── Shape 23: 企业动态（来自 Word 科技方面，按地区分组） ──
    elif shape_id == '23':
        """
        数据来源：Word 科技方面段落（索引 29~31，即"四、科技方面"标题之后的内容）
        美国企业 → 段落中提及美国公司（Anthropic/OpenAI/谷歌等）
        中国及中国香港 → 段落中提及中国公司（阿里/腾讯/字节/华为等）
        每条新闻提炼，不超过3句话。
        """
        def split_sentences(text):
            """将文本拆分为有效句，过滤标题碎句和铺垫引子"""
            import re
            if not text:
                return []
            # 去除句首背景铺垫引子（这些不是有效句子）
            text = re.sub(r'^：?这两天[,，]?\s*', '', text)
            text = re.sub(r'^：?据(CNBC|彭博社|华尔街|路透社|新华社)[，,]\s*', '', text)
            text = re.sub(r'^：?值得关注的是[，,]\s*', '', text)
            text = re.sub(r'^：?公开资料显示?[，,]\s*', '', text)

            # 拆分为句子：按「。」和「；」断句；叹号/问号后若为汉字也断句
            # 先把叹号/问号后有空格接汉字的换成「。」便于统一处理
            text = re.sub(r'([!?])(\s+[\u4e00-\u9fff])', r'。\2', text)
            parts = re.split(r'(?<=[。；])(?!\s*[a-zA-Z0-9%])', text)
            sentences = []
            for p in parts:
                p = p.strip()
                # 过滤：太短（<12字）或纯数字/标题碎句
                if not p or len(p) < 12:
                    continue
                if re.match(r'^[\d万亿]+[亿美元达克]*(?:!|\?)*[^\s。；]*$', p):
                    continue
                sentences.append(p)
            return sentences[:2]

        def summarize_news(text, max_len=120):
            """
            新闻简写：保留事件主体（公司名）、关键动作（融资/收购/上市）、核心数据，
            去掉背景铺垫和冗长条款。
            """
            import re
            if not text:
                return ''
            # 去除开头媒体标注
            text = re.sub(r'^[\s　]*[来源出处记者编者网站][:：\s]*\S*[\s　]*', '', text)
            # 去除句首铺垫词
            skip_prefixes = [
                r'^(?:这两天|这天|今日|昨日|日前|近日、最新|消息人士|公开信息|值得关注)[，,：:\s]+',
                r'^据(?:CNBC|彭博社|路透社|华尔街|SEC|FDA|FTC)[，,：:\s]+',
                r'^公开资料显示?[，,]\s*',
            ]
            for pat in skip_prefixes:
                text = re.sub(pat, '', text)
            # 去除过长括号内容（>40字的解释性括号）
            text = re.sub(r'（[^）]{40,}）', '', text)
            # 去除中间可能出现的铺垫引子句
            text = re.sub(r'[,，]\s*这两天[,，]?\s*AI圈[炸爆]出[重磅巨大]消息', '', text)
            text = re.sub(r'据(?:CNBC|彭博社|路透社)[多方]*证实[,，]?', '', text)
            # 合并多余空格
            text = re.sub(r'\s+', ' ', text).strip()
            # 全局截断（优先在逗号处截断，不在数字/单位中间断开）
            if len(text) > max_len:
                cut = text[:max_len].rfind('，')
                if cut > max_len * 0.5:
                    return text[:cut].rstrip('，。、;：') + '…'
                return text[:max_len].rstrip('，。、;：') + '…'
            return text

        def format_news(text, max_sentences=2):
            """提炼新闻：最多2句，保留核心（主体+动作+数据）"""
            sents = split_sentences(text)
            if not sents:
                return ''
            result_parts = []
            total_len = 0
            for s in sents[:max_sentences]:
                brief = summarize_news(s, max_len=80)
                if not brief:
                    continue
                if total_len + len(brief) > 130:
                    break
                result_parts.append(brief)
                total_len += len(brief)
            result = '。'.join(result_parts)
            return result + ('。' if result and not result.endswith('。') else '')

        def extract_region_news(start_idx, end_idx, country_kws):
            """
            从[start_idx, end_idx)段落范围内，找出含country_kws任一关键词的段落，
            将其内容合并提炼为一个文本（最多3句）。
            """
            candidates = []
            for idx in range(start_idx, min(end_idx, len(word_paragraphs))):
                t = word_paragraphs[idx].strip()
                if t and len(t) > 10 and not t.startswith('四、') and not t.startswith('五、'):
                    if any(kw in t for kw in country_kws):
                        candidates.append(t)
            if not candidates:
                return None
            # 合并所有候选段落的内容（取第一个或拼接前300字）
            merged = ''.join(candidates)[:300]
            return format_news(merged)

        # 美国科技企业关键词
        US_KWS = ['Anthropic', 'OpenAI', '谷歌', 'Google', 'Meta', '微软', 'Microsoft',
                  '亚马逊', 'Amazon', '英伟达', 'Nvidia', '特斯拉', 'Tesla', '苹果', 'Apple',
                  '美国', '美股', '华尔街', '纳斯达克', '纽交所', 'USD', 'OpenAI CFO']
        # 中国科技企业关键词
        CN_KWS = ['阿里', '腾讯', '字节', '百度', '京东', '美团', '小米', '华为',
                  '鸿蒙', '科大讯飞', '商汤', '寒武纪', '海康', '比亚迪', '宁德',
                  '搜狐', '网易', '新浪', '知乎', '哔哩', '港股', 'A股', '中概',
                  '云深处', '望圆科技', '航锦科技', '泰坦科技', '百胜中国', '中信']

        # Word 科技方面内容段落索引：29~31（"四、科技方面"标题在28）
        us_news = extract_region_news(29, 32, US_KWS)
        cn_news = extract_region_news(29, 32, CN_KWS)

        if not us_news:
            us_news = "暂无美国科技企业重大动态。"
        if not cn_news:
            cn_news = "暂无中国科技企业重大动态。"

        ep = [
            ("美国", True, FONT_SIZE_TITLE), (us_news, False, FONT_SIZE_BODY),
            ("中国及中国香港", True, FONT_SIZE_TITLE), (cn_news, False, FONT_SIZE_BODY),
        ]

        rebuild_txBody(sp, ep, title_sf_xml)
        updated += 1
        log.info(f"   ✓ 企业动态 (Shape 23) 已更新（来自 Word 科技方面）")

# 保存
buf = BytesIO()
prs.save(buf)
buf.seek(0)
out = os.path.join(config.OUTPUT_DIR, f"金融市场简报_{config.TODAY.strftime('%Y%m%d')}.pptx")
with open(out, 'wb') as f:
    f.write(buf.read())
log.info(f"\n✅ PPT 简报已保存: {out}")
log.info(f"   更新了 {updated} 个文本框")
