"""PPTX → DOCX (提取文本和表格)

用法:
    python pptx_to_docx.py input.pptx [output.docx]

参数:
    input.pptx  - PPTX 文件路径
    output.docx - 输出文件路径 (默认: output.docx)
"""
from pptx import Presentation
from docx import Document
from docx.shared import Pt
import sys
import os


def main():
    if len(sys.argv) < 2:
        print("用法: python pptx_to_docx.py input.pptx [output.docx]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    out_path = sys.argv[2] if len(sys.argv) > 2 else "output.docx"

    if not os.path.isfile(pptx_path):
        print(f"错误: 文件不存在 - {pptx_path}")
        sys.exit(1)

    prs = Presentation(pptx_path)
    doc = Document()

    for i, slide in enumerate(prs.slides):
        doc.add_heading(f"Slide {i + 1}", level=1)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        p = doc.add_paragraph(text)
                        for run in p.runs:
                            run.font.size = Pt(11)
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    rows.append([cell.text for cell in row.cells])
                if rows:
                    t = doc.add_table(rows=len(rows), cols=len(rows[0]))
                    for r_idx, row_data in enumerate(rows):
                        for c_idx, cell_text in enumerate(row_data):
                            t.cell(r_idx, c_idx).text = cell_text
        doc.add_page_break()

    doc.save(out_path)
    print(f"完成: {len(prs.slides)} 张幻灯片 -> {out_path}")


if __name__ == "__main__":
    main()