"""
PlantUML to SVG/EMF/PPT converter using Inkscape CLI.

Usage:
    python svg2emf.py input.puml [-f ppt|emf|svg|png] [-o output]

Requirements:
    - Inkscape installed
    - Java for PlantUML
    - python-pptx (for PPT output)
"""

import os
import sys
import struct
import subprocess
import shutil

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
INKSCAPE_PATH = os.path.join(SCRIPT_DIR, "inkscape", "bin", "inkscape.exe")
PLANTUML_JAR = os.path.join(SCRIPT_DIR, "plantuml.jar")

EMBEDDED_JAVA_PATH = os.path.join(SCRIPT_DIR, "jdk", "jdk8u482-b08", "bin", "java.exe")

def find_java():
    """Find Java: JAVA_HOME > PATH > embedded JDK"""
    java_home = os.environ.get("JAVA_HOME")
    if java_home:
        java_exe = os.path.join(java_home, "bin", "java.exe")
        if os.path.exists(java_exe):
            return java_exe
    # Try PATH
    from shutil import which
    java_in_path = which("java")
    if java_in_path:
        return java_in_path
    # Fall back to embedded JDK
    if os.path.exists(EMBEDDED_JAVA_PATH):
        return EMBEDDED_JAVA_PATH
    return None

def find_inkscape():
    if os.path.exists(INKSCAPE_PATH):
        return INKSCAPE_PATH
    env_path = os.environ.get("INKSCAPE_PATH")
    if env_path and os.path.exists(env_path):
        return env_path
    common_paths = [
        r"C:\Program Files\Inkscape\bin\inkscape.exe",
        r"C:\Program Files (x86)\Inkscape\bin\inkscape.exe",
    ]
    for path in common_paths:
        if os.path.exists(path):
            return path
    if shutil.which("inkscape"):
        return "inkscape"
    return None


def puml_to_svg(puml_file):
    """PlantUML -> SVG"""
    print("[1/2] Generating SVG via PlantUML...")
    java_path = find_java()
    if not java_path:
        raise RuntimeError(
            "Java not found. Please install JDK 8 and set JAVA_HOME environment variable.\n"
            "Download: https://adoptium.net/zh-CN/temurin/releases?version=8"
        )
    result = subprocess.run(
        [java_path, "-jar", PLANTUML_JAR, "-charset", "UTF-8", "-tsvg", puml_file, "-o", os.path.dirname(puml_file) or "."],
        capture_output=True, text=True
    )
    if result.returncode != 0:
        raise RuntimeError(f"PlantUML failed: {result.stderr}")

    base = os.path.splitext(puml_file)[0]
    svg_out = os.path.join(os.path.dirname(puml_file) or ".", f"{os.path.basename(base)}.svg")
    if not os.path.exists(svg_out):
        raise RuntimeError(f"SVG not created: {svg_out}")

    # Fix numeric character references to UTF-8 characters for proper font rendering
    import re
    with open(svg_out, 'r', encoding='utf-8') as f:
        content = f.read()
    # Replace numeric character references (e.g., &#27979;) with actual UTF-8 chars
    fixed = re.sub(r'&#(\d+);', lambda m: chr(int(m.group(1))), content)
    fixed = fixed.replace('encoding="us-ascii"', 'encoding="UTF-8"')
    with open(svg_out, 'w', encoding='utf-8') as f:
        f.write(fixed)

    return svg_out


def svg_to_emf(svg_file):
    """SVG -> EMF with ungroup via Inkscape"""
    print("[2/2] Converting to EMF via Inkscape (with ungroup)...")
    inkscape = find_inkscape()
    if not inkscape:
        raise RuntimeError("Inkscape not found")

    base = os.path.splitext(svg_file)[0]
    emf_out = f"{base}.emf"

    # Use text-to-path to convert text to vector paths, avoiding font dependency issues
    result = subprocess.run(
        [inkscape, svg_file,
         "--actions", "select-all;selection-ungroup;selection-ungroup;selection-ungroup",
         "--export-text-to-path",
         "--export-filename", emf_out,
         "--export-type", "emf"],
        capture_output=True, text=True
    )
    if result.returncode != 0:
        raise RuntimeError(f"Inkscape failed: {result.stderr}")

    if not os.path.exists(emf_out):
        raise RuntimeError(f"EMF not created: {emf_out}")
    return emf_out


def svg_to_png(svg_file):
    """SVG -> PNG via Inkscape"""
    print("[2/2] Converting to PNG via Inkscape...")
    inkscape = find_inkscape()
    if not inkscape:
        raise RuntimeError("Inkscape not found")

    base = os.path.splitext(svg_file)[0]
    png_out = f"{base}.png"

    result = subprocess.run(
        [inkscape, svg_file,
         "--export-filename", png_out,
         "--export-type", "png"],
        capture_output=True, text=True
    )
    if result.returncode != 0:
        raise RuntimeError(f"Inkscape failed: {result.stderr}")

    if not os.path.exists(png_out):
        raise RuntimeError(f"PNG not created: {png_out}")
    return png_out


def get_emf_bounds_himetric(emf_file):
    """Parse EMF header to get bounds in HIMETRIC units."""
    try:
        with open(emf_file, 'rb') as f:
            data = f.read(88)
            if len(data) < 88:
                return None
            left, top, right, bottom = struct.unpack('<IIII', data[24:40])
            return right - left, bottom - top
    except:
        return None


def create_ppt_with_emf(emf_file, pptx_file):
    """Create a PPT and insert the EMF file with proper sizing."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
    except ImportError:
        raise RuntimeError("python-pptx not installed")

    # Disable PIL decompression bomb check (EMF stores bounds in HIMETRIC, not pixels)
    Image.MAX_IMAGE_PIXELS = None

    print("[3/3] Creating PPT...")

    # Parse EMF bounds to calculate correct size
    # HIMETRIC units: 1 unit = 0.01 mm
    # 1 inch = 25.4 mm = 2540 HIMETRIC units
    himetric_per_inch = 2540

    bounds = get_emf_bounds_himetric(emf_file)
    if bounds:
        width_him, height_him = bounds
        width_inch = width_him / himetric_per_inch
        height_inch = height_him / himetric_per_inch
        print(f"  EMF size: {width_inch:.2f} x {height_inch:.2f} inches")
    else:
        width_inch = 8
        height_inch = 6

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Calculate margins
    margin = Inches(0.5)
    available_width = prs.slide_width - margin * 2
    available_height = prs.slide_height - margin * 2

    # Scale to fit within slide while maintaining aspect ratio
    if bounds:
        aspect_ratio = width_inch / height_inch
        available_aspect = available_width / available_height

        if aspect_ratio > available_aspect:
            emf_width = available_width
            emf_height = int(available_width / aspect_ratio)
        else:
            emf_height = available_height
            emf_width = int(available_height * aspect_ratio)
    else:
        emf_width = available_width
        emf_height = Inches(6)

    left = (prs.slide_width - emf_width) // 2
    top = (prs.slide_height - emf_height) // 2

    slide.shapes.add_picture(emf_file, left, top, width=emf_width)

    prs.save(pptx_file)
    print(f"[OK] PPT saved: {pptx_file}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python svg2emf.py input.puml [-f ppt|emf|svg|png] [-o output]")
        print("  -f: output format, default: ppt")
        print("  -o: output file path (optional)")
        sys.exit(1)

    puml_file = sys.argv[1]
    if not puml_file.lower().endswith(".puml"):
        print("Error: only .puml files supported")
        sys.exit(1)

    # Parse arguments
    fmt = "ppt"
    output_path = None
    i = 2
    while i < len(sys.argv):
        if sys.argv[i] == "-f" and i + 1 < len(sys.argv):
            fmt = sys.argv[i + 1].lower()
            i += 2
        elif sys.argv[i] == "-o" and i + 1 < len(sys.argv):
            output_path = sys.argv[i + 1]
            i += 2
        else:
            i += 1

    if fmt not in ("ppt", "emf", "svg", "png"):
        print("Error: format must be ppt, emf, svg, or png")
        sys.exit(1)

    base = os.path.splitext(puml_file)[0]

    if fmt == "svg":
        svg_file = puml_to_svg(puml_file)
        print(f"[OK] SVG generated: {svg_file}")
    elif fmt == "png":
        svg_file = puml_to_svg(puml_file)
        png_file = svg_to_png(svg_file)
        print(f"[OK] PNG generated: {png_file}")
    elif fmt == "emf":
        svg_file = puml_to_svg(puml_file)
        emf_file = svg_to_emf(svg_file)
        print(f"[OK] EMF generated: {emf_file}")
    else:  # ppt
        svg_file = puml_to_svg(puml_file)
        emf_file = svg_to_emf(svg_file)
        pptx_file = output_path or f"{base}.pptx"
        create_ppt_with_emf(emf_file, pptx_file)
