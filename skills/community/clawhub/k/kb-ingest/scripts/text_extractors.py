from __future__ import annotations

import tarfile
import zipfile
from pathlib import Path


TEXT_SUFFIXES = {
    ".txt", ".md", ".markdown", ".csv", ".json", ".yaml", ".yml", ".xml", ".toml",
    ".py", ".java", ".js", ".ts", ".vue", ".jsx", ".tsx", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".go", ".rs", ".sql", ".sh", ".ps1", ".bat",
}


def extract_text(path: str, max_chars: int = 60000) -> str:
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return p.read_text(encoding="utf-8", errors="ignore")[:max_chars]
    if suffix == ".pdf":
        return extract_pdf(p, max_chars)
    if suffix == ".docx":
        return extract_docx(p, max_chars)
    if suffix == ".pptx":
        return extract_pptx(p, max_chars)
    if suffix in {".xlsx", ".xls"}:
        return extract_spreadsheet(p, max_chars)
    if suffix in {".zip", ".rar", ".7z", ".tar", ".gz"}:
        return inspect_archive(p, max_chars)
    return f"Unsupported binary file: {p.name}"


def extract_pdf(path: Path, max_chars: int) -> str:
    try:
        import fitz
        parts = []
        with fitz.open(str(path)) as doc:
            for page in doc[:80]:
                parts.append(page.get_text())
        return "\n".join(parts)[:max_chars]
    except Exception as exc:
        raise RuntimeError(f"pdf extraction failed: {exc}") from exc


def extract_docx(path: Path, max_chars: int) -> str:
    try:
        import docx
        doc = docx.Document(str(path))
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)[:max_chars]
    except Exception as exc:
        raise RuntimeError(f"docx extraction failed: {exc}") from exc


def extract_pptx(path: Path, max_chars: int) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        lines = []
        for index, slide in enumerate(prs.slides, start=1):
            lines.append(f"# Slide {index}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
        return "\n".join(lines)[:max_chars]
    except Exception as exc:
        raise RuntimeError(f"pptx extraction failed: {exc}") from exc


def extract_spreadsheet(path: Path, max_chars: int) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets[:10]:
            lines.append(f"# Sheet: {sheet.title}")
            for row in sheet.iter_rows(max_row=200, values_only=True):
                lines.append("\t".join("" if cell is None else str(cell) for cell in row))
        return "\n".join(lines)[:max_chars]
    except Exception as exc:
        raise RuntimeError(f"spreadsheet extraction failed: {exc}") from exc


def inspect_archive(path: Path, max_chars: int) -> str:
    try:
        if path.suffix.lower() == ".zip":
            with zipfile.ZipFile(path) as archive:
                names = archive.namelist()[:300]
        elif path.suffix.lower() in {".tar", ".gz"}:
            with tarfile.open(path) as archive:
                names = [member.name for member in archive.getmembers()[:300]]
        else:
            return f"Archive inspection is not available for {path.name}. The file was archived for Agent processing."
        lines = [f"Archive: {path.name}", "", "Contained paths:"]
        lines.extend(f"- {name}" for name in names)
        return "\n".join(lines)[:max_chars]
    except Exception as exc:
        raise RuntimeError(f"archive inspection failed: {exc}") from exc
