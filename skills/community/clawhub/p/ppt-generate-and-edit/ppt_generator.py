"""PPT creation logic for the reusable skill module."""

from __future__ import annotations

from pathlib import Path
from typing import Callable, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from .image_service import generate_image
from .utils import ensure_directory, get_logger, timestamp_str

logger = get_logger(__name__)


def _normalize_content(content: str | List[str]) -> str:
    """Normalize content input into a single text block."""
    if isinstance(content, list):
        return "\n".join(part.strip() for part in content if part and part.strip())
    return content.strip()


def _add_title_and_body(slide, slide_title: str, content: str) -> None:
    """Add title and content text boxes with basic layout."""
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = slide_title
    title_frame.paragraphs[0].font.size = Pt(30)
    title_frame.paragraphs[0].font.bold = True

    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(5.2), Inches(4.8))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.clear()

    for idx, line in enumerate(content.split("\n")):
        paragraph = body_frame.paragraphs[0] if idx == 0 else body_frame.add_paragraph()
        paragraph.text = line.strip()
        paragraph.level = 0
        paragraph.font.size = Pt(18)


def create_ppt(
    title: str,
    slides: List[Dict[str, object]],
    output_dir: str | Path = "ppt_skill/output",
    image_generator: Optional[Callable[[str], str]] = None,
) -> str:
    """Create a PPT based on slide descriptors and return output path.

    Expected slide schema:
    {
      "title": str,
      "content": str | list[str],
      "image_prompt": Optional[str]
    }
    """
    logger.info("Start create_ppt title=%s slide_count=%d", title, len(slides))
    prs = Presentation()
    image_generator = image_generator or generate_image

    # Title page
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(8.0), Inches(1.6))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True

    for index, slide_data in enumerate(slides, start=1):
        logger.info("Rendering slide index=%d data=%s", index, slide_data.get("title", ""))
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_title = str(slide_data.get("title", f"Slide {index}"))
        content = _normalize_content(slide_data.get("content", ""))
        _add_title_and_body(slide, slide_title, content)

        image_prompt = slide_data.get("image_prompt")
        if image_prompt:
            image_path = image_generator(str(image_prompt))
            logger.info("Insert generated image for slide=%d path=%s", index, image_path)
            slide.shapes.add_picture(
                image_path,
                Inches(6.1),
                Inches(1.5),
                width=Inches(3.1),
                height=Inches(3.6),
            )

    output_directory = ensure_directory(output_dir)
    output_path = output_directory / f"{title.replace(' ', '_')}_{timestamp_str()}.pptx"
    prs.save(str(output_path))
    logger.info("create_ppt completed output=%s", output_path)
    return str(output_path)
