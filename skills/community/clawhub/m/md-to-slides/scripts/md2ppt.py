#!/usr/bin/env python3
"""Markdown to PowerPoint converter"""
import sys, os, re

# Check if python-pptx is available
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

def md_to_slides(md):
    """Split markdown into slides by headings"""
    lines = md.split('\n')
    slides = []
    current = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith('# ') and current:
            slides.append('\n'.join(current))
            current = [stripped]
        elif stripped.startswith('## ') and current:
            slides.append('\n'.join(current))
            current = [stripped]
        else:
            current.append(stripped)
    if current:
        slides.append('\n'.join(current))
    return slides

def parse_content(slide_md):
    """Extract title and bullet points from slide markdown"""
    lines = slide_md.split('\n')
    title = ""
    bullets = []
    in_code = False
    code_content = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith('# '):
            title = stripped[2:]
        elif stripped.startswith('## '):
            if not title:
                title = stripped[3:]
        elif stripped.startswith('- ') or stripped.startswith('* '):
            bullets.append(stripped[2:])
        elif stripped.startswith('```'):
            in_code = not in_code
        elif in_code:
            code_content.append(stripped)
        elif stripped and not title:
            if stripped not in ['', ' ']:
                bullets.append(stripped)
    return title, bullets, '\n'.join(code_content) if code_content else None

def create_pptx(slides, output='output.pptx', theme='professional'):
    if not HAS_PPTX:
        # Fallback: create HTML presentation
        html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>Presentation</title>
<style>
body {{ font-family: Arial, sans-serif; margin: 0; }}
.slide {{ width: 100vw; height: 100vh; display: flex; flex-direction: column; justify-content: center; padding: 60px; box-sizing: border-box; page-break-after: always; }}
h1 {{ font-size: 48px; margin-bottom: 40px; color: #1a1a2e; }}
h2 {{ font-size: 36px; margin-bottom: 30px; color: #16213e; }}
li {{ font-size: 28px; margin: 15px 0; color: #333; }}
code {{ background: #f4f4f4; padding: 3px 8px; border-radius: 4px; font-family: monospace; }}
</style></head><body>
"""
        for slide in slides:
            title, bullets, code = parse_content(slide)
            if not title:
                title = "Presentation"
            html += f'<div class="slide"><h1>{title}</h1>\n'
            for b in bullets:
                html += f'<li>{b}</li>\n'
            if code:
                html += f'<pre><code>{code}</code></pre>\n'
            html += '</div>\n'
        html += '</body></html>'
        with open(output.replace('.pptx', '.html'), 'w') as f:
            f.write(html)
        return output.replace('.pptx', '.html')
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    colors = {
        'professional': (26, 26, 46),
        'creative': (41, 128, 185),
        'minimal': (50, 50, 50),
    }
    bg_color = colors.get(theme, colors['professional'])
    
    for slide_md in slides:
        title, bullets, code = parse_content(slide_md)
        if not title:
            title = "Slide"
        
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(*bg_color)
        background.line.fill.background()
        
        txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = txTitle.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        if bullets:
            txBody = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11.5), Inches(5))
            tf = txBody.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets[:8]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(230, 230, 230)
                p.space_before = Pt(12)
        
        if code:
            txCode = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(11.5), Inches(1.5))
            tf = txCode.text_frame
            p = tf.paragraphs[0]
            p.text = code[:200]
            p.font.size = Pt(14)
            p.font.name = "Courier New"
            p.font.color.rgb = RGBColor(150, 255, 150)
    
    prs.save(output)
    return output

def main():
    args = sys.argv[1:]
    md = ""
    output = "output.pptx"
    theme = "professional"
    
    i = 0
    while i < len(args):
        if args[i] == "--file" and i + 1 < len(args):
            with open(args[i+1]) as f:
                md = f.read()
            i += 2
        elif args[i] == "--theme" and i + 1 < len(args):
            theme = args[i+1]
            i += 2
        elif args[i] == "--output" and i + 1 < len(args):
            output = args[i+1]
            i += 2
        else:
            md += args[i] + " "
            i += 1
    
    if not md.strip():
        print("Usage: md2ppt.py [--file <file.md>] [--theme professional|creative|minimal] [--output file.pptx] <markdown>", file=sys.stderr)
        sys.exit(1)
    
    slides = md_to_slides(md)
    result = create_pptx(slides, output, theme)
    print(f"Created: {result}")
    if not HAS_PPTX:
        print("(python-pptx not installed, created HTML instead)")

if __name__ == "__main__":
    main()
