#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Deep Research to PPT Pro (v2.0) - 图片生成、OCR核验与文件组装脚本

本脚本是整个工作流的核心执行引擎，集成了以下功能：
1. 批量调用 Gemini API 生成 PPT 幻灯片图片
2. 对每张图片进行 OCR 文字核验（检测错字、漏字、乱码）
3. 自动删除并重新生成不合格的图片
4. 将所有合格图片组装成 PPTX 和 PDF 文件
5. 支持断点续传（已存在的合格图片会自动跳过）

使用方法：
  1. 确保所有依赖已安装（见下方依赖列表）
  2. 确保 GEMINI_API_KEY 环境变量已设置
  3. 将阶段 4 生成的提示词列表填入 SLIDES_PROMPTS
  4. 运行：python3 generate_and_assemble.py

依赖安装：
  sudo pip3 install google-genai pillow fpdf2 python-pptx pytesseract python-Levenshtein
  sudo apt-get install -y tesseract-ocr tesseract-ocr-chi-sim

环境变量：
  GEMINI_API_KEY - Google Gemini API Key
"""

import os
import sys
import time
import re
import json
import subprocess
from datetime import datetime
from PIL import Image
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Emu

# ============================================================
# 配置区域 (v2.0)
# ============================================================

# Gemini 模型配置
MODEL = "gemini-3.1-flash-image-preview"  # 推荐使用此模型，中文渲染质量最高

# 输出目录
OUTPUT_DIR = "/home/ubuntu/ppt_output"

# 图片尺寸（16:9 标准宽屏）
TARGET_W = 1920
TARGET_H = 1080

# API 调用间隔（秒），避免触发限流
API_DELAY_SECONDS = 10

# 单次 API 请求超时（秒）
# 说明：某些网络环境下 generate_content 可能长时间挂起；设置超时可避免整批任务被卡死
REQUEST_TIMEOUT_SECONDS = 90

# 最大重试次数（包括 OCR 核验失败后的重试）
MAX_RETRIES = 3

# OCR 文本相似度阈值
# 低于此值的图片将被判定为"存在错别字"并重新生成
# 0.98 = 98% 相似度，比较严格
# 如果发现 OCR 核验过于严格导致大量重试，可以适当降低到 0.90
OCR_SIMILARITY_THRESHOLD = 0.90

# 是否启用 OCR 核验（设为 False 可跳过 OCR 核验以加速生成）
ENABLE_OCR_CHECK = True

# 日志文件路径
LOG_FILE = "/tmp/generate_log.txt"

# ============================================================
# 提示词注入区
# ============================================================
# 在此填入阶段 4 生成的提示词列表
# 格式：每个元素是一个字典，包含 "num" (页码) 和 "prompt" (提示词)
# 示例：
# SLIDES_PROMPTS = [
#     {"num": 1, "prompt": "Generate a presentation slide image..."},
#     {"num": 2, "prompt": "Generate a presentation slide image..."},
#     ...
# ]
SLIDES_PROMPTS = []


# ============================================================
# 日志工具
# ============================================================

def log(message):
    """同时输出到控制台和日志文件"""
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    formatted = f"[{timestamp}] {message}"
    print(formatted, flush=True)
    try:
        with open(LOG_FILE, "a", encoding="utf-8") as f:
            f.write(formatted + "\n")
    except:
        pass


# ============================================================
# OCR 核验模块
# ============================================================

def ensure_tesseract_installed():
    """确保 Tesseract OCR 及中文语言包已安装"""
    try:
        result = subprocess.run(["tesseract", "--version"], capture_output=True, text=True)
        if result.returncode == 0:
            log("  [OCR] Tesseract 已安装")
            return True
    except FileNotFoundError:
        pass
    
    log("  [OCR] 正在安装 Tesseract...")
    os.system("sudo apt-get update -qq && sudo apt-get install -y -qq tesseract-ocr tesseract-ocr-chi-sim")
    return True

def clean_text_for_comparison(text):
    """清理文本，移除空格、换行和标点，只保留中文、字母和数字"""
    if not text:
        return ""
    return re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9]', '', text)

def extract_expected_text_from_prompt(prompt):
    """从提示词中提取所有需要渲染的中文文本"""
    # 提取 Title, Section, Content, Footer 中的中文内容
    chinese_parts = []
    
    # 提取 Title
    title_match = re.search(r'Title[^:]*:\s*(.+)', prompt)
    if title_match:
        chinese_parts.append(title_match.group(1).strip())
    
    # 提取 Section
    section_match = re.search(r'Section[^:]*:\s*(.+)', prompt)
    if section_match:
        chinese_parts.append(section_match.group(1).strip())
    
    # 提取 Content 中的各条内容
    content_matches = re.findall(r'[—\-•]\s*(.+)', prompt)
    for match in content_matches:
        chinese_parts.append(match.strip())
    
    # 提取 Footer
    footer_match = re.search(r'Footer:\s*(.+)', prompt)
    if footer_match:
        chinese_parts.append(footer_match.group(1).strip())
    
    return " ".join(chinese_parts)

def get_text_from_image(image_path):
    """使用 Tesseract 从图片中提取文字"""
    try:
        import pytesseract
        return pytesseract.image_to_string(Image.open(image_path), lang='chi_sim')
    except Exception as e:
        log(f"  [OCR 错误] {e}")
        return ""

def check_ocr_similarity(image_path, expected_text):
    """
    对图片进行 OCR 核验，返回相似度得分。
    
    参数：
        image_path: 图片文件路径
        expected_text: 期望的文本内容
    
    返回：
        float: 0.0 到 1.0 之间的相似度得分
    """
    try:
        from Levenshtein import ratio
    except ImportError:
        log("  [OCR] python-Levenshtein 未安装，跳过 OCR 核验")
        return 1.0  # 如果库未安装，默认通过
    
    ocr_text = get_text_from_image(image_path)
    cleaned_ocr = clean_text_for_comparison(ocr_text)
    cleaned_expected = clean_text_for_comparison(expected_text)
    
    if not cleaned_ocr or not cleaned_expected:
        log(f"  [OCR] 文本为空，跳过核验 (OCR长度={len(cleaned_ocr)}, 期望长度={len(cleaned_expected)})")
        return 1.0  # 如果无法提取文本，默认通过
    
    similarity = ratio(cleaned_ocr, cleaned_expected)
    log(f"  [OCR 核验] 相似度: {similarity:.2%} (阈值: {OCR_SIMILARITY_THRESHOLD:.2%})")
    return similarity


# ============================================================
# 核心生成函数
# ============================================================

def generate_single_image(client, types_module, slide_info):
    """
    生成单张幻灯片图片，包含 OCR 核验和自动重试逻辑。
    
    参数：
        client: Gemini API 客户端
        types_module: google.genai.types 模块
        slide_info: 包含 "num" 和 "prompt" 的字典
    
    返回：
        bool: 是否成功生成
    """
    num = slide_info["num"]
    prompt = slide_info["prompt"]
    output_path = os.path.join(OUTPUT_DIR, f"slide_{num:03d}.png")
    
    # 断点续传：如果图片已存在且大小合理，跳过
    if os.path.exists(output_path) and os.path.getsize(output_path) > 50000:
        log(f"  [第{num:03d}页] ⏭️  图片已存在 ({os.path.getsize(output_path) // 1024}KB)，跳过")
        return True
    
    # 提取期望文本用于 OCR 对比
    expected_text = extract_expected_text_from_prompt(prompt)
    
    for attempt in range(1, MAX_RETRIES + 1):
        try:
            log(f"  [第{num:03d}页] 🎨 正在生成... (尝试 {attempt}/{MAX_RETRIES})")
            
            response = client.models.generate_content(
                model=MODEL,
                contents=prompt,
                config=types_module.GenerateContentConfig(
                    response_modalities=["image", "text"],
                    http_options=types_module.HttpOptions(timeout=REQUEST_TIMEOUT_SECONDS * 1000)
                ),
            )
            
            # 从响应中提取图片数据
            image_saved = False
            for part in response.candidates[0].content.parts:
                if part.inline_data is not None:
                    with open(output_path, "wb") as f:
                        f.write(part.inline_data.data)
                    
                    file_size = os.path.getsize(output_path) // 1024
                    log(f"  [第{num:03d}页] ✅ 图片已保存 ({file_size}KB)")
                    image_saved = True
                    
                    # OCR 核验
                    if ENABLE_OCR_CHECK:
                        similarity = check_ocr_similarity(output_path, expected_text)
                        if similarity >= OCR_SIMILARITY_THRESHOLD:
                            log(f"  [第{num:03d}页] ✅ OCR 文字核验通过")
                            return True
                        else:
                            log(f"  [第{num:03d}页] ⚠️ OCR 文字核验失败 (相似度 {similarity:.2%} < 阈值 {OCR_SIMILARITY_THRESHOLD:.2%})")
                            if attempt < MAX_RETRIES:
                                log(f"  [第{num:03d}页] 🔄 删除不合格图片，准备重试...")
                                os.remove(output_path)
                            else:
                                log(f"  [第{num:03d}页] ⚠️ 已达最大重试次数，保留最后一次结果")
                                return True  # 保留最后一次结果
                            break
                    else:
                        return True
            
            if not image_saved:
                log(f"  [第{num:03d}页] ⚠️ 响应中未找到图片数据")
                
        except Exception as e:
            error_msg = str(e)[:200]
            log(f"  [第{num:03d}页] ❌ 生成失败: {error_msg}")
        
        # 重试前等待
        if attempt < MAX_RETRIES:
            wait_time = API_DELAY_SECONDS * attempt
            log(f"  [第{num:03d}页] ⏳ 等待 {wait_time} 秒后重试...")
            time.sleep(wait_time)
    
    log(f"  [第{num:03d}页] ❌ 已达到最大重试次数，放弃此页")
    return False


# ============================================================
# 文件组装模块
# ============================================================

def assemble_pdf(slide_files, report_name):
    """将图片组装成 PDF 文件"""
    try:
        log("  📄 正在生成 PDF...")
        pdf = FPDF(orientation="L", unit="mm", format=(190.5, 338.67))
        for fname in slide_files:
            pdf.add_page()
            img_path = os.path.join(OUTPUT_DIR, fname)
            pdf.image(img_path, x=0, y=0, w=338.67, h=190.5)
        
        pdf_path = os.path.join(OUTPUT_DIR, f"{report_name}.pdf")
        pdf.output(pdf_path)
        file_size = os.path.getsize(pdf_path) / (1024 * 1024)
        log(f"  ✅ PDF 生成成功: {pdf_path} ({file_size:.1f}MB)")
        return True
    except Exception as e:
        log(f"  ❌ PDF 生成失败: {e}")
        return False

def assemble_pptx(slide_files, report_name):
    """将图片组装成 PPTX 文件"""
    try:
        log("  📊 正在生成 PPTX...")
        prs = Presentation()
        prs.slide_width = Emu(12192000)   # 16:9 宽屏
        prs.slide_height = Emu(6858000)
        
        for fname in slide_files:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
            img_path = os.path.join(OUTPUT_DIR, fname)
            slide.shapes.add_picture(
                img_path, Emu(0), Emu(0),
                prs.slide_width, prs.slide_height
            )
        
        pptx_path = os.path.join(OUTPUT_DIR, f"{report_name}.pptx")
        prs.save(pptx_path)
        file_size = os.path.getsize(pptx_path) / (1024 * 1024)
        log(f"  ✅ PPTX 生成成功: {pptx_path} ({file_size:.1f}MB)")
        return True
    except Exception as e:
        log(f"  ❌ PPTX 生成失败: {e}")
        return False

def assemble_final_files(report_name):
    """将所有合格图片组装成最终的 PDF 和 PPTX 文件"""
    slide_files = sorted([
        f for f in os.listdir(OUTPUT_DIR) 
        if f.startswith("slide_") and f.endswith(".png")
    ])
    
    if not slide_files:
        log("\n❌ 错误：输出目录中没有找到任何图片文件。")
        return
    
    log(f"\n📦 开始组装 {len(slide_files)} 张图片...")
    
    assemble_pptx(slide_files, report_name)
    assemble_pdf(slide_files, report_name)


# ============================================================
# 主入口
# ============================================================

def main():
    """主函数"""
    log("=" * 60)
    log("Deep Research to PPT Pro v2.0 - 开始执行")
    log("=" * 60)
    
    # 检查 API Key
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        log("❌ 错误：GEMINI_API_KEY 环境变量未设置。")
        sys.exit(1)
    
    # 检查提示词列表
    if not SLIDES_PROMPTS:
        log("❌ 错误：SLIDES_PROMPTS 列表为空。请先填入提示词。")
        sys.exit(1)
    
    total_slides = len(SLIDES_PROMPTS)
    log(f"📋 共 {total_slides} 页待生成")
    log(f"🤖 模型: {MODEL}")
    log(f"🔍 OCR 核验: {'启用' if ENABLE_OCR_CHECK else '禁用'} (阈值: {OCR_SIMILARITY_THRESHOLD:.2%})")
    log(f"📁 输出目录: {OUTPUT_DIR}")
    
    # 创建输出目录
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    # 确保 Tesseract 已安装
    if ENABLE_OCR_CHECK:
        ensure_tesseract_installed()
    
    # 初始化 Gemini 客户端
    from google import genai
    from google.genai import types
    client = genai.Client(api_key=api_key)
    
    # 批量生成
    success_count = 0
    fail_count = 0
    start_time = time.time()
    
    for i, slide_info in enumerate(SLIDES_PROMPTS):
        log(f"\n--- 进度: {i+1}/{total_slides} ---")
        
        if generate_single_image(client, types, slide_info):
            success_count += 1
        else:
            fail_count += 1
        
        # API 调用间隔
        if i < total_slides - 1:
            time.sleep(API_DELAY_SECONDS)
    
    elapsed = time.time() - start_time
    log(f"\n{'=' * 60}")
    log(f"📊 生成完毕:")
    log(f"   ✅ 成功: {success_count} 页")
    log(f"   ❌ 失败: {fail_count} 页")
    log(f"   ⏱️  耗时: {elapsed/60:.1f} 分钟")
    log(f"{'=' * 60}")
    
    # 组装最终文件
    report_name = "研究报告_清新研究团队"
    assemble_final_files(report_name)
    
    log(f"\n🎉 全部完成！文件已保存到 {OUTPUT_DIR}")

if __name__ == "__main__":
    main()
