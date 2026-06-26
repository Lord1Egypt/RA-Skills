#!/usr/bin/env python3
"""Apply EN edits to a .pptx + write/append the comparison Excel.

Usage:
    python3 apply.py <en_path> <cn_path> <out_pptx> <out_xlsx> <edits.json>

edits.json format:
[
  {
    "slide": 27,
    "loc": "/3/p3",
    "cn": "...",
    "en_optimized": "...",
    "notes": "...",
    "highlight": true   # optional; renders row in red bold
  }
]

For each edit:
  - Walk to the (slide, loc) paragraph in <en_path>; replace its runs' text
    with en_optimized (preserving font/format of first run; subsequent runs
    text-cleared).
  - Append a row to the Excel: page · kind · cn · en_original · en_optimized · notes
    (en_original = pre-edit text)
"""
import sys
import os
import json
import shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openpyxl import load_workbook, Workbook
from openpyxl.styles import Font, Alignment, PatternFill

HEADERS = ["page", "kind", "cn", "en_original", "en_optimized", "notes"]


def lock_path(p):
    d, n = os.path.split(p)
    return os.path.join(d, "~$" + n)


def check_lock(p):
    if os.path.exists(lock_path(p)):
        print(f"⚠ STOP. '{p}' appears open (lock file present). Save+close.")
        sys.exit(2)


def find_paragraph(prs, slide_idx, loc):
    """loc format: /shape_idx/.../[r{r}c{c}/]p{para_idx}"""
    slide = prs.slides[slide_idx - 1]
    parts = [p for p in loc.split("/") if p]
    para_part = parts[-1]
    if not para_part.startswith("p"):
        raise ValueError(f"Loc must end in pN: {loc}")
    pi = int(para_part[1:])
    cur = slide.shapes
    table_coords = None
    for token in parts[:-1]:
        if token.startswith("r") and "c" in token:
            r, c = token[1:].split("c")
            table_coords = (int(r), int(c))
        else:
            idx = int(token)
            shp = list(cur)[idx]
            if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                cur = shp.shapes
                continue
            if table_coords and shp.has_table:
                cell = shp.table.cell(*table_coords)
                return cell.text_frame.paragraphs[pi]
            return shp.text_frame.paragraphs[pi]
    raise ValueError(f"Could not resolve loc {loc}")


def replace_paragraph_text(para, new_text):
    if not para.runs:
        # add a new run
        para.add_run().text = new_text
        return
    para.runs[0].text = new_text
    for r in para.runs[1:]:
        r.text = ""


def open_or_init_excel(path):
    if os.path.exists(path):
        wb = load_workbook(path)
        ws = wb.active
        # Verify header
        hdr = [ws.cell(2, c).value for c in range(2, 2 + len(HEADERS))]
        if hdr != HEADERS:
            # repair
            for ci, h in enumerate(HEADERS, start=2):
                c = ws.cell(2, ci, h)
                c.font = Font(bold=True)
                c.fill = PatternFill("solid", fgColor="DDDDDD")
        return wb, ws
    wb = Workbook()
    ws = wb.active
    ws.title = "comparison"
    for ci, h in enumerate(HEADERS, start=2):
        c = ws.cell(2, ci, h)
        c.font = Font(bold=True)
        c.fill = PatternFill("solid", fgColor="DDDDDD")
    return wb, ws


def main():
    if len(sys.argv) != 6:
        print("Usage: apply.py <en_path> <cn_path> <out_pptx> <out_xlsx> <edits.json>")
        sys.exit(1)
    en_path, cn_path, out_pptx, out_xlsx, edits_json = sys.argv[1:6]
    check_lock(out_pptx)
    check_lock(out_xlsx)

    with open(edits_json) as f:
        edits = json.load(f)
    if "edits" in edits:
        edits = edits["edits"]

    # Copy en to out
    shutil.copyfile(en_path, out_pptx)
    prs = Presentation(out_pptx)

    wb, ws = open_or_init_excel(out_xlsx)
    next_row = ws.max_row + 1 if ws.max_row > 2 else 3

    applied = 0
    for ed in edits:
        try:
            para = find_paragraph(prs, ed["slide"], ed["loc"])
            en_original = "".join(r.text for r in para.runs)
            replace_paragraph_text(para, ed["en_optimized"])
            ws.cell(next_row, 2, ed["slide"])
            ws.cell(next_row, 3, ed.get("kind", "text"))
            ws.cell(next_row, 4, ed.get("cn", ""))
            ws.cell(next_row, 5, en_original)
            ws.cell(next_row, 6, ed["en_optimized"])
            ws.cell(next_row, 7, ed.get("notes", ""))
            if ed.get("highlight"):
                red = Font(color="FF0000", bold=True)
                for c in range(2, 8):
                    ws.cell(next_row, c).font = red
            for c in range(4, 8):
                ws.cell(next_row, c).alignment = Alignment(wrap_text=True, vertical="top")
            applied += 1
            next_row += 1
        except Exception as e:
            print(f"  ⚠ FAILED edit slide={ed.get('slide')} loc={ed.get('loc')}: {e}")

    prs.save(out_pptx)
    wb.save(out_xlsx)

    # Readback
    wb2 = load_workbook(out_xlsx)
    ws2 = wb2.active
    hdr = [ws2.cell(2, c).value for c in range(2, 2 + len(HEADERS))]
    assert hdr == HEADERS, "Header readback failed!"
    print(f"Applied {applied}/{len(edits)} edits")
    print(f"  PPTX → {out_pptx}")
    print(f"  XLSX → {out_xlsx}")


if __name__ == "__main__":
    main()
