"""PDF → PPTX (每页一张幻灯片)

用法:
    python pdf_to_pptx.py input.pdf [output.pptx] [dpi]

参数:
    input.pdf   - PDF 文件路径
    output.pptx - 输出文件路径 (默认: output.pptx)
    dpi         - 渲染分辨率 (默认: 200)
"""
import fitz
import sys
import os
import io
from pptx import Presentation
from pptx.util import Inches


def main():
    if len(sys.argv) < 2:
        print("用法: python pdf_to_pptx.py input.pdf [output.pptx] [dpi]")
        sys.exit(1)

    pdf_path = sys.argv[1]
    out_path = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"
    dpi = int(sys.argv[3]) if len(sys.argv) > 3 else 200

    if not os.path.isfile(pdf_path):
        print(f"错误: 文件不存在 - {pdf_path}")
        sys.exit(1)

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=dpi)
        img_data = pix.tobytes("png")
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_data), Inches(0), Inches(0),
            width=prs.slide_width, height=prs.slide_height,
        )
        print(f"Page {i + 1}/{len(doc)} added")

    prs.save(out_path)
    total = len(doc)
    doc.close()
    print(f"完成: {total} 页 -> {out_path}")


if __name__ == "__main__":
    main()