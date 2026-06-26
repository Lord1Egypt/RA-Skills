#!/usr/bin/env python3
"""Sense Pass — reverse-engineer design DNA from a .pptx.

Usage:
    python3 sense_pass.py <path-to-pptx>

Outputs to stdout:
  - canvas dimensions and aspect
  - fill-color frequency (chrome / banner candidates)
  - font face usage
  - font size clusters with role inference
  - text colors
  - title-zone shape inventory (top ≤ 600K EMU)
  - layout heuristics (top margins, etc.)
  - shape type distribution
"""
import sys
from collections import Counter
from pptx import Presentation


def rgb_hex(color):
    try:
        r, g, b = color.rgb
        return f"#{r:02X}{g:02X}{b:02X}"
    except Exception:
        return None


def main(path):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    print(f"Canvas: {sw} × {sh} EMU = {sw/914400:.2f}\" × {sh/914400:.2f}\"")
    print(f"Aspect: {sw/sh:.4f} ({'16:9' if abs(sw/sh - 16/9) < 0.01 else 'other'})")
    print(f"Slides: {len(prs.slides)}\n")

    fonts = Counter()
    sizes = Counter()
    colors = Counter()
    bold_counter = Counter()
    fill_colors = Counter()
    shape_types = Counter()
    title_zone_shapes = []
    y_positions = []
    shape_widths = []
    shape_heights = []

    def walk_shapes(shapes, top_level=True):
        for shp in shapes:
            shape_types[shp.shape_type] += 1
            try:
                if int(shp.shape_type) == 6:  # GROUP
                    walk_shapes(shp.shapes, top_level=False)
                    continue
            except Exception:
                pass

            try:
                if hasattr(shp, "fill"):
                    f = shp.fill
                    if str(f.type) == "FillType.SOLID":
                        hx = rgb_hex(f.fore_color)
                        if hx:
                            fill_colors[hx] += 1
            except Exception:
                pass

            if top_level and shp.left is not None and shp.top is not None:
                y_positions.append(shp.top)
                if shp.width and shp.height:
                    shape_widths.append(shp.width)
                    shape_heights.append(shp.height)
                if 0 <= shp.top <= 600000 and shp.has_text_frame:
                    t = "".join(r.text for p in shp.text_frame.paragraphs for r in p.runs)
                    if t.strip():
                        title_zone_shapes.append((shp.top, t.strip()[:60]))

            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        if not r.text.strip():
                            continue
                        try:
                            f = r.font
                            name = f.name or "(default)"
                            fonts[name] += 1
                            if f.size is not None:
                                sizes[int(f.size.pt)] += 1
                            if f.bold:
                                bold_counter[name] += 1
                            if f.color and f.color.type is not None:
                                hx = rgb_hex(f.color)
                                if hx:
                                    colors[hx] += 1
                        except Exception:
                            pass

            if shp.has_table:
                for row in shp.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                if not r.text.strip():
                                    continue
                                try:
                                    f = r.font
                                    name = f.name or "(default)"
                                    fonts[name] += 1
                                    if f.size is not None:
                                        sizes[int(f.size.pt)] += 1
                                except Exception:
                                    pass

    for s in prs.slides:
        walk_shapes(s.shapes)

    print("=" * 70)
    print("A.1 — Fill Colors (chrome / band / banner candidates)")
    print("=" * 70)
    for c, n in fill_colors.most_common(15):
        print(f"  {c:10}  {n} occurrences")

    print("\n" + "=" * 70)
    print("A.2 — Font Faces in use")
    print("=" * 70)
    total = sum(fonts.values())
    for f, n in fonts.most_common(15):
        pct = 100 * n / total if total else 0
        bold_n = bold_counter.get(f, 0)
        suffix_flag = " ⚠ suffix-pollution" if any(s in f for s in ["Bold", "Regular", "Italic", "Light"]) else ""
        print(f"  {f:30}  {n:5} runs ({pct:5.1f}%)  bold:{bold_n}{suffix_flag}")

    print("\n" + "=" * 70)
    print("A.3 — Font Sizes (clusters)")
    print("=" * 70)
    total_size = sum(sizes.values())
    canvas_h_pt = sh / 12700
    print(f"Canvas height: {canvas_h_pt:.0f}pt")
    for sz, n in sorted(sizes.items()):
        if n < 5:
            continue
        pct = 100 * n / total_size if total_size else 0
        ratio = sz / canvas_h_pt
        role = ""
        if ratio > 0.10: role = "← hero?"
        elif ratio > 0.04: role = "← H1 title?"
        elif ratio > 0.025: role = "← H2 / data?"
        elif ratio > 0.013: role = "← body"
        elif ratio > 0.010: role = "← caption / fig label"
        else: role = "← chrome / source"
        print(f"  {sz:5}pt  {n:5} runs ({pct:5.1f}%)  ratio={ratio:.4f}  {role}")

    print("\n" + "=" * 70)
    print("A.4 — Text Colors")
    print("=" * 70)
    for c, n in colors.most_common(10):
        print(f"  {c:10}  {n} runs")

    print("\n" + "=" * 70)
    print("A.5 — Title-zone shapes (top ≤ 600000 EMU)")
    print("=" * 70)
    for top, t in title_zone_shapes[:20]:
        print(f"  top={top:8}  {t}")

    print("\n" + "=" * 70)
    print("A.6 — Layout heuristics")
    print("=" * 70)
    if y_positions:
        y_sorted = sorted(y_positions)
        pct5 = y_sorted[len(y_sorted) // 20]
        print(f"  Top-margin estimate (5th percentile): {pct5} EMU = {pct5/sh*100:.1f}% of canvas")
        print(f"  Top-shape y values (first 10): {y_sorted[:10]}")

    print("\n" + "=" * 70)
    print("Shape type distribution")
    print("=" * 70)
    for st, n in shape_types.most_common(10):
        print(f"  type={st}  count={n}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: sense_pass.py <pptx-path>")
        sys.exit(1)
    main(sys.argv[1])
