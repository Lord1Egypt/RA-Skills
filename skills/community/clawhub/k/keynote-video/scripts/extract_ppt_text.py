#!/usr/bin/env python3
"""
PPTX/PDF 文本提取工具

提取 PPTX 或 PDF 文件中每页的标题和正文内容，输出为 Markdown 格式。
供 ppt-video v2.0 的 Phase 1 使用。

用法:
    python3 extract_ppt_text.py <input_file> [> output.md]
"""

import sys
import os

def extract_pptx_text(pptx_path):
    """从 PPTX 文件提取每页文字"""
    try:
        from pptx import Presentation
    except ImportError:
        print("❌ 缺少依赖: pip install python-pptx", file=sys.stderr)
        sys.exit(1)
    
    prs = Presentation(pptx_path)
    slides = []
    
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        title = ""
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # 尝试识别标题（通常是第一个文本框或字体较大）
                        if not title and texts == []:
                            title = text
                        else:
                            texts.append(text)
            
            # 也提取表格中的文字
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            texts.append(text)
        
        slides.append({
            'num': i,
            'title': title or f"第{i}页",
            'body': '\n'.join(texts) if texts else "[仅图片/图表，无文字内容]"
        })
    
    return slides

def extract_pdf_text(pdf_path):
    """从 PDF 文件提取每页文字"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        try:
            from pdfminer.high_level import extract_text
            # pdfminer 不支持逐页标题提取，降级处理
            text = extract_text(pdf_path)
            return [{'num': 1, 'title': 'PDF 内容', 'body': text}]
        except ImportError:
            print("❌ 缺少依赖: pip install PyMuPDF 或 pip install pdfminer.six", file=sys.stderr)
            sys.exit(1)
    
    doc = fitz.open(pdf_path)
    slides = []
    
    for i in range(len(doc)):
        page = doc[i]
        text = page.get_text()
        lines = [l.strip() for l in text.split('\n') if l.strip()]
        
        title = lines[0] if lines else f"第{i+1}页"
        body = '\n'.join(lines[1:]) if len(lines) > 1 else "[仅图片/图表，无文字内容]"
        
        slides.append({
            'num': i + 1,
            'title': title,
            'body': body
        })
    
    doc.close()
    return slides

def format_as_markdown(slides):
    """格式化为 Markdown"""
    output = []
    
    for slide in slides:
        output.append(f"## Slide {slide['num']}")
        output.append(f"[标题]: {slide['title']}")
        output.append(f"[正文]: {slide['body']}")
        output.append("")
    
    return '\n'.join(output)

def main():
    if len(sys.argv) < 2:
        print("用法: python3 extract_ppt_text.py <input_file>", file=sys.stderr)
        sys.exit(1)
    
    input_path = sys.argv[1]
    
    if not os.path.exists(input_path):
        print(f"❌ 文件不存在: {input_path}", file=sys.stderr)
        sys.exit(1)
    
    ext = os.path.splitext(input_path)[1].lower()
    
    if ext in ('.pptx', '.ppt'):
        slides = extract_pptx_text(input_path)
    elif ext == '.pdf':
        slides = extract_pdf_text(input_path)
    else:
        print(f"❌ 不支持的文件格式: {ext}", file=sys.stderr)
        sys.exit(1)
    
    print(format_as_markdown(slides))

if __name__ == '__main__':
    main()
