"""
DeckCraft v5 — Core Drawing Primitives
Low-level helpers for shapes, text, backgrounds, and XML cleanup.
"""

import os
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from .constants import (
    FONT_CN, FONT_EN_TITLE, FONT_EN_BODY,
    MARGIN_LEFT, GAP_SMALL, c
)


def set_bg(slide, color_tuple):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = c(color_tuple)


def add_rect(slide, left, top, width, height, color_tuple, alpha=None):
    """Add a filled rectangle (no border). Returns shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(left), int(top), int(width), int(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = c(color_tuple)
    shape.line.fill.background()
    if alpha is not None:
        _set_shape_alpha(shape, alpha)
    return shape


def add_rounded_rect(slide, left, top, width, height, color_tuple, alpha=None):
    """Add a filled rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        int(left), int(top), int(width), int(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = c(color_tuple)
    shape.line.fill.background()
    if alpha is not None:
        _set_shape_alpha(shape, alpha)
    return shape


def add_textbox(slide, left, top, width, height, text, size=14,
                bold=False, color=(51, 51, 51), align=PP_ALIGN.LEFT,
                font_name=None, word_wrap=True, anchor=MSO_ANCHOR.TOP):
    """Add a text box with formatted text. Returns textbox shape."""
    txBox = slide.shapes.add_textbox(
        int(left), int(top), int(width), int(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    try:
        tf.auto_size = None
    except Exception:
        pass

    # Set vertical anchor
    txBox_elem = txBox._element
    bodyPr = txBox_elem.find(qn('p:txBody')).find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))

    p = tf.paragraphs[0]
    _write_run(p, text, size, bold, color, font_name)
    p.alignment = align
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          size=14, color=(51, 51, 51), bold=False,
                          align=PP_ALIGN.LEFT, font_name=None,
                          bullet=False, spacing_pt=6):
    """Add a text box with multiple paragraphs. Returns textbox shape."""
    txBox = slide.shapes.add_textbox(
        int(left), int(top), int(width), int(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        display_text = f"• {line_text}" if bullet else line_text
        _write_run(p, display_text, size, bold, color, font_name)
        p.alignment = align
        p.space_after = Pt(spacing_pt)

    return txBox


def add_line(slide, left, top, width, color_tuple, thickness=0.04):
    """Add a horizontal decorative line (rectangle, not connector)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(left), int(top), int(width), int(Inches(thickness))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = c(color_tuple)
    shape.line.fill.background()
    return shape


def add_circle(slide, cx, cy, radius, color_tuple):
    """Add a filled circle. cx/cy/radius in Inches-compatible units."""
    d = int(radius * 2)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(cx - radius), int(cy - radius), d, d
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = c(color_tuple)
    shape.line.fill.background()
    return shape


def add_image_safe(slide, image_path, left, top, width=None, height=None):
    """Add image to slide. Returns shape or None if file missing."""
    if not image_path or not os.path.exists(image_path):
        return None
    try:
        kwargs = {"image_file": image_path, "left": int(left), "top": int(top)}
        if width is not None:
            kwargs["width"] = int(width)
        if height is not None:
            kwargs["height"] = int(height)
        return slide.shapes.add_picture(**kwargs)
    except Exception:
        return None


# ── XML Cleanup ───────────────────────────────────────────────────

def clean_shape(shape):
    """Remove shadow, 3D, and other artifacts from a shape."""
    sp = shape._element
    # Remove effectLst (shadows, glows, reflections)
    for tag in ['a:effectLst', 'a:sp3d', 'a:ln', 'p:style']:
        for elem in sp.findall(qn(tag)):
            sp.remove(elem)


def clean_slide(slide):
    """Clean all shapes on a slide."""
    for shape in slide.shapes:
        clean_shape(shape)


def full_cleanup(prs):
    """Clean all shapes on all slides. Call after save for safety."""
    for slide in prs.slides:
        clean_slide(slide)


def set_ea_font(shape, font_name=FONT_CN):
    """Set East Asian font on all text runs in a shape."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is None:
                rPr = run._r.makeelement(qn('a:rPr'), {})
                run._r.insert(0, rPr)
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {})
                rPr.append(ea)
            ea.set('typeface', font_name)


def auto_cjk_font(shape):
    """Auto-detect CJK characters and apply East Asian font."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            text = run.text
            if _has_cjk(text):
                set_ea_font(shape)
                return  # once per shape is enough


# ── Internal Helpers ──────────────────────────────────────────────

def _write_run(paragraph, text, size, bold, color, font_name=None):
    """Clear paragraph and write a single formatted run."""
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    # Ensure size is in EMU (Pt() returns EMU) — if already EMU, pass through
    if isinstance(size, (int, float)) and size > 1000:
        # Already in hundredths of a point (EMU-like) or raw value
        run.font.size = int(size)
    else:
        run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = c(color) if isinstance(color, tuple) else color
    fn = font_name or (_detect_font(text))
    run.font.name = fn


def _detect_font(text):
    """Choose CN or EN font based on content."""
    if _has_cjk(text):
        return FONT_CN
    return FONT_EN_BODY


def _has_cjk(text):
    """Check if text contains CJK characters."""
    for ch in text:
        if '\u4e00' <= ch <= '\u9fff' or '\u3400' <= ch <= '\u4dbf':
            return True
    return False


def _set_shape_alpha(shape, alpha_fraction):
    """Set shape transparency (0=opaque, 1=invisible)."""
    try:
        solidFill = shape.fill._fill
        srgb = solidFill.find(qn('a:solidFill'))
        if srgb is not None:
            clr = srgb.find(qn('a:srgbClr'))
            if clr is not None:
                alpha_elem = clr.makeelement(qn('a:alpha'), {})
                alpha_elem.set('val', str(int(alpha_fraction * 1000)))
                clr.append(alpha_elem)
    except Exception:
        pass
