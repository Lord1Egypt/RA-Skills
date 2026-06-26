#!/usr/bin/env python3
"""
export_pptx.py — Generate a clean scientific presentation PPTX from a JSON input.
Usage:
    python3 export_pptx.py --input paper.json --output slides.pptx
"""

import argparse
import json
import os
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0E, 0x4D, 0x92)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

# ── Slide dimensions: 16:9 widescreen ────────────────────────────────────────
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

# ── Layout constants ──────────────────────────────────────────────────────────
TITLE_BAR_H    = Cm(2.8)
MARGIN         = Cm(1.2)
SLIDE_NUM_W    = Cm(1.5)
SLIDE_NUM_H    = Cm(0.6)
BODY_TOP       = TITLE_BAR_H + Cm(0.3)
BODY_H         = SLIDE_H - BODY_TOP - Cm(1.0)


def _blank_slide(prs: Presentation) -> object:
    """Add a blank slide (layout 6) to the presentation."""
    blank_layout = prs.slide_layouts[6]  # "Blank" layout
    return prs.slides.add_slide(blank_layout)


def _set_background_white(slide) -> None:
    """Ensure slide background is white (already default for blank layout)."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def _add_title_bar(slide, title_text: str) -> None:
    """Draw the navy title bar at the top and place title text inside it."""
    # Navy rectangle
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, TITLE_BAR_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()  # no border

    # Title text box on top of bar
    txb = slide.shapes.add_textbox(MARGIN, Cm(0.2), SLIDE_W - 2 * MARGIN, TITLE_BAR_H - Cm(0.4))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = WHITE


def _add_slide_number(slide, num: int) -> None:
    """Add slide number at bottom right."""
    left = SLIDE_W - MARGIN - SLIDE_NUM_W
    top  = SLIDE_H - Cm(0.8)
    txb = slide.shapes.add_textbox(left, top, SLIDE_NUM_W, SLIDE_NUM_H)
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def _add_body_textbox(slide, text: str, left=None, top=None, width=None, height=None,
                      font_size=Pt(20), align=PP_ALIGN.LEFT, bold=False, color=DARK_GRAY):
    """Generic helper: add a text box with styled content."""
    left   = left   if left   is not None else MARGIN
    top    = top    if top    is not None else BODY_TOP
    width  = width  if width  is not None else SLIDE_W - 2 * MARGIN
    height = height if height is not None else BODY_H

    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for line in text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
        # Indent bullet lines (those starting with • or -)
        if line.strip().startswith(("•", "-", "·")):
            p.level = 1


def _add_figure_placeholder(slide, left, top, width, height, label: str = "Insert Figure here"):
    """Gray rectangle with centered label text."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    txb = slide.shapes.add_textbox(left, top + height // 2 - Cm(0.5), width, Cm(1.0))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run.font.italic = True


# ── Slide builders ────────────────────────────────────────────────────────────

def build_title_slide(prs, data: dict, slide_num: int) -> None:
    slide = _blank_slide(prs)
    _set_background_white(slide)

    # Full-width navy bar (taller for title slide)
    bar_h = Cm(5.5)
    bar = slide.shapes.add_shape(1, 0, Cm(5.5), SLIDE_W, bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Paper title
    txb = slide.shapes.add_textbox(MARGIN, Cm(5.8), SLIDE_W - 2 * MARGIN, Cm(4.5))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = data.get("title", "Untitled")
    run.font.bold = True
    run.font.size = Pt(32)
    run.font.color.rgb = WHITE

    # Authors
    authors = ", ".join(data.get("authors", []))
    txb2 = slide.shapes.add_textbox(MARGIN, Cm(11.5), SLIDE_W - 2 * MARGIN, Cm(1.2))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = authors
    run2.font.size = Pt(18)
    run2.font.color.rgb = DARK_GRAY

    # Journal / Year / DOI
    journal = data.get("journal", "")
    year    = data.get("year", "")
    doi     = data.get("doi", "")
    meta    = f"{journal}  •  {year}"
    if doi:
        meta += f"  |  DOI: {doi}"
    txb3 = slide.shapes.add_textbox(MARGIN, Cm(13.0), SLIDE_W - 2 * MARGIN, Cm(1.0))
    tf3 = txb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = meta
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    run3.font.italic = True

    _add_slide_number(slide, slide_num)


def build_background_slide(prs, slide_data: dict, slide_num: int) -> None:
    slide = _blank_slide(prs)
    _set_background_white(slide)
    title_text = slide_data.get("title", "Background / Introduction")
    _add_title_bar(slide, title_text)
    content = slide_data.get("content", "")
    _add_body_textbox(slide, content, font_size=Pt(20), color=DARK_GRAY)
    _add_slide_number(slide, slide_num)


def build_results_slide(prs, slide_data: dict, slide_num: int, fig_index: int) -> None:
    slide = _blank_slide(prs)
    _set_background_white(slide)
    title_text = slide_data.get("title", "Results")
    _add_title_bar(slide, title_text)

    half_w = (SLIDE_W - 3 * MARGIN) // 2

    # Left: text content
    content = slide_data.get("content", "")
    _add_body_textbox(slide, content,
                      left=MARGIN, top=BODY_TOP,
                      width=half_w, height=BODY_H,
                      font_size=Pt(18), color=DARK_GRAY)

    # Right: figure placeholder
    fig_label = slide_data.get("figure", f"fig{fig_index}.png")
    placeholder_label = f"Insert Figure {fig_index} here\n({fig_label})"
    fig_left = MARGIN + half_w + MARGIN
    fig_top  = BODY_TOP + Cm(0.3)
    fig_h    = BODY_H - Cm(0.6)
    _add_figure_placeholder(slide, fig_left, fig_top, half_w, fig_h, placeholder_label)

    _add_slide_number(slide, slide_num)


def build_conclusion_slide(prs, slide_data: dict, slide_num: int) -> None:
    slide = _blank_slide(prs)
    _set_background_white(slide)
    title_text = slide_data.get("title", "Conclusion")
    _add_title_bar(slide, title_text)
    content = slide_data.get("content", "")
    _add_body_textbox(slide, content, font_size=Pt(20), color=DARK_GRAY)
    _add_slide_number(slide, slide_num)


def build_generic_slide(prs, slide_data: dict, slide_num: int) -> None:
    slide = _blank_slide(prs)
    _set_background_white(slide)
    title_text = slide_data.get("title", slide_data.get("type", "Slide").capitalize())
    _add_title_bar(slide, title_text)
    content = slide_data.get("content", "")
    _add_body_textbox(slide, content, font_size=Pt(20), color=DARK_GRAY)
    _add_slide_number(slide, slide_num)


# ── Main builder ──────────────────────────────────────────────────────────────

def build_presentation(data: dict, output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides_spec = data.get("slides", [])
    # Ensure at least a title slide if none specified
    if not slides_spec:
        slides_spec = [{"type": "title"}]

    slide_num   = 1
    fig_counter = 1

    for spec in slides_spec:
        stype = spec.get("type", "generic").lower()

        if stype == "title":
            build_title_slide(prs, data, slide_num)
        elif stype in ("background", "introduction", "intro"):
            build_background_slide(prs, spec, slide_num)
        elif stype in ("results", "result", "methods", "method", "discussion"):
            build_results_slide(prs, spec, slide_num, fig_counter)
            if spec.get("figure"):
                fig_counter += 1
        elif stype in ("conclusion", "conclusions", "summary"):
            build_conclusion_slide(prs, spec, slide_num)
        else:
            build_generic_slide(prs, spec, slide_num)

        slide_num += 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


# ── CLI ───────────────────────────────────────────────────────────────────────

def parse_args():
    today = datetime.now().strftime("%Y%m%d")
    default_input  = Path.home() / "Downloads" / "nature-paper-draft.json"
    default_output = Path.home() / "Downloads" / f"paper-slides-{today}.pptx"

    parser = argparse.ArgumentParser(
        description="Generate a scientific PPTX presentation from a Nature paper JSON draft."
    )
    parser.add_argument("--input",  default=str(default_input),
                        help=f"Path to input JSON (default: {default_input})")
    parser.add_argument("--output", default=str(default_output),
                        help=f"Path to output PPTX (default: {default_output})")
    return parser.parse_args()


def main():
    args = parse_args()
    input_path  = Path(args.input).expanduser()
    output_path = Path(args.output).expanduser()

    if not input_path.exists():
        print(f"❌ Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    with open(input_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    build_presentation(data, output_path)
    print(f"✅ PPT export successful: {output_path}")


if __name__ == "__main__":
    main()
