#!/usr/bin/env python3
"""Distill a style fingerprint from a reference PDF or .pptx.

Usage:
    python3 style_distill.py <path-to-pdf-or-pptx> [--weight 0.7]

Output: JSON {
  source, weight,
  cadence: {avg_sentence_len, avg_paragraph_len, sentence_len_p90},
  vocab: {signature_phrases[], top_action_verbs[], filler_words_found{}},
  structure: {title_pattern, bullet_pattern, parallel_score},
  tone: {first_person_ratio, hedge_ratio, certainty_ratio},
}

Use the output as input to PROFILE.style_references — it layers on top of
the McKinsey baseline at translation time.

Implementation: PyMuPDF for PDFs, python-pptx for decks. If PyMuPDF missing,
prints install hint and exits.
"""
import sys
import json
import os
import re
from collections import Counter


HEDGES = {"might", "may", "could", "perhaps", "possibly", "likely",
          "tends to", "generally", "often", "usually"}
CERTAIN = {"will", "must", "always", "never", "clearly", "definitively",
           "the answer", "the reason"}
FILLERS = ["in order to", "a number of", "due to the fact that",
           "at this point in time", "it should be noted that",
           "for the purpose of", "with respect to"]


def extract_text_from_pdf(path):
    try:
        import fitz
    except ImportError:
        print("Requires PyMuPDF: pip3 install pymupdf")
        sys.exit(1)
    doc = fitz.open(path)
    text = []
    for page in doc:
        text.append(page.get_text())
    return "\n".join(text)


def extract_text_from_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    def walk(shapes):
        for shp in shapes:
            try:
                if int(shp.shape_type) == 6:
                    walk(shp.shapes); continue
            except Exception: pass
            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs)
                    if t.strip(): parts.append(t.strip())
            if shp.has_table:
                for row in shp.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            t = "".join(r.text for r in p.runs)
                            if t.strip(): parts.append(t.strip())
    for s in prs.slides:
        walk(s.shapes)
    return "\n".join(parts)


def split_sentences(text):
    return [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]


def distill(text, weight):
    sentences = split_sentences(text)
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

    sentence_lens = [len(s.split()) for s in sentences]
    para_lens = [len(p.split()) for p in paragraphs]
    sentence_lens.sort()

    # vocab
    words = re.findall(r'\b[A-Za-z]+\b', text.lower())
    word_count = Counter(words)
    # heuristic action verbs (-ing/-ed/-s common; skip first-letter constraints)
    common_verbs = {"build", "drive", "scale", "deliver", "ship", "launch",
                    "transform", "unlock", "reframe", "rebuild", "accelerate",
                    "compress", "redefine", "reshape", "leverage", "execute"}
    top_action = [v for v in common_verbs if word_count[v] >= 2]

    filler_found = {}
    low = text.lower()
    for fil in FILLERS:
        c = low.count(fil)
        if c > 0:
            filler_found[fil] = c

    # signature phrases: 3-grams appearing ≥3 times
    tokens = re.findall(r'\b\w+\b', text.lower())
    trigrams = Counter()
    for i in range(len(tokens) - 2):
        tg = tuple(tokens[i:i+3])
        if all(len(t) >= 3 for t in tg):
            trigrams[tg] += 1
    sig_phrases = [" ".join(tg) for tg, c in trigrams.most_common(20) if c >= 3]

    # structure: bullet pattern detection
    bullet_lines = [p for p in paragraphs if re.match(r'^[•\-\*■▪◆]', p)]
    parallel_starts = Counter()
    for b in bullet_lines:
        first_word = re.findall(r'\w+', b)
        if first_word:
            parallel_starts[first_word[0].lower()] += 1
    parallel_score = (max(parallel_starts.values()) / len(bullet_lines)
                      if bullet_lines else 0.0)

    # tone
    hedge_count = sum(low.count(h) for h in HEDGES)
    certain_count = sum(low.count(c) for c in CERTAIN)
    first_person = sum(word_count[p] for p in ["we", "our", "us"])
    word_total = sum(word_count.values())

    return {
        "source": os.path.basename(text[:0]) or "<inline>",
        "weight": weight,
        "cadence": {
            "avg_sentence_len": (sum(sentence_lens) / len(sentence_lens)
                                 if sentence_lens else 0),
            "sentence_len_p90": (sentence_lens[int(len(sentence_lens) * 0.9)]
                                 if sentence_lens else 0),
            "avg_paragraph_len": (sum(para_lens) / len(para_lens)
                                  if para_lens else 0),
        },
        "vocab": {
            "signature_phrases": sig_phrases[:10],
            "top_action_verbs": top_action,
            "filler_words_found": filler_found,
        },
        "structure": {
            "bullet_line_count": len(bullet_lines),
            "parallel_score": round(parallel_score, 2),
            "common_bullet_start_words": parallel_starts.most_common(5),
        },
        "tone": {
            "first_person_ratio": (first_person / word_total if word_total else 0),
            "hedge_per_1k": (1000 * hedge_count / word_total if word_total else 0),
            "certainty_per_1k": (1000 * certain_count / word_total if word_total else 0),
        },
    }


def main():
    if len(sys.argv) < 2:
        print("Usage: style_distill.py <path-to-pdf-or-pptx> [--weight 0.7]")
        sys.exit(1)
    path = sys.argv[1]
    weight = 0.7
    if "--weight" in sys.argv:
        weight = float(sys.argv[sys.argv.index("--weight") + 1])

    if path.lower().endswith(".pdf"):
        text = extract_text_from_pdf(path)
    elif path.lower().endswith(".pptx"):
        text = extract_text_from_pptx(path)
    else:
        print("Unsupported format. Use .pdf or .pptx")
        sys.exit(1)

    result = distill(text, weight)
    result["source"] = path
    print(json.dumps(result, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
