"""
yumnb slide deck generator.

Reads a JSON spec describing slides, emits a .pptx (16:9).

Slide types:
  - title       {title, subtitle?}
  - bullets     {title, bullets[], notes?}
  - table       {title, headers[], rows[][]}
  - flow        {title, steps[]}            (horizontal arrow flow, wraps after 4)
  - image       {title?, image_path, caption?}
  - two_column  {title, left, image_path}   (text left, image right)
  - summary     {title, text}
"""
from __future__ import annotations

import argparse
import json
import os
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Pick CJK-friendly fonts that ship with both Windows and most Linux/macOS
# setups. python-pptx just stores font names — the renderer (PowerPoint /
# LibreOffice / Keynote) substitutes if missing.
CN_FONT = "Microsoft YaHei"
EN_FONT = "Segoe UI"

PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_PPTX_NATIVE_EXTS = {".bmp", ".gif", ".jpg", ".jpeg", ".png", ".tif", ".tiff", ".wmf"}
_CONVERTED_CACHE: dict = {}


def _ensure_pptx_image(img_path: str) -> str:
    """Return a path to an image python-pptx can embed.

    Converts WEBP / AVIF / HEIC / SVG / … to PNG via Pillow if available.
    """
    if not img_path or not os.path.isfile(img_path):
        return img_path
    ext = os.path.splitext(img_path)[1].lower()
    if ext in _PPTX_NATIVE_EXTS:
        return img_path
    if img_path in _CONVERTED_CACHE and os.path.isfile(_CONVERTED_CACHE[img_path]):
        return _CONVERTED_CACHE[img_path]
    try:
        from PIL import Image  # type: ignore

        with Image.open(img_path) as im:
            im = im.convert("RGBA") if im.mode in ("RGBA", "LA", "P") else im.convert("RGB")
            fd, out = tempfile.mkstemp(suffix=".png", prefix="yumnb_img_")
            os.close(fd)
            im.save(out, "PNG")
            _CONVERTED_CACHE[img_path] = out
            return out
    except Exception:
        return img_path


def _set_run(run, text, size=18, bold=False, color=None, font=None):
    run.text = text
    run.font.name = font or CN_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                 color=None, align=None, anchor=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    _set_run(run, text, size=size, bold=bold, color=color)
    return tb


def _add_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    _set_run(run, title, size=24, bold=True, color=WHITE)


def slide_title(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = PRIMARY; bg.line.fill.background()
    _add_textbox(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.6),
                 spec.get("title", ""), size=44, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT)
    sub = spec.get("subtitle", "")
    if sub:
        _add_textbox(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.8),
                     sub, size=20, color=WHITE, align=PP_ALIGN.LEFT)
    _add_textbox(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5),
                 "yumnb · auto-generated", size=12, color=WHITE, align=PP_ALIGN.LEFT)


def slide_bullets(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, spec.get("title", ""))
    tb = s.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(spec.get("bullets", []) or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        _set_run(run, f"●  {b}", size=22, color=DARK_TEXT)
        p.space_after = Pt(10)
    notes = spec.get("notes")
    if notes:
        s.notes_slide.notes_text_frame.text = notes


def slide_table(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, spec.get("title", ""))
    headers = spec.get("headers", []) or []
    rows = spec.get("rows", []) or []
    if not headers:
        return
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_w = Inches(12.1)
    tbl_h = Inches(min(0.6 * n_rows, 5.5))
    tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.3), tbl_w, tbl_h)
    table = tbl_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        tf = cell.text_frame; tf.text = ""
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), str(h), size=16, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            if j >= n_cols:
                break
            cell = table.cell(i, j)
            if i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BG
            tf = cell.text_frame; tf.text = ""; tf.word_wrap = True
            p = tf.paragraphs[0]
            _set_run(p.add_run(), str(val), size=14, color=DARK_TEXT)


def slide_flow(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, spec.get("title", ""))
    steps = spec.get("steps", []) or []
    n = len(steps)
    if n == 0:
        return
    per_row = min(n, 4)
    top_start = Inches(1.7)
    row_h = Inches(2.0)
    margin_x = Inches(0.4)
    avail_w = SLIDE_W - margin_x * 2
    arrow_gap = Inches(0.25)
    for idx, step in enumerate(steps):
        r = idx // per_row
        c = idx % per_row
        items_this_row = min(per_row, n - r * per_row)
        actual_w = (avail_w - arrow_gap * (items_this_row - 1)) / items_this_row
        x = margin_x + c * (actual_w + arrow_gap)
        y = top_start + Emu(int(row_h * r))
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   int(x), int(y), int(actual_w), Inches(1.4))
        shape.fill.solid(); shape.fill.fore_color.rgb = PRIMARY
        shape.line.color.rgb = PRIMARY
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), f"{idx+1}. {step}", size=14, bold=True, color=WHITE)
        if c < items_this_row - 1:
            ax = int(x + actual_w)
            ay = int(y) + Inches(0.55)
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay,
                                     int(arrow_gap), Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT
            arr.line.fill.background()


def slide_image(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title = spec.get("title")
    if title:
        _add_title_bar(s, title)
        top = Inches(1.2); avail_h = Inches(5.5)
    else:
        top = Inches(0.4); avail_h = Inches(6.3)
    img_path = spec.get("image_path")
    if img_path and os.path.isfile(img_path):
        try:
            img_path = _ensure_pptx_image(img_path)
            pic = s.shapes.add_picture(img_path, Inches(0.6), top, height=avail_h)
            pic.left = int((SLIDE_W - pic.width) / 2)
            if pic.width > SLIDE_W - Inches(1.2):
                pic.width = SLIDE_W - Inches(1.2)
                pic.left = int(Inches(0.6))
        except Exception as e:
            _add_textbox(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(2),
                         f"[image error: {e}]", size=16, color=ACCENT)
    cap = spec.get("caption")
    if cap:
        _add_textbox(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.5),
                     cap, size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)


def slide_two_column(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(s, spec.get("title", ""))
    left = spec.get("left", "")
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.8))
    tf = tb.text_frame; tf.word_wrap = True
    lines = left.split("\n") if isinstance(left, str) else (left or [])
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        if ln.strip().startswith(("●", "•", "-", "*")):
            _set_run(run, ln.strip(), size=18, color=DARK_TEXT)
        else:
            _set_run(run, f"●  {ln.strip()}" if ln.strip() else "", size=18, color=DARK_TEXT)
        p.space_after = Pt(6)
    img_path = spec.get("image_path")
    if img_path and os.path.isfile(img_path):
        try:
            img_path = _ensure_pptx_image(img_path)
            s.shapes.add_picture(img_path, Inches(7.0), Inches(1.4), width=Inches(5.8))
        except Exception as e:
            _add_textbox(s, Inches(7.0), Inches(3.0), Inches(5.8), Inches(1),
                         f"[image error: {e}]", size=14, color=ACCENT)


def slide_summary(prs, spec):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG; bg.line.fill.background()
    _add_textbox(s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.9),
                 spec.get("title", "Summary"), size=28, bold=True, color=PRIMARY)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5),
                             Inches(2.5), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    _add_textbox(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5),
                 spec.get("text", ""), size=24, color=DARK_TEXT)


BUILDERS = {
    "title": slide_title,
    "bullets": slide_bullets,
    "table": slide_table,
    "flow": slide_flow,
    "image": slide_image,
    "two_column": slide_two_column,
    "summary": slide_summary,
}


def build(spec_path, output_path):
    spec = json.loads(Path(spec_path).read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = spec.get("slides", [])
    if not slides or slides[0].get("type") != "title":
        slide_title(prs, {"title": spec.get("title", "Notes"),
                          "subtitle": spec.get("subtitle", "")})

    for sp in slides:
        fn = BUILDERS.get(sp.get("type"))
        if fn is None:
            print(f"WARN: unknown slide type: {sp.get('type')}")
            continue
        try:
            fn(prs, sp)
        except Exception as e:
            print(f"WARN: slide {sp.get('type')} failed: {e}")

    prs.save(output_path)
    print(f"PPTX: {output_path} ({os.path.getsize(output_path):,} bytes, {len(prs.slides)} slides)")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--json", required=True, help="Path to slide spec JSON")
    ap.add_argument("--output", required=True, help="Output .pptx path")
    args = ap.parse_args()
    build(args.json, args.output)


if __name__ == "__main__":
    main()
